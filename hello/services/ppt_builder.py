from __future__ import annotations

import json
from io import BytesIO
from typing import Iterable

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from hello import models
from hello.services.pptx_update import update_chart as _update_chart
from hello.utils.utils import format_label


# Visual constants
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)  # 16:9
MARGIN_L, MARGIN_T, MARGIN_R, MARGIN_B = Inches(0.6), Inches(0.6), Inches(0.6), Inches(0.6)
TITLE_SIZE = Pt(28)
SUBTITLE_SIZE = Pt(16)
BODY_SIZE = Pt(12)
CODE_SIZE = Pt(10)


def _blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if getattr(layout, "name", "").strip().lower() == "blank":
            return layout
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]


def _textbox(slide, left, top, width, height, text, *, size=BODY_SIZE, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(text) if text is not None else ""
    p.font.size = size
    p.font.bold = bold
    p.alignment = align
    tf.word_wrap = True
    return shape


def _codebox(slide, left, top, width, height, title: str, text: str):
    _textbox(slide, left, top - Inches(0.3), width, Inches(0.28), title, size=SUBTITLE_SIZE, bold=True)
    box = slide.shapes.add_textbox(left, top, width, height)
    # subtle background + border
    try:
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(245, 245, 245)
        box.line.color.rgb = RGBColor(210, 210, 210)
    except Exception:
        pass
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text or ""
    p.font.size = CODE_SIZE
    # p.font.name cannot guarantee availability; rely on default
    tf.word_wrap = True
    return box


def _chart(slide, left, top, width, height, title: str, df: pd.DataFrame, x: str, ys: Iterable[str], kind: str = "line_markers"):
    _textbox(slide, left, top - Inches(0.3), width, Inches(0.28), title, size=SUBTITLE_SIZE, bold=True)
    # seed with basic chart data, then replace using our richer updater
    chart_data = CategoryChartData()
    chart_data.categories = [str(v) for v in df[x]]
    for y in ys:
        series = pd.to_numeric(df[y], errors="coerce").fillna(0).tolist()
        label = str(y)
        # format_label(str(y))
        chart_data.add_series(label, series)
    ct = {
        "line": XL_CHART_TYPE.LINE,
        "line_markers": XL_CHART_TYPE.LINE_MARKERS,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
        "bar_stacked_100": XL_CHART_TYPE.BAR_STACKED_100,
        "area": XL_CHART_TYPE.AREA,
        "pie": XL_CHART_TYPE.PIE,
        "donut": XL_CHART_TYPE.DOUGHNUT,
        "pie_chart": XL_CHART_TYPE.PIE,
        "donut_chart": XL_CHART_TYPE.DOUGHNUT,
    }.get(kind, XL_CHART_TYPE.LINE_MARKERS)
    chart_shape = slide.shapes.add_chart(ct, left, top, width, height, chart_data)

    # Use the presentation-style updater to preserve formatting and scale axes
    try:
        # Build a DataFrame for updater: first column categories, remainder series
        data = pd.DataFrame({x: df[x]})
        for y in ys:
            # Preserve original column but sanitized label is handled when chart data replaced
            data[str(y)] = pd.to_numeric(df[y], errors="coerce")
        _update_chart(chart_shape.chart, data, dynamic_y_axis=True)
    except Exception as _e:
        # Non-fatal; keep the basic chart if anything goes wrong
        pass


def _table(slide, left, top, width, height, title: str, df: pd.DataFrame):
    _textbox(slide, left, top - Inches(0.3), width, Inches(0.28), title, size=SUBTITLE_SIZE, bold=True)
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    # header
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j); cell.text = str(col)
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True; para.font.size = BODY_SIZE
        para.alignment = PP_ALIGN.CENTER
    # body
    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            cell = table.cell(i + 1, j)
            val = df.iat[i, j]
            cell.text = "" if pd.isna(val) else str(val)
            cell.text_frame.paragraphs[0].font.size = BODY_SIZE


def _dummy_chart_df() -> pd.DataFrame:
    return pd.DataFrame({
        "Quarter": ["2023 Q1", "2023 Q2", "2023 Q3", "2023 Q4", "2024 Q1", "2024 Q2"],
        "Vacancy": [18.6, 19.2, 20.1, 20.7, 21.2, 22.1],
        "Absorption": [120, -80, 50, -30, 10, -149],
    })


def _dummy_table_df() -> pd.DataFrame:
    return pd.DataFrame({
        "Metric": ["Vacancy Rate", "Leasing Activity", "Net Absorption"],
        "Current": ["21.9%", "542,721", "50,328"],
        "Prior": ["22.1%", "560,498", "-148,717"],
    })


def _commentary_for_section(report: models.Report, sec: models.ReportSection) -> tuple[str, str]:
    """Return (commentary_text, aggregated_commentary_json_pretty).

    The aggregated JSON mirrors the shape you described, combining:
    - section defaults (name, template, property_type)
    - prompt fields from the commentary element config
    - an array of chart SQLs and table SQLs from other elements
    - plus the actual commentary text if present (under key "commentary")
    """
    text = ""
    default_prompt = ""
    adjust_prompt = ""
    charts_sql: list[str] = []
    tables_sql: list[str] = []

    for el in sorted(sec.elements or [], key=lambda e: e.display_order or 0):
        et = (el.element_type or "").lower()
        cfg = el.config or {}
        if et == "commentary":
            # Strict requirement: use only section_commentary column as the text
            sc = getattr(el, "section_commentary", None)
            text = (sc or "").strip()
            if isinstance(cfg, dict):
                default_prompt = str(cfg.get("prompt_template_body") or "")
                adjust_prompt = str(cfg.get("adjust_prompt") or "")
        elif et == "chart":
            sql = _sql_from_element(el)
            if sql:
                charts_sql.append(sql)
        elif et == "table":
            sql = _sql_from_element(el)
            if sql:
                tables_sql.append(sql)

    data = {
        "section_name": sec.name or sec.key,
        "template_name": report.template_name,
        "property_type": report.property_type,
        "default_prompt": default_prompt,
        "adjust_prompt": adjust_prompt,
        "start_conversation": False,
        "charts": charts_sql,
        "tables": tables_sql,
    }
    if text:
        data["commentary"] = text

    pretty = json.dumps(data, indent=2, ensure_ascii=False)
    return text, pretty


def _sql_from_element(el: models.ReportSectionElement) -> str:
    cfg = el.config or {}
    if isinstance(cfg, dict):
        s = cfg.get("sql")
        if isinstance(s, str) and s.strip():
            return s.strip()
    return "-- SQL not provided"


def _cover_slide(prs: Presentation, report: models.Report):
    slide = prs.slides.add_slide(_blank_layout(prs))
    _textbox(slide, MARGIN_L, MARGIN_T, SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.9), report.name, size=TITLE_SIZE, bold=True)
    meta_lines = []
    if report.template_name:
        meta_lines.append(f"Template: {report.template_name}")
    meta_lines.append(f"Type: {report.report_type or report.template_name or report.property_type}")
    meta_lines.append(f"Property: {report.property_type}")
    if report.quarter:
        meta_lines.append(f"Quarter: {report.quarter}")
    if report.defined_markets:
        meta_lines.append("Markets: " + ", ".join(dict.fromkeys(report.defined_markets)))
    _textbox(slide, MARGIN_L, MARGIN_T + Inches(1.1), SLIDE_W - MARGIN_L - MARGIN_R, Inches(1.2), "\n".join(meta_lines), size=SUBTITLE_SIZE)


def _section_overview_slide(prs: Presentation, report: models.Report, sec: models.ReportSection):
    slide = prs.slides.add_slide(_blank_layout(prs))
    _textbox(slide, MARGIN_L, MARGIN_T, SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.8), sec.name or sec.key or "Section", size=TITLE_SIZE, bold=True)
    text, json_pretty = _commentary_for_section(report, sec)
    # Commentary prose
    _textbox(slide, MARGIN_L, MARGIN_T + Inches(0.9), SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.5), "Commentary", size=SUBTITLE_SIZE, bold=True)
    _textbox(slide, MARGIN_L, MARGIN_T + Inches(1.3), SLIDE_W - MARGIN_L - MARGIN_R, Inches(1.6), text or "—", size=BODY_SIZE)
    # Commentary JSON box
    _codebox(slide, MARGIN_L, MARGIN_T + Inches(3.2), SLIDE_W - MARGIN_L - MARGIN_R, Inches(3.6), "Commentary JSON", json_pretty or "{}")


def _chart_or_table_slide(prs: Presentation, sec_name: str, el: models.ReportSectionElement):
    slide = prs.slides.add_slide(_blank_layout(prs))
    title = f"{sec_name}: {el.label or el.element_type.title()}"
    _textbox(slide, MARGIN_L, MARGIN_T, SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.6), title, size=TITLE_SIZE, bold=True)

    left_main = MARGIN_L
    top_main = MARGIN_T + Inches(0.8)
    width_main = SLIDE_W - MARGIN_L - MARGIN_R - Inches(4.5)
    height_main = SLIDE_H - top_main - MARGIN_B

    left_code = left_main + width_main + Inches(0.3)
    width_code = SLIDE_W - left_code - MARGIN_R
    height_code = height_main

    et = (el.element_type or "").lower()
    if et == "chart":
        df = _dummy_chart_df()
        x_col = "Quarter"
        y_cols = [c for c in df.columns if c != x_col]
        kind = (el.config or {}).get("chart_type") if isinstance(el.config, dict) else None
        _chart(slide, left_main, top_main, width_main, height_main, el.label or "Chart", df, x_col, y_cols, str(kind or "line_markers"))
    elif et == "table":
        df = _dummy_table_df()
        _table(slide, left_main, top_main, width_main, height_main, el.label or "Table", df)
    else:
        # fallback text
        _textbox(slide, left_main, top_main, width_main, height_main, "(No visual)", size=BODY_SIZE)

    _codebox(slide, left_code, top_main, width_code, height_code, "SQL", _sql_from_element(el))


def build_report_ppt_bytes(report: models.Report) -> bytes:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    # Cover
    _cover_slide(prs, report)

    # Sections
    for sec in sorted(report.sections or [], key=lambda s: s.display_order or 0):
        if sec.selected is False:
            continue
        _section_overview_slide(prs, report, sec)
        for el in sorted(sec.elements or [], key=lambda e: e.display_order or 0):
            if (el.element_type or '').lower() in {"chart", "table"}:
                _chart_or_table_slide(prs, sec.name or sec.key or "Section", el)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_report_ppt_from_context(context: dict, report: models.Report | None = None) -> bytes:
    """Temporary adapter that accepts a fully-expanded report context dict.

    - Prints the inbound context to stdout (for debugging/instrumentation)
    - Delegates to the legacy builder using the provided SQLAlchemy `report`
      instance so that the current PPT still renders as before.
    - If `report` is not provided, returns an empty PPT with a single slide.
    """
    try:
        # Pretty-print (ensure ASCII safe) to avoid control chars in logs
        print("[PPT] build_report_ppt_from_context INPUT:")
        print(json.dumps(context, indent=2, ensure_ascii=False, default=str))
    except Exception as e:  # Never fail due to logging
        print("[PPT] failed to log context:", repr(e))

    if report is not None:
        return build_report_ppt_bytes(report)

    # Fallback: minimal empty deck
    prs = Presentation()
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()

"""
Slide Orchestrator - Intelligent Layout Engine for PowerPoint Presentations

This module provides a pure, deterministic layout engine that:
- Arranges content blocks in adaptive grids (2×2 default)
- Estimates content fit and adapts layout accordingly
- Handles overflow with graceful pagination
- Maintains consistency within sections
- Generates layout specifications for renderers

Author: AI Assistant
Date: 2025-10-01
"""

from dataclasses import dataclass, field, asdict
from typing import List, Dict, Optional, Tuple, Callable, Set
from enum import Enum
import json
import math
import os

from hello.utils.ppt_helpers_utils.services.template_config import (
    SlideConstraints,
    SlideLayoutConfig,
    get_allowed_layout_types,
    get_layout_preference_config,
    get_slide_layout_config,
    get_element_dimensions,
    get_layout_threshold_config,
    LayoutPreferenceConfig,
    LayoutPreferenceRule,
    determine_layout_type_from_criteria,
)


# ============================================================================
# CONFIGURATION & CONSTANTS
# ============================================================================


class LayoutType(Enum):
    """Available layout types for slides."""

    GRID_2x2 = "grid_2x2"  # Four blocks in 2×2 grid
    FULL_WIDTH = "full_width"  # Full width vertical stack layout
    BASE_SLIDE = (
        "base_slide"  # Uses base slide template with KPIs, populates remaining space
    )
    HYBRID_GRID = "hybrid_grid"  # DEPRECATED: Kept for backwards compatibility only


class ContentType(Enum):
    """Types of content blocks."""

    TEXT = "text"
    CHART = "chart"
    TABLE = "table"
    IMAGE = "image"


DEFAULT_SLIDE_CONSTRAINTS = SlideConstraints()


@dataclass
class CellFrame:
    """Physical dimensions of a cell on the slide."""

    left: float  # inches from left edge
    top: float  # inches from top edge
    width: float  # inches
    height: float  # inches

    @property
    def aspect_ratio(self) -> float:
        """Width to height ratio."""
        return self.width / self.height if self.height > 0 else 1.0

    @property
    def area(self) -> float:
        """Area in square inches."""
        return self.width * self.height

    def overlaps_with(self, other: "CellFrame", tolerance: float = 0.1) -> bool:
        """
        Check if this cell overlaps with another cell.

        Args:
            other: Another CellFrame to check against
            tolerance: Minimum overlap in inches to consider (default 0.1")

        Returns:
            True if cells overlap significantly
        """
        # Calculate boundaries
        self_right = self.left + self.width
        self_bottom = self.top + self.height
        other_right = other.left + other.width
        other_bottom = other.top + other.height

        # Check for overlap
        horizontal_overlap = not (
            self_right <= other.left + tolerance or self.left >= other_right - tolerance
        )
        vertical_overlap = not (
            self_bottom <= other.top + tolerance or self.top >= other_bottom - tolerance
        )

        return horizontal_overlap and vertical_overlap


# ============================================================================
# CONTENT BLOCKS
# ============================================================================


@dataclass
class ContentBlock:
    """Base class for content blocks."""

    id: str = ""
    type: ContentType = ContentType.TEXT
    priority: int = 0  # Higher priority = prefer larger space
    label: str = ""  # Optional label/heading for the content
    section_id: Optional[str] = None  # Track original section for boundary constraints
    display_order: Optional[int] = None  # Order within section (0 = first element)

    def estimate_intrinsic_size(
        self, constraints: SlideConstraints
    ) -> Tuple[float, float]:
        """
        Estimate the intrinsic (preferred) size for this block.
        Returns: (width, height) in inches
        """
        raise NotImplementedError("Subclasses must implement")

    def can_fit_in_cell(self, cell: CellFrame, constraints: SlideConstraints) -> bool:
        """Check if content can fit in the given cell while maintaining readability."""
        raise NotImplementedError("Subclasses must implement")

    def estimate_scale_factor(
        self, cell: CellFrame, constraints: SlideConstraints
    ) -> float:
        """
        Calculate scale factor needed to fit content in cell.
        Returns: 1.0 if fits perfectly, <1.0 if needs shrinking, >1.0 if room to grow
        """
        intrinsic_w, intrinsic_h = self.estimate_intrinsic_size(constraints)
        scale_w = cell.width / intrinsic_w if intrinsic_w > 0 else 1.0
        scale_h = cell.height / intrinsic_h if intrinsic_h > 0 else 1.0
        return min(scale_w, scale_h)

    def get_aspect_ratio(self) -> Optional[float]:
        """
        Get the aspect ratio (width/height) for this block if it has a fixed aspect ratio.
        Returns None if the block doesn't have a fixed aspect ratio.
        """
        return None

    def get_actual_bounds(
        self, cell: CellFrame, constraints: SlideConstraints
    ) -> CellFrame:
        """
        Get the actual bounds this content will occupy when rendered.
        Accounts for content that may exceed its allocated cell.
        """
        intrinsic_w, intrinsic_h = self.estimate_intrinsic_size(constraints)

        # Content may extend beyond cell if it doesn't fit
        actual_w = max(intrinsic_w, cell.width)
        actual_h = max(intrinsic_h, cell.height)

        return CellFrame(left=cell.left, top=cell.top, width=actual_w, height=actual_h)


@dataclass
class TextBlock(ContentBlock):
    """Text content block (commentary)."""

    text: str = ""
    bullet_points: List[str] = field(default_factory=list)
    font_size: float = 14.0
    line_spacing: float = 1.2
    estimated_chars_per_line: int = 80

    # Cached height calculations (calculated once, used everywhere)
    _cached_height: Optional[float] = field(default=None, repr=False)
    _cached_width: Optional[float] = field(default=None, repr=False)
    _cached_for_constraints_width: Optional[float] = field(default=None, repr=False)

    def __post_init__(self):
        self.type = ContentType.TEXT

    def calculate_content_based_height(
        self, constraints_width: float, force_recalculate: bool = False
    ) -> Tuple[float, float]:
        """
        Calculate commentary height based on actual content.
        This is the SINGLE SOURCE OF TRUTH for text/commentary height calculation.

        Results are cached and reused across all modules.

        Args:
            constraints_width: Available content width in inches
            force_recalculate: If True, ignore cache and recalculate

        Returns:
            Tuple of (width, height) in inches
        """
        # Return cached value if available
        if (
            not force_recalculate
            and self._cached_height is not None
            and self._cached_width is not None
            and self._cached_for_constraints_width is not None
            and abs(self._cached_for_constraints_width - constraints_width) < 0.01
        ):
            return (self._cached_width, self._cached_height)

        # Get element dimensions from config
        element_dims = get_element_dimensions()
        commentary_chars_per_line = element_dims.commentary_chars_per_line

        # Calculate total content
        title_chars = len(self.text) if self.text else 0
        bullet_chars = sum(len(bullet) for bullet in self.bullet_points)

        # Estimate lines more accurately
        chars_per_line = max(
            40,
            (
                commentary_chars_per_line
                if commentary_chars_per_line
                else self.estimated_chars_per_line
            )
            // 2,
        )

        # Title lines
        title_lines = max(1, title_chars // chars_per_line) if title_chars > 0 else 0

        # Bullet lines (each bullet gets its own line minimum)
        bullet_lines = 0
        for bullet in self.bullet_points:
            bullet_lines += max(1, len(bullet) // chars_per_line)

        total_lines = title_lines + bullet_lines

        # Add spacing
        if self.text and self.bullet_points:
            total_lines += 1  # Space between title and bullets
        if len(self.bullet_points) > 1:
            total_lines += len(self.bullet_points) * 0.2  # Inter-bullet spacing

        # Height calculation
        commentary_line_height = element_dims.commentary_line_height
        if commentary_line_height:
            height = total_lines * commentary_line_height
        else:
            line_height_inches = (self.font_size * self.line_spacing) / 72.0
            height = total_lines * line_height_inches + 0.3

        # Width calculation
        if len(self.bullet_points) > 2 or (self.text and len(self.text) > 30):
            width = constraints_width * 0.45
        else:
            width = constraints_width * 0.35

        # Cache results
        self._cached_height = height
        self._cached_width = width
        self._cached_for_constraints_width = constraints_width

        return (width, height)

    def get_cached_height(self) -> Optional[float]:
        """Get cached height calculation if available."""
        return self._cached_height

    def estimate_intrinsic_size(
        self, constraints: SlideConstraints
    ) -> Tuple[float, float]:
        """Estimate size based on text length and formatting.

        Uses cached calculation from calculate_content_based_height() - single source of truth.
        """
        return self.calculate_content_based_height(constraints.content_width)

    def can_fit_in_cell(self, cell: CellFrame, constraints: SlideConstraints) -> bool:
        """Check if text can fit without going below minimum font size."""
        _, intrinsic_height = self.estimate_intrinsic_size(constraints)
        scale = cell.height / intrinsic_height if intrinsic_height > 0 else 1.0

        scaled_font = self.font_size * scale
        return scaled_font >= constraints.min_font_size

    def get_actual_bounds(
        self, cell: CellFrame, constraints: SlideConstraints
    ) -> CellFrame:
        """
        Get the actual bounds this text will occupy when rendered.
        For text blocks, we constrain to the allocated cell since text wraps and truncates.
        """
        # Text should fit within its allocated cell, not expand beyond it
        # This prevents overlaps in grid layouts
        return CellFrame(
            left=cell.left,
            top=cell.top,
            width=cell.width,  # Use allocated width
            height=cell.height,  # Use allocated height - text will wrap/truncate
        )


@dataclass
class ChartBlock(ContentBlock):
    """Chart content block."""

    template_path: str = ""
    chart_type: str = "bar"  # bar, line, pie, etc.
    intrinsic_aspect_ratio: float = 16.0 / 9.0  # width/height
    min_width: float = 3.0  # inches
    min_height: float = 2.0
    data_points: int = 10  # affects readability
    data: List[Dict] = field(default_factory=list)  # Chart data
    figure_number: Optional[int] = None
    figure_label: Optional[str] = None
    figure_name: Optional[str] = None
    figure_source: Optional[str] = None
    # Axis title configuration (from axisConfig)
    primary_y_axis_title: Optional[str] = None  # Left Y-axis title (isPrimary: true)
    secondary_y_axis_title: Optional[str] = (
        None  # Right Y-axis title (isPrimary: false)
    )
    x_axis_title: Optional[str] = None  # X-axis (category axis) title
    # Y-axis keys for column selection (ordered: primary first, then secondary)
    y_axis_keys: List[str] = field(default_factory=list)
    # Multi-axis flag: when False, secondary axis is removed and all series plot on primary
    is_multi_axis: bool = True
    # Y-axis number formatting (Excel-style format codes from frontend axisConfig)
    primary_y_axis_format_code: Optional[str] = None
    secondary_y_axis_format_code: Optional[str] = None

    # Cached height calculations (calculated once, used everywhere)
    _cached_content_height: Optional[float] = field(default=None, repr=False)
    _cached_total_height: Optional[float] = field(default=None, repr=False)
    _cached_for_width: Optional[float] = field(default=None, repr=False)

    def __post_init__(self):
        self.type = ContentType.CHART

    def _estimate_legend_height(self, data: List[Dict]) -> float:
        """
        Estimate legend height based on number of series in data.

        Args:
            data: List of data dictionaries

        Returns:
            Legend height in inches
        """
        if not data or len(data) == 0:
            # No data, use default legend height
            return 0.15

        # Analyze data structure to count number of series
        # For chart data, series are typically the value columns (non-category columns)
        if isinstance(data[0], dict):
            # Count columns that are likely numeric (series)
            # First column is usually category, rest are series
            columns = list(data[0].keys())
            if len(columns) > 1:
                # Estimate series count (all columns except first, or all numeric columns)
                series_count = len(columns) - 1  # Assume first is category
            else:
                series_count = 1
        else:
            series_count = 1

        # Estimate legend height based on series count
        # More series = taller legend (may wrap to multiple lines)
        if series_count <= 2:
            return 0.15  # 1-2 series: single line legend
        elif series_count <= 4:
            return 0.20  # 3-4 series: may wrap to 2 lines
        else:
            return 0.25  # 5+ series: likely 2+ lines

    def estimate_intrinsic_size(
        self, constraints: SlideConstraints, available_width: Optional[float] = None
    ) -> Tuple[float, float]:
        """
        Estimate using content-based calculation to match slide assigner.

        Uses the SAME calculation as the slide assigner to ensure consistency
        and prevent wasted space due to height mismatches.

        Height is calculated ONCE via calculate_content_based_height() and cached.
        """
        # Get element dimensions from config
        element_dims = get_element_dimensions()

        # Use full available width (for full-width layouts) or content width
        width = (
            available_width
            if available_width is not None
            else constraints.content_width
        )

        # Use cached height calculation (single source of truth)
        chart_content_height = self.calculate_content_based_height(width)

        # Apply dynamic layout constraints (same as slide assigner)
        min_height = width * element_dims.dynamic_layout_min_height_ratio
        max_height = width * element_dims.dynamic_layout_max_height_ratio
        chart_content_height = max(min_height, min(chart_content_height, max_height))

        # Cache the constrained height for use by renderer
        self._cached_total_height = chart_content_height

        # Add label spaces (figure label and source - these are included in intrinsic size)
        figure_label_space = element_dims.get_figure_label_total_height()
        source_label_space = element_dims.source_gap + element_dims.source_label_height
        total_label_space = figure_label_space + source_label_space
        height = chart_content_height + total_label_space

        return (width, height)

    def get_aspect_ratio(self) -> Optional[float]:
        """Return the intrinsic aspect ratio for this chart."""
        return self.intrinsic_aspect_ratio

    def can_fit_in_cell(self, cell: CellFrame, constraints: SlideConstraints) -> bool:
        """
        Check if chart maintains minimum size and aspect ratio.

        Uses dynamic layout ratios: min height = 20% of width (from config).
        """
        # Get element dimensions from config
        element_dims = get_element_dimensions()

        # Check minimum dimensions based on dynamic layout ratios
        min_height = cell.width * element_dims.dynamic_layout_min_height_ratio
        if cell.height < min_height:
            return False

        # Check if aspect ratio distortion is acceptable (within 30%)
        intrinsic_w, intrinsic_h = self.estimate_intrinsic_size(constraints)
        intrinsic_ratio = intrinsic_w / intrinsic_h
        cell_ratio = cell.aspect_ratio

        distortion = abs(intrinsic_ratio - cell_ratio) / intrinsic_ratio
        return distortion < 0.3

    def get_actual_bounds(
        self, cell: CellFrame, constraints: SlideConstraints
    ) -> CellFrame:
        """
        Get the actual bounds this chart will occupy when rendered.
        Charts are scaled to fit within their allocated cell, so actual bounds = cell bounds.
        """
        # Charts are scaled to fit within their cell, so they don't extend beyond cell bounds
        return CellFrame(
            left=cell.left,
            top=cell.top,
            width=cell.width,  # Use cell width (chart is scaled to fit)
            height=cell.height,  # Use cell height (chart is scaled to fit)
        )

    def calculate_content_based_height(
        self, chart_width: float, force_recalculate: bool = False
    ) -> float:
        """
        Calculate chart height based on actual content (number of series, data points, etc.).
        This is the SINGLE SOURCE OF TRUTH for chart height calculation.

        Results are cached and reused across all modules (orchestrator, renderer, etc.).

        Args:
            chart_width: Available width for the chart in inches
            force_recalculate: If True, ignore cache and recalculate

        Returns:
            Total chart content height in inches (chart + legend, no labels/source)
        """
        # Return cached value if available and width matches
        if (
            not force_recalculate
            and self._cached_content_height is not None
            and self._cached_for_width is not None
            and abs(self._cached_for_width - chart_width) < 0.01
        ):
            return self._cached_content_height

        from hello.utils.ppt_helpers_utils.ppt_helpers.content_height_calculator import (
            calculate_chart_content_height,
        )
        from hello.utils.ppt_helpers_utils.services.template_config import (
            get_element_dimensions,
        )

        element_dims = get_element_dimensions()

        content_height = calculate_chart_content_height(
            chart_data=self.data,
            chart_width=chart_width,
            element_dims=element_dims,
            chart_type=self.chart_type,
        )

        # Cache results
        self._cached_content_height = content_height
        self._cached_for_width = chart_width

        return content_height

    def get_cached_height(self) -> Optional[float]:
        """Get cached height calculation if available.

        Returns:
            Cached content height if available, None otherwise
        """
        return self._cached_content_height


@dataclass
class TableBlock(ContentBlock):
    """Table content block."""

    rows: int = 1
    columns: int = 1
    has_header: bool = True
    min_row_height: float = 0.22  # inches (reduced from 0.25 to fit more rows)
    min_col_width: float = 0.75  # inches (reduced from 0.8 to fit more columns)
    cell_padding: float = 0.01  # reduced from 0.1 for more compact tables
    can_split_rows: bool = True
    template_path: Optional[str] = None
    data: List[Dict] = field(default_factory=list)  # Table data
    figure_number: Optional[int] = None
    figure_label: Optional[str] = None
    figure_name: Optional[str] = None
    table_source: Optional[str] = None

    # Template column count - read from template file for accurate height calculation
    _template_columns: Optional[int] = field(default=None, repr=False)

    # Cached height calculations (calculated once, used everywhere)
    _cached_total_height: Optional[float] = field(default=None, repr=False)
    _cached_row_heights: Optional[List[float]] = field(default=None, repr=False)
    _cached_column_widths: Optional[List[float]] = field(default=None, repr=False)
    _cached_for_width: Optional[float] = field(default=None, repr=False)

    # Continuation table attributes (set by renderer for overflow handling)
    _continuation_data: Optional[List[Dict]] = field(default=None, repr=False)
    _continuation_id: Optional[str] = field(default=None, repr=False)
    _continuation_label: Optional[str] = field(default=None, repr=False)
    _continuation_section: Optional[str] = field(default=None, repr=False)
    _skip_source: bool = field(default=False, repr=False)

    def __post_init__(self):
        self.type = ContentType.TABLE
        # Read template column count if template_path is set and columns not yet cached
        if self.template_path and self._template_columns is None:
            self._template_columns = self._read_template_columns()

    def _read_template_columns(self) -> int:
        """Read the column count from the template file."""
        if not self.template_path or not os.path.exists(self.template_path):
            return 11  # Default fallback
        try:
            from pptx import Presentation

            prs = Presentation(self.template_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_table:
                        return len(shape.table.columns)
            return 11  # Fallback if no table found
        except Exception:
            return 11  # Fallback on error

    @property
    def template_columns(self) -> int:
        """Get the template column count."""
        if self._template_columns is None:
            self._template_columns = self._read_template_columns()
        return self._template_columns

    def _calculate_actual_dimensions(self, data: List[Dict]) -> Tuple[int, int]:
        """
        Calculate actual dimensions from data.

        Args:
            data: List of data dictionaries

        Returns:
            Tuple of (actual_rows, actual_columns)
        """
        if not data or len(data) == 0:
            # No data available, return fallback dimensions
            return (self.rows, self.columns)

        # Calculate actual rows: data rows + header row
        actual_rows = len(data) + (1 if self.has_header else 0)

        # Calculate actual columns from first data row keys
        if isinstance(data[0], dict):
            actual_columns = len(data[0].keys())
        else:
            # Fallback if data structure is unexpected
            actual_columns = self.columns

        return (actual_rows, actual_columns)

    def calculate_content_based_height(
        self, table_width: float, force_recalculate: bool = False
    ) -> Tuple[float, List[float]]:
        """
        Calculate table height based on actual content and text wrapping.

        Delegates to the shared calculate_table_content_height() function
        which is the SINGLE SOURCE OF TRUTH for all table height calculations.

        Results are CACHED for efficiency - all downstream code uses the same values.
        Use force_recalculate=True to bypass cache (e.g., after data changes).

        Args:
            table_width: Available width for the table in inches
            force_recalculate: If True, bypass cache and recalculate

        Returns:
            Tuple of (total_table_height, list_of_row_heights)
            - total_table_height: Total height including all rows (no labels)
            - list_of_row_heights: List of heights for each row [header, row1, row2, ...]
        
        Note: Column widths are also calculated and cached (access via get_cached_column_widths())
        """
        # Return cached values if available and width hasn't changed
        if (
            not force_recalculate
            and self._cached_total_height is not None
            and self._cached_row_heights is not None
            and self._cached_for_width is not None
            and abs(self._cached_for_width - table_width) < 0.01
        ):
            return (self._cached_total_height, self._cached_row_heights)

        from hello.utils.ppt_helpers_utils.services.template_config import (
            get_element_dimensions,
        )
        from hello.utils.ppt_helpers_utils.ppt_helpers.content_height_calculator import (
            calculate_table_content_height,
        )

        element_dims = get_element_dimensions()

        # Get table data
        data = self.data if hasattr(self, "data") and self.data else []

        # Check if table has source (source row is added INSIDE the table)
        has_table_source = hasattr(self, "table_source") and self.table_source
        
        # Delegate to the shared height calculation function (single source of truth)
        # Pass template_columns to limit columns used for height calculation
        # Include source row height in the calculation since it's part of the table
        # Returns: (total_height, row_heights, column_widths)
        total_height, row_heights, column_widths = calculate_table_content_height(
            data=data,
            table_width=table_width,
            element_dims=element_dims,
            has_header=self.has_header,
            max_columns=self.template_columns,
            has_source=has_table_source,
        )

        # Handle no data case - use self.rows for fallback count
        if not data:
            rows_to_use = self.rows
            row_heights = [element_dims.table_min_row_height] * rows_to_use
            total_height = sum(row_heights)
            column_widths = []

        # Log result (calculation details are logged in calculate_table_content_height)
        print(
            f'    📐 Table height calculation: {len(row_heights)} rows, width={table_width:.2f}", total height = {total_height:.3f}"'
        )

        # Cache results (including column widths for PPT rendering)
        self._cached_total_height = total_height
        self._cached_row_heights = row_heights
        self._cached_column_widths = column_widths
        self._cached_for_width = table_width
        
        return (total_height, row_heights)

    def get_cached_height(self) -> Optional[Tuple[float, List[float]]]:
        """Get cached height calculation if available.

        Returns:
            Tuple of (total_height, row_heights) if cached, None otherwise
        """
        if (
            self._cached_total_height is not None
            and self._cached_row_heights is not None
        ):
            return (self._cached_total_height, self._cached_row_heights)
        return None

    def get_cached_column_widths(self) -> Optional[List[float]]:
        """Get cached column widths if available.
        
        Column widths are calculated during height calculation and cached.
        First column is expanded to fit content without wrapping,
        remaining columns share the leftover space equally.

        Returns:
            List of column widths in inches, or None if not cached
        """
        return self._cached_column_widths

    def get_max_rows_for_available_height(
        self, available_height: float, table_width: float
    ) -> int:
        """Calculate how many rows fit in available height using cached row heights.

        This is the SINGLE SOURCE OF TRUTH for row fitting calculations.
        Uses actual content-based row heights (sum of individual rows, not average).

        Args:
            available_height: Available height in inches (excluding labels)
            table_width: Table width for calculation (used if cache miss)

        Returns:
            Number of DATA rows that fit (excluding header and source row)
        """
        from hello.utils.ppt_helpers_utils.services.template_config import (
            get_element_dimensions,
        )
        element_dims = get_element_dimensions()
        
        # Ensure we have calculated heights (uses cache if available)
        _, row_heights = self.calculate_content_based_height(table_width)

        if not row_heights:
            return 0

        # Check if table has source row (source is ALWAYS the LAST element in row_heights if present)
        # BUT: Check _skip_source flag - first table of a split doesn't render source
        has_source = hasattr(self, "table_source") and self.table_source
        skip_source = getattr(self, "_skip_source", False)
        will_render_source = has_source and not skip_source
        
        # Subtract border overhead (same as in content_height_calculator)
        remaining = available_height - element_dims.table_border_overhead
        row_gap_padding = element_dims.table_row_gap_padding

        # If table WILL render source row, subtract its height + gap from available space
        # Source row is at the END of row_heights and must always fit
        if will_render_source and len(row_heights) > 0:
            source_height = row_heights[-1]
            remaining -= source_height
            remaining -= row_gap_padding  # Gap before source row
        
        # Determine where data rows start and end in row_heights
        # Structure: [header?, data1, data2, ..., dataN, source?]
        data_start_idx = 1 if self.has_header else 0
        # Only exclude source from iteration if it's in row_heights (has_source, not will_render_source)
        data_end_idx = len(row_heights) - 1 if has_source else len(row_heights)

        # Header row (if present) - use actual calculated height + gap after
        if self.has_header and len(row_heights) > 0:
            header_height = row_heights[0]
            remaining -= header_height
            remaining -= row_gap_padding  # Gap after header

        if remaining <= 0:
            return 0

        # Count DATA rows that fit (excluding header and source)
        data_rows_that_fit = 0
        for idx in range(data_start_idx, data_end_idx):
            row_height = row_heights[idx]
            # Include gap padding (except for first data row after header)
            height_needed = row_height
            if data_rows_that_fit > 0:  # Add gap before this row
                height_needed += row_gap_padding
            
            if remaining >= height_needed:
                data_rows_that_fit += 1
                remaining -= height_needed
            else:
                break

        return data_rows_that_fit

    def _calculate_max_rows_for_height(
        self, available_height: float, element_dims
    ) -> int:
        """Calculate maximum number of rows that can fit in available height.

        Args:
            available_height: Total available height in inches for the table and its labels
            element_dims: ElementDimensions configuration object

        Returns:
            Maximum number of rows (including header) that can fit
        """
        # Account for labels (heading + source)
        label_space = 0.0

        # Top space for heading/label
        has_heading = (hasattr(self, "figure_name") and self.figure_name) or (
            hasattr(self, "label") and self.label
        )
        if has_heading:
            # Use table-specific gap for tables (larger than chart gap to prevent overlapping)
            label_space += element_dims.get_table_label_total_height()

        # Source is now a row INSIDE the table - use config value for height
        # Subtract from available space so we calculate how many DATA rows fit
        has_source = hasattr(self, "table_source") and self.table_source
        source_row_space = element_dims.table_source_row_height if has_source else 0.0

        # Available space for actual table (data rows + header, not source)
        # Subtract fixed border overhead (top + bottom borders) - matches content_height_calculator
        table_height = (
            available_height - label_space - source_row_space - element_dims.table_border_overhead
        )

        # Calculate rows that fit using same formula as content_height_calculator
        # Each row has: row_height + row_gap_padding (except last row has no trailing gap)
        # Formula: table_height = (num_rows * row_height) + ((num_rows - 1) * row_gap_padding)
        fixed_row_height = (
            element_dims.table_min_row_height * element_dims.table_height_safety_margin
        )
        row_gap_padding = element_dims.table_row_gap_padding

        # Solve for num_rows: 
        # table_height = num_rows * row_height + (num_rows - 1) * gap
        # table_height = num_rows * (row_height + gap) - gap
        # num_rows = (table_height + gap) / (row_height + gap)
        max_rows = int(
            (table_height + row_gap_padding) / (fixed_row_height + row_gap_padding)
        )

        # At least header row
        return max(1, max_rows)

    def estimate_intrinsic_size(
        self, constraints: SlideConstraints, available_width: Optional[float] = None
    ) -> Tuple[float, float]:
        """Estimate based on row/column count, using actual data when available.

        Args:
            constraints: Slide constraints
            available_width: Optional available width in inches. If provided, uses this width
                           for calculating height with text wrapping. Otherwise calculates width
                           from column count.
        """

        # Get element dimensions from config
        element_dims = get_element_dimensions()

        # Use element dimensions if available, otherwise fall back to instance defaults
        min_col_width = element_dims.table_min_col_width
        cell_padding = element_dims.table_cell_padding
        height_safety_margin = element_dims.table_height_safety_margin

        # Use actual data dimensions if available, otherwise fall back to self.rows/self.columns
        if self.data and len(self.data) > 0:
            actual_rows, actual_columns = self._calculate_actual_dimensions(self.data)
            columns_to_use = actual_columns
        else:
            # Fall back to estimated dimensions
            columns_to_use = self.columns

        # Determine width to use for height calculation
        if available_width is not None:
            # Use provided width (e.g., full content width for full_width layout)
            width = available_width
        else:
            # Calculate table width from column count
            # Width: columns * min_col_width * safety_margin + padding
            width = columns_to_use * min_col_width * height_safety_margin
            width += (columns_to_use + 1) * cell_padding

        # Calculate height using content-based calculation
        # This accounts for text wrapping and actual cell content
        table_height, row_heights = self.calculate_content_based_height(width)

        print(
            f'    📐 Table height calculation: {len(row_heights)} rows, width={width:.2f}", total height = {table_height:.3f}"'
        )

        # Always reserve space for labels within cell (show/hide based on config)
        # This must match the label space calculation in the renderer for consistency.
        #
        # Note: renderer adds "Figure X" labels for charts/tables even if `label` is empty.
        # We gate this with a config flag to keep behavior config-driven.
        reserve_figure_label_space = (
            getattr(element_dims, "reserve_figure_label_space_for_charts_and_tables", True)
        )
        has_heading = reserve_figure_label_space or (
            (hasattr(self, "figure_name") and self.figure_name)
            or (hasattr(self, "label") and self.label)
        )
        has_table_source = hasattr(self, "table_source") and self.table_source

        # Calculate top space (table heading/label only)
        # NOTE: Section title space is NOT included here because the orchestrator
        # will add it separately for the first block in a section.
        top_space = 0.0

        # Table label/heading (top space)
        if has_heading:
            # Use table-specific gap for tables (larger than chart gap to prevent overlapping)
            top_space += element_dims.get_table_label_total_height()

        # Source row height is now included in table_height from calculate_content_based_height
        # (it's added as a row inside the table, so it's part of the table content height)
        # Cell height = table content (including source row) + top labels
        height = table_height + top_space

        print(
            f'    📏 estimate_intrinsic_size: table={table_height:.2f}" (includes source row), top_space={top_space:.2f}", total={height:.2f}"'
        )

        # Return the width we used for calculation (already accounts for available_width)
        return (width, height)

    def estimate_scale_factor(
        self, cell: CellFrame, constraints: SlideConstraints
    ) -> float:
        """
        Calculate scale factor needed to fit table in cell.
        
        Override to use cell.width as available_width for accurate height calculation.
        This ensures the intrinsic height is calculated with the ACTUAL cell width,
        not a fallback based on column count.
        """
        intrinsic_w, intrinsic_h = self.estimate_intrinsic_size(
            constraints, available_width=cell.width
        )
        scale_w = cell.width / intrinsic_w if intrinsic_w > 0 else 1.0
        scale_h = cell.height / intrinsic_h if intrinsic_h > 0 else 1.0
        return min(scale_w, scale_h)

    def can_fit_in_cell(self, cell: CellFrame, constraints: SlideConstraints) -> bool:
        """Check if table fits without truncating columns."""
        # Get element dimensions from config
        element_dims = get_element_dimensions()
        min_col_width = element_dims.table_min_col_width
        cell_padding = element_dims.table_cell_padding

        # Use actual data dimensions if available, otherwise fall back to self.columns
        if self.data and len(self.data) > 0:
            _, actual_columns = self._calculate_actual_dimensions(self.data)
            columns_to_use = actual_columns
        else:
            columns_to_use = self.columns

        # Columns must NEVER be truncated
        min_width = columns_to_use * min_col_width + (columns_to_use + 1) * cell_padding
        if cell.width < min_width:
            return False

        # Rows can be split if enabled
        if self.can_split_rows:
            # Get element dimensions for row height
            min_row_height = element_dims.table_min_row_height
            # Need space for at least header + 1 data row
            min_rows = 2 if self.has_header else 1
            min_height = min_rows * min_row_height + (min_rows + 1) * cell_padding
            return cell.height >= min_height
        else:
            # All rows must fit - use cell.width for accurate calculation
            _, intrinsic_height = self.estimate_intrinsic_size(
                constraints, available_width=cell.width
            )
            return cell.height >= intrinsic_height

    def split_rows(self, available_height: float) -> List["TableBlock"]:
        """Split table into multiple blocks by rows."""
        if not self.can_split_rows:
            return [self]

        # Get element dimensions from config
        element_dims = get_element_dimensions()
        min_row_height = element_dims.table_min_row_height
        cell_padding = element_dims.table_cell_padding

        # Calculate how many rows fit in available height
        row_height = min_row_height + cell_padding
        header_rows = 1 if self.has_header else 0
        header_height = header_rows * row_height

        available_for_data = available_height - header_height - 0.2  # margin
        rows_per_page = max(1, int(available_for_data / row_height))

        data_rows = self.rows - header_rows
        if rows_per_page >= data_rows:
            return [self]

        # Split into multiple table blocks
        splits = []
        remaining_rows = data_rows

        while remaining_rows > 0:
            rows_this_page = min(rows_per_page, remaining_rows)
            # Each split includes header
            # Get element dimensions for split block
            element_dims = get_element_dimensions()
            split_block = TableBlock(
                id=f"{self.id}_split_{len(splits)}",
                type=self.type,
                rows=rows_this_page + header_rows,
                columns=self.columns,
                has_header=self.has_header,
                min_row_height=element_dims.table_min_row_height,
                min_col_width=element_dims.table_min_col_width,
                cell_padding=element_dims.table_cell_padding,
                can_split_rows=False,  # Already split
                template_path=self.template_path,
                priority=self.priority,
            )
            splits.append(split_block)
            remaining_rows -= rows_this_page

        return splits

    def get_actual_bounds(
        self, cell: CellFrame, constraints: SlideConstraints
    ) -> CellFrame:
        """
        Get the actual bounds this table will occupy when rendered.
        For tables, we constrain to the allocated cell to prevent overlaps in table rows layout.
        """
        # Tables should fit within their allocated cell, not expand beyond it
        # This prevents overlaps in table rows layout
        return CellFrame(
            left=cell.left,
            top=cell.top,
            width=cell.width,  # Use allocated width, don't expand
            height=cell.height,  # Use allocated height, don't expand
        )


# ============================================================================
# SECTION & LAYOUT SPECIFICATIONS
# ============================================================================


@dataclass
class Section:
    """A logical section containing multiple content blocks."""

    id: str
    title: str
    blocks: List[ContentBlock]
    style: Dict = field(default_factory=dict)  # fonts, colors, etc.
    constraints: SlideConstraints = field(default_factory=SlideConstraints)
    layout_preference: Optional[str] = None  # User's preferred layout from UI


@dataclass
class SlideLayout:
    """Complete layout specification for a single slide."""

    slide_number: int
    section_id: str
    layout_type: LayoutType
    cell_frames: List[CellFrame]
    assigned_blocks: List[Tuple[int, str]]  # (cell_index, block_id)
    scale_factors: Dict[str, float]  # block_id -> scale_factor
    overflow_notes: List[str] = field(default_factory=list)
    metadata: Dict = field(default_factory=dict)

    def to_dict(self) -> Dict:
        """Convert to dictionary for JSON serialization."""
        return {
            "slide_number": self.slide_number,
            "section_id": self.section_id,
            "layout_type": self.layout_type.value,
            "cell_frames": [asdict(cf) for cf in self.cell_frames],
            "assigned_blocks": self.assigned_blocks,
            "scale_factors": self.scale_factors,
            "overflow_notes": self.overflow_notes,
            "metadata": self.metadata,
        }


# ============================================================================
# LAYOUT ENGINE
# ============================================================================


class SlideOrchestrator:
    """
    Pure layout engine that generates slide specifications.

    This is a deterministic, side-effect-free function that takes
    content blocks and produces layout specifications.
    """

    def __init__(
        self,
        constraints: Optional[SlideConstraints] = None,
        property_sub_type: Optional[str] = None,
    ):
        """Initialize orchestrator with layout constraints and subtype policy."""
        self._custom_constraints_supplied = not (
            constraints is None or constraints == DEFAULT_SLIDE_CONSTRAINTS
        )
        self.set_property_sub_type(property_sub_type)
        self.constraints = self._resolve_constraints(constraints)

    def set_property_sub_type(self, property_sub_type: Optional[str]) -> None:
        """Configure layout policy for a property sub-type."""
        self.property_sub_type = self._normalize_property_sub_type(property_sub_type)
        if hasattr(self, "constraints") and not self._custom_constraints_supplied:
            self.constraints = self._build_constraints_for_property_sub_type()

    def _normalize_property_sub_type(self, property_sub_type: Optional[str]) -> str:
        """Return canonical property sub-type label."""
        if not property_sub_type:
            return "figures"
        return property_sub_type.strip().lower()

    def _resolve_constraints(
        self, provided_constraints: Optional[SlideConstraints]
    ) -> SlideConstraints:
        """
        Determine which constraints should be used for the orchestrator or section.

        - If no constraints (or defaults) are provided, use property_sub_type-specific constraints.
        - Otherwise, respect the explicit overrides.
        """
        if (
            provided_constraints is None
            or provided_constraints == DEFAULT_SLIDE_CONSTRAINTS
        ):
            return self._build_constraints_for_property_sub_type()
        return provided_constraints

    def _build_constraints_for_property_sub_type(self) -> SlideConstraints:
        """Return slide constraints tuned to the current property_sub_type."""
        layout_config = get_slide_layout_config(self.property_sub_type)
        # Use base constraints for regular slides (orchestrator default)
        return layout_config.get_constraints(is_first_slide=False)

    def _is_layout_allowed(self, layout_type: LayoutType) -> bool:
        """Check whether a layout type is allowed for the current property sub-type."""
        allowed_values = get_allowed_layout_types(self.property_sub_type)
        return layout_type.value in allowed_values

    def _first_slide_capacity(self) -> Optional[int]:
        """Maximum number of elements allowed on the first slide per sub-type."""
        layout_config = get_slide_layout_config(self.property_sub_type)
        return layout_config.first_slide_max_elements

    @staticmethod
    def _limit_blocks(
        blocks: List[ContentBlock], max_elements: Optional[int]
    ) -> List[ContentBlock]:
        """Return blocks trimmed to the configured capacity, if any."""
        if max_elements is None:
            return blocks
        return blocks[:max_elements]

    def _get_strategy_sequence(self, blocks: List[ContentBlock]) -> List[Callable]:
        """Determine which layout strategies should be attempted."""
        # Check if blocks from different sections have mixed layout requirements
        has_mixed_layouts = self._has_mixed_layout_requirements(blocks)

        if has_mixed_layouts:
            # When blocks have mixed layout requirements, prioritize hybrid grid
            print("    🔀 Detected mixed layout requirements, prioritizing hybrid grid")
            return [
                self._try_hybrid_grid,
                self._try_grid_2x2,
                self._try_full_width_stack,
            ]

        preference_config = get_layout_preference_config(self.property_sub_type)
        allowed_layouts = get_allowed_layout_types(self.property_sub_type)
        layout_sequence = self._resolve_layout_sequence(preference_config, blocks)
        strategies = self._build_strategy_sequence(layout_sequence, allowed_layouts)
        if strategies:
            return strategies

        # Fallback: ensure we always have at least one strategy based on allowed layouts
        fallback_sequence = (
            LayoutType.GRID_2x2.value,
            LayoutType.FULL_WIDTH.value,
            LayoutType.BASE_SLIDE.value,
        )
        fallback_strategies = self._build_strategy_sequence(
            fallback_sequence, allowed_layouts
        )
        return fallback_strategies or [self._try_grid_2x2]

    def _has_mixed_layout_requirements(self, blocks: List[ContentBlock]) -> bool:
        """
        Check if blocks from different sections have mixed layout requirements.

        Groups blocks by section_id and determines the layout for each section.
        Returns True if different sections require different layouts (e.g., one needs grid_2x2, another needs full_width).
        """
        if not blocks or len(blocks) < 2:
            return False

        # Group blocks by section_id
        from collections import defaultdict

        blocks_by_section = defaultdict(list)
        for block in blocks:
            section_id = getattr(block, "section_id", None) or "unknown"
            blocks_by_section[section_id].append(block)

        # If all blocks are from the same section, no mixed requirements
        if len(blocks_by_section) < 2:
            print("       All blocks from same section - no mixed layout requirements")
            return False

        # Determine layout for each section
        section_layouts = {}
        for section_id, section_blocks in blocks_by_section.items():
            # First, check if layout is in block metadata
            layout_from_metadata = None
            for block in section_blocks:
                if hasattr(block, "metadata") and isinstance(block.metadata, dict):
                    layout_from_metadata = block.metadata.get("layout")
                    if layout_from_metadata:
                        break

            if layout_from_metadata:
                section_layouts[section_id] = layout_from_metadata
                print(
                    f"       Section {section_id}: layout={layout_from_metadata} (from metadata)"
                )
            else:
                # Determine layout using criteria
                layout = determine_layout_type_from_criteria(
                    property_sub_type=self.property_sub_type,
                    is_first_slide=False,  # Mixed layouts only apply to middle slides
                    normalized_preference=None,
                    blocks=section_blocks,
                )
                section_layouts[section_id] = layout
                print(
                    f"       Section {section_id}: layout={layout} (determined from criteria)"
                )

        # Check if different sections have different layouts
        unique_layouts = set(section_layouts.values())
        is_mixed = len(unique_layouts) > 1

        if is_mixed:
            print(f"       🔀 Mixed layout detected: sections require {unique_layouts}")
        else:
            print(
                f"       All sections use same layout: {unique_layouts.pop() if unique_layouts else 'unknown'}"
            )

        return is_mixed

    @staticmethod
    def _resolve_layout_sequence(
        preference_config: LayoutPreferenceConfig,
        blocks: List[ContentBlock],
    ) -> Tuple[str, ...]:
        """Return the layout sequence after evaluating conditional rules."""
        for rule in preference_config.rules:
            if SlideOrchestrator._matches_layout_rule(rule, blocks):
                return rule.layout_sequence
        return preference_config.default_sequence

    @staticmethod
    def _matches_layout_rule(
        rule: LayoutPreferenceRule, blocks: List[ContentBlock]
    ) -> bool:
        """Evaluate whether a conditional rule applies for the given blocks."""
        if not rule.condition:
            return True
        if rule.condition == "all_tables":
            return bool(blocks) and all(
                block.type == ContentType.TABLE for block in blocks
            )
        return False

    def _build_strategy_sequence(
        self,
        layout_sequence: Tuple[str, ...],
        allowed_layouts: Set[str],
    ) -> List[Callable]:
        """Map layout labels to callable strategies while respecting allowed layouts."""
        strategy_map = {
            LayoutType.GRID_2x2.value: self._try_grid_2x2,
            LayoutType.FULL_WIDTH.value: self._try_full_width_stack,
            LayoutType.BASE_SLIDE.value: self._try_base_slide,
            LayoutType.HYBRID_GRID.value: self._try_hybrid_grid,
        }

        strategies: List[Callable] = []
        for layout_label in layout_sequence:
            if layout_label not in allowed_layouts:
                continue
            strategy = strategy_map.get(layout_label)
            if strategy and strategy not in strategies:
                strategies.append(strategy)
        return strategies

    def orchestrate_section(
        self,
        section: Section,
        total_slides: Optional[int] = None,
        start_slide_number: int = 1,
    ) -> List[SlideLayout]:
        """
        Main entry point: orchestrate a complete section into slides.

        NOTE: This method now expects blocks to have slide_number and layout pre-assigned.
        It only computes positions (CellFrame) for the pre-assigned blocks, not assignment logic.

        Args:
            section: Section with content blocks (with pre-assigned slide_number and layout)
            total_slides: Total number of slides across all sections (None if unknown)
            start_slide_number: Starting slide number for this section (default 1)

        Returns:
            List of SlideLayout specifications (deterministic)
        """
        section.constraints = self._resolve_constraints(section.constraints)

        # Group blocks by their pre-assigned slide number
        blocks_by_slide = {}
        for block in section.blocks:
            slide_num = None
            layout_type = None

            if hasattr(block, "metadata") and isinstance(block.metadata, dict):
                slide_num = block.metadata.get("slide_number")
                layout_type = block.metadata.get("layout")

            if slide_num is None:
                slide_num = start_slide_number  # Fallback

            if slide_num not in blocks_by_slide:
                blocks_by_slide[slide_num] = {"blocks": [], "layouts": []}

            blocks_by_slide[slide_num]["blocks"].append(block)
            # Track all unique layouts on this slide
            if layout_type and layout_type not in blocks_by_slide[slide_num]["layouts"]:
                blocks_by_slide[slide_num]["layouts"].append(layout_type)

        # Process each slide
        layouts = []
        for slide_num in sorted(blocks_by_slide.keys()):
            slide_info = blocks_by_slide[slide_num]
            blocks = slide_info["blocks"]
            unique_layouts = slide_info["layouts"]

            # Detect mixed layout: if blocks have different layout assignments
            # This should not happen after slide_number_assigner changes enforce layout consistency
            if len(unique_layouts) > 1:
                import logging

                logging.warning(
                    f"Mixed layouts detected on slide {slide_num}: {unique_layouts}. "
                    "This may indicate a bug in slide_number_assigner - layouts should be homogeneous per slide. "
                    "Using first layout type."
                )
                assigned_layout = unique_layouts[0]
                print(
                    f"    ⚠️ Processing slide {slide_num}: {len(blocks)} blocks, MIXED LAYOUTS {unique_layouts} → using '{assigned_layout}' (hybrid disabled)"
                )
            else:
                assigned_layout = unique_layouts[0] if unique_layouts else "full_width"
                print(
                    f"    📋 Processing slide {slide_num}: {len(blocks)} blocks, layout='{assigned_layout}'"
                )

            # Compute positions for all pre-assigned blocks
            layout = self._compute_positions_for_assigned_blocks(
                section=section,
                blocks=blocks,
                slide_number=slide_num,
                layout_type=assigned_layout,
                total_slides=total_slides,
            )

            layouts.append(layout)

        return layouts

    def _compute_positions_for_assigned_blocks(
        self,
        section: Section,
        blocks: List[ContentBlock],
        slide_number: int,
        layout_type: Optional[str],
        total_slides: Optional[int] = None,
    ) -> SlideLayout:
        """
        Compute cell positions for pre-assigned blocks.

        This method does NOT make assignment decisions - it only calculates positions.
        All blocks MUST be placed, no skipping.

        Args:
            section: Section with constraints
            blocks: Pre-assigned blocks for this slide
            slide_number: Slide number (1-indexed)
            layout_type: Pre-assigned layout type ('full_width', 'grid_2x2', 'base_slide')
            total_slides: Total slides (for first/last detection)

        Returns:
            SlideLayout with positions computed for ALL blocks
        """
        if not blocks:
            return SlideLayout(
                slide_number=slide_number,
                section_id=section.id,
                layout_type=LayoutType.FULL_WIDTH,
                cell_frames=[],
                assigned_blocks=[],
                scale_factors={},
                overflow_notes=[],
            )

        constraints = section.constraints
        is_first_slide = slide_number == 1

        # Normalize layout type
        if layout_type == "full_width":
            layout_enum = LayoutType.FULL_WIDTH
        elif layout_type == "grid_2x2":
            layout_enum = LayoutType.GRID_2x2
        elif layout_type == "base_slide":
            layout_enum = LayoutType.BASE_SLIDE
        elif layout_type == "hybrid_grid":
            # DEPRECATED: hybrid_grid is no longer used as slides now have homogeneous layouts
            # Fall back to full_width for backwards compatibility
            import logging

            logging.warning(
                f"hybrid_grid layout requested for slide {slide_number} but is deprecated. "
                "Falling back to full_width layout."
            )
            layout_enum = LayoutType.FULL_WIDTH
        else:
            # Default to full_width
            layout_enum = LayoutType.FULL_WIDTH

        print(
            f"       Computing positions for {len(blocks)} blocks using {layout_enum.value} layout"
        )

        # Compute positions based on layout type
        if layout_enum == LayoutType.FULL_WIDTH:
            return self._compute_full_width_positions(
                section, blocks, slide_number, is_first_slide
            )
        elif layout_enum == LayoutType.GRID_2x2:
            return self._compute_grid_positions(
                section, blocks, slide_number, is_first_slide
            )
        elif layout_enum == LayoutType.BASE_SLIDE:
            return self._compute_base_slide_positions(section, blocks, slide_number)
        else:
            # Fallback (hybrid_grid case is already handled above)
            return self._compute_full_width_positions(
                section, blocks, slide_number, is_first_slide
            )

    def _compute_full_width_positions(
        self,
        section: Section,
        blocks: List[ContentBlock],
        slide_number: int,
        is_first_slide: bool,
    ) -> SlideLayout:
        """
        Compute positions for full-width stacked layout - places ALL blocks.

        For full_width layout (no fixed rows×cols):
        - Charts/commentary: min height = 20% of width, max height = 50% of width
        - Gutter is considered in available space calculations
        """
        constraints = section.constraints
        layout_config = get_slide_layout_config(self.property_sub_type)
        element_dims = get_element_dimensions()

        # Determine starting position and margins from config
        if is_first_slide:
            current_top = layout_config.first_slide_start_top
            margin_left = (
                layout_config.first_slide_margin_left
                or layout_config.base_constraints.margin_left
            )
            margin_right = (
                layout_config.first_slide_margin_right
                or layout_config.base_constraints.margin_right
            )
            margin_bottom = (
                layout_config.first_slide_margin_bottom
                or layout_config.base_constraints.margin_bottom
            )
        else:
            current_top = constraints.margin_top
            margin_left = constraints.margin_left
            margin_right = constraints.margin_right
            margin_bottom = constraints.margin_bottom

        content_width = constraints.slide_width - margin_left - margin_right
        max_bottom = constraints.slide_height - margin_bottom
        gutter = layout_config.get_full_width_gutter(is_first_slide=is_first_slide)
        table_gutter = layout_config.get_full_width_table_gutter()

        cell_frames = []
        assignments = []
        scale_factors = {}

        # Calculate total gutter space needed
        total_gutter_space = gutter * max(0, len(blocks) - 1)
        available_height = max_bottom - current_top - total_gutter_space

        for idx, block in enumerate(blocks):
            # Get intrinsic size (pass content width for accurate calculations)
            if block.type == ContentType.TABLE or block.type == ContentType.CHART:
                intrinsic_w, intrinsic_h = block.estimate_intrinsic_size(
                    constraints, available_width=content_width
                )
            else:
                intrinsic_w, intrinsic_h = block.estimate_intrinsic_size(constraints)

            # GENERIC: Add section title space if this is the ABSOLUTE first element of its
            # logical section (display_order == 0) AND layout shows section titles.
            # Section title should only appear ONCE per section, not on every slide.
            block_display_order = getattr(block, "display_order", None)
            is_first_in_section_overall = block_display_order == 0
            show_section_titles = (
                layout_config.show_section_titles
                if hasattr(layout_config, "show_section_titles")
                else True
            )

            if is_first_in_section_overall and show_section_titles:
                section_title_space = element_dims.get_section_title_total_height()
                intrinsic_h += section_title_space
                print(
                    f'    📋 Section title for first {block.type.value} (display_order=0): +{section_title_space:.2f}"'
                )

            # Set minimum heights and apply constraints based on element type
            if block.type == ContentType.CHART or block.type == ContentType.TEXT:
                min_height = (
                    content_width * element_dims.dynamic_layout_min_height_ratio
                )
                max_height = (
                    content_width * element_dims.dynamic_layout_max_height_ratio
                )
                block_height = max(min_height, min(intrinsic_h, max_height))
            else:
                # Tables: use full intrinsic height (no max cap for accurate sizing)
                min_height = element_dims.table_min_row_height * 3
                block_height = max(min_height, intrinsic_h)

            # Check if element fits within slide bounds
            # NOTE: Source space is now included in block_height via estimate_intrinsic_size()
            remaining = max_bottom - current_top
            if block_height > remaining:
                print(
                    f'    ⚠️ {block.type.value.title()} ({block_height:.2f}") exceeds remaining ({remaining:.2f}")'
                )
                block_height = max(min_height, remaining)

            # Create cell
            print(
                f'    📦 Creating cell for block {block.id} ({block.type.value}): height={block_height:.2f}", intrinsic_h={intrinsic_h:.2f}"'
            )
            print(
                f'       Cell: top={current_top:.2f}", bottom={current_top + block_height:.2f}" (margin_bottom={margin_bottom:.2f}")'
            )
            cell = CellFrame(
                left=margin_left,
                top=current_top,
                width=content_width,
                height=block_height,
            )

            cell_frames.append(cell)
            assignments.append((idx, block.id))
            scale_factors[block.id] = block.estimate_scale_factor(cell, constraints)

            # Move to next position
            # NOTE: Source space is now included in block_height via estimate_intrinsic_size()
            current_top += block_height

            if idx < len(blocks) - 1:  # Add gutter between blocks
                next_block = blocks[idx + 1]
                # Use reduced gutter between consecutive tables
                is_current_table = block.type == ContentType.TABLE
                is_next_table = next_block.type == ContentType.TABLE
                effective_gutter = table_gutter if (is_current_table and is_next_table) else gutter
                print(
                    f'    ↕️  Adding gutter: +{effective_gutter:.2f}" (next element will start at {current_top + effective_gutter:.2f}")'
                    + (f' [table→table reduced gutter]' if (is_current_table and is_next_table) else '')
                )
                current_top += effective_gutter

        return SlideLayout(
            slide_number=slide_number,
            section_id=section.id,
            layout_type=LayoutType.FULL_WIDTH,
            cell_frames=cell_frames,
            assigned_blocks=assignments,
            scale_factors=scale_factors,
            overflow_notes=[],
            metadata={"all_blocks_placed": True},
        )

    def _compute_grid_positions(
        self,
        section: Section,
        blocks: List[ContentBlock],
        slide_number: int,
        is_first_slide: bool,
    ) -> SlideLayout:
        """
        Compute positions for 2x2 grid layout - places ALL blocks.

        For grid layout with fixed rows×cols (2×2):
        - Cell dimensions calculated from margins and gutter
        - Min = Max = cell dimensions (fixed layout constraint)
        - Elements must fit within their assigned cell
        - NEW: Enforces section boundary constraint - new sections must start from leftmost cell
        """
        constraints = section.constraints
        layout_config = get_slide_layout_config(self.property_sub_type)

        # Get grid dimensions from config (2x2 for standard grid)
        rows = 2
        cols = 2

        # Get margins and gutters from config
        if (
            is_first_slide
            and layout_config.first_slide_rows
            and layout_config.first_slide_cols
        ):
            # Use first slide specific configuration if available
            rows = layout_config.first_slide_rows
            cols = layout_config.first_slide_cols
            margin_left = (
                layout_config.first_slide_margin_left
                or layout_config.base_constraints.margin_left
            )
            margin_top = layout_config.first_slide_start_top
            margin_right = (
                layout_config.first_slide_margin_right
                or layout_config.base_constraints.margin_right
            )
            margin_bottom = (
                layout_config.first_slide_margin_bottom
                or layout_config.base_constraints.margin_bottom
            )
            gutter_h = (
                layout_config.first_slide_gutter_horizontal
                or layout_config.base_constraints.gutter_horizontal
            )
            gutter_v = (
                layout_config.first_slide_gutter_vertical
                or layout_config.base_constraints.gutter_vertical
            )
        else:
            # Use standard constraints
            margin_left = constraints.margin_left
            margin_top = constraints.margin_top
            margin_right = constraints.margin_right
            margin_bottom = constraints.margin_bottom
            gutter_h = constraints.gutter_horizontal
            gutter_v = constraints.gutter_vertical

        # Calculate content area
        content_width = constraints.slide_width - margin_left - margin_right
        content_height = constraints.slide_height - margin_top - margin_bottom

        print(
            f"       Grid layout for slide {slide_number}: {rows}×{cols}, first_slide={is_first_slide}"
        )

        # Calculate total gutter space
        total_h_gutter = gutter_h * (cols - 1)
        total_v_gutter = gutter_v * (rows - 1)

        # Calculate cell dimensions based on margins and gutter
        # For fixed layouts: min = max = cell dimensions
        cell_w = (content_width - total_h_gutter) / cols
        cell_h = (content_height - total_v_gutter) / rows

        cell_frames = []
        assignments = []
        scale_factors = {}
        current_section_id = None  # Track current section for boundary checking
        cell_idx = 0  # Track actual cell index for placement

        for block in blocks:
            # Calculate which row and column this cell would be at
            row = cell_idx // cols
            col = cell_idx % cols

            # SECTION BOUNDARY CONSTRAINT: New sections must start from leftmost cell (col == 0)
            if current_section_id is not None and block.section_id and block.section_id != current_section_id:
                # Section changed - check if we're at a leftmost cell
                if col != 0:
                    # Not at leftmost cell - skip to next row's leftmost cell
                    print(
                        f"       Section boundary: Section '{block.section_id}' must start from left, "
                        f"skipping from col {col} to next row"
                    )
                    # Move to the start of the next row
                    cell_idx = (row + 1) * cols
                    row = cell_idx // cols
                    col = cell_idx % cols

            # Check if we've exceeded the grid capacity
            if cell_idx >= rows * cols:
                print(
                    f"       Grid capacity exceeded at block {block.id}, stopping placement"
                )
                break

            # Fixed layout: all cells have same dimensions (min = max)
            cell = CellFrame(
                left=margin_left + col * (cell_w + gutter_h),
                top=margin_top + row * (cell_h + gutter_v),
                width=cell_w,
                height=cell_h,
            )

            cell_frames.append(cell)
            assignments.append((len(cell_frames) - 1, block.id))
            scale_factors[block.id] = block.estimate_scale_factor(cell, constraints)

            # Update current section tracking
            if block.section_id:
                current_section_id = block.section_id

            cell_idx += 1

        return SlideLayout(
            slide_number=slide_number,
            section_id=section.id,
            layout_type=LayoutType.GRID_2x2,
            cell_frames=cell_frames,
            assigned_blocks=assignments,
            scale_factors=scale_factors,
            overflow_notes=[],
            metadata={
                "all_blocks_placed": len(assignments) == len(blocks),
                "rows": rows,
                "cols": cols,
                "fixed_layout": True,
            },
        )

    def _compute_hybrid_grid_positions(
        self,
        section: Section,
        blocks: List[ContentBlock],
        slide_number: int,
        is_first_slide: bool,
    ) -> SlideLayout:
        """
        Compute positions for hybrid grid layout - mixed grid and full-width elements.

        Processes blocks in their original order:
        - Consecutive grid elements are grouped and placed in grid layout (N×cols)
        - Full-width elements are placed individually in full-width layout
        - Supports any combination in any order (e.g., grid→full→grid→full or full→grid→full)
        """
        constraints = section.constraints
        layout_config = get_slide_layout_config(self.property_sub_type)
        element_dims = get_element_dimensions()

        # Get margins and gutters
        margin_left = constraints.margin_left
        margin_top = constraints.margin_top
        margin_right = constraints.margin_right
        margin_bottom = constraints.margin_bottom
        gutter_h = constraints.gutter_horizontal
        gutter_v = constraints.gutter_vertical

        content_width = constraints.slide_width - margin_left - margin_right
        max_bottom = constraints.slide_height - margin_bottom

        # Tag blocks with their layout type while preserving order
        blocks_with_layout = []
        grid_count = 0
        full_width_count = 0

        for block in blocks:
            layout = None
            if hasattr(block, "metadata") and isinstance(block.metadata, dict):
                layout = block.metadata.get("layout")

            if layout == "grid_2x2":
                blocks_with_layout.append((block, "grid_2x2"))
                grid_count += 1
            elif layout == "full_width":
                blocks_with_layout.append((block, "full_width"))
                full_width_count += 1
            else:
                # Default to grid if not specified
                blocks_with_layout.append((block, "grid_2x2"))
                grid_count += 1

        print(
            f"       Hybrid layout: {grid_count} grid blocks, {full_width_count} full-width blocks (preserving order)"
        )

        cell_frames = []
        assignments = []
        scale_factors = {}
        current_top = margin_top

        # Process blocks in original order, grouping consecutive grid blocks
        idx = 0
        while idx < len(blocks_with_layout):
            block, layout_type = blocks_with_layout[idx]

            if layout_type == "grid_2x2":
                # Collect consecutive grid blocks
                grid_group = [block]
                idx += 1
                while (
                    idx < len(blocks_with_layout)
                    and blocks_with_layout[idx][1] == "grid_2x2"
                ):
                    grid_group.append(blocks_with_layout[idx][0])
                    idx += 1

                # Place grid group using grid layout
                cols = layout_config.hybrid_grid_cols
                # Calculate actual rows needed based on number of blocks (don't waste space on empty rows)
                rows = math.ceil(len(grid_group) / cols)
                total_h_gutter = gutter_h * (cols - 1)
                # Use smaller vertical gutter for hybrid layouts to maximize space
                hybrid_gutter_v = layout_config.get_hybrid_gutter_vertical(
                    is_first_slide=is_first_slide
                )
                total_v_gutter = hybrid_gutter_v * (rows - 1)

                # Calculate required height for grid blocks based on actual content
                max_grid_block_height = 0.0
                for block in grid_group:
                    _, intrinsic_h = block.estimate_intrinsic_size(constraints)
                    max_grid_block_height = max(max_grid_block_height, intrinsic_h)

                # Available space for content
                available_height = max_bottom - current_top

                # Use intrinsic size if it fits, otherwise allocate proportionally
                grid_height_needed = (max_grid_block_height * rows) + total_v_gutter
                if grid_height_needed <= available_height:
                    # Enough space - use full intrinsic heights
                    grid_height = grid_height_needed
                else:
                    # Not enough space - use proportional allocation
                    # For multi-section hybrid layouts, give reasonable space based on available height
                    grid_height = min(
                        available_height * layout_config.hybrid_grid_height_ratio,
                        available_height * 0.8,
                    )

                cell_w = (content_width - total_h_gutter) / cols
                cell_h = (grid_height - total_v_gutter) / rows

                for grid_idx, grid_block in enumerate(grid_group):
                    row = grid_idx // cols
                    col = grid_idx % cols

                    cell = CellFrame(
                        left=margin_left + col * (cell_w + gutter_h),
                        top=current_top + row * (cell_h + hybrid_gutter_v),
                        width=cell_w,
                        height=cell_h,
                    )

                    cell_frames.append(cell)
                    assignments.append((blocks.index(grid_block), grid_block.id))
                    scale_factors[grid_block.id] = grid_block.estimate_scale_factor(
                        cell, constraints
                    )

                # Move past grid area - use configured spacing for hybrid layouts
                hybrid_spacing = layout_config.hybrid_section_spacing
                current_top += grid_height + hybrid_spacing

            elif layout_type == "full_width":
                # Place single full-width block
                gutter = layout_config.get_full_width_gutter(
                    is_first_slide=is_first_slide
                )

                # For full-width elements, pass the actual available width for accurate height calculation
                if block.type == ContentType.TABLE or block.type == ContentType.CHART:
                    intrinsic_w, intrinsic_h = block.estimate_intrinsic_size(
                        constraints, available_width=content_width
                    )
                else:
                    intrinsic_w, intrinsic_h = block.estimate_intrinsic_size(
                        constraints
                    )

                # GENERIC: Add section title space if first element AND layout shows section titles
                # Applies to ALL element types
                is_first_in_section = idx == 0
                show_section_titles = (
                    layout_config.show_section_titles
                    if hasattr(layout_config, "show_section_titles")
                    else True
                )

                if is_first_in_section and show_section_titles:
                    section_title_space = element_dims.get_section_title_total_height()
                    intrinsic_h += section_title_space
                    print(
                        f'    📋 Section title for first {block.type.value}: +{section_title_space:.2f}"'
                    )

                # Set minimum heights and apply constraints based on element type
                if block.type == ContentType.CHART or block.type == ContentType.TEXT:
                    min_height = (
                        content_width * element_dims.dynamic_layout_min_height_ratio
                    )
                    max_height = (
                        content_width * element_dims.dynamic_layout_max_height_ratio
                    )
                    block_height = max(min_height, min(intrinsic_h, max_height))
                else:
                    # Tables: use full intrinsic height (no max cap for accurate sizing)
                    min_height = element_dims.table_min_row_height * 3
                    block_height = max(min_height, intrinsic_h)

                # GENERIC: Check if element + source fits within slide bounds
                # For tables, source is now a row INSIDE the table, so don't add external source_space
                # For charts, source is still external, so add source_space
                has_source = any(
                    [
                        hasattr(block, "table_source") and block.table_source,
                        hasattr(block, "figure_source") and block.figure_source,
                        hasattr(block, "source") and block.source,
                    ]
                )

                remaining = max_bottom - current_top
                # Tables: source is inside table (already in block_height via estimate_intrinsic_size)
                # Charts/others: source is external, need to add source_space
                is_table = block.type == ContentType.TABLE
                if has_source and not is_table:
                    source_space = (
                        element_dims.source_gap + element_dims.source_label_height
                    )
                    total_space_needed = block_height + source_space

                    if total_space_needed > remaining:
                        print(
                            f'    ⚠️ {block.type.value.title()} + source ({total_space_needed:.2f}") exceeds remaining ({remaining:.2f}")'
                        )
                        # Reduce block height to fit (source space is mandatory)
                        block_height = max(min_height, remaining - source_space)
                else:
                    # No source - just check element height
                    if block_height > remaining:
                        print(
                            f'    ⚠️ {block.type.value.title()} ({block_height:.2f}") exceeds remaining ({remaining:.2f}")'
                        )
                        block_height = max(min_height, remaining)

                cell = CellFrame(
                    left=margin_left,
                    top=current_top,
                    width=content_width,
                    height=block_height,
                )

                cell_frames.append(cell)
                assignments.append((blocks.index(block), block.id))
                scale_factors[block.id] = block.estimate_scale_factor(cell, constraints)

                # NOTE: Source space is now included in block_height via estimate_intrinsic_size()
                # For tables, the source row is INSIDE the table, so it's part of block_height
                # but doesn't create extra visual gap between tables (only gutter creates gap)
                current_top += block_height

                # Add gutter after full-width block (will be overridden if next is a grid group)
                if idx < len(blocks_with_layout):
                    current_top += gutter

                idx += 1

        return SlideLayout(
            slide_number=slide_number,
            section_id=section.id,
            layout_type=LayoutType.HYBRID_GRID,
            cell_frames=cell_frames,
            assigned_blocks=assignments,
            scale_factors=scale_factors,
            overflow_notes=[],
            metadata={
                "all_blocks_placed": True,
                "hybrid": True,
                "grid_blocks": grid_count,
                "full_width_blocks": full_width_count,
            },
        )

    def _compute_base_slide_positions(
        self, section: Section, blocks: List[ContentBlock], slide_number: int
    ) -> SlideLayout:
        """
        Compute positions for base slide layout - places ALL blocks.

        Base slide layout positions elements horizontally below KPIs.
        Uses config-based margins, gutters, and KPI section height.
        """
        constraints = section.constraints
        layout_config = get_slide_layout_config(self.property_sub_type)

        # Get KPI section height and spacing from config
        kpi_height = layout_config.first_slide_kpi_section_height
        if kpi_height is None:
            # Fallback if not configured
            kpi_height = layout_config.first_slide_start_top

        # Get margins from config
        margin_left = (
            layout_config.first_slide_margin_left
            or layout_config.base_constraints.margin_left
        )
        margin_right = (
            layout_config.first_slide_margin_right
            or layout_config.base_constraints.margin_right
        )
        margin_bottom = (
            layout_config.first_slide_margin_bottom
            or layout_config.base_constraints.margin_bottom
        )

        # Calculate content dimensions
        content_width = constraints.slide_width - margin_left - margin_right

        # Get gutter from config
        gutter = (
            layout_config.first_slide_gutter_horizontal
            or layout_config.base_constraints.gutter_horizontal
        )

        # Start position below KPIs (use config-based spacing)
        start_top = kpi_height

        # Calculate available height
        available_height = constraints.slide_height - margin_bottom - start_top

        # Distribute width among blocks with gutter
        num_blocks = len(blocks)
        total_gutter = gutter * max(0, num_blocks - 1)
        block_width = (
            (content_width - total_gutter) / num_blocks
            if num_blocks > 0
            else content_width
        )

        cell_frames = []
        assignments = []
        scale_factors = {}

        current_left = margin_left
        for idx, block in enumerate(blocks):
            cell = CellFrame(
                left=current_left,
                top=start_top,
                width=block_width,
                height=available_height,
            )

            cell_frames.append(cell)
            assignments.append((idx, block.id))
            scale_factors[block.id] = block.estimate_scale_factor(cell, constraints)

            current_left += block_width + gutter

        return SlideLayout(
            slide_number=slide_number,
            section_id=section.id,
            layout_type=LayoutType.BASE_SLIDE,
            cell_frames=cell_frames,
            assigned_blocks=assignments,
            scale_factors=scale_factors,
            overflow_notes=[],
            metadata={"all_blocks_placed": True, "kpi_height": kpi_height},
        )

    def _create_slide_layout(
        self,
        section: Section,
        blocks: List[ContentBlock],
        slide_number: int,
        total_slides: Optional[int] = None,
    ) -> Tuple[SlideLayout, List[str]]:
        """
        Create layout for a single slide.

        Args:
            section: Section with content blocks
            blocks: List of content blocks to place
            slide_number: Current slide number (1-indexed)
            total_slides: Total number of slides (None if unknown)

        Returns:
            (SlideLayout, list of consumed block IDs)
        """
        # CRITICAL: Filter blocks to only those assigned to this slide number
        # This prevents blocks assigned to other slides from being processed on the wrong slide
        filtered_blocks = []
        blocks_without_slide_num = []
        blocks_wrong_slide = []
        for block in blocks:
            block_slide_num = None
            if hasattr(block, "metadata") and isinstance(block.metadata, dict):
                block_slide_num = block.metadata.get("slide_number")

            # Include block if it's assigned to this slide, or if it has no slide number (fallback)
            if block_slide_num == slide_number:
                filtered_blocks.append(block)
            elif block_slide_num is None:
                # Only include blocks without slide number if we're in a context where that makes sense
                # (e.g., when processing a section that hasn't been pre-assigned)
                filtered_blocks.append(block)
                blocks_without_slide_num.append(block.id)
            else:
                blocks_wrong_slide.append((block.id, block_slide_num))

        if blocks_without_slide_num:
            print(
                f"    ⚠ {len(blocks_without_slide_num)} blocks without slide number included (fallback): {blocks_without_slide_num}"
            )
        if blocks_wrong_slide:
            print(
                f"    ⚠ {len(blocks_wrong_slide)} blocks excluded (wrong slide): {blocks_wrong_slide}"
            )
        print(
            f"    📋 Filtered to {len(filtered_blocks)} blocks for slide {slide_number} (from {len(blocks)} total)"
        )

        # Use filtered blocks for all subsequent processing
        blocks = filtered_blocks

        # If no blocks after filtering, return empty layout immediately
        if not blocks:
            print(
                f"    ⚠ No blocks to process for slide {slide_number} after filtering"
            )
            empty_layout = SlideLayout(
                slide_number=slide_number,
                section_id=section.id,
                layout_type=LayoutType.GRID_2x2,
                cell_frames=[],
                assigned_blocks=[],
                scale_factors={},
                overflow_notes=[f"No blocks assigned to slide {slide_number}"],
            )
            return (empty_layout, [])

        # Layout preference is not allowed for 1st and last slides (only applies to middle slides)
        is_first_slide = slide_number == 1
        is_last_slide = total_slides is not None and slide_number == total_slides

        # Check if user has a specific layout preference
        # Ignore layout_preference for first and last slides
        if section.layout_preference and not is_first_slide and not is_last_slide:
            preferred_strategy, is_known_preference = self._get_preferred_strategy(
                section.layout_preference
            )
            if preferred_strategy:
                print(f"    🎯 Using preferred layout: {section.layout_preference}")
                result = preferred_strategy(section, blocks, slide_number)
                if result and self._is_layout_allowed(result[0].layout_type):
                    return result
                else:
                    print(
                        "    ⚠ Preferred layout failed, falling back to auto selection"
                    )
            elif not is_known_preference:
                print(
                    f"    ⚠ Preferred layout '{section.layout_preference}' not supported for {self.property_sub_type}"
                )
        elif section.layout_preference and (is_first_slide or is_last_slide):
            slide_type = "first" if is_first_slide else "last"
            print(
                f"    ℹ️  Layout preference '{section.layout_preference}' ignored for {slide_type} slide (only applies to middle slides)"
            )

        # For first slide: enforce rules based on property_sub_type from config
        if is_first_slide:
            first_slide_result = self._try_base_slide(section, blocks, slide_number)
            if first_slide_result:
                return first_slide_result
            # If first slide layout fails, fall through to default strategies

        # Try different layout strategies in order of preference
        # For middle slides: use default layout sequence when no preference, or use preference when valid
        strategies = self._get_strategy_sequence(blocks)

        for strategy in strategies:
            result = strategy(section, blocks, slide_number)
            if result and self._is_layout_allowed(result[0].layout_type):
                return result

        # Fallback: create empty slide with overflow note
        empty_layout = SlideLayout(
            slide_number=slide_number,
            section_id=section.id,
            layout_type=LayoutType.GRID_2x2,
            cell_frames=[],
            assigned_blocks=[],
            scale_factors={},
            overflow_notes=["Could not fit any content - check constraints"],
        )
        return (empty_layout, [])

    def _get_preferred_strategy(
        self, layout_preference: str
    ) -> Tuple[Optional[Callable], bool]:
        """Map UI layout preference to orchestrator strategy method."""
        preference_map = {
            "First Slide (Base with KPIs)": self._try_base_slide,
            "Content (2x2 Grid)": self._try_grid_2x2,
            "Full Width": self._try_full_width_stack,
            "Auto (Smart Layout)": None,  # Use default strategy order
        }
        if layout_preference in preference_map:
            return preference_map[layout_preference], True
        return None, False

    def _optimize_table_layout(
        self,
        table_block: TableBlock,
        constraints: SlideConstraints,
        current_y: float = 0.0,
        slide_number: int = 1,
    ) -> Tuple[str, CellFrame, float]:
        """
        Optimize table layout by evaluating potential layouts and choosing the best fit.

        Steps:
        1. Identify potential layouts based on property_sub_type and constraints
        2. For each potential layout, calculate table width/height
        3. Choose the layout that optimizes space usage and satisfies all constraints
        4. Account for gutters in calculations

        Args:
            table_block: TableBlock to optimize
            constraints: Slide constraints
            current_y: Current vertical position (for full_width calculations)
            slide_number: Current slide number

        Returns:
            Tuple of (best_layout_type, best_cell_frame, best_scale_factor)
        """
        from hello.utils.ppt_helpers_utils.services.template_config import (
            get_allowed_layout_types,
            get_element_dimensions,
            get_layout_threshold_config,
        )

        # Get potential layouts based on property_sub_type
        allowed_layouts = get_allowed_layout_types(self.property_sub_type)
        element_dims = get_element_dimensions()
        threshold_config = get_layout_threshold_config(self.property_sub_type)

        # Calculate intrinsic table size
        intrinsic_w, intrinsic_h = table_block.estimate_intrinsic_size(constraints)

        best_layout = None
        best_cell = None
        best_score = -1.0  # Score: higher is better (space utilization)
        best_scale = 1.0

        # Evaluate each potential layout
        for layout_type in allowed_layouts:
            cell = None

            if layout_type == "grid_2x2":
                # Grid cell: use _calculate_grid_cells for consistency
                # This properly accounts for gutters and uses chart dimensions
                grid_cells = self._calculate_grid_cells(
                    rows=2, cols=2, constraints=constraints
                )

                # Use the first cell (top-left) as reference for optimization
                # In actual assignment, all 4 cells would be available
                if grid_cells:
                    cell = grid_cells[0]
                else:
                    continue

            elif layout_type == "full_width":
                # Full width: account for table margins and gutters
                # Check if this is a submarket table
                is_submarket_table = (
                    hasattr(table_block, "template_path")
                    and table_block.template_path
                    and "table_market_submarket"
                    in str(table_block.template_path).lower()
                )

                if is_submarket_table:
                    # Submarket table: use constraints margins
                    table_margin_left = constraints.margin_left
                    table_margin_right = constraints.margin_right
                    table_margin_top = constraints.margin_top

                    available_width = (
                        constraints.slide_width - table_margin_left - table_margin_right
                    )
                    available_height = (
                        constraints.slide_height
                        - constraints.margin_bottom
                        - table_margin_top
                    )

                    cell = CellFrame(
                        left=table_margin_left,
                        top=table_margin_top,
                        width=available_width,
                        height=min(intrinsic_h, available_height),
                    )
                else:
                    # Regular table: use constraints margins
                    table_margin_left = constraints.margin_left
                    table_margin_right = constraints.margin_right

                    available_width = (
                        constraints.slide_width - table_margin_left - table_margin_right
                    )
                    available_height = (
                        constraints.slide_height - constraints.margin_bottom - current_y
                    )

                    # Cap height to reasonable maximum using element dimensions
                    max_reasonable_height = min(
                        constraints.content_height * element_dims.max_height_ratio,
                        element_dims.max_height_absolute,
                    )
                    capped_intrinsic_h = min(intrinsic_h, max_reasonable_height)
                    row_height = min(capped_intrinsic_h, available_height)

                    cell = CellFrame(
                        left=table_margin_left,
                        top=current_y,
                        width=available_width,
                        height=row_height,
                    )

            if cell is None:
                continue

            # Check if table can fit in this cell
            if not table_block.can_fit_in_cell(cell, constraints):
                continue

            # Calculate scale factor
            scale = table_block.estimate_scale_factor(cell, constraints)

            # Use threshold from config (already loaded at method start)
            min_scale_threshold = threshold_config.min_scale_threshold_normal

            if scale < min_scale_threshold:
                continue  # Doesn't meet minimum scale requirement

            # Calculate score: space utilization (area used / area available)
            # Higher score = better utilization
            table_area = intrinsic_w * intrinsic_h * scale
            cell_area = cell.width * cell.height
            utilization = table_area / cell_area if cell_area > 0 else 0.0

            # Prefer layouts with better utilization and higher scale
            score = utilization * scale

            if score > best_score:
                best_score = score
                best_layout = layout_type
                best_cell = cell
                best_scale = scale

        # Fallback to full_width if no layout found
        if best_layout is None:
            best_layout = "full_width"
            table_margin_left = constraints.margin_left
            table_margin_right = constraints.margin_right
            available_width = (
                constraints.slide_width - table_margin_left - table_margin_right
            )
            available_height = (
                constraints.slide_height - constraints.margin_bottom - current_y
            )

            # Cap height to reasonable maximum
            max_reasonable_height = min(
                constraints.content_height * element_dims.max_height_ratio,
                element_dims.max_height_absolute,
            )
            capped_intrinsic_h = min(intrinsic_h, max_reasonable_height)
            row_height = min(capped_intrinsic_h, available_height)

            best_cell = CellFrame(
                left=table_margin_left,
                top=current_y,
                width=available_width,
                height=row_height,
            )
            best_scale = table_block.estimate_scale_factor(best_cell, constraints)

        return (best_layout, best_cell, best_scale)

    def _assign_blocks_to_cells(
        self,
        blocks: List[ContentBlock],
        cell_frames: List[CellFrame],
        constraints: SlideConstraints,
        allow_aggressive_scaling: bool = True,
        respect_section_boundaries: bool = True,
    ) -> Tuple[List[Tuple[int, str]], Dict[str, float], List[str], List[str]]:
        """
        Assign blocks to cells using skip-and-backfill logic.
        Preserves display_order while maximizing cell utilization.

        For each cell, tries all unassigned blocks in their original order
        and assigns the first one that fits. This ensures:
        - Display order is preserved (blocks appear in original sequence)
        - Grid cells are maximally utilized (no empty cells when content could fit)
        - Mixed content types can share slides
        - NEW: Sections start from leftmost cells (column 1) when respect_section_boundaries=True

        Args:
            blocks: List of content blocks in display order
            cell_frames: List of available cells on the slide
            constraints: Layout constraints
            allow_aggressive_scaling: If True, allow scaling down to 40%, else 60%
            respect_section_boundaries: If True, ensure new sections start from leftmost cell

        Returns:
            Tuple of (assignments, scale_factors, overflow_notes, consumed_block_ids)
            - assignments: List of (cell_index, block_id) pairs
            - scale_factors: Dict mapping block_id to scale factor
            - overflow_notes: List of warning messages
            - consumed_block_ids: List of block IDs in display order
        """
        assignments = []
        scale_factors = {}
        overflow_notes = []
        assigned_block_ids = set()
        current_section_id = None  # Track current section for boundary checking
        section_boundary_blocked = (
            False  # Flag to stop processing when section needs fresh slide
        )

        # Get threshold config
        threshold_config = get_layout_threshold_config(self.property_sub_type)
        min_scale_threshold = (
            threshold_config.min_scale_threshold_aggressive
            if allow_aggressive_scaling
            else threshold_config.min_scale_threshold_normal
        )

        # Pre-identify leftmost cell indices (cells in the leftmost column)
        leftmost_cell_indices = set()
        if cell_frames:
            min_left = min(c.left for c in cell_frames)
            tolerance = 0.1
            for idx, cell in enumerate(cell_frames):
                if abs(cell.left - min_left) < tolerance:
                    leftmost_cell_indices.add(idx)
            print(f"       Leftmost cell indices: {sorted(leftmost_cell_indices)}")

        # For each cell, try to find the first unassigned block (in order) that fits
        for cell_idx, cell in enumerate(cell_frames):
            assigned_to_cell = False

            # Try each unassigned block in display order
            for block in blocks:
                if block.id in assigned_block_ids:
                    continue  # Already assigned

                # NEW: Check section boundary constraint
                # Only enforce when switching to a TRULY DIFFERENT section
                # Allow same section to continue even if not at leftmost
                if respect_section_boundaries and current_section_id is not None:
                    # If block belongs to a different section
                    if block.section_id and block.section_id != current_section_id:
                        # Section changed - new section needs to start at leftmost cell
                        if not self._is_leftmost_cell(cell, cell_frames):
                            # Current cell is NOT leftmost
                            # Check if there are any leftmost cells available ahead
                            available_leftmost_ahead = any(
                                idx > cell_idx and idx in leftmost_cell_indices
                                for idx in leftmost_cell_indices
                            )

                            if available_leftmost_ahead:
                                # There's a leftmost cell ahead, skip to it
                                print(
                                    f"       Section boundary: Section '{block.section_id}' needs leftmost cell, skipping cell {cell_idx}"
                                )
                                break  # Break inner loop to try next cell
                            else:
                                # NO leftmost cells available ahead
                                # This means new section would share slide with previous section
                                # Stop processing entirely to let new section overflow to fresh slide
                                print(
                                    f"       Section boundary BLOCKED: Section '{block.section_id}' needs fresh slide (no leftmost cells available)"
                                )
                                section_boundary_blocked = True
                                break  # Break inner loop
                        # else: Current cell IS leftmost
                        # Check if this is cell 0 (top-left) with existing assignments, or cell 2 (bottom-left)
                        elif assignments and cell_idx == 0:
                            # New section trying to start at cell 0, but previous section already has assignments
                            # This should not happen - previous section should have its own slide
                            print(
                                f"       Section boundary: Section '{block.section_id}' needs fresh slide (previous section has {len(assignments)} assignments at cell 0)"
                            )
                            section_boundary_blocked = True
                            break  # Break inner loop
                        # If current cell is cell 2 (bottom-left) and assignments exist in cells 0-1, that's OK - sections can share
                    # If same section continues, allow it regardless of cell position
                    # This handles cases where sections were split across slides by assignment logic

                # Check if block fits in this cell
                can_fit = False
                scale = 1.0

                # Be more lenient for text blocks - they can scale down more
                if block.type == ContentType.TEXT:
                    scale = block.estimate_scale_factor(cell, constraints)
                    can_fit = scale >= min_scale_threshold
                else:
                    # For charts and tables, check fit and scale factor
                    scale = block.estimate_scale_factor(cell, constraints)
                    can_fit = scale >= min_scale_threshold

                if can_fit:
                    # This block fits! Assign it to this cell
                    assignments.append((cell_idx, block.id))
                    scale_factors[block.id] = scale
                    assigned_block_ids.add(block.id)
                    assigned_to_cell = True

                    # Track current section for boundary checking
                    if block.section_id:
                        current_section_id = block.section_id

                    # Add overflow note if heavily scaled
                    if scale < 0.7:
                        overflow_notes.append(
                            f"Block {block.id} scaled to {scale:.1%} in cell {cell_idx}"
                        )

                    print(
                        f"       Cell {cell_idx}: Assigned {block.id} (scale={scale:.1%})"
                    )
                    break  # Move to next cell

            # Check if section boundary blocked - exit outer loop
            if section_boundary_blocked:
                print(
                    "       ⚠️  Stopping cell assignment - remaining blocks will overflow to next slide"
                )
                break  # Exit outer loop

            if not assigned_to_cell:
                print(f"       Cell {cell_idx}: Empty (no suitable blocks found)")

        # Build consumed_block_ids list in display order
        consumed_block_ids = [
            block.id for block in blocks if block.id in assigned_block_ids
        ]

        return (assignments, scale_factors, overflow_notes, consumed_block_ids)

    def _is_leftmost_cell(self, cell: CellFrame, cell_frames: List[CellFrame]) -> bool:
        """
        Check if a cell is in the leftmost column (column 1).

        Args:
            cell: The cell to check
            cell_frames: All cell frames in the layout

        Returns:
            True if cell is in leftmost column, False otherwise
        """
        if not cell_frames:
            return True

        # Find the minimum left position (leftmost column)
        min_left = min(c.left for c in cell_frames)

        # Cell is leftmost if within tolerance of minimum
        tolerance = 0.1  # 0.1" tolerance for floating point comparison
        return abs(cell.left - min_left) < tolerance

    def _try_base_slide(
        self, section: Section, blocks: List[ContentBlock], slide_number: int
    ) -> Optional[Tuple[SlideLayout, List[str]]]:
        """
        Try property-sub-type-specific first slide layouts.
        """
        if slide_number != 1:
            return None

        layout_config = get_slide_layout_config(self.property_sub_type)
        print(
            f"    🎛  First slide profile '{layout_config.first_slide_style}' for sub-type '{self.property_sub_type}'"
        )

        if layout_config.first_slide_style == "kpi_row":
            return self._layout_first_slide_kpi(
                section, blocks, slide_number, layout_config
            )
        if layout_config.first_slide_style == "full_width":
            return self._layout_first_slide_full_width(
                section, blocks, slide_number, layout_config
            )
        if layout_config.first_slide_style == "grid":
            return self._layout_first_slide_grid(
                section, blocks, slide_number, layout_config
            )

        print(
            f"    ⚠ Unknown first slide style '{layout_config.first_slide_style}', skipping"
        )
        return None

    def _layout_first_slide_kpi(
        self,
        section: Section,
        blocks: List[ContentBlock],
        slide_number: int,
        layout_config: SlideLayoutConfig,
    ) -> Optional[Tuple[SlideLayout, List[str]]]:
        """Figures-style first slide layout (existing KPI row logic)."""
        suitable_blocks = [
            block
            for block in blocks
            if block.type in [ContentType.CHART, ContentType.TEXT]
        ]

        if not suitable_blocks:
            print("    ⚠ Base slide: No suitable blocks (need CHART or TEXT)")
            return None

        constraints = section.constraints
        slide_width = constraints.slide_width
        slide_height = constraints.slide_height
        margin_bottom = (
            layout_config.first_slide_margin_bottom
            or layout_config.base_constraints.margin_bottom
        )
        available_height = (
            slide_height - margin_bottom - layout_config.first_slide_start_top
        )

        if available_height < 2.0:
            return None

        margin_left = (
            layout_config.first_slide_margin_left
            or layout_config.base_constraints.margin_left
        )
        margin_right = (
            layout_config.first_slide_margin_right
            or layout_config.base_constraints.margin_right
        )
        content_area = CellFrame(
            left=margin_left,
            top=layout_config.first_slide_start_top,
            width=slide_width - margin_left - margin_right,
            height=available_height,
        )

        capacity = layout_config.first_slide_max_elements
        blocks_to_fit = self._limit_blocks(suitable_blocks, capacity)
        num_blocks = len(blocks_to_fit)

        print(
            f"    📊 Base slide: Fitting {num_blocks} blocks (from {len(suitable_blocks)} suitable, {len(blocks)} total)"
        )
        capacity_label = f"max {capacity}" if capacity is not None else "unbounded"
        print(
            f"    📊 First slide capacity ({self.property_sub_type}): {capacity_label} elements below KPIs"
        )

        # Calculate dimensions dynamically based on content area and number of blocks
        cells: List[CellFrame] = []
        if num_blocks == 1 and blocks_to_fit:
            # Single block: use full content area dimensions
            cells.append(
                CellFrame(
                    left=content_area.left,
                    top=content_area.top,
                    width=content_area.width,
                    height=content_area.height,
                )
            )
        elif num_blocks > 1:
            # Multiple blocks: distribute width evenly with gutters
            gutter = (
                layout_config.first_slide_gutter_horizontal
                or layout_config.base_constraints.gutter_horizontal
            )
            total_gutter = gutter * (num_blocks - 1)
            available_width = content_area.width - total_gutter
            cell_width = available_width / max(1, num_blocks)
            current_left = content_area.left
            for block in blocks_to_fit:
                cells.append(
                    CellFrame(
                        left=current_left,
                        top=content_area.top,
                        width=cell_width,
                        height=content_area.height,
                    )
                )
                current_left += cell_width + gutter

        print(f"    📐 Base slide created {len(cells)} cells for {num_blocks} blocks:")
        for i, cell in enumerate(cells):
            print(
                f'       Cell {i}: left={cell.left:.2f}", top={cell.top:.2f}", w={cell.width:.2f}", h={cell.height:.2f}"'
            )

        assignments: List[Tuple[int, str]] = []
        scale_factors: Dict[str, float] = {}
        consumed_ids: List[str] = []

        for i, block in enumerate(blocks_to_fit):
            if i >= len(cells):
                break

            cell = cells[i]
            if block.can_fit_in_cell(cell, constraints):
                assignments.append((i, block.id))
                scale_factors[block.id] = block.estimate_scale_factor(cell, constraints)
                consumed_ids.append(block.id)
            else:
                scale = block.estimate_scale_factor(cell, constraints)
                threshold_config = get_layout_threshold_config(self.property_sub_type)
                if scale >= threshold_config.min_scale_threshold_normal:
                    assignments.append((i, block.id))
                    scale_factors[block.id] = scale
                    consumed_ids.append(block.id)
                else:
                    print(
                        f"       Block {i} rejected: scale factor {scale:.2f} < 0.6 threshold"
                    )

        if not assignments:
            print("    ⚠ Base slide: No blocks could be assigned to cells")
            return None

        layout = SlideLayout(
            slide_number=slide_number,
            section_id=section.id,
            layout_type=LayoutType.BASE_SLIDE,
            cell_frames=cells[: len(assignments)],
            assigned_blocks=assignments,
            scale_factors=scale_factors,
            overflow_notes=[],
            metadata={
                "blocks_fitted": len(assignments),
                "layout_style": "base_slide_with_kpis",
                "kpi_section_height": layout_config.first_slide_kpi_section_height,
                "available_space_used": True,
            },
        )

        if not self._validate_layout_no_overlaps(layout, section, blocks):
            return None

        return (layout, consumed_ids)

    def _layout_first_slide_full_width(
        self,
        section: Section,
        blocks: List[ContentBlock],
        slide_number: int,
        layout_config: SlideLayoutConfig,
    ) -> Optional[Tuple[SlideLayout, List[str]]]:
        """Snapshot-style first slide layout (full width stack)."""
        if not blocks:
            return None

        blocks_to_fit = self._limit_blocks(
            blocks, layout_config.first_slide_max_elements
        )
        constraints = section.constraints
        current_top = layout_config.first_slide_start_top
        margin_bottom = (
            layout_config.first_slide_margin_bottom
            or layout_config.base_constraints.margin_bottom
        )
        max_top = constraints.slide_height - margin_bottom

        # Get configurable gutter for first slide
        gutter = layout_config.get_full_width_gutter(is_first_slide=True)

        cell_frames: List[CellFrame] = []
        assignments: List[Tuple[int, str]] = []
        scale_factors: Dict[str, float] = {}
        consumed_ids: List[str] = []
        overflow_notes: List[str] = []

        # Get threshold config for scale factor checks
        threshold_config = get_layout_threshold_config(self.property_sub_type)
        min_scale_threshold = threshold_config.min_scale_threshold_normal

        print(
            f'    📐 First slide full-width layout: starting at current_top={current_top:.2f}", max_top={max_top:.2f}", gutter={gutter:.2f}"'
        )

        # Track block index to determine if more blocks exist
        for block_idx, block in enumerate(blocks_to_fit):
            remaining = max_top - current_top
            if remaining <= 0:
                overflow_notes.append(
                    f'No vertical space remaining on slide (current_top={current_top:.2f}", max_top={max_top:.2f}")'
                )
                print(
                    f'    ⚠️  No space remaining: current_top={current_top:.2f}", max_top={max_top:.2f}", breaking'
                )
                break

            # Always use full width for full_width layout - never cap or reduce width
            margin_left = (
                layout_config.first_slide_margin_left
                or layout_config.base_constraints.margin_left
            )
            margin_right = (
                layout_config.first_slide_margin_right
                or layout_config.base_constraints.margin_right
            )
            cell_width = constraints.slide_width - margin_left - margin_right

            # Only cap height at reasonable maximum: 50% of content height or 4 inches, whichever is smaller
            # Width remains full width regardless of height cap
            # Get element dimensions for height capping
            element_dims = get_element_dimensions()
            max_reasonable_height = min(
                constraints.content_height * element_dims.max_height_ratio,
                element_dims.max_height_absolute,
            )

            # Get intrinsic height for reference
            # For tables, pass the actual available width for accurate height calculation
            if block.type == ContentType.TABLE:
                _, intrinsic_h = block.estimate_intrinsic_size(
                    constraints, available_width=cell_width
                )
            else:
                _, intrinsic_h = block.estimate_intrinsic_size(constraints)

            # Apply height cap
            capped_height = min(intrinsic_h, max_reasonable_height)

            # Check if there are more blocks after this one (to determine if gutter is needed)
            is_last_block = block_idx == len(blocks_to_fit) - 1
            has_more_blocks = not is_last_block

            # Calculate space needed: block height + gutter (if more blocks exist)
            space_needed = capped_height
            if has_more_blocks:
                space_needed += gutter

            # Check if block can fit in remaining space (accounting for gutter if needed)
            if space_needed > remaining:
                # Block doesn't fit even with height cap - try scaling down
                # Create a temporary cell with available space to check if scaling can help
                temp_cell_height = remaining
                if has_more_blocks and remaining > gutter:
                    # Reserve space for gutter if more blocks exist
                    temp_cell_height = remaining - gutter
                elif has_more_blocks:
                    # Not enough space even for gutter - can't fit
                    temp_cell_height = 0

                if temp_cell_height > 0:
                    # Try to see if block can fit scaled down
                    temp_cell = CellFrame(
                        left=margin_left,
                        top=current_top,
                        width=cell_width,
                        height=temp_cell_height,
                    )
                    scale_factor = block.estimate_scale_factor(temp_cell, constraints)

                    # Check if scaled block meets minimum threshold
                    # For full_width, we prioritize scale factor over strict can_fit_in_cell checks
                    if scale_factor >= min_scale_threshold:
                        # Block can fit scaled down - use the scaled height
                        print(
                            f'    🔄 Block {block.id} scaled down to fit: scale={scale_factor:.1%}, height={temp_cell_height:.2f}" (was {capped_height:.2f}")'
                        )
                        capped_height = temp_cell_height
                        space_needed = capped_height
                        if has_more_blocks:
                            space_needed += gutter
                    else:
                        # Block can't fit even scaled down - skip it
                        overflow_notes.append(
                            f'{block.id} skipped (too large: needs {capped_height:.2f}", available {remaining:.2f}", scale={scale_factor:.1%} < {min_scale_threshold:.1%})'
                        )
                        print(
                            f'    ❌ Block {block.id} skipped: too large even scaled (scale={scale_factor:.1%} < {min_scale_threshold:.1%}, needs {capped_height:.2f}", remaining={remaining:.2f}")'
                        )
                        continue
                else:
                    # Not enough space even for minimum height - skip this block
                    if assignments:
                        overflow_notes.append(
                            f'{block.id} moved to next slide (needs {capped_height:.2f}" + {gutter:.2f}" gutter but only {remaining:.2f}" left, current_top={current_top:.2f}")'
                        )
                        print(
                            f'    ⚠️  Block {block.id} overflow: needs {capped_height:.2f}" + {gutter:.2f}" gutter, remaining={remaining:.2f}", current_top={current_top:.2f}", will start at margin_top on next slide'
                        )
                    else:
                        # First block doesn't fit - try scaling or skip
                        temp_cell = CellFrame(
                            left=margin_left,
                            top=current_top,
                            width=cell_width,
                            height=remaining,
                        )
                        scale_factor = block.estimate_scale_factor(
                            temp_cell, constraints
                        )
                        # For full_width, we prioritize scale factor over strict can_fit_in_cell checks
                        if scale_factor >= min_scale_threshold:
                            # Can fit scaled - use remaining space
                            print(
                                f'    🔄 First block {block.id} scaled to fit: scale={scale_factor:.1%}, height={remaining:.2f}"'
                            )
                            capped_height = remaining
                            space_needed = capped_height
                        else:
                            overflow_notes.append(
                                f'{block.id} skipped (first block too large: needs {capped_height:.2f}", available {remaining:.2f}", scale={scale_factor:.1%} < {min_scale_threshold:.1%})'
                            )
                            print(
                                f"    ❌ First block {block.id} skipped: too large even scaled (scale={scale_factor:.1%} < {min_scale_threshold:.1%})"
                            )
                            continue
                    break

            # Use the smaller of capped height or remaining space
            cell_height = min(capped_height, remaining)

            if intrinsic_h > max_reasonable_height:
                # Height was capped
                overflow_notes.append(
                    f'{block.id} height capped from {intrinsic_h:.2f}" to {cell_height:.2f}"'
                )

            cell = CellFrame(
                left=margin_left,
                top=current_top,
                width=cell_width,
                height=cell_height,
            )

            # Check if block can fit in cell (may need scaling)
            scale_factor = block.estimate_scale_factor(cell, constraints)

            # For full_width layout, we're more lenient: if scale factor is acceptable, allow it
            # The can_fit_in_cell check is too strict for scaled-down blocks in full_width
            # We prioritize scale factor over strict dimension checks when scaling is acceptable
            if scale_factor < min_scale_threshold:
                # Block can't fit even when scaled down - skip it
                overflow_notes.append(
                    f"{block.id} skipped (scale too low: {scale_factor:.1%} < {min_scale_threshold:.1%})"
                )
                print(
                    f"    ❌ Block {block.id} skipped: scale too low (scale={scale_factor:.1%} < {min_scale_threshold:.1%})"
                )
                continue

            # Optional: Log if can_fit_in_cell would reject but we're allowing due to acceptable scale
            can_fit = block.can_fit_in_cell(cell, constraints)
            if not can_fit:
                print(
                    f"    ⚠️  Block {block.id} doesn't meet strict fit requirements but scale={scale_factor:.1%} is acceptable, allowing"
                )

            print(
                f'    📍 Placing block {block.id} at top={current_top:.2f}", height={cell_height:.2f}", remaining={remaining:.2f}", scale={scale_factor:.1%}'
            )
            if scale_factor < 1.0:
                print(f"    🔄 Block {block.id} will be scaled to {scale_factor:.1%}")

            cell_idx = len(cell_frames)
            cell_frames.append(cell)
            assignments.append((cell_idx, block.id))
            scale_factors[block.id] = scale_factor
            consumed_ids.append(block.id)

            # Labels are now part of the cell, so use cell_height for spacing
            # Update current_top for next block, but only add gutter if there are more blocks
            if has_more_blocks:
                new_current_top = current_top + cell_height + gutter
                print(
                    f'    📐 Updated current_top to {new_current_top:.2f}" (added {cell_height:.2f}" + {gutter:.2f}" gutter for next block)'
                )
            else:
                new_current_top = current_top + cell_height
                print(
                    f'    📐 Updated current_top to {new_current_top:.2f}" (last block, no gutter)'
                )

            # Validate new_current_top doesn't exceed bounds
            if new_current_top > max_top:
                print(
                    f'    ⚠️  Warning: next position {new_current_top:.2f}" would exceed max_top={max_top:.2f}", capping to max_top'
                )
                new_current_top = max_top
            current_top = new_current_top

        if not assignments:
            return None

        layout = SlideLayout(
            slide_number=slide_number,
            section_id=section.id,
            layout_type=LayoutType.FULL_WIDTH,
            cell_frames=cell_frames,
            assigned_blocks=assignments,
            scale_factors=scale_factors,
            overflow_notes=overflow_notes,
            metadata={
                "blocks_fitted": len(assignments),
                "layout_style": "first_slide_full_width",
            },
        )

        return (layout, consumed_ids)

    def _layout_first_slide_grid(
        self,
        section: Section,
        blocks: List[ContentBlock],
        slide_number: int,
        layout_config: SlideLayoutConfig,
    ) -> Optional[Tuple[SlideLayout, List[str]]]:
        """Submarket-style first slide (2x3 grid)."""
        if (
            not blocks
            or not layout_config.first_slide_rows
            or not layout_config.first_slide_cols
        ):
            return None

        constraints = section.constraints
        margin_left = (
            layout_config.first_slide_margin_left
            or layout_config.base_constraints.margin_left
        )
        margin_right = (
            layout_config.first_slide_margin_right
            or layout_config.base_constraints.margin_right
        )
        margin_bottom = (
            layout_config.first_slide_margin_bottom
            or layout_config.base_constraints.margin_bottom
        )
        content_width = constraints.slide_width - margin_left - margin_right
        content_height = (
            constraints.slide_height
            - margin_bottom
            - layout_config.first_slide_start_top
        )

        if content_width <= 0 or content_height <= 0:
            return None

        gutter_h = (
            layout_config.first_slide_gutter_horizontal
            or layout_config.base_constraints.gutter_horizontal
        )
        gutter_v = (
            layout_config.first_slide_gutter_vertical
            or layout_config.base_constraints.gutter_vertical
        )
        total_h_gutter = gutter_h * (layout_config.first_slide_cols - 1)
        total_v_gutter = gutter_v * (layout_config.first_slide_rows - 1)

        cell_w = (content_width - total_h_gutter) / layout_config.first_slide_cols
        cell_h = (content_height - total_v_gutter) / layout_config.first_slide_rows

        cell_frames: List[CellFrame] = []
        for row in range(layout_config.first_slide_rows):
            for col in range(layout_config.first_slide_cols):
                cell_frames.append(
                    CellFrame(
                        left=margin_left + col * (cell_w + gutter_h),
                        top=layout_config.first_slide_start_top
                        + row * (cell_h + gutter_v),
                        width=cell_w,
                        height=cell_h,
                    )
                )

        blocks_to_fit = self._limit_blocks(
            blocks, layout_config.first_slide_max_elements
        )
        assignments, scale_factors, overflow_notes, consumed_ids = (
            self._assign_blocks_to_cells(
                blocks_to_fit,
                cell_frames,
                constraints,
                allow_aggressive_scaling=True,
                respect_section_boundaries=True,
            )
        )

        if not assignments:
            return None

        layout = SlideLayout(
            slide_number=slide_number,
            section_id=section.id,
            layout_type=LayoutType.GRID_2x2,
            cell_frames=cell_frames,
            assigned_blocks=assignments,
            scale_factors=scale_factors,
            overflow_notes=overflow_notes,
            metadata={
                "blocks_fitted": len(assignments),
                "layout_style": f"first_slide_grid_{layout_config.first_slide_rows}x{layout_config.first_slide_cols}",
            },
        )

        return (layout, consumed_ids)

    def _validate_layout_no_overlaps(
        self, layout: SlideLayout, section: Section, blocks: List[ContentBlock]
    ) -> bool:
        """
        Validate that a layout won't cause overlaps when rendered.

        Checks actual content bounds vs allocated cell frames.
        """
        block_dict = {b.id: b for b in blocks}
        actual_bounds = []

        # Get actual bounds for each block
        for cell_idx, block_id in layout.assigned_blocks:
            block = block_dict.get(block_id)
            if not block:
                continue

            cell = layout.cell_frames[cell_idx]
            bounds = block.get_actual_bounds(cell, section.constraints)
            actual_bounds.append((block_id, bounds))

        # Check for overlaps between actual bounds
        for i, (id1, bounds1) in enumerate(actual_bounds):
            for j, (id2, bounds2) in enumerate(actual_bounds):
                if i >= j:
                    continue

                if bounds1.overlaps_with(bounds2):
                    print(
                        f"       ⚠ Overlap detected: Block {id1} overlaps with Block {id2}"
                    )
                    print(
                        f"          {id1}: left={bounds1.left:.2f}, top={bounds1.top:.2f}, right={bounds1.left + bounds1.width:.2f}, bottom={bounds1.top + bounds1.height:.2f}"
                    )
                    print(
                        f"          {id2}: left={bounds2.left:.2f}, top={bounds2.top:.2f}, right={bounds2.left + bounds2.width:.2f}, bottom={bounds2.top + bounds2.height:.2f}"
                    )
                    layout.overflow_notes.append(
                        f"⚠ Overlap detected between {id1} and {id2}"
                    )
                    return False

        return True

    def _try_cover_layout(
        self, section: Section, blocks: List[ContentBlock], slide_number: int
    ) -> Optional[Tuple[SlideLayout, List[str]]]:
        """Cover slide: Row 1 full width + Row 2 split (summary left + chart right)"""

        # Only use for slides with title or KPIs
        has_title = any(hasattr(b, "font_size") and b.font_size >= 20.0 for b in blocks)
        has_kpis = any("►" in getattr(b, "text", "") for b in blocks)

        if not (has_title or has_kpis):
            return None

        constraints = section.constraints

        # Row 1: Full width (title + KPIs)
        row1_cell = CellFrame(
            left=constraints.margin_left,
            top=constraints.margin_top,
            width=constraints.content_width,
            height=2.2,  # Space for title and KPI row
        )

        # Row 2: Split layout
        row2_top = constraints.margin_top + 2.5
        row2_height = constraints.content_height - 2.5

        summary_cell = CellFrame(
            left=constraints.margin_left,
            top=row2_top,
            width=constraints.content_width * 0.4,  # Left 40%
            height=row2_height,
        )

        chart_cell = CellFrame(
            left=constraints.margin_left + constraints.content_width * 0.4 + 0.5,
            top=row2_top,
            width=constraints.content_width * 0.6 - 0.5,  # Right 60% minus gap
            height=row2_height,
        )

        # Assign blocks
        cell_frames = [row1_cell, summary_cell, chart_cell]
        assignments = []
        scale_factors = {}
        consumed_ids = []

        # Put title and KPIs in row 1
        for block in blocks:
            if (
                hasattr(block, "font_size") and block.font_size >= 16.0
            ) or "►" in getattr(block, "text", ""):
                assignments.append((0, block.id))  # Row 1
                scale_factors[block.id] = 1.0
                consumed_ids.append(block.id)

        # Put summary and chart in row 2
        remaining = [b for b in blocks if b.id not in consumed_ids]
        for i, block in enumerate(remaining[:2]):
            if block.type == ContentType.TEXT:
                assignments.append((1, block.id))  # Summary cell
            elif block.type == ContentType.CHART:
                assignments.append((2, block.id))  # Chart cell
            scale_factors[block.id] = 1.0
            consumed_ids.append(block.id)

        if len(assignments) < 2:
            return None

        layout = SlideLayout(
            slide_number=slide_number,
            section_id=section.id,
            layout_type=LayoutType.HYBRID_GRID,
            cell_frames=cell_frames,
            assigned_blocks=assignments,
            scale_factors=scale_factors,
            overflow_notes=[],
            metadata={"layout": "cover", "blocks_fitted": len(assignments)},
        )

        return (layout, consumed_ids)

    def _try_hybrid_grid(
        self, section: Section, blocks: List[ContentBlock], slide_number: int
    ) -> Optional[Tuple[SlideLayout, List[str]]]:
        # Import here to avoid circular dependencies
        from hello.utils.ppt_helpers_utils.services.template_config import (
            get_layout_threshold_config,
            get_element_dimensions,
        )

        """
        Try flexible 2×2 grid where large items get full rows.
        
        This creates a smart grid that:
        - Fits small items in 2×2 cells
        - Gives large items a full row (1×2)
        - Maximizes content per slide
        - PRESERVES display_order (blocks appear in original sequence)
        """
        if len(blocks) < 1:
            return None

        constraints = self.constraints
        content_w = constraints.content_width
        content_h = constraints.content_height

        # Use gutter from constraints for hybrid grid
        # This accounts for spacing between columns in the grid
        half_width = (content_w - constraints.gutter_horizontal) / 2

        # Analyze which blocks need full width vs can fit in half width
        # KEEP IN DISPLAY ORDER - do NOT sort
        block_info = []
        print(f"    🔍 Hybrid grid: Analyzing {len(blocks)} blocks")
        for block in blocks[:6]:  # Try up to 6 blocks
            intrinsic_w, intrinsic_h = block.estimate_intrinsic_size(constraints)

            # Log data availability for debugging
            if block.type == ContentType.TABLE:
                data_count = (
                    len(block.data) if hasattr(block, "data") and block.data else 0
                )
                print(
                    f'       Block {block.id} (table): {data_count} data rows, intrinsic size: {intrinsic_w:.2f}" × {intrinsic_h:.2f}"'
                )
            elif block.type == ContentType.CHART:
                data_count = (
                    len(block.data) if hasattr(block, "data") and block.data else 0
                )
                print(
                    f'       Block {block.id} (chart): {data_count} data rows, intrinsic size: {intrinsic_w:.2f}" × {intrinsic_h:.2f}"'
                )

            # Check if block has explicit layout requirement in metadata
            explicit_layout = None
            if hasattr(block, "metadata") and isinstance(block.metadata, dict):
                explicit_layout = block.metadata.get("layout")

            # Determine if block needs full width:
            # 1. If explicit layout is "full_width", use that
            # 2. Otherwise, check if block would fit comfortably in half width
            if explicit_layout == "full_width":
                needs_full = True
                print(
                    f"       Block {block.id} ({block.type.value}): explicit layout=full_width → needs_full=True"
                )
            elif explicit_layout == "grid_2x2":
                needs_full = False
                print(
                    f"       Block {block.id} ({block.type.value}): explicit layout=grid_2x2 → needs_full=False"
                )
            else:
                # Use threshold from config
                threshold_config = get_layout_threshold_config(self.property_sub_type)
                threshold = threshold_config.full_width_detection_threshold
                needs_full = intrinsic_w > (half_width * threshold)
                print(
                    f"       Block {block.id} ({block.type.value}): no explicit layout, intrinsic_w={intrinsic_w:.2f} > {half_width * threshold:.2f} → needs_full={needs_full}"
                )

            block_info.append((block, needs_full, intrinsic_w, intrinsic_h))

        # Build flexible layout, preserving display order
        cell_frames = []
        assignments = []
        scale_factors = {}
        # Use margin_top from constraints (accounts for header space)
        current_y = constraints.margin_top
        blocks_consumed = []
        assigned_block_ids = set()
        current_section_id = None  # Track current section for boundary checking

        i = 0
        while (
            i < len(block_info)
            and current_y < constraints.slide_height - constraints.margin_bottom
        ):
            block, needs_full, intrinsic_w, intrinsic_h = block_info[i]

            if block.id in assigned_block_ids:
                i += 1
                continue

            # Check section boundary: if section changes and we're not at start of new row, stop
            if (
                current_section_id is not None
                and block.section_id
                and block.section_id != current_section_id
            ):
                # Section changed - this should start at a new row (leftmost)
                # In hybrid grid, each row naturally starts at leftmost, so we're good
                # But we need to ensure we're starting a new row, not continuing a partial row
                print(
                    f"       Section boundary: Section '{block.section_id}' starts at new row"
                )
                current_section_id = block.section_id

            if needs_full:
                # Give this block the entire row
                print(
                    f"       Placing block {block.id} ({block.type.value}) in FULL WIDTH row at y={current_y:.2f}"
                )

                # For tables, optimize layout selection before assignment
                use_optimized_cell = False
                optimized_cell = None
                if block.type == ContentType.TABLE:
                    optimized_layout, optimized_cell, optimized_scale = (
                        self._optimize_table_layout(
                            table_block=block,
                            constraints=constraints,
                            current_y=current_y,
                            slide_number=slide_number,
                        )
                    )
                    print(
                        f"       📊 Table {block.id} optimized layout: {optimized_layout}, scale: {optimized_scale:.2%}"
                    )

                    # Use optimized cell if it's full_width or provides better fit
                    # Get threshold config to determine acceptable scale
                    threshold_config = get_layout_threshold_config(
                        self.property_sub_type
                    )
                    good_scale_threshold = threshold_config.min_scale_threshold_normal

                    if optimized_layout == "full_width":
                        use_optimized_cell = True
                    elif optimized_scale >= good_scale_threshold:
                        # If grid layout has acceptable scale, consider it but respect needs_full
                        print(
                            f"       ℹ️  Optimizer suggests {optimized_layout} with scale {optimized_scale:.2%}, but using full_width as determined by needs_full"
                        )

                # Check if this is a SUBMARKET table that needs special positioning
                is_submarket_table = (
                    block.type == ContentType.TABLE
                    and hasattr(block, "template_path")
                    and block.template_path
                    and "table_market_submarket" in str(block.template_path).lower()
                )

                if use_optimized_cell and optimized_cell:
                    # Use the optimized cell from layout optimization
                    cell = optimized_cell
                    print(
                        f'       ✅ Using optimized cell: {cell.width:.2f}" × {cell.height:.2f}" at ({cell.left:.2f}", {cell.top:.2f}")'
                    )
                elif is_submarket_table:
                    # SUBMARKET TABLE: Use constraints margins
                    table_margin_left = constraints.margin_left
                    table_margin_right = constraints.margin_right
                    table_content_width = (
                        constraints.slide_width - table_margin_left - table_margin_right
                    )
                    table_margin_top = constraints.margin_top

                    print("🔍 DEBUG BIG SUBMARKET TABLE in hybrid_grid:")
                    print(
                        f'🔍   left={table_margin_left}", top={table_margin_top}", width={table_content_width}"'
                    )

                    # Calculate available space from the forced-up position
                    available_height = (
                        constraints.slide_height
                        - constraints.margin_bottom
                        - table_margin_top
                    )
                    row_height = min(intrinsic_h, available_height)

                    cell = CellFrame(
                        left=table_margin_left,
                        top=table_margin_top,
                        width=table_content_width,
                        height=row_height,
                    )
                else:
                    # REGULAR FULL-WIDTH BLOCK: Use normal positioning
                    # Calculate available height from current_y to bottom of slide
                    available_height = (
                        constraints.slide_height - constraints.margin_bottom - current_y
                    )

                    # Get element dimensions for height capping (reuse from _try_full_width_stack)
                    element_dims = get_element_dimensions()
                    max_reasonable_height = min(
                        constraints.content_height * element_dims.max_height_ratio,
                        element_dims.max_height_absolute,
                    )

                    # Cap intrinsic height to reasonable maximum
                    capped_intrinsic_h = min(intrinsic_h, max_reasonable_height)

                    # Use the smaller of capped height or available space
                    row_height = min(capped_intrinsic_h, available_height)

                    # Check if it fits in remaining space
                    if (
                        row_height <= 0
                        or current_y + row_height
                        > constraints.slide_height - constraints.margin_bottom
                    ):
                        print(
                            f'       ❌ Block {block.id} cannot fit in full-width cell (needs {row_height:.2f}", available: {available_height:.2f}"), stopping'
                        )
                        break  # No more room

                    # Use constraints margins for regular tables
                    table_margin_left = constraints.margin_left
                    table_margin_right = constraints.margin_right
                    table_content_width = (
                        constraints.slide_width - table_margin_left - table_margin_right
                    )

                    cell = CellFrame(
                        left=table_margin_left,
                        top=current_y,
                        width=table_content_width,
                        height=row_height,
                    )

                if block.can_fit_in_cell(cell, constraints):
                    cell_frames.append(cell)
                    assignments.append((len(cell_frames) - 1, block.id))
                    scale_factors[block.id] = block.estimate_scale_factor(
                        cell, constraints
                    )
                    blocks_consumed.append(block)
                    assigned_block_ids.add(block.id)

                    print(
                        f'       ✅ Block {block.id} placed: left={cell.left:.2f}", top={cell.top:.2f}", width={cell.width:.2f}", height={cell.height:.2f}"'
                    )

                    # Track current section for boundary checking
                    if block.section_id:
                        current_section_id = block.section_id

                    # Update current_y for next block
                    # Use configurable gutter for full-width layouts
                    layout_config = get_slide_layout_config(self.property_sub_type)
                    is_first_slide = slide_number == 1
                    full_width_gutter = layout_config.get_full_width_gutter(
                        is_first_slide
                    )

                    if is_submarket_table:
                        current_y = table_margin_top + row_height + full_width_gutter
                    else:
                        current_y += row_height + full_width_gutter

                    i += 1
                else:
                    print(
                        f"       ❌ Block {block.id} cannot fit in full-width cell, stopping"
                    )
                    break  # Can't fit, stop here
            else:
                # Try to fit blocks side by side in this row using skip-and-backfill
                # Look ahead for up to 2 blocks that can fit side-by-side
                row_blocks = []
                max_row_height = intrinsic_h
                row_section_id = None  # Track section for this row

                # Try to find up to 2 blocks that can fit in this row (in display order)
                # SECTION BOUNDARY: Only include blocks from the same section in a row
                for j in range(
                    i, min(len(block_info), i + 4)
                ):  # Look ahead max 4 blocks
                    if len(row_blocks) >= 2:
                        break

                    b, needs_f, w, h = block_info[j]
                    if b.id in assigned_block_ids:
                        continue

                    # Only add blocks that don't need full width
                    if not needs_f:
                        # SECTION BOUNDARY CHECK: Only add blocks from the same section
                        # This ensures new sections always start from the leftmost cell
                        if row_section_id is None:
                            row_section_id = b.section_id
                        elif b.section_id and b.section_id != row_section_id:
                            # Different section - don't add to this row
                            # New section will start from left in next row
                            print(
                                f"       Section boundary in row_blocks: Block {b.id} (section '{b.section_id}') "
                                f"skipped to ensure new section starts from left"
                            )
                            break

                        row_blocks.append((j, b, w, h))
                        max_row_height = max(max_row_height, h)

                if not row_blocks:
                    i += 1
                    continue

                # Cap row height to allow content to fit
                row_height = min(max_row_height, content_h * 0.42)

                # Check if row fits
                if (
                    current_y + row_height
                    > constraints.slide_height - constraints.margin_bottom
                ):
                    break

                # Create cells for this row and try to assign blocks using skip-and-backfill
                row_cell_frames = []
                row_start_cell_index = len(
                    cell_frames
                )  # Track where this row starts in cell_frames
                for col_idx in range(min(2, len(row_blocks))):
                    cell = CellFrame(
                        left=constraints.margin_left
                        + col_idx * (half_width + constraints.gutter_horizontal),
                        top=current_y,
                        width=half_width,
                        height=row_height,
                    )
                    row_cell_frames.append(cell)

                # Assign blocks to row cells
                row_assigned = False
                row_section_boundary_hit = (
                    False  # Track if we hit section boundary mid-row
                )
                blocks_assigned_this_row = []  # Track which blocks were assigned in this row

                for cell_idx, cell in enumerate(row_cell_frames):
                    # Try each available block for this cell (skip-and-backfill)
                    block_fitted = False
                    for idx, blk, w, h in row_blocks:
                        if blk.id in assigned_block_ids:
                            continue

                        # Check section boundary in the middle of a row
                        if (
                            current_section_id is not None
                            and blk.section_id
                            and blk.section_id != current_section_id
                        ):
                            # Section changed mid-row
                            if cell_idx > 0:  # Not the first cell in row
                                # Stop here - new section needs to start at leftmost (first cell of new row)
                                print(
                                    f"       Section boundary in row: Section '{blk.section_id}' needs new row"
                                )
                                row_section_boundary_hit = True
                                break

                        # Check if block can fit using scale factor (reuse logic from _assign_blocks_to_cells)
                        from hello.utils.ppt_helpers_utils.services.template_config import (
                            get_layout_threshold_config,
                        )

                        threshold_config = get_layout_threshold_config(
                            self.property_sub_type
                        )
                        min_scale_threshold = (
                            threshold_config.min_scale_threshold_aggressive
                        )  # 40% minimum

                        scale = blk.estimate_scale_factor(cell, constraints)
                        can_fit = scale >= min_scale_threshold

                        if not can_fit:
                            print(
                                f'       ⚠️  Block {blk.id} ({blk.type.value}) cannot fit in cell {cell_idx}: scale={scale:.2%} < {min_scale_threshold:.2%}, cell={cell.width:.2f}"×{cell.height:.2f}", intrinsic={w:.2f}"×{h:.2f}"'
                            )

                        if can_fit:
                            cell_frames.append(cell)
                            assignments.append((len(cell_frames) - 1, blk.id))
                            scale_factors[blk.id] = blk.estimate_scale_factor(
                                cell, constraints
                            )
                            blocks_consumed.append(blk)
                            assigned_block_ids.add(blk.id)
                            blocks_assigned_this_row.append(blk.id)
                            row_assigned = True
                            block_fitted = True

                            # Track current section for boundary checking
                            if blk.section_id:
                                current_section_id = blk.section_id

                            print(
                                f"       ✅ Block {blk.id} ({blk.type.value}) placed in cell {cell_idx} of row"
                            )
                            break  # Move to next cell

                    if row_section_boundary_hit:
                        break  # Stop filling this row

                    if not block_fitted and cell_idx == 0:
                        # No block could fit in the first cell - this row won't work
                        print(
                            "       ⚠️  No blocks could fit in first cell of row, skipping row"
                        )
                        break

                # Check if section boundary hit - stop main loop too
                if row_section_boundary_hit:
                    print(
                        "       ⚠️  Section boundary hit mid-row - stopping hybrid grid placement"
                    )
                    break  # Exit main while loop to let new section overflow

                if row_assigned:
                    # Calculate actual bottom of all cells in this row (not just row_height)
                    # This prevents overlaps when blocks are scaled or extend beyond row_height
                    # Use the actual cells from cell_frames (the ones that were appended), not row_cell_frames
                    max_row_bottom = current_y  # Start with row top
                    for cell_idx in range(len(row_cell_frames)):
                        cell_index_in_full_list = row_start_cell_index + cell_idx
                        # Check if this cell was assigned to a block
                        cell_assigned = any(
                            (cell_index_in_full_list, blk_id) in assignments
                            for blk_id in blocks_assigned_this_row
                        )
                        if cell_assigned and cell_index_in_full_list < len(cell_frames):
                            # Use actual cell from cell_frames (the one that was appended)
                            actual_cell = cell_frames[cell_index_in_full_list]
                            # Use actual cell bottom (top + height)
                            cell_bottom = actual_cell.top + actual_cell.height
                            max_row_bottom = max(max_row_bottom, cell_bottom)

                    # Update current_y to the maximum bottom of all cells in the row, plus gutter
                    current_y = max_row_bottom + constraints.gutter_vertical
                    print(
                        f'       📐 Row completed at y={current_y:.2f}" (max cell bottom: {max_row_bottom:.2f}"), assigned blocks: {blocks_assigned_this_row}'
                    )
                else:
                    # No blocks were assigned in this row - try next block as full-width
                    print(
                        "       ⚠️  No blocks assigned in row, will try next block as full-width"
                    )

                # Only move past blocks that were actually assigned
                # If a block wasn't assigned, we'll try it again in the next iteration
                # But we need to avoid infinite loops, so track which blocks we've tried
                blocks_tried_this_iteration = [b.id for _, b, _, _ in row_blocks]
                if not row_assigned and blocks_tried_this_iteration:
                    # None of the blocks in row_blocks could be placed
                    # Move past the first unassigned block to avoid infinite loop
                    for idx, blk, _, _ in row_blocks:
                        if blk.id not in assigned_block_ids:
                            i = idx + 1  # Move past this block
                            print(
                                f"       ⚠️  Block {blk.id} could not be placed, moving past it"
                            )
                            break
                else:
                    # Move past the blocks we considered (they were either assigned or we'll try them again)
                    i += 1

        if not assignments:
            return None

        layout = SlideLayout(
            slide_number=slide_number,
            section_id=section.id,
            layout_type=LayoutType.HYBRID_GRID,
            cell_frames=cell_frames,
            assigned_blocks=assignments,
            scale_factors=scale_factors,
            overflow_notes=[],
            metadata={"blocks_fitted": len(assignments), "rows": "hybrid"},
        )

        # Validate no overlaps
        # Note: get_actual_bounds now returns cell bounds for charts/text (scaled to fit)
        # and cell bounds for tables (constrained to cell), so overlap detection should be accurate
        if not self._validate_layout_no_overlaps(layout, section, blocks):
            print("    ❌ Hybrid grid failed: overlap validation failed")
            return None

        # Return consumed blocks in display order
        consumed_ids = [b.id for b in blocks if b.id in assigned_block_ids]
        print(f"    ✅ Hybrid grid succeeded with {len(consumed_ids)} blocks")
        return (layout, consumed_ids)

    def _try_grid_2x2(
        self, section: Section, blocks: List[ContentBlock], slide_number: int
    ) -> Optional[Tuple[SlideLayout, List[str]]]:
        """
        Try to fit up to 4 blocks in 2×2 grid using skip-and-backfill logic.

        This method now uses intelligent placement:
        - For each cell, tries all unassigned blocks in display order
        - Assigns the first block that fits
        - Preserves display order in final output
        - Maximizes grid utilization
        """
        if len(blocks) < 2:  # Need at least 2 blocks for grid
            return None

        # Calculate cell frames for 2×2 grid
        cell_frames = self._calculate_grid_cells(2, 2, section.constraints)

        print("    📊 Grid 2x2: Trying to fit blocks with skip-and-backfill")

        # Use skip-and-backfill helper to assign blocks to cells
        # Consider more blocks than cells to enable backfill
        candidates = blocks[
            : min(8, len(blocks))
        ]  # Consider up to 8 blocks for 4 cells

        assignments, scale_factors, overflow_notes, consumed_ids = (
            self._assign_blocks_to_cells(
                blocks=candidates,
                cell_frames=cell_frames,
                constraints=section.constraints,
                allow_aggressive_scaling=True,  # Allow scaling down to 40%
                respect_section_boundaries=True,  # Enforce section leftmost cell constraint
            )
        )

        # Debug: show what got assigned
        print(f"    📊 Grid 2x2: Assigned {len(assignments)}/4 blocks to cells")

        # Require at least one block to succeed
        if len(assignments) < 1:
            print("    ⚠ Grid 2x2 failed: no blocks fit")
            return None

        layout = SlideLayout(
            slide_number=slide_number,
            section_id=section.id,
            layout_type=LayoutType.GRID_2x2,
            cell_frames=cell_frames,
            assigned_blocks=assignments,
            scale_factors=scale_factors,
            overflow_notes=overflow_notes,
            metadata={"blocks_fitted": len(assignments)},
        )

        # Validate no overlaps
        if not self._validate_layout_no_overlaps(layout, section, blocks):
            print("    ⚠ Grid 2x2 failed: overlap validation failed")
            return None

        print(f"    ✅ Grid 2x2 succeeded with {len(consumed_ids)} blocks")
        return (layout, consumed_ids)

    def _try_full_width_stack(
        self, section: Section, blocks: List[ContentBlock], slide_number: int
    ) -> Optional[Tuple[SlideLayout, List[str]]]:
        """Full-width layout: stack blocks vertically using the entire slide width."""
        if not blocks:
            return None

        constraints = section.constraints
        margin_left = constraints.margin_left
        margin_right = constraints.margin_right
        margin_top = constraints.margin_top
        margin_bottom = constraints.margin_bottom

        # Get configurable gutter based on property_sub_type and slide number
        layout_config = get_slide_layout_config(self.property_sub_type)
        is_first_slide = slide_number == 1
        gutter = layout_config.get_full_width_gutter(is_first_slide)

        cell_frames: List[CellFrame] = []
        assignments: List[Tuple[int, str]] = []
        scale_factors: Dict[str, float] = {}
        consumed_ids: List[str] = []
        overflow_notes: List[str] = []

        # CRITICAL: Always start at margin_top for each slide to ensure proper positioning
        # This ensures overflow blocks from previous slide start at the top of new slide
        current_top = margin_top

        # Validation: Ensure current_top is properly initialized
        max_top = constraints.slide_height - margin_bottom
        if current_top < margin_top or current_top > max_top:
            print(
                f'    ⚠️  Warning: current_top={current_top:.2f}" is out of bounds, resetting to margin_top={margin_top:.2f}"'
            )
            current_top = margin_top

        print(
            f'    📐 Full-width layout slide {slide_number}: starting at current_top={current_top:.2f}", max_top={max_top:.2f}", gutter={gutter:.2f}"'
        )

        # Get threshold config for scale factor checks
        threshold_config = get_layout_threshold_config(self.property_sub_type)
        min_scale_threshold = threshold_config.min_scale_threshold_normal

        # Track block index to determine if more blocks exist
        for block_idx, content_block in enumerate(blocks):
            # Calculate remaining space from current position to bottom of slide
            remaining = max_top - current_top
            if remaining <= 0:
                overflow_notes.append(
                    f'No vertical space remaining on slide (current_top={current_top:.2f}", max_top={max_top:.2f}")'
                )
                print(
                    f'    ⚠️  No space remaining: current_top={current_top:.2f}", max_top={max_top:.2f}", breaking'
                )
                break

            # Always use full width for full_width layout - never cap or reduce width
            cell_width = constraints.slide_width - margin_left - margin_right

            # Height capping policy:
            # - For tables in full_width with overflow enabled: NEVER cap here.
            #   If a table doesn't fit, it should consume the rest of the slide and continue
            #   on the next slide(s). This preserves report ordering (a table's continuation
            #   should appear before later elements).
            # - For other elements: keep the existing reasonable cap behavior.
            # Get element dimensions for height capping / overflow policy
            element_dims = get_element_dimensions()
            max_reasonable_height = min(
                constraints.content_height * element_dims.max_height_ratio,
                element_dims.max_height_absolute,
            )

            # Get intrinsic height for reference
            _, intrinsic_h = content_block.estimate_intrinsic_size(constraints)

            # Apply height cap (non-tables only; tables may overflow)
            is_table = getattr(content_block, "type", None) == ContentType.TABLE
            if is_table and element_dims.allow_full_width_overflow:
                capped_height = intrinsic_h  # true intrinsic height; renderer will split rows if needed
            else:
                capped_height = min(intrinsic_h, max_reasonable_height)

            # Check if there are more blocks after this one (to determine if gutter is needed)
            is_last_block = block_idx == len(blocks) - 1
            has_more_blocks = not is_last_block

            # Calculate space needed: block height + gutter (if more blocks exist)
            space_needed = capped_height
            if has_more_blocks:
                space_needed += gutter

            # Check if block can fit in remaining space (accounting for gutter if needed)
            if space_needed > remaining:
                # Special handling: tables should not be scaled/capped to "fit" alongside subsequent blocks.
                # If a table would overflow, force a page break to preserve ordering.
                if is_table and element_dims.allow_full_width_overflow:
                    if assignments:
                        overflow_notes.append(
                            f'{content_block.id} moved to next slide (table overflow; preserve ordering)'
                        )
                        print(
                            f'    📄 Table {content_block.id} would overflow after prior blocks; moving to next slide to preserve ordering'
                        )
                        break
                    # First block on slide: allocate remaining space and stop placing more blocks.
                    print(
                        f'    📄 Table {content_block.id} overflows slide; placing first page in remaining {remaining:.2f}" and forcing page break'
                    )
                    capped_height = remaining
                    cell_height = remaining
                    cell = CellFrame(
                        left=margin_left,
                        top=current_top,
                        width=cell_width,
                        height=cell_height,
                    )
                    cell_index = len(cell_frames)
                    cell_frames.append(cell)
                    assignments.append((cell_index, content_block.id))
                    scale_factors[content_block.id] = 1.0
                    consumed_ids.append(content_block.id)
                    # End slide after placing an overflowing table.
                    break

                # Block doesn't fit even with height cap - try scaling down
                # Create a temporary cell with available space to check if scaling can help
                temp_cell_height = remaining
                if has_more_blocks and remaining > gutter:
                    # Reserve space for gutter if more blocks exist
                    temp_cell_height = remaining - gutter
                elif has_more_blocks:
                    # Not enough space even for gutter - can't fit
                    temp_cell_height = 0

                if temp_cell_height > 0:
                    # Try to see if block can fit scaled down
                    temp_cell = CellFrame(
                        left=margin_left,
                        top=current_top,
                        width=cell_width,
                        height=temp_cell_height,
                    )
                    scale_factor = content_block.estimate_scale_factor(
                        temp_cell, constraints
                    )

                    # Check if scaled block meets minimum threshold
                    # For full_width, we prioritize scale factor over strict can_fit_in_cell checks
                    if scale_factor >= min_scale_threshold:
                        # Block can fit scaled down - use the scaled height
                        # Note: After this update, space_needed will equal remaining (temp_cell_height + gutter = remaining),
                        # so we'll fall through to cell creation below
                        print(
                            f'    🔄 Block {content_block.id} scaled down to fit: scale={scale_factor:.1%}, height={temp_cell_height:.2f}" (was {capped_height:.2f}")'
                        )
                        capped_height = temp_cell_height
                        space_needed = capped_height
                        if has_more_blocks:
                            space_needed += gutter
                        # space_needed now equals remaining, so we'll proceed to cell creation
                    else:
                        # Block can't fit even scaled down - skip it
                        overflow_notes.append(
                            f'{content_block.id} skipped (too large: needs {capped_height:.2f}", available {remaining:.2f}", scale={scale_factor:.1%} < {min_scale_threshold:.1%})'
                        )
                        print(
                            f'    ❌ Block {content_block.id} skipped: too large even scaled (scale={scale_factor:.1%} < {min_scale_threshold:.1%}, needs {capped_height:.2f}", remaining={remaining:.2f}")'
                        )
                        # Continue to next block instead of breaking
                        continue
                else:
                    # Not enough space even for minimum height - skip this block
                    if assignments:
                        overflow_notes.append(
                            f'{content_block.id} moved to next slide (needs {capped_height:.2f}" + {gutter:.2f}" gutter but only {remaining:.2f}" left, current_top={current_top:.2f}")'
                        )
                        print(
                            f'    ⚠️  Block {content_block.id} overflow: needs {capped_height:.2f}" + {gutter:.2f}" gutter, remaining={remaining:.2f}", current_top={current_top:.2f}", will start at margin_top on next slide'
                        )
                    else:
                        # First block doesn't fit - try scaling or skip
                        temp_cell = CellFrame(
                            left=margin_left,
                            top=current_top,
                            width=cell_width,
                            height=remaining,
                        )
                        scale_factor = content_block.estimate_scale_factor(
                            temp_cell, constraints
                        )
                        # For full_width, we prioritize scale factor over strict can_fit_in_cell checks
                        if scale_factor >= min_scale_threshold:
                            # Can fit scaled - use remaining space
                            # Note: After this update, space_needed equals remaining, so we'll fall through to cell creation
                            print(
                                f'    🔄 First block {content_block.id} scaled to fit: scale={scale_factor:.1%}, height={remaining:.2f}"'
                            )
                            capped_height = remaining
                            space_needed = capped_height
                            # space_needed now equals remaining, so we'll proceed to cell creation
                        else:
                            overflow_notes.append(
                                f'{content_block.id} skipped (first block too large: needs {capped_height:.2f}", available {remaining:.2f}", scale={scale_factor:.1%} < {min_scale_threshold:.1%})'
                            )
                            print(
                                f"    ❌ First block {content_block.id} skipped: too large even scaled (scale={scale_factor:.1%} < {min_scale_threshold:.1%})"
                            )
                            continue
                    break

            # Use the smaller of capped height or remaining space
            cell_height = min(capped_height, remaining)

            if intrinsic_h > max_reasonable_height:
                # Height was capped
                overflow_notes.append(
                    f'{content_block.id} height capped from {intrinsic_h:.2f}" to {cell_height:.2f}"'
                )

            # Validate current_top is within bounds before creating cell
            if current_top < margin_top:
                print(
                    f'    ⚠️  Warning: current_top={current_top:.2f}" < margin_top={margin_top:.2f}", resetting'
                )
                current_top = margin_top
            elif current_top > max_top:
                print(
                    f'    ⚠️  Warning: current_top={current_top:.2f}" > max_top={max_top:.2f}", cannot place block'
                )
                overflow_notes.append(
                    f'{content_block.id} cannot be placed (current_top={current_top:.2f}" exceeds max_top={max_top:.2f}")'
                )
                break

            cell = CellFrame(
                left=margin_left, top=current_top, width=cell_width, height=cell_height
            )

            # Check if block can fit in cell (may need scaling)
            scale_factor = content_block.estimate_scale_factor(cell, constraints)

            # For full_width layout, we're more lenient: if scale factor is acceptable, allow it
            # The can_fit_in_cell check is too strict for scaled-down blocks in full_width
            # We prioritize scale factor over strict dimension checks when scaling is acceptable
            if scale_factor < min_scale_threshold:
                # Block can't fit even when scaled down - skip it
                overflow_notes.append(
                    f"{content_block.id} skipped (scale too low: {scale_factor:.1%} < {min_scale_threshold:.1%})"
                )
                print(
                    f"    ❌ Block {content_block.id} skipped: scale too low (scale={scale_factor:.1%} < {min_scale_threshold:.1%})"
                )
                continue

            # Optional: Log if can_fit_in_cell would reject but we're allowing due to acceptable scale
            can_fit = content_block.can_fit_in_cell(cell, constraints)
            if not can_fit:
                print(
                    f"    ⚠️  Block {content_block.id} doesn't meet strict fit requirements but scale={scale_factor:.1%} is acceptable, allowing"
                )

            print(
                f'    📍 Placing block {content_block.id} at top={current_top:.2f}", height={cell_height:.2f}", remaining={remaining:.2f}", scale={scale_factor:.1%}'
            )
            if scale_factor < 1.0:
                print(
                    f"    🔄 Block {content_block.id} will be scaled to {scale_factor:.1%}"
                )

            cell_index = len(cell_frames)
            cell_frames.append(cell)
            assignments.append((cell_index, content_block.id))
            scale_factors[content_block.id] = scale_factor
            consumed_ids.append(content_block.id)

            # Labels are now part of the cell, so use cell_height for spacing
            # Update current_top for next block, but only add gutter if there are more blocks
            if has_more_blocks:
                new_current_top = current_top + cell_height + gutter
                print(
                    f'    📐 Updated current_top to {new_current_top:.2f}" (added {cell_height:.2f}" + {gutter:.2f}" gutter for next block)'
                )
            else:
                new_current_top = current_top + cell_height
                print(
                    f'    📐 Updated current_top to {new_current_top:.2f}" (last block, no gutter)'
                )

            # Validate new_current_top doesn't exceed bounds
            if new_current_top > max_top:
                print(
                    f'    ⚠️  Warning: next position {new_current_top:.2f}" would exceed max_top={max_top:.2f}", capping to max_top'
                )
                new_current_top = max_top
            current_top = new_current_top

        if not assignments:
            return None

        layout = SlideLayout(
            slide_number=slide_number,
            section_id=section.id,
            layout_type=LayoutType.FULL_WIDTH,
            cell_frames=cell_frames,
            assigned_blocks=assignments,
            scale_factors=scale_factors,
            overflow_notes=overflow_notes,
            metadata={"blocks_fitted": len(assignments), "layout_style": "full_width"},
        )

        return (layout, consumed_ids)

    def _calculate_grid_cells(
        self, rows: int, cols: int, constraints: SlideConstraints
    ) -> List[CellFrame]:
        """Calculate cell frames for a regular grid using the PERFECT settings from json_processor."""
        slide_width = constraints.slide_width
        slide_height = constraints.slide_height
        margin_left = constraints.margin_left
        margin_right = constraints.margin_right
        margin_top = constraints.margin_top
        margin_bottom = constraints.margin_bottom
        h_gutter = constraints.gutter_horizontal
        v_gutter = constraints.gutter_vertical

        content_w = slide_width - margin_left - margin_right
        content_h = slide_height - margin_top - margin_bottom

        total_h_gutter = (cols - 1) * h_gutter
        total_v_gutter = (rows - 1) * v_gutter

        # For 2x2 grid (standard content layout), use specific chart dimensions
        if rows == 2 and cols == 2:
            max_cell_w = (content_w - total_h_gutter) / cols if cols else content_w
            max_cell_h = (content_h - total_v_gutter) / rows if rows else content_h
            cell_w = min(constraints.second_slide_chart_width, max_cell_w)
            cell_h = min(constraints.second_slide_chart_height, max_cell_h)
        else:
            # For other grid configurations, calculate proportionally
            cell_w = (content_w - total_h_gutter) / cols
            cell_h = (content_h - total_v_gutter) / rows

        cells = []
        for row in range(rows):
            for col in range(cols):
                left = margin_left + col * (cell_w + h_gutter)
                top = margin_top + row * (cell_h + v_gutter)

                cells.append(CellFrame(left=left, top=top, width=cell_w, height=cell_h))

        return cells


# ============================================================================
# UTILITIES
# ============================================================================


def export_layouts_to_json(layouts: List[SlideLayout], output_path: str):
    """Export layout specifications to JSON file."""
    data = {
        "version": "1.0",
        "total_slides": len(layouts),
        "layouts": [layout.to_dict() for layout in layouts],
    }

    with open(output_path, "w") as f:
        json.dump(data, f, indent=2)


def print_layout_summary(layouts: List[SlideLayout]):
    """Print human-readable summary of layouts."""
    print(f"\n{'=' * 60}")
    print(f"LAYOUT SUMMARY: {len(layouts)} slides")
    print(f"{'=' * 60}\n")

    for layout in layouts:
        print(f"Slide {layout.slide_number}: {layout.layout_type.value}")
        print(f"  Blocks: {len(layout.assigned_blocks)}")
        for cell_idx, block_id in layout.assigned_blocks:
            scale = layout.scale_factors.get(block_id, 1.0)
            print(f"    Cell {cell_idx}: {block_id} (scale: {scale:.1%})")

        if layout.overflow_notes:
            print("  Overflow notes:")
            for note in layout.overflow_notes:
                print(f"    - {note}")
        print()

"""
Data Populator for Charts and Tables

This module handles populating PowerPoint charts and tables with data from JSON.
It extracts embedded Excel data from charts, modifies it, and updates table cells.
"""

import math
import os
import tempfile
import zipfile
import shutil
from typing import Any, Dict, List, Optional, Sequence, Tuple

import lxml.etree as etree

from hello.utils.utils import (
    is_total_label,
    total_display_text,
)
from hello.utils.ppt_helpers_utils.ppt_helpers.chart_formatting import (
    apply_axis_format_codes_to_chart_space,
)


def write_xml_with_office_declaration(tree_or_root, file_path: str) -> None:
    """
    Write XML with double-quote XML declaration that Office expects.

    lxml's tree.write(xml_declaration=True) produces single quotes:
        <?xml version='1.0' encoding='UTF-8' standalone='yes'?>

    But Office expects double quotes:
        <?xml version="1.0" encoding="UTF-8" standalone="yes"?>

    This helper serializes without declaration and prepends the correct one.
    """
    # Get the root element (handle both ElementTree and Element)
    if hasattr(tree_or_root, "getroot"):
        root = tree_or_root.getroot()
    else:
        root = tree_or_root

    # Serialize without declaration
    xml_bytes = etree.tostring(root, encoding="UTF-8")

    # Prepend correct double-quote declaration
    xml_decl = b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'

    with open(file_path, "wb") as f:
        f.write(xml_decl + xml_bytes)


def create_pptx_from_dir(extract_dir: str, output_path: str) -> None:
    """
    Repackage an extracted PPTX directory back into a .pptx file.

    Prefer system `zip` for Office compatibility; fall back to Python zipfile if unavailable.
    """
    zip_cli = shutil.which("zip")

    # Prefer OS zip for better interoperability with Office.
    if zip_cli:
        import subprocess

        abs_output_path = os.path.abspath(output_path)
        if os.path.exists(abs_output_path):
            os.remove(abs_output_path)

        # -r: recurse, -X: strip extra file attrs, -q: quiet
        # -n .xlsx: store embedded workbooks without compression (they are already zip files)
        # -D: do not add directory entries (Office files typically don't have them)
        proc = subprocess.run(
            [zip_cli, "-q", "-r", "-X", "-D", "-n", ".xlsx", abs_output_path, "."],
            cwd=extract_dir,
            capture_output=True,
            text=True,
        )
        if proc.returncode == 0 and os.path.exists(abs_output_path):
            return

        # rc=15 often means "nothing to do" - try without -n flag as fallback
        if proc.returncode == 15:
            proc2 = subprocess.run(
                [zip_cli, "-q", "-r", "-X", "-D", abs_output_path, "."],
                cwd=extract_dir,
                capture_output=True,
                text=True,
            )
            if proc2.returncode == 0 and os.path.exists(abs_output_path):
                return

        stderr = (proc.stderr or "").strip()
        print(
            f"  ⚠️  zip CLI repack failed (rc={proc.returncode}); falling back to python zipfile. {stderr}"
        )

    # Fallback: Python zipfile (best-effort)
    with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zipf:
        for root_dir, _dirs, files in os.walk(extract_dir):
            for file in files:
                file_path = os.path.join(root_dir, file)
                arcname = os.path.relpath(file_path, extract_dir)
                zipf.write(file_path, arcname)


class ChartDataPopulator:
    """
    Populates chart data by modifying embedded Excel workbooks
    """

    def __init__(self):
        self.namespaces = {
            "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }
        # Track pending Excel updates to be applied via python-pptx
        self._pending_excel_updates: List[Dict[str, Any]] = []
        self._pending_excel_update: Optional[Dict[str, Any]] = None
        # Store legend dimensions for plot area calculation
        self._combo_legend_height: float = 0.08
        self._combo_legend_y: float = 0.92

    @staticmethod
    def _write_xml_with_office_declaration(tree_or_root, file_path: str) -> None:
        # Kept as a class method for backward-compat; shared implementation is module-level.
        write_xml_with_office_declaration(tree_or_root, file_path)

    def _rewrite_zip_members_in_place(
        self, zip_path: str, replacements: Dict[str, bytes]
    ) -> None:
        """
        Rewrite selected members in a ZIP file in-place (atomic swap) using pure Python.

        This is a fallback for environments where the system `zip` CLI isn't available.
        We preserve existing ZipInfo metadata where possible.

        Note: Python's zipfile ignores flag_bits for deflated entries when writing,
        so we patch them afterward.
        """
        tmp_path = f"{zip_path}.tmp"
        modified_members = []
        with zipfile.ZipFile(zip_path, "r") as zin:
            with zipfile.ZipFile(tmp_path, "w") as zout:
                seen = set()
                for info in zin.infolist():
                    name = info.filename
                    is_modified = name in replacements
                    data = replacements[name] if is_modified else zin.read(name)
                    # Preserve metadata where possible
                    new_info = zipfile.ZipInfo(
                        filename=info.filename, date_time=info.date_time
                    )
                    new_info.compress_type = info.compress_type
                    new_info.create_system = info.create_system
                    new_info.create_version = info.create_version
                    new_info.extract_version = info.extract_version
                    new_info.flag_bits = info.flag_bits
                    new_info.external_attr = info.external_attr
                    new_info.comment = info.comment
                    new_info.extra = info.extra
                    new_info.internal_attr = info.internal_attr
                    zout.writestr(new_info, data)
                    seen.add(name)
                    if is_modified and info.compress_type == zipfile.ZIP_DEFLATED:
                        modified_members.append(name)

                # Add any new files not present in original archive
                for name, data in replacements.items():
                    if name in seen:
                        continue
                    zout.writestr(name, data)

        shutil.move(tmp_path, zip_path)

        # Python's zipfile ignores flag_bits for deflated entries. Patch them to 6
        # for Office compatibility.
        if modified_members:
            self._patch_zip_flag_bits(zip_path, modified_members, flag_bits=6)

    def _patch_zip_flag_bits(
        self, zip_path: str, members: List[str], flag_bits: int = 6
    ) -> None:
        """
        Patch flag_bits (general purpose bit flag) for specific members in a ZIP file.

        The system `zip` CLI resets flag_bits to 0 when updating members. PowerPoint/Excel
        expect deflated members in embedded xlsx files to have flag_bits=6 (compression
        options). This method patches both the local file header and central directory
        entry for each specified member.

        Args:
            zip_path: Path to the ZIP file
            members: List of member names to patch
            flag_bits: The flag_bits value to set (default 6 for Office compatibility)
        """
        import struct

        members_set = set(members)

        with open(zip_path, "r+b") as f:
            # First, find and patch local file headers
            # Local file header signature: 0x04034b50
            local_sig = b"PK\x03\x04"
            f.seek(0)
            content = f.read()

            # Patch local file headers
            pos = 0
            while True:
                pos = content.find(local_sig, pos)
                if pos == -1:
                    break
                # Local file header structure:
                # 0-3: signature (4 bytes)
                # 4-5: version needed (2 bytes)
                # 6-7: flag_bits (2 bytes)
                # 8-9: compression method (2 bytes)
                # ...
                # 26-27: file name length (2 bytes)
                # 28-29: extra field length (2 bytes)
                # 30+: file name
                try:
                    name_len = struct.unpack_from("<H", content, pos + 26)[0]
                    name = content[pos + 30 : pos + 30 + name_len].decode(
                        "utf-8", errors="replace"
                    )
                    if name in members_set:
                        # Patch flag_bits at offset 6-7
                        f.seek(pos + 6)
                        f.write(struct.pack("<H", flag_bits))
                except Exception:
                    pass
                pos += 4

            # Now patch central directory entries
            # Central directory signature: 0x02014b50
            central_sig = b"PK\x01\x02"
            pos = 0
            while True:
                pos = content.find(central_sig, pos)
                if pos == -1:
                    break
                # Central directory structure:
                # 0-3: signature (4 bytes)
                # 4-5: version made by (2 bytes)
                # 6-7: version needed (2 bytes)
                # 8-9: flag_bits (2 bytes)
                # ...
                # 28-29: file name length (2 bytes)
                # 30-31: extra field length (2 bytes)
                # 32-33: comment length (2 bytes)
                # ...
                # 46+: file name
                try:
                    name_len = struct.unpack_from("<H", content, pos + 28)[0]
                    name = content[pos + 46 : pos + 46 + name_len].decode(
                        "utf-8", errors="replace"
                    )
                    if name in members_set:
                        # Patch flag_bits at offset 8-9
                        f.seek(pos + 8)
                        f.write(struct.pack("<H", flag_bits))
                except Exception:
                    pass
                pos += 4

    def _update_xlsx_members_preserving_powerpoint_compat(
        self, *, xlsx_path: str, replacements: Dict[str, bytes]
    ) -> bool:
        """
        Update XLSX members in a way PowerPoint/Excel-in-PowerPoint accepts.

        PowerPoint is sensitive to ZIP metadata (e.g. deflate flag_bits) for embedded
        XLSX packages. Python's `zipfile` writer always resets `flag_bits` to 0 for
        deflated entries, which can trigger "linked file isn't available" when
        editing chart data in Excel.

        If the system `zip` CLI is available, we use it to "freshen" files in-place,
        then patch flag_bits to 6 (what Office expects). If not available, we fall back
        to pure-python rewrite (best-effort, but less compatible).
        """
        zip_cli = shutil.which("zip")
        if not zip_cli:
            self._rewrite_zip_members_in_place(xlsx_path, replacements)
            return False

        with tempfile.TemporaryDirectory() as td:
            # Write replacement files to temp dir with correct relative paths
            for rel_path, content in replacements.items():
                full_path = os.path.join(td, rel_path)
                os.makedirs(os.path.dirname(full_path), exist_ok=True)
                with open(full_path, "wb") as f:
                    f.write(content)

            # Update members using the system `zip` CLI (more PowerPoint-compatible than
            # pure-python zipfile). We update by relative path from the temp dir so the
            # member name matches exactly (e.g. xl/worksheets/sheet1.xml).
            import subprocess

            for rel_path in replacements.keys():
                proc = subprocess.run(
                    [zip_cli, "-q", xlsx_path, rel_path],
                    cwd=td,
                    capture_output=True,
                    text=True,
                )
                if proc.returncode != 0:
                    self._rewrite_zip_members_in_place(xlsx_path, replacements)
                    return False

        # The system `zip` CLI resets flag_bits to 0. Patch them back to 6 for
        # Office compatibility (deflate compression options flag).
        self._patch_zip_flag_bits(xlsx_path, list(replacements.keys()), flag_bits=6)

        return True

    def _resolve_xlsx_worksheet_path(
        self, original_files: Dict[str, bytes], sheet_name: str
    ) -> str:
        """
        Resolve an Excel worksheet part path (e.g. 'xl/worksheets/sheet2.xml') by sheet name.

        Many embedded workbooks don't store chart data in 'sheet1.xml'. PowerPoint charts
        reference a sheet name via formulas (e.g. "Sheet2!$A$2:$A$8"), so we must update
        the worksheet part that corresponds to that sheet name.
        """
        xlsx_main_ns = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
        rels_ns = "http://schemas.openxmlformats.org/package/2006/relationships"

        # Normalize sheet name: strip outer quotes used for names with spaces
        normalized = sheet_name.strip()
        if normalized.startswith("'") and normalized.endswith("'") and len(normalized) >= 2:
            normalized = normalized[1:-1]

        workbook_xml = original_files.get("xl/workbook.xml")
        workbook_rels_xml = original_files.get("xl/_rels/workbook.xml.rels")
        if workbook_xml and workbook_rels_xml:
            try:
                wb_root = etree.fromstring(workbook_xml)
                sheets = wb_root.findall(f".//{{{xlsx_main_ns}}}sheet")
                rid = None
                for sh in sheets:
                    name = sh.get("name")
                    if name == normalized:
                        rid = sh.get(
                            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
                        )
                        break
                if rid:
                    rels_root = etree.fromstring(workbook_rels_xml)
                    target = None
                    for rel in rels_root.findall(f".//{{{rels_ns}}}Relationship"):
                        if rel.get("Id") == rid:
                            target = rel.get("Target")
                            break
                    if target:
                        # Targets are relative to xl/
                        return "xl/" + target.lstrip("/")
            except Exception:
                # Fall back below
                pass

        # Fallback: pick sheet1.xml if present, else the first worksheet part found
        if "xl/worksheets/sheet1.xml" in original_files:
            return "xl/worksheets/sheet1.xml"
        worksheet_parts = sorted(
            p
            for p in original_files.keys()
            if p.startswith("xl/worksheets/sheet") and p.endswith(".xml")
        )
        if worksheet_parts:
            return worksheet_parts[0]
        return "xl/worksheets/sheet1.xml"

    def _get_chart_sheet_name(self, chart_path: str) -> Optional[str]:
        """Extract the sheet name used by chart formulas, if present."""
        try:
            ns = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}
            tree = etree.parse(chart_path)
            root = tree.getroot()
            for f_elem in root.findall(".//c:f", namespaces=ns):
                if f_elem is not None and f_elem.text and "!" in f_elem.text:
                    return self._extract_sheet_name(f_elem.text)
        except Exception:
            return None
        return None

    def _calculate_legend_dimensions(
        self,
        series_names: List[str],
        chart_width_inches: Optional[float] = None,
        chart_type: Optional[str] = None,
    ) -> dict:
        """
        Dynamically calculate legend dimensions based on series names and available width.
        Uses chart-specific config for consistent styling per chart group.

        Args:
            series_names: List of series names to display in legend
            chart_width_inches: Approximate chart width in inches (uses config default if None)
            chart_type: Chart type string for config lookup (e.g., 'bar', 'line', 'combo')

        Returns:
            dict with 'height', 'y', 'num_rows', 'entries_per_row'
        """
        from hello.utils.ppt_helpers_utils.services.template_config import (
            get_chart_layout_config,
            get_chart_group,
        )

        # Get chart-specific config
        config = get_chart_layout_config(chart_type)
        chart_group = get_chart_group(chart_type)

        # If legend should not be shown for this chart type, return minimal dimensions
        if not config.show_legend:
            return {"height": 0.0, "y": 1.0, "num_rows": 0, "entries_per_row": []}

        result = config.calculate_legend_dimensions(series_names, chart_width_inches)

        # Log for debugging
        entry_widths = [
            config.calculate_legend_entry_width(len(name)) for name in series_names
        ]
        available_width = (
            chart_width_inches or config.default_chart_width_inches
        ) * config.legend_width

        print(
            f"    [LEGEND] {chart_group}: {len(series_names)} series, {result['num_rows']} row(s) needed"
        )
        print(f'    [LEGEND] Entry widths: {[f"{w:.2f}" for w in entry_widths]}')
        print(
            f'    [LEGEND] Available width: {available_width:.2f}", entries per row: {result["entries_per_row"]}'
        )

        return result

    def _calculate_plot_area_dimensions(
        self,
        legend_y: float,
        chart_type: Optional[str] = None,
    ) -> dict:
        """
        Calculate plot area dimensions based on legend position.
        Uses chart-specific config for consistent spacing per chart group.

        Args:
            legend_y: Legend Y position (factor 0.0-1.0)
            chart_type: Chart type string for config lookup (e.g., 'bar', 'line', 'combo')

        Returns:
            dict with 'x', 'y', 'width', 'height', 'x_axis_space'
        """
        from hello.utils.ppt_helpers_utils.services.template_config import (
            get_chart_layout_config,
            get_chart_group,
        )

        # Get chart-specific config
        config = get_chart_layout_config(chart_type)
        chart_group = get_chart_group(chart_type)

        result = config.calculate_plot_area_dimensions(legend_y)

        print(
            f"    [PLOT AREA] {chart_group}: x={result['x']:.2f}, y={result['y']:.2f}, "
            f"w={result['width']:.2f}, h={result['height']:.2f}"
        )

        return result

    def _get_chart_layout_config(self, chart_type: Optional[str] = None):
        """
        Get the chart layout configuration for a given chart type.

        Args:
            chart_type: Chart type string (e.g., 'bar', 'line', 'pie', 'combo')

        Returns:
            ChartLayoutConfig for the chart's group
        """
        from hello.utils.ppt_helpers_utils.services.template_config import (
            get_chart_layout_config,
        )
        return get_chart_layout_config(chart_type)

    def update_embedded_excel_via_pptx(
        self,
        pptx_path: str,
        slide_index: int,
        chart_index: int,
        categories: List[str],
        series_data: List[tuple],
        primary_y_axis_format_code: Optional[str] = None,
        secondary_y_axis_format_code: Optional[str] = None,
    ) -> bool:
        """
        Update the embedded Excel workbook using python-pptx's replace_data().

        This method uses python-pptx's native chart.replace_data() method which
        properly updates both the chart XML cache and the embedded Excel workbook
        in a format that PowerPoint accepts.

        Args:
            pptx_path: Path to the PPTX file
            slide_index: Slide index (0-based)
            chart_index: Chart index on the slide (0-based)
            categories: List of category labels
            series_data: List of (series_name, values) tuples

        Returns:
            True if successful, False otherwise
        """
        try:
            from pptx import Presentation
            from pptx.chart.data import CategoryChartData

            prs = Presentation(pptx_path)
            slide = prs.slides[slide_index]

            # Find the chart
            chart_count = 0
            target_chart = None
            for shape in slide.shapes:
                if shape.has_chart:
                    if chart_count == chart_index:
                        target_chart = shape.chart  # type: ignore[attr-defined]
                        break
                    chart_count += 1

            if target_chart is None:
                print(
                    f"    ⚠️  Chart {chart_index} not found on slide {slide_index + 1}"
                )
                return False

            # DEBUG: Check chart type before replace_data
            print(f"    [DEBUG] Chart type (python-pptx): {target_chart.chart_type}")
            print(
                f"    [DEBUG] Number of series before replace_data: {len(target_chart.series)}"
            )
            for i, ser in enumerate(target_chart.series):
                print(f"    [DEBUG] Pre-replace Series {i}: name='{ser.name}'")

            # Create chart data using python-pptx's CategoryChartData
            chart_data = CategoryChartData()
            chart_data.categories = categories

            print(
                f"    [DEBUG] replace_data: Adding {len(series_data)} series to CategoryChartData"
            )
            for i, (series_name, values) in enumerate(series_data):
                print(
                    f"    [DEBUG] replace_data: Series {i}: '{series_name}' with {len(values)} values"
                )
                chart_data.add_series(series_name, values)

            # Use python-pptx's native replace_data method
            # This properly updates both chart XML and embedded Excel
            target_chart.replace_data(chart_data)

            # Apply frontend-driven axis/data-label formatting (Excel-style format codes).
            apply_axis_format_codes_to_chart_space(
                getattr(target_chart, "_chartSpace", None),
                primary_y_axis_format_code=primary_y_axis_format_code,
                secondary_y_axis_format_code=secondary_y_axis_format_code,
                apply_to_data_labels=True,
            )

            # Save the presentation
            prs.save(pptx_path)
            # Ensure embedded workbooks remain PowerPoint-editable after python-pptx save.
            self._normalize_pptx_embedded_workbooks(pptx_path)

            print(
                f"    ✓ Updated chart data via python-pptx replace_data(): {len(categories)} rows, {len(series_data)} series"
            )
            return True

        except Exception as e:
            print(f"    ⚠️  Error updating chart via python-pptx: {e}")
            import traceback

            traceback.print_exc()
            return False

    def update_chart_via_pptx(
        self,
        pptx_path: str,
        slide_index: int,
        chart_index: int,
        series_names: List[str],
        series_data: Optional[Sequence[Tuple[str, Sequence[Any]]]] = None,
        primary_y_axis_title: Optional[str] = None,
        secondary_y_axis_title: Optional[str] = None,
        primary_y_axis_format_code: Optional[str] = None,
        secondary_y_axis_format_code: Optional[str] = None,
    ) -> bool:
        """
        Update chart axis titles and series names using python-pptx's native API.
        This preserves the combo chart structure unlike replace_data().

        Args:
            pptx_path: Path to the PPTX file
            slide_index: Slide index (0-based)
            chart_index: Chart index on the slide (0-based)
            series_names: List of series names to set
            primary_y_axis_title: Title for primary (left) Y-axis
            secondary_y_axis_title: Title for secondary (right) Y-axis
        """
        try:
            from pptx import Presentation

            prs = Presentation(pptx_path)
            slide = prs.slides[slide_index]

            # Find the chart
            chart_count = 0
            target_chart = None
            for shape in slide.shapes:
                if shape.has_chart:
                    if chart_count == chart_index:
                        target_chart = shape.chart  # type: ignore[attr-defined]
                        break
                    chart_count += 1

            if target_chart is None:
                print(
                    f"    ⚠️  Chart {chart_index} not found on slide {slide_index + 1}"
                )
                return False

            print("    [PPTX] Updating chart via python-pptx native API")

            # Update series names
            for i, ser in enumerate(target_chart.series):
                if i < len(series_names):
                    # Access the series name through the underlying XML
                    # python-pptx's series.name is read-only, so we use XML
                    tx = ser._element.find(
                        ".//{http://schemas.openxmlformats.org/drawingml/2006/chart}tx"
                    )
                    if tx is not None:
                        v = tx.find(
                            ".//{http://schemas.openxmlformats.org/drawingml/2006/chart}v"
                        )
                        if v is not None:
                            v.text = series_names[i]
                            print(
                                f"    [PPTX] Set series {i} name to: '{series_names[i]}'"
                            )

            # Update axis titles using the chart's XML elements directly
            # python-pptx's chart.value_axis only returns the primary axis
            # For combo charts, we need to find both axes and update separately

            ns = {
                "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            }

            # Access the chart XML element
            chart_elem = target_chart._chartSpace

            # Find all value axes
            val_axes = chart_elem.findall(".//c:valAx", namespaces=ns)
            print(f"    [PPTX] Found {len(val_axes)} value axes")

            for i, va in enumerate(val_axes):
                # Get axis position (l=left/primary, r=right/secondary)
                ax_pos = va.find("./c:axPos", namespaces=ns)
                pos = ax_pos.get("val") if ax_pos is not None else None
                ax_id = va.find("./c:axId", namespaces=ns)
                ax_id_val = ax_id.get("val") if ax_id is not None else "N/A"

                print(f"    [PPTX] ValAx {i}: position={pos}, id={ax_id_val}")

                # Determine which title to set
                title_to_set = None
                if pos == "l" and primary_y_axis_title:
                    title_to_set = primary_y_axis_title
                    print(f"    [PPTX] Setting PRIMARY title on ValAx {i}")
                elif pos == "r" and secondary_y_axis_title:
                    title_to_set = secondary_y_axis_title
                    print(f"    [PPTX] Setting SECONDARY title on ValAx {i}")

                if title_to_set:
                    # Find title element and update text
                    title_elem = va.find("./c:title", namespaces=ns)
                    if title_elem is not None:
                        t_elem = title_elem.find(".//a:t", namespaces=ns)
                        if t_elem is not None:
                            old_text = t_elem.text
                            t_elem.text = title_to_set
                            print(
                                f"    [PPTX] Updated ValAx {i} title: '{old_text}' -> '{title_to_set}'"
                            )

            # Apply frontend-driven axis/data-label formatting (Excel-style format codes).
            apply_axis_format_codes_to_chart_space(
                getattr(target_chart, "_chartSpace", None),
                primary_y_axis_format_code=primary_y_axis_format_code,
                secondary_y_axis_format_code=secondary_y_axis_format_code,
                apply_to_data_labels=True,
            )

            # Save
            prs.save(pptx_path)
            # Ensure embedded workbooks remain PowerPoint-editable after python-pptx save.
            self._normalize_pptx_embedded_workbooks(pptx_path)
            print("    [PPTX] Chart updated successfully")
            return True

        except Exception as e:
            print(f"    ⚠️  Error updating chart via python-pptx: {e}")
            import traceback

            traceback.print_exc()
            return False

    def populate_chart_data(
        self,
        pptx_path: str,
        slide_index: int,
        chart_index: int,
        data: List[Dict[str, Any]],
        output_path: Optional[str] = None,
        primary_y_axis_title: Optional[str] = None,
        secondary_y_axis_title: Optional[str] = None,
        x_axis_title: Optional[str] = None,
        y_axis_keys: Optional[List[str]] = None,
        is_multi_axis: bool = True,
        primary_y_axis_format_code: Optional[str] = None,
        secondary_y_axis_format_code: Optional[str] = None,
    ) -> str:
        """
        Populate chart with new data by modifying chart XML directly

        Args:
            pptx_path: Path to PowerPoint file
            slide_index: Slide number (0-based)
            chart_index: Chart index on slide (0-based)
            data: List of data dictionaries
            output_path: Optional output path
            primary_y_axis_title: Optional title for primary (left) Y-axis
            secondary_y_axis_title: Optional title for secondary (right) Y-axis
            x_axis_title: Optional title for X-axis (category axis)
            y_axis_keys: Optional ordered list of column keys from yAxis config
                         (primary columns first, then secondary columns)
            is_multi_axis: When False, removes secondary axis and plots all series on primary

        Returns:
            Path to updated PowerPoint file
        """
        if output_path is None:
            output_path = pptx_path

        if not data or len(data) == 0:
            print("  No data to populate")
            return pptx_path

        # Create temp directory for extraction
        with tempfile.TemporaryDirectory() as temp_dir:
            extract_dir = os.path.join(temp_dir, "pptx_extract")

            # Extract PPTX
            with zipfile.ZipFile(pptx_path, "r") as zip_ref:
                zip_ref.extractall(extract_dir)

            # Find chart file
            slide_path = f"ppt/slides/slide{slide_index + 1}.xml"
            slide_rels_path = f"ppt/slides/_rels/slide{slide_index + 1}.xml.rels"

            slide_file = os.path.join(extract_dir, slide_path.replace("/", os.sep))
            rels_file = os.path.join(extract_dir, slide_rels_path.replace("/", os.sep))

            if not os.path.exists(slide_file):
                print(f"  Warning: Slide {slide_index + 1} not found")
                return pptx_path

            # Find chart relationships
            chart_rid = self._find_chart_rid(slide_file, chart_index)
            if not chart_rid:
                print(
                    f"  Warning: Chart {chart_index} not found on slide {slide_index + 1}"
                )
                return pptx_path

            # Get chart file path from relationships
            chart_file = self._get_chart_file(rels_file, chart_rid)
            if not chart_file:
                print(f"  Warning: Chart file not found for rId {chart_rid}")
                return pptx_path

            chart_path = os.path.join(
                extract_dir, "ppt", chart_file.replace("/", os.sep)
            )

            # Update chart XML directly
            success = self._update_chart_xml(
                chart_path,
                data,
                primary_y_axis_title=primary_y_axis_title,
                secondary_y_axis_title=secondary_y_axis_title,
                x_axis_title=x_axis_title,
                y_axis_keys=y_axis_keys,
                is_multi_axis=is_multi_axis,
                primary_y_axis_format_code=primary_y_axis_format_code,
                secondary_y_axis_format_code=secondary_y_axis_format_code,
            )

            if success:
                # Update embedded Excel workbook contents for *all* charts, including combo charts.
                # This prevents "Edit Data in Excel" from snapping back to template values.
                if self._pending_excel_update:
                    from hello.utils.ppt_helpers_utils.services.template_config import (
                        get_element_dimensions,
                    )

                    mode = getattr(
                        get_element_dimensions(), "embedded_excel_update_mode", "auto"
                    )
                    chart_type = self._pending_excel_update.get("chart_type", "")
                    is_combo = "combo" in str(chart_type).lower()
                    if mode in ("openxml", "auto") and (
                        is_combo or mode == "openxml"
                    ):
                        try:
                            self._update_embedded_excel(
                                chart_path=chart_path,
                                extract_dir=extract_dir,
                                categories=self._pending_excel_update["categories"],
                                series_data=self._pending_excel_update["series_data"],
                            )
                        except Exception as e:
                            print(f"    ⚠️  Embedded Excel update failed: {e}")

                # Repackage PPTX
                self._create_pptx(extract_dir, output_path)
                print(f"  ✓ Updated chart data: {len(data)} rows")

                # Now update embedded Excel via python-pptx
                # This uses python-pptx's OPC package handling which produces
                # Excel files that PowerPoint accepts
                #
                # IMPORTANT: Skip replace_data() for combo charts!
                # python-pptx sees combo charts as COLUMN_CLUSTERED and replace_data()
                # destroys the combo chart structure (converts lineChart to barChart)
                if self._pending_excel_update:
                    # Determine strategy (config-driven)
                    from hello.utils.ppt_helpers_utils.services.template_config import (
                        get_element_dimensions,
                    )

                    mode = getattr(
                        get_element_dimensions(), "embedded_excel_update_mode", "auto"
                    )

                    # Check if this is a combo chart that should skip replace_data
                    chart_type = self._pending_excel_update.get("chart_type", "")
                    is_combo = "combo" in str(chart_type).lower()

                    if is_combo:
                        print(
                            "    ⚠️  Skipping replace_data() for combo chart - would destroy chart structure"
                        )
                        print(
                            "    ℹ️  Using python-pptx native API for axis titles and series names"
                        )

                        # Use python-pptx native API to set axis titles and series names
                        series_names = [
                            name
                            for name, _ in self._pending_excel_update["series_data"]
                        ]
                        self.update_chart_via_pptx(
                            pptx_path=output_path,
                            slide_index=slide_index,
                            chart_index=chart_index,
                            series_names=series_names,
                            series_data=self._pending_excel_update["series_data"],
                            primary_y_axis_title=primary_y_axis_title,
                            secondary_y_axis_title=secondary_y_axis_title,
                            primary_y_axis_format_code=primary_y_axis_format_code,
                            secondary_y_axis_format_code=secondary_y_axis_format_code,
                        )
                    elif mode in ("python_pptx", "auto"):
                        self.update_embedded_excel_via_pptx(
                            pptx_path=output_path,
                            slide_index=slide_index,
                            chart_index=chart_index,
                            categories=self._pending_excel_update["categories"],
                            series_data=self._pending_excel_update["series_data"],
                            primary_y_axis_format_code=primary_y_axis_format_code,
                            secondary_y_axis_format_code=secondary_y_axis_format_code,
                        )
                    self._pending_excel_update = None
            else:
                print("  Warning: Could not update chart data")

        return output_path

    def _find_chart_rid(self, slide_file: str, chart_index: int) -> Optional[str]:
        """Find chart relationship ID in slide XML"""
        try:
            tree = etree.parse(slide_file)
            root = tree.getroot()

            # Find all chart references
            charts = root.findall(".//c:chart", namespaces=self.namespaces)

            if chart_index < len(charts):
                chart_elem = charts[chart_index]
                rid = chart_elem.get(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
                )
                return str(rid) if rid is not None else None

            return None
        except Exception as e:
            print(f"Error finding chart: {e}")
            return None

    def _get_chart_file(self, rels_file: str, rid: str) -> Optional[str]:
        """Get chart file path from relationships"""
        try:
            tree = etree.parse(rels_file)
            root = tree.getroot()

            for rel in root.findall(
                ".//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"
            ):
                if rel.get("Id") == rid:
                    target = rel.get("Target")
                    if target is not None:
                        # Remove ../ from path
                        target = str(target).replace("../", "")
                        return target

            return None
        except Exception as e:
            print(f"Error getting chart file: {e}")
            return None

    def _get_excel_file(self, chart_path: str, extract_dir: str) -> Optional[str]:
        """Get embedded Excel file path from chart relationships."""
        try:
            # Get PPTX directory names from config
            from hello.utils.ppt_helpers_utils.services.template_config import (
                get_element_dimensions,
            )

            element_dims = get_element_dimensions()
            ppt_dir = element_dims.pptx_ppt_dir
            embeddings_dir = element_dims.pptx_embeddings_dir

            # Read chart relationships
            chart_dir = os.path.dirname(chart_path)
            chart_name = os.path.basename(chart_path)
            rels_path = os.path.join(chart_dir, "_rels", f"{chart_name}.rels")

            if not os.path.exists(rels_path):
                print(f"  Chart rels not found: {rels_path}")
                return None

            tree = etree.parse(rels_path)
            root = tree.getroot()

            # Find embedded Excel file (usually rId1)
            for rel in root.findall(
                ".//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship"
            ):
                target = rel.get("Target")
                if target and target.endswith(".xlsx"):
                    # Properly resolve relative path from chart directory
                    # Target is like "../embeddings/Microsoft_Excel_Worksheet1.xlsx"
                    # We need to resolve this relative to chart_dir
                    target_os = target.replace("/", os.sep)
                    excel_path = os.path.normpath(os.path.join(chart_dir, target_os))

                    # Check if file exists
                    if os.path.exists(excel_path):
                        print(
                            f"    ✓ Found embedded Excel: {os.path.basename(excel_path)}"
                        )
                        return excel_path

                    # Try alternate path: directly under ppt/embeddings
                    excel_filename = os.path.basename(target)
                    excel_path_alt = os.path.join(
                        extract_dir, ppt_dir, embeddings_dir, excel_filename
                    )
                    if os.path.exists(excel_path_alt):
                        print(
                            f"    ✓ Found embedded Excel (alt): {os.path.basename(excel_path_alt)}"
                        )
                        return excel_path_alt

                    # Try another alternate: relative to extract_dir/ppt
                    target_clean = target.lstrip("../").replace("/", os.sep)
                    excel_path_alt2 = os.path.join(extract_dir, ppt_dir, target_clean)
                    if os.path.exists(excel_path_alt2):
                        print(
                            f"    ✓ Found embedded Excel (alt2): {os.path.basename(excel_path_alt2)}"
                        )
                        return excel_path_alt2

                    print(f"    ⚠️  Excel not found at: {excel_path}")
                    print(f"       Also tried: {excel_path_alt}")
                    print(f"       Also tried: {excel_path_alt2}")

            return None
        except Exception as e:
            print(f"Error getting Excel file: {e}")
            import traceback

            traceback.print_exc()
            return None

    def _analyze_data_types(self, data: List[Dict[str, Any]]) -> Dict[str, str]:
        """
        Analyze data to determine which columns are categorical vs numerical.

        Args:
            data: List of data dictionaries

        Returns:
            Dict mapping column names to types ('categorical' or 'numerical')
        """
        if not data:
            return {}

        columns = list(data[0].keys())
        column_types = {}

        for column in columns:
            # Sample values from this column (up to 10 non-null values)
            sample_values = []
            for row in data:
                value = row.get(column)
                if value is not None and value != "":
                    sample_values.append(value)
                    if len(sample_values) >= 10:
                        break

            if not sample_values:
                column_types[column] = "categorical"  # Default for empty columns
                continue

            # Check if values are numeric
            numeric_count = 0
            for value in sample_values:
                try:
                    float(value)
                    numeric_count += 1
                except (ValueError, TypeError):
                    pass

            # If most values (>= 70%) are numeric, consider it numerical
            if numeric_count / len(sample_values) >= 0.7:
                column_types[column] = "numerical"
            else:
                column_types[column] = "categorical"

        return column_types

    def _select_chart_columns(
        self,
        data: List[Dict[str, Any]],
        chart_type: str = "single_line",
        y_axis_keys: Optional[List[str]] = None,
        template_series_count: int = 0,
    ) -> tuple:
        """
        Select appropriate columns for chart based on data types, chart type, and yAxis config.

        Args:
            data: List of data dictionaries
            chart_type: Type of chart with specific handling for various chart types
            y_axis_keys: Optional ordered list of column keys from yAxis config
                         (primary columns first, then secondary columns)
            template_series_count: Number of series in the chart template (used to limit columns)

        Supported chart types:
        - 'single_line': Single line chart - 1 categorical (X) + 1 numerical (Y)
        - 'multi_line': Multi-line chart - 1 categorical (X) + multiple numerical (Y)
        - 'single_bar'/'single_column': Single bar/column - 1 categorical (X) + 1 numerical (Y)
        - 'multi_bar'/'multi_column': Multi-bar/column - 1 categorical (X) + multiple numerical (Y)
        - 'horizontal_bar': Horizontal bar - 1 categorical (Y) + 1 numerical (X) - AXES SWAPPED
        - 'combo_bar_line': Combo chart - 1 categorical (X) + 2 numerical (dual Y-axes)
        - 'combo_double_bar_line': Combo chart - 1 categorical (X) + 3 numerical (dual Y-axes)
        - 'combo_stacked_bar_line': Combo chart - 1 categorical (X) + 2+ numerical (stacked bars + line)
        - 'combo_area_bar': Combo chart - 1 categorical (X) + 2 numerical (dual Y-axes)
        - 'line_single_axis': Line chart single axis - 1 categorical (X) + 1 numerical (Y)
        - 'line_multi_axis': Line chart multi axis - 1 categorical (X) + multiple numerical (Y)
        - 'single_column_stacked': Single column stacked - 1 categorical (X) + multiple numerical (Y) stacked

        Returns:
            Tuple of (category_column, value_columns)
        """
        column_types = self._analyze_data_types(data)
        columns = list(data[0].keys())

        # Separate categorical and numerical columns
        categorical_columns = [
            col for col, dtype in column_types.items() if dtype == "categorical"
        ]
        numerical_columns = [
            col for col, dtype in column_types.items() if dtype == "numerical"
        ]

        print(f"    📊 Data analysis for {chart_type}:")
        print(f"       Categorical columns: {categorical_columns}")
        print(f"       Numerical columns: {numerical_columns}")
        if y_axis_keys:
            print(f"       Y-axis keys from config: {y_axis_keys}")
        if template_series_count > 0:
            print(f"       Template series count: {template_series_count}")

        # If y_axis_keys provided, use them to select/order value columns
        if y_axis_keys and len(y_axis_keys) > 0:
            # Filter y_axis_keys to only include columns that exist in the data
            available_keys = [key for key in y_axis_keys if key in columns]

            if available_keys:
                print("    🔑 Using y_axis_keys config to select columns")

                # Limit to template_series_count if provided
                if (
                    template_series_count > 0
                    and len(available_keys) > template_series_count
                ):
                    print(
                        f"    ✂️  Limiting to {template_series_count} columns (template has {template_series_count} series)"
                    )
                    available_keys = available_keys[:template_series_count]

                # Determine category column (first categorical, or first column not in y_axis_keys)
                if categorical_columns:
                    category_column = categorical_columns[0]
                else:
                    # Use first column that's not in y_axis_keys
                    non_value_cols = [col for col in columns if col not in y_axis_keys]
                    category_column = (
                        non_value_cols[0] if non_value_cols else columns[0]
                    )
                    print(
                        f"    ⚠️  No categorical columns found, using '{category_column}' as category"
                    )

                value_columns = available_keys

                print("    📈 Selected using y_axis_keys config:")
                print(f"       Category (X-axis): {category_column}")
                print(f"       Values (Y-axis): {value_columns}")

                return category_column, value_columns
            else:
                print(
                    "    ⚠️  y_axis_keys provided but none found in data, falling back to auto-detection"
                )

        # Fallback to original auto-detection logic
        # Special handling for pie and donut charts (single value column)
        if chart_type in ["pie", "donut"]:
            # Pie/donut charts: 1 categorical (labels) + 1 numerical (values)
            if categorical_columns and numerical_columns:
                category_column = categorical_columns[0]  # Slice labels
                value_columns = [numerical_columns[0]]  # Slice values
                print(
                    f"    🥧 {chart_type.title()} chart: Category '{category_column}' → labels, Value '{value_columns[0]}' → values"
                )
            else:
                # Fallback
                category_column = columns[0]
                value_columns = [columns[1]] if len(columns) > 1 else []
                print(
                    f"    ⚠️  {chart_type.title()} chart fallback: using '{category_column}' and '{value_columns}'"
                )
            print(f"    📈 Selected for {chart_type}:")
            print(f"       Category (labels): {category_column}")
            print(f"       Values (sizes): {value_columns}")
            return category_column, value_columns

        # Special handling for horizontal bar charts (axes are swapped)
        if chart_type == "horizontal_bar":
            # For horizontal bars: categorical goes to Y-axis, numerical to X-axis
            # But we still return them in the same order for processing
            if categorical_columns and numerical_columns:
                category_column = categorical_columns[0]  # This will be the Y-axis
                value_columns = [numerical_columns[0]]  # This will be the X-axis
                print(
                    f"    🔄 Horizontal bar: Category '{category_column}' → Y-axis, Value '{value_columns[0]}' → X-axis"
                )
            else:
                # Fallback
                category_column = columns[0]
                value_columns = [columns[1]] if len(columns) > 1 else []
                print(
                    f"    ⚠️  Horizontal bar fallback: using '{category_column}' and '{value_columns}'"
                )
        else:
            # Standard charts: categorical to X-axis, numerical to Y-axis
            if categorical_columns:
                category_column = categorical_columns[0]
            else:
                # Fallback: use first column if no categorical found
                category_column = columns[0]
                print(
                    f"    ⚠️  No categorical columns found, using first column '{category_column}' as category"
                )

            # Select value columns based on chart type
            if chart_type in [
                "single_line",
                "single_bar",
                "single_column",
                "line_single_axis",
            ]:
                # Single series charts: use first numerical column
                if numerical_columns:
                    value_columns = [numerical_columns[0]]
                else:
                    # Fallback: use second column if available
                    value_columns = [columns[1]] if len(columns) > 1 else []
                    print(
                        f"    ⚠️  No numerical columns found, using '{value_columns[0] if value_columns else 'none'}' as value"
                    )

            elif chart_type in ["combo_bar_line", "combo_area_bar"]:
                # Combo charts with 2 series: use first 2 numerical columns
                if len(numerical_columns) >= 2:
                    value_columns = numerical_columns[:2]
                elif len(numerical_columns) == 1:
                    value_columns = numerical_columns
                    print("    ⚠️  Only 1 numerical column for combo chart, expected 2")
                else:
                    # Fallback: use next available columns
                    available_cols = [col for col in columns if col != category_column]
                    value_columns = available_cols[:2]
                    print(
                        "    ⚠️  No numerical columns for combo chart, using fallback columns"
                    )

            elif chart_type == "combo_double_bar_line":
                # Double combo charts with 3 series: use first 3 numerical columns
                if len(numerical_columns) >= 3:
                    value_columns = numerical_columns[:3]
                elif len(numerical_columns) >= 1:
                    value_columns = numerical_columns
                    print(
                        f"    ⚠️  Only {len(numerical_columns)} numerical columns for double combo chart, expected 3"
                    )
                else:
                    # Fallback: use next available columns
                    available_cols = [col for col in columns if col != category_column]
                    value_columns = available_cols[:3]
                    print(
                        "    ⚠️  No numerical columns for double combo chart, using fallback columns"
                    )

            elif chart_type == "combo_stacked_bar_line":
                # Stacked bar + line combo charts: use all available numerical columns
                # (first N-1 for stacked bars, last for line)
                if len(numerical_columns) >= 2:
                    value_columns = numerical_columns
                elif len(numerical_columns) >= 1:
                    value_columns = numerical_columns
                    print(
                        f"    ⚠️  Only {len(numerical_columns)} numerical column for stacked bar + line combo chart, expected at least 2"
                    )
                else:
                    # Fallback: use next available columns
                    available_cols = [col for col in columns if col != category_column]
                    value_columns = available_cols
                    print(
                        "    ⚠️  No numerical columns for stacked bar + line combo chart, using fallback columns"
                    )

            elif chart_type in [
                "single_column_stacked",
                "single_column_stacked_chart",
                "Single_column_stacked_chart",
            ]:
                # Single column stacked chart: use all numerical columns as stacked segments
                # The chart shows one vertical column with multiple stacked series
                if numerical_columns:
                    value_columns = numerical_columns
                    print(
                        f"    📊 Single column stacked chart: using {len(value_columns)} numerical columns as stacked segments"
                    )
                else:
                    # Fallback: use all columns except the category column
                    available_cols = [col for col in columns if col != category_column]
                    value_columns = available_cols
                    print(
                        "    ⚠️  No numerical columns for single column stacked chart, using fallback columns"
                    )

            else:
                # Multi-series charts: use all numerical columns
                if numerical_columns:
                    value_columns = numerical_columns
                else:
                    # Fallback: use all columns except the category column
                    value_columns = [col for col in columns if col != category_column]

                # Limit to template_series_count if provided
                if (
                    template_series_count > 0
                    and len(value_columns) > template_series_count
                ):
                    print(
                        f"    ✂️  Limiting to {template_series_count} columns (template has {template_series_count} series)"
                    )
                    value_columns = value_columns[:template_series_count]

        print(f"    📈 Selected for {chart_type}:")
        print(f"       Category (X-axis): {category_column}")
        print(f"       Values (Y-axis): {value_columns}")

        return category_column, value_columns

    def _detect_chart_type(self, root, namespaces, series_count: int) -> str:
        """
        Detect the specific chart type from PowerPoint XML structure.

        Args:
            root: XML root element
            namespaces: XML namespaces
            series_count: Number of series in the chart

        Returns:
            String representing the chart type
        """
        ns = namespaces

        # Check for combo charts first (they have multiple chart types)
        has_line = root.find(".//c:lineChart", namespaces=ns) is not None
        has_bar = root.find(".//c:barChart", namespaces=ns) is not None
        has_column = root.find(".//c:columnChart", namespaces=ns) is not None
        has_area = root.find(".//c:areaChart", namespaces=ns) is not None
        has_pie = root.find(".//c:pieChart", namespaces=ns) is not None
        has_donut = root.find(".//c:doughnutChart", namespaces=ns) is not None

        # Pie and donut charts are standalone - return early if found
        if has_donut:
            return "donut"
        if has_pie:
            return "pie"

        chart_types_found = []
        if has_line:
            chart_types_found.append("line")
        if has_bar:
            chart_types_found.append("bar")
        if has_column:
            chart_types_found.append("column")
        if has_area:
            chart_types_found.append("area")

        # Combo chart detection
        if len(chart_types_found) > 1:
            if "bar" in chart_types_found and "line" in chart_types_found:
                # Check if bar is stacked
                bar_chart = root.find(".//c:barChart", namespaces=ns)
                is_stacked = False
                if bar_chart is not None:
                    grouping = bar_chart.find(".//c:grouping", namespaces=ns)
                    is_stacked = grouping is not None and grouping.get("val") in (
                        "stacked",
                        "percentStacked",
                    )

                if is_stacked:
                    return "combo_stacked_bar_line"
                elif series_count >= 3:
                    return "combo_double_bar_line"
                else:
                    return "combo_bar_line"
            elif "area" in chart_types_found and "bar" in chart_types_found:
                return "combo_area_bar"
            elif "column" in chart_types_found and "line" in chart_types_found:
                if series_count >= 3:
                    return "combo_double_bar_line"  # Treat column+line same as bar+line
                else:
                    return "combo_bar_line"
            else:
                # Other combo types - treat as multi-series
                return "multi_series"

        # Single chart type detection
        if has_line:
            # Check for multi-axis by looking for secondary axis
            # Use a two-step approach to avoid XPath predicate issues
            has_secondary_axis = (
                root.find('.//c:valAx[@val="1"]', namespaces=ns) is not None
            )

            # Alternative approach: find axId with val="2" then check if parent is valAx
            axid_elem = root.find('.//c:axId[@val="2"]', namespaces=ns)
            if axid_elem is not None:
                parent = axid_elem.getparent()
                if parent is not None and parent.tag.endswith("}valAx"):
                    has_secondary_axis = True

            if has_secondary_axis or series_count > 1:
                return "line_multi_axis"
            else:
                return "line_single_axis"

        elif has_bar:
            bar_chart = root.find(".//c:barChart", namespaces=ns)
            if bar_chart is not None:
                # Check for stacked grouping
                grouping = bar_chart.find(".//c:grouping", namespaces=ns)
                is_stacked = grouping is not None and grouping.get("val") in (
                    "stacked",
                    "percentStacked",
                )

                # Check if it's horizontal bar by looking at barDir attribute
                bar_dir = bar_chart.find(".//c:barDir", namespaces=ns)
                is_horizontal = bar_dir is not None and bar_dir.get("val") == "bar"

                if is_stacked:
                    if is_horizontal:
                        return "stacked_horizontal_bar"
                    return "stacked_bar"
                elif is_horizontal:
                    return "horizontal_bar"

            # Regular vertical bar
            if series_count == 1:
                return "single_bar"
            else:
                return "multi_bar"

        elif has_column:
            if series_count == 1:
                return "single_column"
            else:
                return "multi_column"

        elif has_area:
            if series_count == 1:
                return "single_area"
            else:
                return "multi_area"

        # Default fallback
        print("    ⚠️  Unknown chart type, using multi_series default")
        return "multi_series"

    def _update_chart_xml(
        self,
        chart_path: str,
        data: List[Dict[str, Any]],
        primary_y_axis_title: Optional[str] = None,
        secondary_y_axis_title: Optional[str] = None,
        x_axis_title: Optional[str] = None,
        y_axis_keys: Optional[List[str]] = None,
        is_multi_axis: bool = True,
        primary_y_axis_format_code: Optional[str] = None,
        secondary_y_axis_format_code: Optional[str] = None,
    ) -> bool:
        """Update chart data directly in XML, update axis titles, and ensure legend is present for all series.

        Args:
            chart_path: Path to chart XML file
            data: List of data dictionaries
            primary_y_axis_title: Title for primary (left) Y-axis
            secondary_y_axis_title: Title for secondary (right) Y-axis
            x_axis_title: Title for X-axis
            y_axis_keys: Ordered list of Y-axis column keys
            is_multi_axis: When False, removes secondary axis and plots all series on primary
        """
        try:
            # Parse chart XML
            tree = etree.parse(chart_path)
            root = tree.getroot()

            # Find chart data elements
            ns = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}

            # Apply frontend-driven axis/data-label formatting (Excel-style format codes).
            apply_axis_format_codes_to_chart_space(
                root,
                primary_y_axis_format_code=primary_y_axis_format_code,
                secondary_y_axis_format_code=secondary_y_axis_format_code,
                apply_to_data_labels=True,
            )

            # Find all series in the chart
            series_list = root.findall(".//c:ser", namespaces=ns)

            if not series_list:
                print("    No chart series found in XML")
                return False

            # Determine chart type from XML structure with enhanced detection
            chart_type = self._detect_chart_type(root, ns, len(series_list))

            # Intelligently select columns based on data types, chart type, and y_axis_keys config
            category_column, value_columns = self._select_chart_columns(
                data,
                chart_type,
                y_axis_keys=y_axis_keys,
                template_series_count=len(series_list),
            )

            # Extract category values
            categories = [str(row[category_column]) for row in data]

            print(
                f"    📊 Chart has {len(series_list)} series, data has {len(value_columns)} value columns"
            )
            print(f"    [DEBUG] Template series count: {len(series_list)}")
            print(f"    [DEBUG] Data value columns: {value_columns}")

            # REMOVE extra series that don't have corresponding data
            if len(series_list) > len(value_columns):
                series_to_remove = len(series_list) - len(value_columns)
                print(f"    🗑️  Removing {series_to_remove} extra series from chart...")
                for i, series in enumerate(series_list[len(value_columns) :]):
                    # Get series name before removing
                    ser_name = series.find(".//c:tx//c:v", namespaces=ns)
                    ser_name_text = (
                        ser_name.text
                        if ser_name is not None
                        else f"Series {len(value_columns) + i}"
                    )
                    print(f"    [DEBUG] Removing series: '{ser_name_text}'")
                    parent = series.getparent()
                    if parent is not None:
                        parent.remove(series)
                        print("✓ Removed extra series")
                series_list = series_list[: len(value_columns)]

            print(f"    [DEBUG] Final series count after removal: {len(series_list)}")

            # Debug: Print each series name/index
            for idx, series in enumerate(series_list):
                ser_idx = series.find("./c:idx", namespaces=ns)
                ser_order = series.find("./c:order", namespaces=ns)
                ser_tx = series.find("./c:tx", namespaces=ns)
                ser_name = (
                    ser_tx.find(".//c:v", namespaces=ns) if ser_tx is not None else None
                )
                idx_val = ser_idx.get("val") if ser_idx is not None else "N/A"
                order_val = ser_order.get("val") if ser_order is not None else "N/A"
                name_val = ser_name.text if ser_name is not None else "N/A"
                # Check which chart type this series belongs to (bar, line, etc.)
                parent = series.getparent()
                parent_tag = parent.tag.split("}")[-1] if parent is not None else "N/A"
                print(
                    f"    [DEBUG] Series {idx}: idx={idx_val}, order={order_val}, name='{name_val}', parent={parent_tag}"
                )

            # Collect series data for embedded Excel update
            series_data_for_excel = []

            # Update each remaining series
            for series_idx, series in enumerate(series_list):
                if series_idx >= len(value_columns):
                    break

                value_column = value_columns[series_idx]
                values = [
                    float(row.get(value_column, 0))
                    if isinstance(row.get(value_column, 0), (int, float))
                    else 0
                    for row in data
                ]

                # Force unique series name for legend
                # Format series legend name: replace underscores with spaces, title case
                raw_name = str(value_column)
                series_name = raw_name
                # format_label(raw_name, fallback="Series")
                if len(value_columns) > 1:
                    # Keep distinct names while still sanitized
                    series_name = (
                        series_name  # placeholder if future disambiguation needed
                    )

                # Collect series data for embedded Excel update
                series_data_for_excel.append((series_name, values))

                # Ensure <c:idx> and <c:order> are unique and sequential
                idx_elem = series.find("./c:idx", namespaces=ns)
                if idx_elem is not None:
                    idx_elem.set("val", str(series_idx))
                else:
                    idx_elem = etree.Element(
                        "{http://schemas.openxmlformats.org/drawingml/2006/chart}idx"
                    )
                    idx_elem.set("val", str(series_idx))
                    series.insert(0, idx_elem)
                order_elem = series.find("./c:order", namespaces=ns)
                if order_elem is not None:
                    order_elem.set("val", str(series_idx))
                else:
                    order_elem = etree.Element(
                        "{http://schemas.openxmlformats.org/drawingml/2006/chart}order"
                    )
                    order_elem.set("val", str(series_idx))
                    series.insert(1, order_elem)

                tx = series.find(".//c:tx", namespaces=ns)
                if tx is not None:
                    v = tx.find(".//c:v", namespaces=ns)
                    if v is not None:
                        v.text = series_name
                        print(
                            f"    [DEBUG] Updated series {series_idx} name to: '{series_name}'"
                        )
                else:
                    tx = etree.SubElement(
                        series,
                        "{http://schemas.openxmlformats.org/drawingml/2006/chart}tx",
                    )
                    v = etree.SubElement(
                        tx, "{http://schemas.openxmlformats.org/drawingml/2006/chart}v"
                    )
                    v.text = series_name
                    print(
                        f"    [DEBUG] Created series {series_idx} name: '{series_name}'"
                    )

                # ...existing code for category axis and values...
                cat = series.find(".//c:cat", namespaces=ns)
                if cat is not None:
                    strRef = cat.find(".//c:strRef", namespaces=ns)
                    if strRef is not None:
                        strCache = strRef.find(".//c:strCache", namespaces=ns)
                        if strCache is not None:
                            ptCount = strCache.find(".//c:ptCount", namespaces=ns)
                            if ptCount is not None:
                                ptCount.set("val", str(len(categories)))
                            for pt in strCache.findall(".//c:pt", namespaces=ns):
                                strCache.remove(pt)
                            for idx, cat_value in enumerate(categories):
                                pt = etree.SubElement(
                                    strCache,
                                    "{http://schemas.openxmlformats.org/drawingml/2006/chart}pt",
                                )
                                pt.set("idx", str(idx))
                                v_elem = etree.SubElement(
                                    pt,
                                    "{http://schemas.openxmlformats.org/drawingml/2006/chart}v",
                                )
                                v_elem.text = cat_value

                val = series.find(".//c:val", namespaces=ns)
                if val is not None:
                    numRef = val.find(".//c:numRef", namespaces=ns)
                    if numRef is not None:
                        numCache = numRef.find(".//c:numCache", namespaces=ns)
                        if numCache is not None:
                            ptCount = numCache.find(".//c:ptCount", namespaces=ns)
                            if ptCount is not None:
                                ptCount.set("val", str(len(values)))
                            for pt in numCache.findall(".//c:pt", namespaces=ns):
                                numCache.remove(pt)
                            for idx, val_value in enumerate(values):
                                pt = etree.SubElement(
                                    numCache,
                                    "{http://schemas.openxmlformats.org/drawingml/2006/chart}pt",
                                )
                                pt.set("idx", str(idx))
                                v_elem = etree.SubElement(
                                    pt,
                                    "{http://schemas.openxmlformats.org/drawingml/2006/chart}v",
                                )
                                v_elem.text = str(val_value)

            # Handle legend using chart-specific configuration
            chart_elem = root.find(".//c:chart", namespaces=ns)
            legend_elem = None

            # Get chart-specific layout config
            from hello.utils.ppt_helpers_utils.services.template_config import (
                get_chart_layout_config,
                get_chart_group,
            )
            chart_config = get_chart_layout_config(chart_type)
            chart_group = get_chart_group(chart_type)

            # Determine series count for single-series legend hiding
            series_count = len(value_columns) if value_columns else 1

            # Disable legends for single-series line and bar charts
            # Single-series charts don't need a legend as there's only one data series
            is_single_series_line_or_bar = (
                series_count == 1 and chart_group in ("LINE", "VERTICAL_BAR")
            )

            # Use config to determine if legend should be shown
            is_horizontal_bar = chart_group == "HORIZONTAL_BAR"
            skip_legend = not chart_config.show_legend or is_single_series_line_or_bar

            print(
                f"    [DEBUG] Legend handling: group={chart_group}, show_legend={chart_config.show_legend}, chart_type={chart_type}, series_count={series_count}"
            )
            if is_single_series_line_or_bar:
                print(
                    f"    ℹ️  Single-series {chart_group} chart: Legend disabled (only 1 series)"
                )

            if chart_elem is not None:
                # For combo charts, preserve existing legend but adjust width for smaller charts
                existing_legends = [
                    child
                    for child in list(chart_elem)
                    if child.tag
                    == "{http://schemas.openxmlformats.org/drawingml/2006/chart}legend"
                ]
                print(
                    f"    [DEBUG] Found {len(existing_legends)} existing legend element(s)"
                )

                is_combo = chart_group == "COMBO"
                is_line = chart_group == "LINE"

                # Use template's legend presence to determine if legend space is needed
                has_legend_in_template = len(existing_legends) > 0

                # Determine if legend should be hidden (config disabled OR single-series line/bar)
                should_hide_legend = not chart_config.show_legend or is_single_series_line_or_bar

                if not has_legend_in_template or should_hide_legend:
                    # No legend in template or legend should be hidden - no legend space needed
                    self._combo_legend_height = 0.0
                    self._combo_legend_y = 1.0  # No legend = no space reserved
                    if should_hide_legend:
                        if is_single_series_line_or_bar:
                            print(f"    ℹ️  Single-series {chart_group} chart: Legend disabled")
                        else:
                            print(f"    ℹ️  {chart_group} config: Legend disabled")
                        # Remove existing legend from template when legend should be hidden
                        if has_legend_in_template:
                            for legend in existing_legends:
                                legend_parent = legend.getparent()
                                if legend_parent is not None:
                                    legend_parent.remove(legend)
                                    print(f"    🗑️  Removed existing legend from {chart_group} chart")
                    else:
                        print("    ℹ️  No legend in template - no legend space needed")

                if chart_config.show_legend and not is_single_series_line_or_bar and has_legend_in_template and (is_combo or is_line):
                    # Dynamically calculate legend dimensions using chart-specific config
                    legend = existing_legends[0]
                    layout = legend.find("./c:layout", namespaces=ns)
                    if layout is not None:
                        manual = layout.find("./c:manualLayout", namespaces=ns)
                        if manual is not None:
                            # Get series names for dimension calculation
                            series_names = [
                                str(col)
                                for col in (value_columns or [])
                            ]

                            # Calculate legend dimensions using chart-specific config
                            legend_dims = self._calculate_legend_dimensions(
                                series_names, chart_type=chart_type
                            )

                            # Apply legend dimensions from chart-specific config
                            w_elem = manual.find("./c:w", namespaces=ns)
                            h_elem = manual.find("./c:h", namespaces=ns)
                            x_elem = manual.find("./c:x", namespaces=ns)
                            y_elem = manual.find("./c:y", namespaces=ns)

                            if w_elem is not None:
                                w_elem.set("val", f"{chart_config.legend_width:.4f}")
                            if x_elem is not None:
                                x_elem.set("val", f"{chart_config.legend_x:.4f}")
                            if h_elem is not None:
                                h_elem.set("val", f"{legend_dims['height']:.4f}")
                            if y_elem is not None:
                                y_elem.set("val", f"{legend_dims['y']:.4f}")

                            # Store legend info for plot area adjustment
                            self._combo_legend_height = legend_dims["height"]
                            self._combo_legend_y = legend_dims["y"]

                            print(
                                f"    📐 {chart_group} chart: Legend needs {legend_dims['num_rows']} row(s) - height={legend_dims['height'] * 100:.0f}%, y={legend_dims['y']:.2f}"
                            )
                    skip_legend = True  # Skip recreation, keep adjusted template legend
                elif chart_config.show_legend and not is_single_series_line_or_bar and has_legend_in_template:
                    # Skip legend recreation - just use the template's legend
                    skip_legend = True

                if skip_legend:
                    # Skip legend creation based on config
                    print(f"    ⏭️  Skipping legend recreation for {chart_group} ({chart_type})")
                else:
                    # ---------------- Legend (Bottom, Non‑Overlapping) ---------------- #
                    # We recreate the legend and explicitly reserve vertical space so that
                    # category axis tick labels never collide with the legend. In some
                    # PowerPoint builds the default auto-layout still causes partial
                    # overlap when the chart shape is short or category labels wrap.
                    legend_elem = etree.SubElement(
                        chart_elem,
                        "{http://schemas.openxmlformats.org/drawingml/2006/chart}legend",
                    )
                    legend_pos = etree.SubElement(
                        legend_elem,
                        "{http://schemas.openxmlformats.org/drawingml/2006/chart}legendPos",
                    )
                    # Use chart-specific legend position
                    legend_pos_val = "b"  # default bottom
                    if chart_config.legend_position == "right":
                        legend_pos_val = "r"
                    elif chart_config.legend_position == "top":
                        legend_pos_val = "t"
                    legend_pos.set("val", legend_pos_val)
                    overlay = etree.SubElement(
                        legend_elem,
                        "{http://schemas.openxmlformats.org/drawingml/2006/chart}overlay",
                    )
                    overlay.set("val", "0")  # NOT overlay -> allocate its own space

                    # Provide a manual layout with DYNAMIC height based on series count
                    # (The units are factors of the chart container. y+h must be <= 1)
                    # Legend is positioned based on chart-specific config
                    try:
                        # Calculate dynamic legend height using chart-specific config
                        series_count = len(value_columns) if value_columns else 1
                        legend_height = min(
                            chart_config.legend_height_max,
                            series_count * chart_config.legend_row_height_factor + chart_config.legend_padding,
                        )
                        # Position legend at bottom (tight positioning, no gap)
                        legend_y = 1.0 - legend_height - chart_config.legend_bottom_margin

                        layout = etree.SubElement(
                            legend_elem,
                            "{http://schemas.openxmlformats.org/drawingml/2006/chart}layout",
                        )
                        manual = etree.SubElement(
                            layout,
                            "{http://schemas.openxmlformats.org/drawingml/2006/chart}manualLayout",
                        )
                        for tag, val in [
                            ("xMode", "factor"),
                            ("yMode", "factor"),
                            ("wMode", "factor"),
                            ("hMode", "factor"),
                        ]:
                            t = etree.SubElement(
                                manual,
                                f"{{http://schemas.openxmlformats.org/drawingml/2006/chart}}{tag}",
                            )
                            t.set("val", val)
                        x = etree.SubElement(
                            manual,
                            "{http://schemas.openxmlformats.org/drawingml/2006/chart}x",
                        )
                        x.set("val", f"{chart_config.legend_x:.4f}")
                        y = etree.SubElement(
                            manual,
                            "{http://schemas.openxmlformats.org/drawingml/2006/chart}y",
                        )
                        y.set("val", f"{legend_y:.4f}")
                        w = etree.SubElement(
                            manual,
                            "{http://schemas.openxmlformats.org/drawingml/2006/chart}w",
                        )
                        w.set("val", f"{chart_config.legend_width:.4f}")
                        h = etree.SubElement(
                            manual,
                            "{http://schemas.openxmlformats.org/drawingml/2006/chart}h",
                        )
                        h.set("val", f"{legend_height:.4f}")

                        print(f"    📐 {chart_group} legend: height={legend_height:.2f}, y={legend_y:.2f}")
                    except Exception:
                        # Non-fatal; if PPT ignores manual layout it will fallback to auto.
                        pass

                # ---------------- Plot Area (MAXIMIZED DIMENSIONS) ---------------- #
                # Plot area is maximized to fill container from top to legend (no gaps)
                # Uses chart-specific configuration for each chart group

                # Use chart group to determine behavior
                is_combo_chart = chart_group == "COMBO"
                is_line_chart = chart_group == "LINE"

                # For combo/line charts, we need to adjust plot area to accommodate legend
                combo_legend_y = getattr(self, "_combo_legend_y", 0.92)

                try:
                    plot_area = root.find(".//c:plotArea", namespaces=ns)
                    if plot_area is not None and (is_combo_chart or is_line_chart):
                        # For combo/line charts: adjust plot area using chart-specific dimensions
                        layout = plot_area.find("./c:layout", namespaces=ns)
                        if layout is not None:
                            manual = layout.find("./c:manualLayout", namespaces=ns)
                            if manual is not None:
                                # Calculate plot area dimensions using chart-specific config
                                plot_dims = self._calculate_plot_area_dimensions(
                                    combo_legend_y, chart_type=chart_type
                                )

                                # Apply plot area dimensions
                                h_elem = manual.find("./c:h", namespaces=ns)
                                y_elem = manual.find("./c:y", namespaces=ns)
                                if h_elem is not None:
                                    h_elem.set("val", f"{plot_dims['height']:.4f}")
                                if y_elem is not None:
                                    y_elem.set("val", f"{plot_dims['y']:.4f}")

                                print(
                                    f"    📐 {chart_group} chart: Plot area y={plot_dims['y']:.0%}, h={plot_dims['height']:.0%}, x_axis_space={plot_dims['x_axis_space']:.0%} (legend at y={combo_legend_y:.2f})"
                                )
                        else:
                            # No existing layout - preserve template defaults
                            print(
                                f"    📐 {chart_group} chart: Preserving template plot area (no manual layout)"
                            )

                    elif (
                        plot_area is not None
                        and not is_combo_chart
                        and not is_line_chart
                    ):
                        layout = plot_area.find("./c:layout", namespaces=ns)
                        if layout is None:
                            layout = etree.SubElement(
                                plot_area,
                                "{http://schemas.openxmlformats.org/drawingml/2006/chart}layout",
                            )
                        # Remove prior manualLayout definitions to avoid duplication
                        for ml in layout.findall("./c:manualLayout", namespaces=ns):
                            layout.remove(ml)
                        manual = etree.SubElement(
                            layout,
                            "{http://schemas.openxmlformats.org/drawingml/2006/chart}manualLayout",
                        )
                        for tag, val in [
                            ("xMode", "factor"),
                            ("yMode", "factor"),
                            ("wMode", "factor"),
                            ("hMode", "factor"),
                        ]:
                            t = etree.SubElement(
                                manual,
                                f"{{http://schemas.openxmlformats.org/drawingml/2006/chart}}{tag}",
                            )
                            t.set("val", val)

                        # Use chart-specific config for plot area dimensions
                        if is_horizontal_bar:
                            # Horizontal bar: use config values (typically full area)
                            plot_x = chart_config.plot_area_x
                            plot_y = chart_config.get_plot_area_y()
                            plot_width = chart_config.plot_area_width
                            plot_height = chart_config.plot_area_height or 1.0
                            print(
                                f"    📐 {chart_group}: Plot area at ({plot_x:.2f},{plot_y:.2f}) with dimensions ({plot_width:.2f}x{plot_height:.2f})"
                            )
                        else:
                            # Calculate dynamic legend height using chart-specific config
                            series_count = len(value_columns) if value_columns else 1
                            
                            # Check if legend is hidden for single-series line/bar charts
                            is_single_series_no_legend = (
                                series_count == 1 and chart_group in ("LINE", "VERTICAL_BAR")
                            )
                            
                            if is_single_series_no_legend:
                                # No legend space needed for single-series charts
                                actual_legend_height = 0.0
                            else:
                                actual_legend_height = min(
                                    chart_config.legend_height_max,
                                    series_count * chart_config.legend_row_height_factor + chart_config.legend_padding,
                                )

                            # Use chart-specific plot area margins
                            plot_x = chart_config.plot_area_x
                            plot_y = chart_config.plot_area_top_margin
                            plot_width = chart_config.plot_area_width
                            # Plot height fills to legend - no gap between plot and legend
                            plot_height = 1.0 - plot_y - actual_legend_height

                            print(
                                f"    📐 {chart_group} chart: Plot area y={plot_y:.0%}, h={plot_height:.0%} (legend={actual_legend_height:.0%})"
                            )

                        x = etree.SubElement(
                            manual,
                            "{http://schemas.openxmlformats.org/drawingml/2006/chart}x",
                        )
                        x.set("val", f"{plot_x:.4f}")
                        y = etree.SubElement(
                            manual,
                            "{http://schemas.openxmlformats.org/drawingml/2006/chart}y",
                        )
                        y.set("val", f"{plot_y:.4f}")
                        w = etree.SubElement(
                            manual,
                            "{http://schemas.openxmlformats.org/drawingml/2006/chart}w",
                        )
                        w.set("val", f"{plot_width:.4f}")
                        h = etree.SubElement(
                            manual,
                            "{http://schemas.openxmlformats.org/drawingml/2006/chart}h",
                        )
                        h.set("val", f"{plot_height:.4f}")
                except Exception:
                    pass

                # ---------------- Axis Label Font Scaling (optional) ---------------- #
                # If many categories, reduce font size to mitigate vertical spillover.
                try:
                    cat_axis = root.find(".//c:catAx", namespaces=ns)
                    if cat_axis is not None:
                        # Set tick label position to "low" to ensure labels appear below axis
                        # This prevents labels from appearing inside the plot area when charts are scaled
                        tickLblPos = cat_axis.find("./c:tickLblPos", namespaces=ns)
                        if tickLblPos is None:
                            # Insert tickLblPos after basic axis elements (before txPr if it exists)
                            # Find insertion point (before txPr or at end)
                            txPr_elem = cat_axis.find("./c:txPr", namespaces=ns)
                            if txPr_elem is not None:
                                insert_idx = list(cat_axis).index(txPr_elem)
                                tickLblPos = etree.Element(
                                    "{http://schemas.openxmlformats.org/drawingml/2006/chart}tickLblPos"
                                )
                                tickLblPos.set("val", "low")
                                cat_axis.insert(insert_idx, tickLblPos)
                            else:
                                tickLblPos = etree.SubElement(
                                    cat_axis,
                                    "{http://schemas.openxmlformats.org/drawingml/2006/chart}tickLblPos",
                                )
                                tickLblPos.set("val", "low")
                        else:
                            # Update existing tickLblPos to "low"
                            tickLblPos.set("val", "low")

                        txPr = cat_axis.find("./c:txPr", namespaces=ns)
                        if txPr is None:
                            txPr = etree.SubElement(
                                cat_axis,
                                "{http://schemas.openxmlformats.org/drawingml/2006/chart}txPr",
                            )
                            etree.SubElement(
                                txPr,
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr",
                            )
                            etree.SubElement(
                                txPr,
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle",
                            )
                            p = etree.SubElement(
                                txPr,
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}p",
                            )
                        else:
                            p = txPr.find(
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}p"
                            )
                            if p is None:
                                p = etree.SubElement(
                                    txPr,
                                    "{http://schemas.openxmlformats.org/drawingml/2006/main}p",
                                )
                        pPr = p.find(
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
                        )
                        if pPr is None:
                            pPr = etree.SubElement(
                                p,
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr",
                            )
                        defRPr = pPr.find(
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr"
                        )
                        if defRPr is None:
                            defRPr = etree.SubElement(
                                pPr,
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr",
                            )
                        # Font size attribute is in 1/100 points. Use 900 (=9pt) when crowded.
                        if len(categories) > 10:
                            defRPr.set("sz", "900")
                        elif len(categories) > 7:
                            defRPr.set("sz", "1000")  # 10pt
                except Exception:
                    pass

                # ---------------- Value Axis Tick Label Position ---------------- #
                # Set tick label position based on axis position:
                # - Left axis (axPos="l"): tickLblPos="low" (labels on left)
                # - Right axis (axPos="r"): tickLblPos="high" (labels on right)
                try:
                    value_axes = root.findall(".//c:valAx", namespaces=ns)
                    for val_axis in value_axes:
                        # Determine axis position
                        ax_pos_elem = val_axis.find("./c:axPos", namespaces=ns)
                        axis_position = (
                            ax_pos_elem.get("val", "l")
                            if ax_pos_elem is not None
                            else "l"
                        )

                        # Set appropriate tickLblPos based on axis position
                        tick_pos_value = "high" if axis_position == "r" else "low"

                        tickLblPos = val_axis.find("./c:tickLblPos", namespaces=ns)
                        if tickLblPos is None:
                            # Insert tickLblPos at appropriate position
                            txPr_elem = val_axis.find("./c:txPr", namespaces=ns)
                            if txPr_elem is not None:
                                insert_idx = list(val_axis).index(txPr_elem)
                                tickLblPos = etree.Element(
                                    "{http://schemas.openxmlformats.org/drawingml/2006/chart}tickLblPos"
                                )
                                tickLblPos.set("val", tick_pos_value)
                                val_axis.insert(insert_idx, tickLblPos)
                            else:
                                tickLblPos = etree.SubElement(
                                    val_axis,
                                    "{http://schemas.openxmlformats.org/drawingml/2006/chart}tickLblPos",
                                )
                                tickLblPos.set("val", tick_pos_value)
                        else:
                            # Update existing tickLblPos based on axis position
                            tickLblPos.set("val", tick_pos_value)
                except Exception:
                    pass

            # Remove all <c:legendEntry> elements anywhere in the tree (paranoia)
            for elem in root.findall(".//c:legendEntry", namespaces=ns):
                parent = elem.getparent()
                if parent is not None:
                    parent.remove(elem)

            # Handle multi-axis configuration
            # When is_multi_axis=False, remove secondary axis and plot all series on primary
            if not is_multi_axis:
                self._remove_secondary_axis(root, ns)
                print(
                    "    📊 Single-axis mode: Removed secondary Y-axis (isMultiAxis=false)"
                )

            # Enable auto-scaling for all value axes with data-driven min/max
            self._enable_auto_scaling(root, ns, data, value_columns, chart_type)

            # Use chart-specific config to determine if axis titles should be shown
            # Charts with show_axis_titles=False (e.g., horizontal bar, pie/donut) skip titles
            skip_axis_titles = not chart_config.show_axis_titles

            # Update axis titles if provided (unless skipped by config)
            if not skip_axis_titles:
                # When is_multi_axis=False, don't set secondary axis title
                effective_secondary_title = (
                    secondary_y_axis_title if is_multi_axis else None
                )
                self._update_axis_titles(
                    root,
                    ns,
                    primary_y_axis_title=primary_y_axis_title,
                    secondary_y_axis_title=effective_secondary_title,
                    x_axis_title=x_axis_title,
                    chart_type=chart_type,
                )
            else:
                # Remove any existing axis titles from the template
                self._remove_axis_titles(root, ns)
                print(f"    ⏭️  Removed axis titles for {chart_type} chart")

            # Update formula references in chart XML to match new data range
            # This ensures formulas point to the correct Excel cells
            self._update_formula_references(
                root,
                ns,
                num_categories=len(categories),
                num_series=len(series_data_for_excel),
            )

            # Update embedded Excel workbook with chart data
            # This fixes the "Edit Data in Excel" issue where old template data was shown
            if series_data_for_excel:
                # Store data for embedded Excel update (will be done via python-pptx after save)
                # The actual Excel update happens in _update_embedded_excel_via_pptx()
                # after the PPTX is reassembled
                self._pending_excel_update = {
                    "categories": categories,
                    "series_data": series_data_for_excel,
                    "slide_index": None,  # Will be set by caller
                    "chart_index": None,  # Will be set by caller
                    "chart_type": chart_type,  # Store chart type to check for combo charts
                }

            # Save modified XML with correct double-quote declaration for Office
            self._write_xml_with_office_declaration(tree, chart_path)
            return True

        except Exception as e:
            print(f"    Error updating chart XML: {e}")
            import traceback

            traceback.print_exc()
            return False

    def _round_to_nice_number(
        self, value: float, direction: str = "round", force_integer: bool = False
    ) -> float:
        """
        Round a value to a "nice" number for axis scaling.

        Uses floor for min values and ceil for max values to ensure
        all data points are visible within the axis range.

        Args:
            value: The value to round
            direction: "floor" for min values, "ceil" for max values, "round" for nearest
            force_integer: If True, always returns integers. If False, allows decimals.

        Returns:
            A nicely rounded number (float or int depending on force_integer)
        """
        if value == 0:
            return 0 if force_integer else 0.0

        abs_value = abs(value)
        step: float

        if force_integer:
            # Integer-only steps: 1, 2, 5, 10, 50, 100, 500, 1000, etc.
            if abs_value >= 10000:
                step = 1000 if abs_value >= 50000 else 500
            elif abs_value >= 1000:
                step = 100 if abs_value >= 5000 else 50
            elif abs_value >= 100:
                step = 10 if abs_value >= 500 else 5
            elif abs_value >= 10:
                step = 5 if abs_value >= 50 else 2
            else:
                step = 1
        else:
            # Allow decimal steps for better precision
            if abs_value >= 10000:
                step = 1000 if abs_value >= 50000 else 500
            elif abs_value >= 1000:
                step = 100 if abs_value >= 5000 else 50
            elif abs_value >= 100:
                step = 10 if abs_value >= 500 else 5
            elif abs_value >= 10:
                step = 1.0 if abs_value >= 50 else 0.5
            elif abs_value >= 1:
                step = 0.5 if abs_value >= 5 else 0.1
            else:
                # For values < 1, use smaller steps
                if abs_value >= 0.1:
                    step = 0.05
                elif abs_value >= 0.01:
                    step = 0.005
                else:
                    step = 0.001

        # Apply floor or ceil based on direction
        if direction == "floor":
            result = math.floor(value / step) * step
        elif direction == "ceil":
            result = math.ceil(value / step) * step
        else:
            result = round(value / step) * step

        if force_integer:
            return int(result)
        return result

    def _get_axis_series_data(
        self,
        root,
        ns: dict,
        axis_id: str,
        data: List[Dict[str, Any]],
        value_columns: List[str],
    ) -> List[float]:
        """
        Get numeric values for all series that belong to a specific axis.

        Args:
            root: XML root element
            ns: XML namespaces dict
            axis_id: The axis ID to filter series by
            data: List of data dictionaries
            value_columns: List of value column names

        Returns:
            List of numeric values for series belonging to this axis
        """
        numeric_values: List[float] = []
        series_list = root.findall(".//c:ser", namespaces=ns)

        for series_idx, series in enumerate(series_list):
            # Check if this series belongs to the given axis
            ax_id_refs = series.findall(".//c:axId", namespaces=ns)
            belongs_to_axis = False

            for ax_id_ref in ax_id_refs:
                if ax_id_ref.get("val") == axis_id:
                    belongs_to_axis = True
                    break

            if belongs_to_axis and series_idx < len(value_columns):
                col = value_columns[series_idx]
                for row in data:
                    value = row.get(col)
                    if value is not None:
                        try:
                            numeric_val = float(value)
                            # Skip NaN and infinite values
                            if not (
                                numeric_val != numeric_val
                                or numeric_val == float("inf")
                                or numeric_val == float("-inf")
                            ):
                                numeric_values.append(numeric_val)
                        except (ValueError, TypeError):
                            pass

        return numeric_values

    def _get_columns_for_axis_position(
        self,
        root,
        ns: dict,
        axis_position: str,
        value_columns: List[str],
    ) -> List[str]:
        """
        Get column names for an axis based on its position (left/right).

        For combo charts:
        - Left axis (primary): All columns except the last one (bars/area)
        - Right axis (secondary): The last column only (line)

        Args:
            root: XML root element
            ns: XML namespaces dict
            axis_position: 'l' for left/primary, 'r' for right/secondary
            value_columns: List of all value column names

        Returns:
            List of column names for this axis
        """
        if not value_columns:
            return []

        # Count series in different chart types to determine split
        bar_series = root.findall(".//c:barChart//c:ser", namespaces=ns)
        line_series = root.findall(".//c:lineChart//c:ser", namespaces=ns)
        area_series = root.findall(".//c:areaChart//c:ser", namespaces=ns)

        bar_count = len(bar_series)
        line_count = len(line_series)
        area_count = len(area_series)

        # For combo charts: bars/area on primary (left), line on secondary (right)
        if line_count > 0 and (bar_count > 0 or area_count > 0):
            primary_count = bar_count + area_count
            if axis_position == "l":
                # Primary axis: first N columns (where N = bar + area series)
                return value_columns[:primary_count]
            else:
                # Secondary axis: remaining columns (line series)
                return value_columns[primary_count:]

        # Default: all columns for any axis
        return value_columns

    def _enable_auto_scaling(
        self,
        root,
        ns,
        data: Optional[List[Dict[str, Any]]] = None,
        value_columns: Optional[List[str]] = None,
        chart_type: Optional[str] = None,
    ):
        """
        Enable automatic axis scaling by setting min/max values based on data.
        This allows the chart Y-axis to start from near the minimum data value
        instead of always starting from 0.

        For multi-axis charts, each axis is scaled independently based on the
        series data that belongs to that axis.

        For stacked bar charts, calculates the stacked totals per category
        to determine the correct axis max.

        Uses floor for minimum values and ceil for maximum values to ensure
        all data points are visible and axes have nice round numbers.

        Args:
            root: XML root element
            ns: XML namespaces dict
            data: List of data dictionaries containing chart values
            value_columns: List of column names that contain numeric values
            chart_type: The detected chart type (e.g., 'stacked_bar', 'combo_bar_line')
        """
        try:
            # Find all value axes (both primary and secondary)
            value_axes = root.findall(".//c:valAx", namespaces=ns)
            c_ns = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"

            if not value_axes:
                return

            # Determine if this is a multi-axis chart
            is_multi_axis = len(value_axes) > 1

            if is_multi_axis and data and value_columns:
                # Multi-axis chart: scale each axis independently
                print(f"    📊 Multi-axis chart detected: {len(value_axes)} axes")

                # Check if this is a combo stacked bar + line chart
                is_combo_stacked = chart_type == "combo_stacked_bar_line"

                # Collect axis scales for synchronization
                axis_scales: list[tuple[float, float]] = []

                for axis_index, axis in enumerate(value_axes):
                    # Get axis ID
                    ax_id_elem = axis.find("./c:axId", namespaces=ns)
                    axis_id = ax_id_elem.get("val") if ax_id_elem is not None else None

                    # Get axis position for logging and column selection
                    ax_pos_elem = axis.find("./c:axPos", namespaces=ns)
                    axis_pos = ax_pos_elem.get("val") if ax_pos_elem is not None else "?"

                    if axis_id is None:
                        continue

                    # Get values for series belonging to this axis
                    axis_values = self._get_axis_series_data(
                        root, ns, axis_id, data, value_columns
                    )

                    if not axis_values:
                        # Fallback: determine columns based on axis position (left=primary, right=secondary)
                        # For combo charts: bars on primary (left), line on secondary (right)
                        axis_columns = self._get_columns_for_axis_position(
                            root, ns, axis_pos, value_columns
                        )
                        axis_label = "primary" if axis_pos == "l" else "secondary"
                        print(
                            f"    📊 Using position-based column selection for {axis_label} axis: {axis_columns}"
                        )

                        # For combo stacked bar + line, primary axis (left) needs stacked totals
                        if is_combo_stacked and axis_pos == "l" and len(axis_columns) > 1:
                            # Calculate stacked totals per row for the primary axis (stacked bars)
                            # Track positive and negative sums separately for proper axis range
                            print("    📊 Combo stacked bar: calculating stacked totals for primary axis (with negative support)")
                            axis_values = []
                            for row in data:
                                positive_sum = 0.0
                                negative_sum = 0.0
                                for col in axis_columns:
                                    value = row.get(col)
                                    if value is not None:
                                        try:
                                            numeric_val = float(value)
                                            if not (
                                                numeric_val != numeric_val
                                                or numeric_val == float("inf")
                                                or numeric_val == float("-inf")
                                            ):
                                                if numeric_val >= 0:
                                                    positive_sum += numeric_val
                                                else:
                                                    negative_sum += numeric_val
                                        except (ValueError, TypeError):
                                            pass
                                # Add positive stack total for max calculation
                                if positive_sum != 0.0:
                                    axis_values.append(positive_sum)
                                # Add negative stack total for min calculation
                                if negative_sum != 0.0:
                                    axis_values.append(negative_sum)
                        else:
                            # Standard behavior: use individual values
                            axis_values = []
                            for row in data:
                                for col in axis_columns:
                                    value = row.get(col)
                                    if value is not None:
                                        try:
                                            numeric_val = float(value)
                                            if not (
                                                numeric_val != numeric_val
                                                or numeric_val == float("inf")
                                                or numeric_val == float("-inf")
                                            ):
                                                axis_values.append(numeric_val)
                                        except (ValueError, TypeError):
                                            pass

                    if axis_values:
                        scale_result = self._set_axis_scale(
                            axis, ns, c_ns, axis_values, axis_pos, axis_index
                        )
                        axis_scales.append(scale_result)

                # Synchronize divisions for dual-axis charts (exactly 2 axes)
                if len(value_axes) == 2 and len(axis_scales) == 2:
                    self._synchronize_dual_axis_divisions(
                        value_axes, ns, c_ns, axis_scales
                    )

            else:
                # Single-axis chart or no data: calculate global min/max
                all_numeric_values: List[float] = []

                # Check if this is a stacked bar chart type
                is_stacked_chart = chart_type in (
                    "stacked_bar",
                    "stacked_horizontal_bar",
                    "single_column_stacked",
                    "single_column_stacked_chart",
                    "Single_column_stacked_chart",
                )
                # Combo stacked bar + line needs special handling - only sum bar columns, not line
                is_combo_stacked_single_axis = chart_type == "combo_stacked_bar_line"

                if data and value_columns:
                    if is_stacked_chart:
                        # For stacked bar charts, calculate per-row positive and negative sums
                        # The axis should span from min(negative_sums) to max(positive_sums)
                        for row in data:
                            positive_sum = 0.0
                            negative_sum = 0.0
                            for col in value_columns:
                                value = row.get(col)
                                if value is not None:
                                    try:
                                        numeric_val = float(value)
                                        # Skip NaN and infinite values
                                        if not (
                                            numeric_val != numeric_val
                                            or numeric_val == float("inf")
                                            or numeric_val == float("-inf")
                                        ):
                                            if numeric_val >= 0:
                                                positive_sum += numeric_val
                                            else:
                                                negative_sum += numeric_val
                                    except (ValueError, TypeError):
                                        pass
                            # Add positive stack total for max calculation
                            if positive_sum != 0.0:
                                all_numeric_values.append(positive_sum)
                            # Add negative stack total for min calculation
                            if negative_sum != 0.0:
                                all_numeric_values.append(negative_sum)
                        print("    📊 Stacked bar chart: calculating stacked totals per category (with negative support)")
                    elif is_combo_stacked_single_axis and len(value_columns) >= 2:
                        # For combo stacked bar + line in single-axis mode:
                        # Sum all columns EXCEPT the last one (the line series)
                        # Track positive and negative sums separately for proper axis range
                        bar_columns = value_columns[:-1]  # All but last (bar series)
                        print(f"    📊 Combo stacked bar (single-axis): stacking columns {bar_columns} (with negative support)")
                        for row in data:
                            positive_sum = 0.0
                            negative_sum = 0.0
                            for col in bar_columns:
                                value = row.get(col)
                                if value is not None:
                                    try:
                                        numeric_val = float(value)
                                        if not (
                                            numeric_val != numeric_val
                                            or numeric_val == float("inf")
                                            or numeric_val == float("-inf")
                                        ):
                                            if numeric_val >= 0:
                                                positive_sum += numeric_val
                                            else:
                                                negative_sum += numeric_val
                                    except (ValueError, TypeError):
                                        pass
                            # Add positive stack total for max calculation
                            if positive_sum != 0.0:
                                all_numeric_values.append(positive_sum)
                            # Add negative stack total for min calculation
                            if negative_sum != 0.0:
                                all_numeric_values.append(negative_sum)
                    else:
                        # For non-stacked charts, use individual values
                        for row in data:
                            for col in value_columns:
                                value = row.get(col)
                                if value is not None:
                                    try:
                                        numeric_val = float(value)
                                        # Skip NaN and infinite values
                                        if not (
                                            numeric_val != numeric_val
                                            or numeric_val == float("inf")
                                            or numeric_val == float("-inf")
                                        ):
                                            all_numeric_values.append(numeric_val)
                                    except (ValueError, TypeError):
                                        pass

                if all_numeric_values:
                    print(
                        f"    📊 Data range: min={min(all_numeric_values):.2f}, "
                        f"max={max(all_numeric_values):.2f}"
                    )

                for axis_index, axis in enumerate(value_axes):
                    ax_pos_elem = axis.find("./c:axPos", namespaces=ns)
                    axis_pos = ax_pos_elem.get("val") if ax_pos_elem is not None else "?"

                    if all_numeric_values:
                        # For stacked charts, force minimum to 0 since stacking starts from 0
                        self._set_axis_scale(
                            axis, ns, c_ns, all_numeric_values, axis_pos, axis_index,
                            force_min_zero=is_stacked_chart or is_combo_stacked_single_axis
                        )
                    else:
                        # Just ensure orientation is set
                        scaling = axis.find(".//c:scaling", namespaces=ns)
                        if scaling is not None:
                            orientation = scaling.find(".//c:orientation", namespaces=ns)
                            if orientation is None:
                                orientation = etree.SubElement(
                                    scaling, f"{c_ns}orientation"
                                )
                                orientation.set("val", "minMax")

            if value_axes:
                print(
                    f"    ✓ Enabled auto-scaling for {len(value_axes)} value axis(es)"
                )

        except Exception as e:
            print(f"    Warning: Could not enable auto-scaling: {e}")

    def _set_axis_scale(
        self,
        axis,
        ns: dict,
        c_ns: str,
        numeric_values: List[float],
        axis_pos: str = "?",
        axis_index: int = 0,
        force_min_zero: bool = False,
    ) -> tuple[float, float]:
        """
        Set the min/max scale for a single axis based on numeric values.

        Uses floor for min and ceil for max to ensure nice numbers.
        Primary axis (left) allows decimals, secondary axis (right) uses integers.
        If all data is non-negative, minimum starts from 0.
        Adds padding on top of the chart.

        Args:
            axis: The axis XML element
            ns: XML namespaces dict
            c_ns: Chart namespace string
            numeric_values: List of numeric values for this axis
            axis_pos: Axis position ('l' for left, 'r' for right)
            axis_index: Index of the axis for logging
            force_min_zero: If True, always set minimum to 0 (for stacked charts)

        Returns:
            Tuple of (scaled_min, scaled_max) for use in dual-axis synchronization
        """
        if not numeric_values:
            return (0.0, 0.0)

        scaling = axis.find(".//c:scaling", namespaces=ns)
        if scaling is None:
            return (0.0, 0.0)

        # Remove existing min/max elements first
        min_elem = scaling.find(".//c:min", namespaces=ns)
        if min_elem is not None:
            scaling.remove(min_elem)

        max_elem = scaling.find(".//c:max", namespaces=ns)
        if max_elem is not None:
            scaling.remove(max_elem)

        data_min = min(numeric_values)
        data_max = max(numeric_values)

        # Calculate padding (10% of the range, or 10% of the value if range is 0)
        data_range = data_max - data_min
        if data_range == 0:
            # If all values are the same, use 10% of the value
            padding = abs(data_min) * 0.1 if data_min != 0 else 1.0
        else:
            padding = data_range * 0.1

        # Primary axis (left) uses decimals, secondary axis (right) uses integers
        force_integer = (axis_pos == "r")

        # Determine scaled minimum with dynamic scaling
        scaled_min_raw = data_min - padding
        scaled_min = self._round_to_nice_number(
            scaled_min_raw, direction="floor", force_integer=force_integer
        )

        # Ensure min is not greater than data_min
        if scaled_min > data_min:
            scaled_min = self._round_to_nice_number(
                data_min, direction="floor", force_integer=force_integer
            )

        # If all data is non-negative, don't let the axis go negative
        if data_min >= 0 and scaled_min < 0:
            scaled_min = 0 if force_integer else 0.0

        # For stacked charts, only force minimum to 0 if all data is non-negative
        # If data contains negative values, allow the axis to extend into negative range
        if force_min_zero and data_min >= 0:
            scaled_min = 0 if force_integer else 0.0

        # Calculate scaled max with padding
        scaled_max_raw = data_max + padding
        scaled_max = self._round_to_nice_number(
            scaled_max_raw, direction="ceil", force_integer=force_integer
        )

        # Ensure max is not less than data_max
        if scaled_max < data_max:
            scaled_max = self._round_to_nice_number(
                data_max, direction="ceil", force_integer=force_integer
            )

        # Add padding on top of the chart (1 unit for integers, proportional for decimals)
        if force_integer:
            scaled_max = scaled_max + 1
        else:
            # Add ~5% more padding for decimal scales
            scaled_max = scaled_max + (scaled_max - scaled_min) * 0.05

        # Create and add min element
        new_min = etree.SubElement(scaling, f"{c_ns}min")
        new_min.set("val", str(scaled_min))

        # Create and add max element
        new_max = etree.SubElement(scaling, f"{c_ns}max")
        new_max.set("val", str(scaled_max))

        axis_label = "primary" if axis_pos == "l" else (
            "secondary" if axis_pos == "r" else f"axis-{axis_index}"
        )
        number_type = "integers" if force_integer else "decimals"
        print(
            f"    ✓ Set {axis_label} axis scale: "
            f"min={scaled_min}, max={scaled_max} ({number_type}) "
            f"[data: {data_min:.2f}-{data_max:.2f}]"
        )

        # Ensure orientation is set (required element)
        orientation = scaling.find(".//c:orientation", namespaces=ns)
        if orientation is None:
            orientation = etree.SubElement(scaling, f"{c_ns}orientation")
            orientation.set("val", "minMax")

        return (scaled_min, scaled_max)

    def _synchronize_dual_axis_divisions(
        self,
        value_axes: list,
        ns: dict,
        c_ns: str,
        axis_scales: list[tuple[float, float]],
        num_divisions: int = 5,
    ) -> None:
        """
        Synchronize the number of divisions (major units) between dual axes.

        For dual-axis charts, both axes should have the same number of gridlines/ticks
        for visual consistency. This method calculates the majorUnit for each axis
        and adjusts the max value to ensure exactly num_divisions divisions.

        Primary axis (left) uses decimal major units, secondary axis (right) uses integers.

        Args:
            value_axes: List of value axis XML elements
            ns: XML namespaces dict
            c_ns: Chart namespace string
            axis_scales: List of (min, max) tuples for each axis
            num_divisions: Target number of divisions (default: 5)
        """
        if len(value_axes) != 2 or len(axis_scales) != 2:
            return

        try:
            for axis_index, (axis, (axis_min, axis_max)) in enumerate(
                zip(value_axes, axis_scales)
            ):
                # Get axis position to determine if we use integers or decimals
                ax_pos_elem = axis.find("./c:axPos", namespaces=ns)
                axis_pos = ax_pos_elem.get("val") if ax_pos_elem is not None else "?"
                force_integer = (axis_pos == "r")  # Secondary (right) uses integers

                # Calculate the range for this axis
                axis_range = axis_max - axis_min

                if axis_range <= 0:
                    continue

                # Calculate major unit to achieve desired number of divisions
                raw_major_unit = axis_range / num_divisions

                # Round major unit to a nice number (integer for secondary, decimal for primary)
                major_unit = self._round_to_nice_number(
                    raw_major_unit, direction="ceil", force_integer=force_integer
                )

                # Ensure major unit is at least a minimum value
                min_major_unit = 1 if force_integer else 0.1
                if major_unit < min_major_unit:
                    major_unit = min_major_unit

                # Adjust the max value to ensure exactly num_divisions divisions
                adjusted_max = axis_min + (major_unit * num_divisions)

                # Update the axis max in the scaling element
                scaling = axis.find(".//c:scaling", namespaces=ns)
                if scaling is not None:
                    max_elem = scaling.find("./c:max", namespaces=ns)
                    if max_elem is not None:
                        max_elem.set("val", str(adjusted_max))
                    else:
                        max_elem = etree.SubElement(scaling, f"{c_ns}max")
                        max_elem.set("val", str(adjusted_max))

                # Remove existing majorUnit element if present
                existing_major = axis.find("./c:majorUnit", namespaces=ns)
                if existing_major is not None:
                    axis.remove(existing_major)

                # Add majorUnit element
                major_unit_elem = etree.SubElement(axis, f"{c_ns}majorUnit")
                major_unit_elem.set("val", str(major_unit))

                axis_label = "primary" if axis_pos == "l" else "secondary"
                number_type = "integers" if force_integer else "decimals"

                print(
                    f"    ✓ Set {axis_label} axis: majorUnit={major_unit} ({number_type}), "
                    f"max adjusted {axis_max}→{adjusted_max} "
                    f"({num_divisions} divisions, {num_divisions + 1} tick marks)"
                )

        except Exception as e:
            print(f"    Warning: Could not synchronize dual axis divisions: {e}")

    def _update_axis_titles(
        self,
        root,
        ns: dict,
        primary_y_axis_title: Optional[str] = None,
        secondary_y_axis_title: Optional[str] = None,
        x_axis_title: Optional[str] = None,
        chart_type: Optional[str] = None,
    ) -> None:
        """
        Update axis titles in chart XML.

        Args:
            root: XML root element
            ns: XML namespaces dict
            primary_y_axis_title: Title for primary (left) Y-axis
            secondary_y_axis_title: Title for secondary (right) Y-axis
            x_axis_title: Title for X-axis (category axis)
            chart_type: Chart type for group-specific positioning
        """
        # Skip if no titles provided
        if not any([primary_y_axis_title, secondary_y_axis_title, x_axis_title]):
            return

        try:
            from hello.utils.ppt_helpers_utils.services.template_config import (
                get_element_dimensions,
            )

            element_dims = get_element_dimensions()
            max_title_chars = element_dims.chart_axis_title_max_chars

            # Add 'a' namespace for drawing elements if not present
            ns_extended = ns.copy()
            ns_extended["a"] = "http://schemas.openxmlformats.org/drawingml/2006/main"

            # Update Y-axis titles (value axes)
            value_axes = root.findall(".//c:valAx", namespaces=ns)

            # Identify axes by position: primary (left) vs secondary (right)
            left_axis = None
            right_axis = None

            for axis in value_axes:
                # Check axis position: c:axPos val="l" (left) or "r" (right)
                ax_pos_elem = axis.find("./c:axPos", namespaces=ns)
                if ax_pos_elem is not None:
                    pos = ax_pos_elem.get("val", "")
                    if pos == "l":
                        left_axis = axis
                    elif pos == "r":
                        right_axis = axis

            # Fallback: if no position found, use axId ordering (lower = primary/left)
            if left_axis is None and right_axis is None and len(value_axes) >= 1:
                axes_with_id = []
                for axis in value_axes:
                    ax_id_elem = axis.find("./c:axId", namespaces=ns)
                    ax_id = (
                        int(ax_id_elem.get("val", "0")) if ax_id_elem is not None else 0
                    )
                    axes_with_id.append((ax_id, axis))
                axes_with_id.sort(key=lambda x: x[0])
                left_axis = axes_with_id[0][1]
                if len(axes_with_id) > 1:
                    right_axis = axes_with_id[1][1]

            # Update primary (left) Y-axis title
            if primary_y_axis_title and left_axis is not None:
                self._set_axis_title(
                    left_axis, primary_y_axis_title, ns_extended, is_secondary=False, chart_type=chart_type
                )
                # Ensure axis is positioned on left
                ax_pos = left_axis.find("./c:axPos", namespaces=ns)
                if ax_pos is not None:
                    ax_pos.set("val", "l")
                truncated_title = (
                    primary_y_axis_title[:max_title_chars]
                    if len(primary_y_axis_title) > max_title_chars
                    else primary_y_axis_title
                )
                print(f"    ✓ Updated primary Y-axis (left) title: '{truncated_title}'")

            # Update secondary (right) Y-axis title
            if secondary_y_axis_title and right_axis is not None:
                self._set_axis_title(
                    right_axis, secondary_y_axis_title, ns_extended, is_secondary=True, chart_type=chart_type
                )
                # Ensure axis is positioned on right
                ax_pos = right_axis.find("./c:axPos", namespaces=ns)
                if ax_pos is not None:
                    ax_pos.set("val", "r")
                truncated_title = (
                    secondary_y_axis_title[:max_title_chars]
                    if len(secondary_y_axis_title) > max_title_chars
                    else secondary_y_axis_title
                )
                print(
                    f"    ✓ Updated secondary Y-axis (right) title: '{truncated_title}'"
                )

            # Update X-axis title (category axis)
            if x_axis_title:
                cat_axis = root.find(".//c:catAx", namespaces=ns)
                if cat_axis is not None:
                    self._set_axis_title(cat_axis, x_axis_title, ns_extended)
                    truncated_title = (
                        x_axis_title[:max_title_chars]
                        if len(x_axis_title) > max_title_chars
                        else x_axis_title
                    )
                    print(f"    ✓ Updated X-axis title: '{truncated_title}'")
                else:
                    # Some charts use dateAx instead of catAx
                    date_axis = root.find(".//c:dateAx", namespaces=ns)
                    if date_axis is not None:
                        self._set_axis_title(date_axis, x_axis_title, ns_extended)
                        truncated_title = (
                            x_axis_title[:max_title_chars]
                            if len(x_axis_title) > max_title_chars
                            else x_axis_title
                        )
                        print(f"    ✓ Updated X-axis (date) title: '{truncated_title}'")

        except Exception as e:
            print(f"    Warning: Could not update axis titles: {e}")

    def _set_axis_title(
        self, axis, title_text: str, ns: dict, is_secondary: bool = False, chart_type: Optional[str] = None
    ) -> None:
        """
        Set or update the title text for an axis.

        Args:
            axis: XML axis element (c:valAx, c:catAx, or c:dateAx)
            title_text: New title text (will be truncated using config-driven max chars)
            ns: XML namespaces dict
            is_secondary: Whether this is a secondary (right) Y-axis
            chart_type: Chart type for group-specific positioning
        """
        axis_type = "secondary" if is_secondary else "primary"
        print(
            f"    [DEBUG] _set_axis_title: Setting {axis_type} axis title to '{title_text}'"
        )
        from hello.utils.ppt_helpers_utils.services.template_config import (
            get_element_dimensions,
        )

        element_dims = get_element_dimensions()
        max_title_chars = element_dims.chart_axis_title_max_chars

        # Truncate title to config-driven max characters
        if title_text and len(title_text) > max_title_chars:
            title_text = title_text[:max_title_chars]

        # Determine axis type from element tag
        axis_tag = axis.tag.split("}")[-1] if "}" in axis.tag else axis.tag
        is_y_axis = axis_tag == "valAx"

        # Find existing title element
        title_elem = axis.find("./c:title", namespaces=ns)

        # DEBUG: Check existing title and layout
        if title_elem is not None:
            existing_layout = title_elem.find("./c:layout", namespaces=ns)
            if existing_layout is not None:
                manual = existing_layout.find("./c:manualLayout", namespaces=ns)
                if manual is not None:
                    x_elem = manual.find("./c:x", namespaces=ns)
                    y_elem = manual.find("./c:y", namespaces=ns)
                    x_val = x_elem.get("val") if x_elem is not None else "N/A"
                    y_val = y_elem.get("val") if y_elem is not None else "N/A"
                    print(
                        f"    [DEBUG] Existing {axis_type} axis title layout: x={x_val}, y={y_val}"
                    )
                else:
                    print(
                        f"    [DEBUG] Existing {axis_type} axis title has layout but no manualLayout"
                    )
            else:
                print(
                    f"    [DEBUG] Existing {axis_type} axis title has NO layout element"
                )

            # Get existing title text
            existing_text = title_elem.find(".//a:t", namespaces=ns)
            if existing_text is not None:
                print(
                    f"    [DEBUG] Existing {axis_type} axis title text: '{existing_text.text}'"
                )

        if title_elem is not None:
            # Track if we successfully updated the text
            text_updated = False
            
            # Find the rich text container and simplify it to just one text run
            rich_elem = title_elem.find(".//c:rich", namespaces=ns)
            if rich_elem is not None:
                # Find all paragraphs
                paragraphs = rich_elem.findall("./a:p", namespaces=ns)
                if paragraphs:
                    # Keep only the first paragraph, remove others
                    for p in paragraphs[1:]:
                        rich_elem.remove(p)

                    # In the first paragraph, find all runs
                    first_p = paragraphs[0]
                    runs = first_p.findall("./a:r", namespaces=ns)
                    if runs:
                        # Update first run's text ONLY - preserve template's font properties
                        first_run = runs[0]
                        text_elem = first_run.find("./a:t", namespaces=ns)
                        if text_elem is not None:
                            text_elem.text = title_text
                            text_updated = True
                        # DON'T modify font properties - preserve template styling
                        # This prevents PowerPoint from recalculating title position
                        # Remove extra runs if any
                        for run in runs[1:]:
                            first_p.remove(run)
                    else:
                        # No runs found, try to find text element directly
                        text_elem = first_p.find(".//a:t", namespaces=ns)
                        if text_elem is not None:
                            text_elem.text = title_text
                            text_updated = True
            
            # If rich_elem not found or text not updated, try alternate structures
            if not text_updated:
                # Try alternate structure: c:title/c:tx/c:strRef/c:strCache/c:pt/c:v
                v_elem = title_elem.find(".//c:v", namespaces=ns)
                if v_elem is not None:
                    v_elem.text = title_text
                    text_updated = True
                else:
                    # Fallback: find any text element
                    text_elem = title_elem.find(".//a:t", namespaces=ns)
                    if text_elem is not None:
                        text_elem.text = title_text
                        text_updated = True
                    else:
                        # Create rich text structure if nothing exists
                        # _create_title_structure now preserves the c:layout element
                        self._create_title_structure(
                            title_elem, title_text, ns, is_y_axis, is_secondary, chart_type=chart_type
                        )
                        text_updated = True

            # PRESERVE the template's existing layout for axis titles
            # The template has them positioned correctly - don't remove or modify
            pass
        else:
            # Create new title element
            title_elem = etree.SubElement(
                axis, "{http://schemas.openxmlformats.org/drawingml/2006/chart}title"
            )
            self._create_title_structure(
                title_elem, title_text, ns, is_y_axis, is_secondary, chart_type=chart_type
            )

    def _add_axis_title_layout(
        self, title_elem, ns: dict, is_secondary: bool = False, chart_type: Optional[str] = None
    ) -> None:
        """
        Add manual layout positioning to an axis title to place it at the plot area top.

        Both primary and secondary Y-axis titles use the same Y position to ensure
        they are aligned on the same horizontal line (important for dual-axis charts).

        Args:
            title_elem: c:title element
            ns: XML namespaces dict
            is_secondary: Whether this is a secondary (right) Y-axis
            chart_type: Chart type for group-specific positioning
        """
        from hello.utils.ppt_helpers_utils.services.template_config import (
            get_element_dimensions,
            get_chart_layout_config,
        )

        element_dims = get_element_dimensions()
        
        # Get chart-specific layout config for positioning
        chart_config = get_chart_layout_config(chart_type)

        c_ns = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"

        # Remove existing layout if present
        existing_layout = title_elem.find("./c:layout", namespaces=ns)
        if existing_layout is not None:
            title_elem.remove(existing_layout)

        # Create layout with manual positioning
        layout = etree.SubElement(title_elem, f"{c_ns}layout")
        manual = etree.SubElement(layout, f"{c_ns}manualLayout")

        # Set positioning modes to 'edge' (relative to chart edge)
        x_mode = etree.SubElement(manual, f"{c_ns}xMode")
        x_mode.set("val", "edge")
        y_mode = etree.SubElement(manual, f"{c_ns}yMode")
        y_mode.set("val", "edge")

        # Position the title ABOVE the plot area using chart-specific config
        # X position comes from chart group config (left or right edge)
        if is_secondary:
            x_pos = chart_config.secondary_y_axis_title_x
        else:
            x_pos = chart_config.primary_y_axis_title_x

        # Y position from chart group config (both titles use same Y for alignment)
        y_pos = chart_config.y_axis_title_y

        x_elem = etree.SubElement(manual, f"{c_ns}x")
        x_elem.set("val", f"{x_pos:.4f}")
        y_elem = etree.SubElement(manual, f"{c_ns}y")
        y_elem.set("val", f"{y_pos:.4f}")

    def _create_title_structure(
        self,
        title_elem,
        title_text: str,
        ns: dict,
        is_y_axis: bool = True,
        is_secondary: bool = False,
        chart_type: Optional[str] = None,
    ) -> None:
        """
        Create the rich text structure for an axis title.

        Structure: c:title/c:tx/c:rich/a:p/a:r/a:rPr/a:t

        Args:
            title_elem: c:title element
            title_text: Title text
            ns: XML namespaces dict
            is_y_axis: Whether this is a Y-axis (value axis)
            is_secondary: Whether this is a secondary (right) Y-axis
            chart_type: Chart type for group-specific positioning
        """
        from hello.utils.ppt_helpers_utils.services.template_config import (
            get_element_dimensions,
            get_chart_layout_config,
        )
        import copy

        element_dims = get_element_dimensions()

        c_ns = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
        a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

        # IMPORTANT: Preserve the existing c:layout element before clearing content
        # The layout element controls the position of the title (horizontal at top vs vertical on side)
        # Without it, PowerPoint defaults to vertical positioning for Y-axis titles
        ns_for_search = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
                         "a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        
        existing_layout = title_elem.find("./c:layout", namespaces=ns_for_search)
        layout_copy = None
        if existing_layout is not None:
            layout_copy = copy.deepcopy(existing_layout)
            print("    [DEBUG] Preserving existing axis title layout element")
        else:
            print("    [DEBUG] No existing c:layout element found in title")
        
        # Also preserve the bodyPr rotation if it exists (controls horizontal vs vertical text)
        existing_bodyPr = title_elem.find(".//a:bodyPr", namespaces=ns_for_search)
        bodyPr_rot = None
        if existing_bodyPr is not None:
            bodyPr_rot = existing_bodyPr.get("rot")
            print(f"    [DEBUG] Existing bodyPr rotation: {bodyPr_rot}")
        
        # Clear existing content
        for child in list(title_elem):
            title_elem.remove(child)

        # Create structure: c:tx/c:rich/a:bodyPr, a:lstStyle, a:p/a:r/a:t
        tx = etree.SubElement(title_elem, f"{c_ns}tx")
        rich = etree.SubElement(tx, f"{c_ns}rich")

        # Add body properties - preserve rotation if it was horizontal (rot="0" or None for horizontal Y-axis title)
        # Default Y-axis title rotation is -5400000 (vertical), so we need to set rot="0" for horizontal
        bodyPr = etree.SubElement(rich, f"{a_ns}bodyPr")
        if bodyPr_rot is not None:
            # Preserve existing rotation
            bodyPr.set("rot", bodyPr_rot)
            print(f"    [DEBUG] Preserved bodyPr rotation: {bodyPr_rot}")
        elif is_y_axis:
            # For Y-axis titles, set horizontal rotation (rot="0") to avoid default vertical
            # This ensures the title appears horizontally at the top instead of vertically on the side
            bodyPr.set("rot", "0")
            print("    [DEBUG] Set bodyPr rotation to 0 (horizontal) for Y-axis title")
        
        etree.SubElement(rich, f"{a_ns}lstStyle")

        # Create paragraph with text run
        p = etree.SubElement(rich, f"{a_ns}p")
        r = etree.SubElement(p, f"{a_ns}r")

        # Add run properties with font size and Calibri font
        rPr = etree.SubElement(r, f"{a_ns}rPr")
        rPr.set(
            "sz", str(element_dims.chart_axis_title_font_size)
        )  # Font size in hundredths of a point (9pt = 900)
        rPr.set("b", "0")  # Not bold

        # Set font to Calibri
        latin = etree.SubElement(rPr, f"{a_ns}latin")
        latin.set("typeface", "Calibre (Body)")

        t = etree.SubElement(r, f"{a_ns}t")
        t.text = title_text

        # Restore the preserved layout element OR create a new one for horizontal positioning
        # This is critical for stacked bar charts where the Y-axis title should appear
        # horizontally at the top instead of vertically on the side
        if layout_copy is not None:
            # Restore preserved layout - insert BEFORE overlay (correct OOXML order: tx, layout, overlay)
            title_elem.append(layout_copy)
            print("    [DEBUG] Restored axis title layout element")
        elif is_y_axis:
            # Get chart-specific layout config for positioning
            chart_config = get_chart_layout_config(chart_type)
            
            # Create a layout element to position Y-axis title horizontally at top
            # Without explicit layout, PowerPoint defaults to vertical positioning
            layout = etree.SubElement(title_elem, f"{c_ns}layout")
            manual_layout = etree.SubElement(layout, f"{c_ns}manualLayout")
            
            # Set positioning mode to "edge" (relative to chart edge)
            xMode = etree.SubElement(manual_layout, f"{c_ns}xMode")
            xMode.set("val", "edge")
            yMode = etree.SubElement(manual_layout, f"{c_ns}yMode")
            yMode.set("val", "edge")
            
            # Position using chart-specific config (top-left for primary, top-right for secondary)
            x_pos = chart_config.secondary_y_axis_title_x if is_secondary else chart_config.primary_y_axis_title_x
            y_pos = chart_config.y_axis_title_y
            
            x_elem = etree.SubElement(manual_layout, f"{c_ns}x")
            x_elem.set("val", str(x_pos))
            y_elem = etree.SubElement(manual_layout, f"{c_ns}y")
            y_elem.set("val", str(y_pos))
            
            print(f"    [DEBUG] Created new layout for Y-axis title at x={x_pos}, y={y_pos}")

        # Add overlay element (required by PowerPoint)
        overlay = etree.SubElement(title_elem, f"{c_ns}overlay")
        overlay.set("val", "0")

    def _remove_axis_titles(self, root, ns: dict) -> None:
        """
        Remove all axis titles from chart XML.
        Used for horizontal bar charts where axis titles don't fit well due to swapped axes.

        Args:
            root: XML root element
            ns: XML namespaces dict
        """
        try:
            # Find all axis types and remove their titles
            axis_types = ["c:valAx", "c:catAx", "c:dateAx"]
            removed_count = 0

            for axis_type in axis_types:
                axes = root.findall(f".//{axis_type}", namespaces=ns)
                for axis in axes:
                    title_elem = axis.find("./c:title", namespaces=ns)
                    if title_elem is not None:
                        axis.remove(title_elem)
                        removed_count += 1

            if removed_count > 0:
                print(f"    ✓ Removed {removed_count} axis title(s)")
        except Exception as e:
            print(f"    Warning: Could not remove axis titles: {e}")

    def _remove_secondary_axis(self, root, ns: dict) -> Optional[str]:
        """
        Handle single-axis mode by hiding the secondary Y-axis and redirecting all series.
        Used when isMultiAxis=false to ensure all series plot on primary axis.

        Instead of deleting the secondary axis (which can corrupt complex combo charts),
        this method:
        - Hides the secondary axis (c:delete val="1")
        - Removes axis title and labels
        - Updates all series to reference the primary axis

        Args:
            root: XML root element
            ns: XML namespaces dict

        Returns:
            Primary axis ID string if found, None otherwise
        """
        try:
            # Find all value axes
            value_axes = root.findall(".//c:valAx", namespaces=ns)

            # Safety check: if only one value axis, nothing to do
            if len(value_axes) <= 1:
                print(
                    f"    ℹ️  Only {len(value_axes)} value axis found, nothing to remove"
                )
                if value_axes:
                    ax_id_elem = value_axes[0].find("./c:axId", namespaces=ns)
                    return ax_id_elem.get("val") if ax_id_elem is not None else None
                return None

            primary_axis_id: Optional[str] = None
            secondary_axis_id: Optional[str] = None
            secondary_axis_elem = None

            for axis in value_axes:
                # Get axis position
                ax_pos_elem = axis.find("./c:axPos", namespaces=ns)
                ax_id_elem = axis.find("./c:axId", namespaces=ns)

                if ax_id_elem is None:
                    continue

                pos = ax_pos_elem.get("val", "l") if ax_pos_elem is not None else "l"
                axis_id = ax_id_elem.get("val", "")

                if pos == "l":  # Left = primary
                    primary_axis_id = axis_id
                elif pos == "r":  # Right = secondary
                    secondary_axis_id = axis_id
                    secondary_axis_elem = axis

            # Fallback: if no position found or no primary identified, use axId ordering (lower = primary)
            if primary_axis_id is None or secondary_axis_id is None:
                axes_with_id = []
                for axis in value_axes:
                    ax_id_elem = axis.find("./c:axId", namespaces=ns)
                    ax_id = (
                        int(ax_id_elem.get("val", "0")) if ax_id_elem is not None else 0
                    )
                    axes_with_id.append((ax_id, axis))
                axes_with_id.sort(key=lambda x: x[0])

                if len(axes_with_id) >= 1 and primary_axis_id is None:
                    primary_axis_id = str(axes_with_id[0][0])
                if len(axes_with_id) >= 2 and secondary_axis_id is None:
                    secondary_axis_id = str(axes_with_id[1][0])
                    secondary_axis_elem = axes_with_id[1][1]

            if secondary_axis_elem is None or secondary_axis_id is None:
                print("    ℹ️  No secondary axis found to hide")
                return primary_axis_id

            if primary_axis_id is None:
                print("    ⚠️  No primary axis found, cannot proceed")
                return None

            print(
                f"    🔧 Hiding secondary axis (id={secondary_axis_id}), keeping primary (id={primary_axis_id})"
            )

            # Step 1: Update ALL series axId references to point to primary axis
            # This ensures all data series plot on the primary (left) axis
            series_list = root.findall(".//c:ser", namespaces=ns)
            updated_series = 0
            for series in series_list:
                ax_id_refs = series.findall(".//c:axId", namespaces=ns)
                for ax_id_ref in ax_id_refs:
                    if ax_id_ref.get("val") == secondary_axis_id:
                        ax_id_ref.set("val", primary_axis_id)
                        updated_series += 1

            if updated_series > 0:
                print(
                    f"    ✓ Redirected {updated_series} series reference(s) to primary axis"
                )

            # Step 2: Hide the secondary axis instead of deleting it
            # This preserves the XML structure while making the axis invisible
            c_ns = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"

            # Set c:delete to hide the axis
            delete_elem = secondary_axis_elem.find("./c:delete", namespaces=ns)
            if delete_elem is None:
                delete_elem = etree.SubElement(secondary_axis_elem, f"{c_ns}delete")
            delete_elem.set("val", "1")

            # Remove axis title if present
            title_elem = secondary_axis_elem.find("./c:title", namespaces=ns)
            if title_elem is not None:
                secondary_axis_elem.remove(title_elem)

            # Hide tick labels
            tick_lbl_pos = secondary_axis_elem.find("./c:tickLblPos", namespaces=ns)
            if tick_lbl_pos is None:
                tick_lbl_pos = etree.SubElement(
                    secondary_axis_elem, f"{c_ns}tickLblPos"
                )
            tick_lbl_pos.set("val", "none")

            # Hide major/minor tick marks
            major_tick = secondary_axis_elem.find("./c:majorTickMark", namespaces=ns)
            if major_tick is None:
                major_tick = etree.SubElement(
                    secondary_axis_elem, f"{c_ns}majorTickMark"
                )
            major_tick.set("val", "none")

            minor_tick = secondary_axis_elem.find("./c:minorTickMark", namespaces=ns)
            if minor_tick is None:
                minor_tick = etree.SubElement(
                    secondary_axis_elem, f"{c_ns}minorTickMark"
                )
            minor_tick.set("val", "none")

            print("    ✓ Hidden secondary Y-axis (labels, ticks, title removed)")

            return primary_axis_id

        except Exception as e:
            print(f"    Warning: Could not hide secondary axis: {e}")
            import traceback

            traceback.print_exc()
            return None

    def _update_embedded_excel(
        self,
        chart_path: str,
        extract_dir: str,
        categories: List[str],
        series_data: List[tuple],  # List of (series_name, values_list) tuples
    ) -> bool:
        """
        Update the embedded Excel workbook with chart data by directly manipulating
        the XLSX internal XML files. This preserves PowerPoint-specific structures
        that openpyxl would otherwise corrupt.

        This ensures that when users right-click the chart and choose
        "Edit Data in Excel", they see the actual chart data instead of
        the template data.

        Args:
            chart_path: Path to the chart XML file
            extract_dir: Path to the extracted PPTX directory
            categories: List of category labels (X-axis values)
            series_data: List of tuples (series_name, values_list) for each series

        Returns:
            True if successful, False otherwise
        """
        try:
            # Get configuration for Excel data layout
            from hello.utils.ppt_helpers_utils.services.template_config import (
                get_element_dimensions,
            )

            element_dims = get_element_dimensions()

            header_row = element_dims.excel_header_row
            data_start_row = element_dims.excel_data_start_row
            category_col = element_dims.excel_category_column
            series_start_col = element_dims.excel_series_start_column

            # Find the embedded Excel file using the existing method
            excel_path = self._get_excel_file(chart_path, extract_dir)

            if excel_path is None:
                print(
                    "    ⚠️  No embedded Excel file found - chart may use external data source"
                )
                return False

            if not os.path.exists(excel_path):
                print(f"    ⚠️  Excel file path does not exist: {excel_path}")
                return False

            # Update the Excel file by directly modifying its XML
            sheet_name = self._get_chart_sheet_name(chart_path)
            success = self._update_excel_xml_directly(
                excel_path,
                categories,
                series_data,
                header_row,
                data_start_row,
                category_col,
                series_start_col,
                sheet_name=sheet_name,
            )

            if success:
                print(
                    f"    ✓ Updated embedded Excel: {len(categories)} rows, {len(series_data)} series"
                )
            return success

        except Exception as e:
            print(f"    ⚠️  Error updating embedded Excel: {e}")
            import traceback

            traceback.print_exc()
            return False

    def _update_excel_xml_directly(
        self,
        excel_path: str,
        categories: List[str],
        series_data: List[tuple],
        header_row: int,
        data_start_row: int,
        category_col: int,
        series_start_col: int,
        sheet_name: Optional[str] = None,
    ) -> bool:
        """
        Update Excel file by modifying specific files within the ZIP archive
        without fully extracting and repackaging. This preserves the original
        ZIP structure that PowerPoint expects.

        Args:
            excel_path: Path to the XLSX file
            categories: List of category labels (X-axis values)
            series_data: List of (series_name, values) tuples
            header_row: Row number for headers (1-based)
            data_start_row: Row number where data starts (1-based)
            category_col: Column number for categories (1-based, A=1)
            series_start_col: Column number where series data starts (1-based, B=2)

        Returns:
            True if successful, False otherwise
        """
        try:
            xlsx_main_ns = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"

            # Read the original ZIP file contents
            original_files = {}
            with zipfile.ZipFile(excel_path, "r") as xlsx_zip:
                for name in xlsx_zip.namelist():
                    original_files[name] = xlsx_zip.read(name)
            had_shared_strings = "xl/sharedStrings.xml" in original_files

            # Resolve which worksheet actually holds chart data
            from hello.utils.ppt_helpers_utils.services.template_config import (
                get_element_dimensions,
            )

            element_dims = get_element_dimensions()
            target_sheet_name = sheet_name or element_dims.excel_default_sheet_name
            worksheet_part = self._resolve_xlsx_worksheet_path(
                original_files, target_sheet_name
            )

            # Parse worksheet XML
            sheet_xml = original_files.get(worksheet_part)
            if not sheet_xml:
                print(f"    ⚠️  Worksheet not found in Excel: {worksheet_part}")
                return False

            parser = etree.XMLParser(remove_blank_text=False)
            root = etree.fromstring(sheet_xml, parser)

            sheet_data = root.find(f".//{{{xlsx_main_ns}}}sheetData")
            if sheet_data is None:
                print("    ⚠️  sheetData element not found")
                return False

            # Build maps of existing cells and rows
            existing_cells = {}
            existing_rows = {}
            for row_elem in sheet_data.findall(f"{{{xlsx_main_ns}}}row"):
                row_num = row_elem.get("r")
                if row_num:
                    existing_rows[int(row_num)] = row_elem
                    for cell_elem in row_elem.findall(f"{{{xlsx_main_ns}}}c"):
                        cell_ref = cell_elem.get("r")
                        if cell_ref:
                            existing_cells[cell_ref] = cell_elem

            # Capture original template max row BEFORE we add any new data rows
            # This is needed to detect when we have MORE rows than the template
            original_template_max_row = max(existing_rows.keys()) if existing_rows else 1

            # Read existing sharedStrings
            existing_strings: List[str] = []
            string_to_index: Dict[str, int] = {}
            ss_xml = original_files.get("xl/sharedStrings.xml")
            if ss_xml:
                ss_root = etree.fromstring(ss_xml)
                for si_elem in ss_root.findall(f"{{{xlsx_main_ns}}}si"):
                    t_elem = si_elem.find(f"{{{xlsx_main_ns}}}t")
                    text = t_elem.text if t_elem is not None and t_elem.text else ""
                    string_to_index[text] = len(existing_strings)
                    existing_strings.append(text)
            
            # Capture original shared strings BEFORE we modify them (for formula updates)
            original_shared_strings = list(existing_strings)

            def get_string_index(text: str) -> int:
                if text not in string_to_index:
                    string_to_index[text] = len(existing_strings)
                    existing_strings.append(text)
                return string_to_index[text]

            def col_to_letter(col_num: int) -> str:
                result = ""
                while col_num > 0:
                    col_num, remainder = divmod(col_num - 1, 26)
                    result = chr(65 + remainder) + result
                return result

            def update_or_create_cell(
                row_num: int, col_num: int, value, is_string: bool = False
            ):
                cell_ref = f"{col_to_letter(col_num)}{row_num}"

                if row_num in existing_rows:
                    row_elem = existing_rows[row_num]
                else:
                    row_elem = etree.SubElement(sheet_data, f"{{{xlsx_main_ns}}}row")
                    row_elem.set("r", str(row_num))
                    existing_rows[row_num] = row_elem

                if cell_ref in existing_cells:
                    cell_elem = existing_cells[cell_ref]
                else:
                    cell_elem = etree.SubElement(row_elem, f"{{{xlsx_main_ns}}}c")
                    cell_elem.set("r", cell_ref)
                    existing_cells[cell_ref] = cell_elem

                # Remove existing value/formula
                for child in list(cell_elem):
                    if etree.QName(child.tag).localname in ("v", "f"):
                        cell_elem.remove(child)

                if is_string:
                    cell_elem.set("t", "s")
                    v_elem = etree.SubElement(cell_elem, f"{{{xlsx_main_ns}}}v")
                    v_elem.text = str(get_string_index(str(value)))
                elif isinstance(value, (int, float)) and value is not None:
                    if "t" in cell_elem.attrib:
                        del cell_elem.attrib["t"]
                    v_elem = etree.SubElement(cell_elem, f"{{{xlsx_main_ns}}}v")
                    v_elem.text = str(value)
                elif value is not None:
                    cell_elem.set("t", "s")
                    v_elem = etree.SubElement(cell_elem, f"{{{xlsx_main_ns}}}v")
                    v_elem.text = str(get_string_index(str(value)))

            # Capture original header names BEFORE updating (for formula reference updates)
            # This is needed when templates have formulas that reference column names.
            original_headers: Dict[str, str] = {}
            new_headers = [" "] + [name for name, _ in series_data]
            for col_offset in range(len(new_headers)):
                col_num = category_col + col_offset
                cell_ref = f"{col_to_letter(col_num)}{header_row}"
                if cell_ref in existing_cells:
                    cell = existing_cells[cell_ref]
                    v_elem = cell.find(f"{{{xlsx_main_ns}}}v")
                    if v_elem is not None and v_elem.text:
                        if cell.get("t") == "s":
                            # String index - look up original string
                            try:
                                orig_idx = int(v_elem.text)
                                if orig_idx < len(original_shared_strings):
                                    orig_name = original_shared_strings[orig_idx]
                                    if col_offset < len(new_headers) and orig_name != new_headers[col_offset]:
                                        original_headers[orig_name] = new_headers[col_offset]
                            except (ValueError, IndexError):
                                pass

            # Update header row
            update_or_create_cell(header_row, category_col, " ", is_string=True)
            for col_offset, (series_name, _) in enumerate(series_data):
                update_or_create_cell(
                    header_row,
                    series_start_col + col_offset,
                    series_name,
                    is_string=True,
                )

            # Update data rows
            for row_offset, category in enumerate(categories):
                current_row_num = data_start_row + row_offset
                update_or_create_cell(
                    current_row_num, category_col, category, is_string=True
                )
                for col_offset, (_, values) in enumerate(series_data):
                    if row_offset < len(values):
                        val = values[row_offset]
                        is_numeric = isinstance(val, (int, float))
                        update_or_create_cell(
                            current_row_num,
                            series_start_col + col_offset,
                            val,
                            is_string=not is_numeric,
                        )

            # Clear excess data cells beyond our new data range.
            # This ensures the xlsx data matches what the chart formulas reference.
            # We don't delete rows (which could break tables), but we clear values
            # AND formulas in cells that would otherwise contain old template data.
            last_data_row = data_start_row + len(categories) - 1
            num_series_cols = len(series_data)

            # Find the max row in existing data (to know how many excess rows to clear)
            max_existing_row = max(existing_rows.keys()) if existing_rows else last_data_row

            # =========================================================================
            # GENERIC FORMULA COLUMN HANDLING
            # =========================================================================
            # Templates may have formula columns (e.g., SUM, AVERAGE, calculated fields)
            # that exist outside our data columns. These need special handling:
            #   1. When we have FEWER rows than template: clear formulas from excess rows
            #   2. When we have MORE rows than template: copy formula pattern to new rows
            #   3. When column names change: update formula references (handled later)
            #
            # This logic is generic and works for any chart template with formula columns.
            # =========================================================================
            
            formula_columns: Dict[int, str] = {}  # col_num -> formula template
            last_series_col = series_start_col + num_series_cols - 1
            
            # Helper to parse cell reference into (col_letter, row_num)
            import re
            def parse_cell_ref(ref: str) -> tuple:
                match = re.match(r'^([A-Z]+)(\d+)$', ref)
                if match:
                    return match.group(1), int(match.group(2))
                return None, None
            
            # Helper to convert column letter to number
            def letter_to_col(col_letter: str) -> int:
                col_num = 0
                for i, c in enumerate(reversed(col_letter)):
                    col_num += (ord(c.upper()) - ord('A') + 1) * (26 ** i)
                return col_num
            
            # Scan first data row to detect ALL formula columns outside our data range.
            # This includes columns BEFORE category column (rare) and AFTER series columns.
            # We only care about formula columns we're NOT overwriting with data.
            data_columns = set([category_col] + [series_start_col + i for i in range(num_series_cols)])
            
            for cell_ref, cell in existing_cells.items():
                col_letter, row_num = parse_cell_ref(cell_ref)
                if col_letter is None or row_num != data_start_row:
                    continue
                
                col_num = letter_to_col(col_letter)
                
                # Skip columns we're writing data to
                if col_num in data_columns:
                    continue
                
                # Check if this cell has a formula
                formula_elem = cell.find(f"{{{xlsx_main_ns}}}f")
                if formula_elem is not None and formula_elem.text:
                    formula_columns[col_num] = formula_elem.text

            for excess_row_num in range(last_data_row + 1, max_existing_row + 1):
                # Clear category column
                cat_ref = f"{col_to_letter(category_col)}{excess_row_num}"
                if cat_ref in existing_cells:
                    cell = existing_cells[cat_ref]
                    # Remove value and formula elements
                    for child in list(cell):
                        if etree.QName(child.tag).localname in ("v", "f"):
                            cell.remove(child)
                    if "t" in cell.attrib:
                        del cell.attrib["t"]

                # Clear series columns
                for col_offset in range(num_series_cols):
                    series_col = series_start_col + col_offset
                    cell_ref = f"{col_to_letter(series_col)}{excess_row_num}"
                    if cell_ref in existing_cells:
                        cell = existing_cells[cell_ref]
                        # Remove value and formula elements
                        for child in list(cell):
                            if etree.QName(child.tag).localname in ("v", "f"):
                                cell.remove(child)
                        if "t" in cell.attrib:
                            del cell.attrib["t"]
                
                # Clear formula columns (e.g., column E with SUM formulas)
                for formula_col in formula_columns.keys():
                    cell_ref = f"{col_to_letter(formula_col)}{excess_row_num}"
                    if cell_ref in existing_cells:
                        cell = existing_cells[cell_ref]
                        # Remove value and formula elements
                        for child in list(cell):
                            if etree.QName(child.tag).localname in ("v", "f"):
                                cell.remove(child)
                        if "t" in cell.attrib:
                            del cell.attrib["t"]

            # For MORE rows case: Add formula columns to new rows beyond the template
            # (rows that didn't exist in the template but we're adding data to)
            # Use original_template_max_row (captured before we added data) to detect this case
            if formula_columns and last_data_row > original_template_max_row:
                for new_row_num in range(original_template_max_row + 1, last_data_row + 1):
                    for formula_col, formula_template in formula_columns.items():
                        cell_ref = f"{col_to_letter(formula_col)}{new_row_num}"
                        
                        # Ensure row exists
                        if new_row_num not in existing_rows:
                            row_elem = etree.SubElement(sheet_data, f"{{{xlsx_main_ns}}}row")
                            row_elem.set("r", str(new_row_num))
                            existing_rows[new_row_num] = row_elem
                        
                        # Create cell with formula
                        row_elem = existing_rows[new_row_num]
                        cell_elem = etree.SubElement(row_elem, f"{{{xlsx_main_ns}}}c")
                        cell_elem.set("r", cell_ref)
                        
                        # Add formula element
                        formula_elem = etree.SubElement(cell_elem, f"{{{xlsx_main_ns}}}f")
                        formula_elem.text = formula_template
                        
                        existing_cells[cell_ref] = cell_elem

            # IMPORTANT: Do NOT modify worksheet dimension or table ref.
            # PowerPoint embedded xlsx files are sensitive to these references.
            # Keeping them at their original values (e.g., A1:D20) even when we have
            # fewer data rows prevents the "linked file isn't available" error.
            # Empty rows beyond our data will be handled by clearing cell values.

            # =========================================================================
            # GENERIC FORMULA REFERENCE UPDATES
            # =========================================================================
            # When column names change (e.g., "Class A" -> "BAR_STACK_A"), we need to
            # update ALL formulas that reference the old names using Excel's structured
            # reference syntax: [column_name]
            #
            # This handles formulas like:
            #   - SUM(Table1[[#This Row],[Class A]:[Class B]])
            #   - AVERAGE(Table1[Column1])
            #   - Any formula using [old_column_name] syntax
            #
            # This is generic and works for any template with any formula types.
            # =========================================================================
            if original_headers:
                for cell in root.findall(f".//{{{xlsx_main_ns}}}c"):
                    formula_elem = cell.find(f"{{{xlsx_main_ns}}}f")
                    if formula_elem is not None and formula_elem.text:
                        updated_formula = formula_elem.text
                        for old_name, new_name in original_headers.items():
                            # Replace [old_name] with [new_name] in structured references
                            # This handles all Excel structured reference patterns
                            updated_formula = updated_formula.replace(
                                f"[{old_name}]", f"[{new_name}]"
                            )
                        formula_elem.text = updated_formula

            # Serialize updated sheet XML with proper XML declaration (double quotes)
            updated_sheet_xml = etree.tostring(root, encoding="UTF-8")
            # Replace lxml's single-quote XML declaration with double-quote version
            xml_decl = b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
            if updated_sheet_xml.startswith(b"<?xml"):
                # Remove existing declaration
                updated_sheet_xml = updated_sheet_xml.split(b"?>", 1)[1].lstrip()
            original_files[worksheet_part] = xml_decl + updated_sheet_xml

            # Create updated sharedStrings.xml with proper XML declaration
            sst = etree.Element(f"{{{xlsx_main_ns}}}sst", nsmap={None: xlsx_main_ns})
            sst.set("count", str(len(existing_strings)))
            sst.set("uniqueCount", str(len(existing_strings)))
            for text in existing_strings:
                si = etree.SubElement(sst, f"{{{xlsx_main_ns}}}si")
                t = etree.SubElement(si, f"{{{xlsx_main_ns}}}t")
                t_val = str(text) if text else ""
                # IMPORTANT: preserve whitespace exactly for Office compatibility.
                # The template often uses a single space " " as the first header cell.
                # OOXML requires xml:space="preserve" to keep leading/trailing spaces.
                if t_val.startswith(" ") or t_val.endswith(" "):
                    t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
                t.text = t_val
            sst_xml = etree.tostring(sst, encoding="UTF-8")
            xml_decl = b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
            original_files["xl/sharedStrings.xml"] = xml_decl + sst_xml

            # =========================================================================
            # GENERIC TABLE DEFINITION HANDLING
            # =========================================================================
            # Some embedded chart workbooks include Excel Table(s) (xl/tables/table*.xml).
            # These tables have:
            #   - A ref attribute (e.g., A1:E20) that defines the table range
            #   - tableColumn elements with name attributes for each column
            #   - Optional calculatedColumnFormula elements for computed columns
            #
            # CRITICAL: Do NOT modify the table ref attribute. Keeping it at its
            # original value prevents the "linked file isn't available" error.
            #
            # What we DO update:
            #   1. tableColumn name attributes to match our new header names
            #   2. calculatedColumnFormula references to use new column names
            #
            # This is generic and works for any table with any number of columns.
            # =========================================================================
            
            # Process all table files (table1.xml, table2.xml, etc.)
            table_files = [f for f in original_files.keys() if f.startswith("xl/tables/table") and f.endswith(".xml")]
            
            for table_file in table_files:
                table_xml = original_files.get(table_file)
                if not table_xml:
                    continue
                    
                try:
                    table_root = etree.fromstring(table_xml)
                    
                    # DO NOT change table ref - keep original value
                    # DO NOT change autoFilter ref - keep original value
                    
                    # Build mapping of old column names -> new column names
                    table_columns = table_root.find(
                        f".//{{{xlsx_main_ns}}}tableColumns"
                    )
                    column_name_mapping: Dict[str, str] = {}
                    if table_columns is not None:
                        existing_cols = list(
                            table_columns.findall(f"{{{xlsx_main_ns}}}tableColumn")
                        )
                        # Column A header is a space per our layout
                        new_names = [" "] + [name for name, _ in series_data]
                        for i, new_name in enumerate(new_names):
                            if i < len(existing_cols):
                                old_name = existing_cols[i].get("name", "")
                                # Only update the display name; preserve ids/uids/count/ref.
                                existing_cols[i].set("name", new_name)
                                if old_name and old_name != new_name:
                                    column_name_mapping[old_name] = new_name
                    
                    # Update any calculatedColumnFormula references to use new column names
                    # This handles ALL formulas in the table, regardless of type
                    if column_name_mapping:
                        for formula_elem in table_root.findall(
                            f".//{{{xlsx_main_ns}}}calculatedColumnFormula"
                        ):
                            if formula_elem.text:
                                updated_formula = formula_elem.text
                                for old_name, new_name in column_name_mapping.items():
                                    # Replace [old_name] with [new_name] in structured references
                                    updated_formula = updated_formula.replace(
                                        f"[{old_name}]", f"[{new_name}]"
                                    )
                                formula_elem.text = updated_formula

                    table_xml_out = etree.tostring(table_root, encoding="UTF-8")
                    xml_decl = (
                        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
                    )
                    original_files[table_file] = xml_decl + table_xml_out
                    
                except Exception as e:
                    print(
                        f"    ⚠️  Warning: could not update table column names in {table_file}: {e}"
                    )

            # Rewrite the XLSX archive in-place (pure Python)
            replacements: Dict[str, bytes] = {
                worksheet_part: original_files[worksheet_part],
                "xl/sharedStrings.xml": original_files["xl/sharedStrings.xml"],
            }
            # Include all table files that were processed
            for table_file in table_files:
                if original_files.get(table_file):
                    replacements[table_file] = original_files[table_file]

            # If sharedStrings.xml was not present originally, ensure the workbook registers it.
            if not had_shared_strings:
                content_types = original_files.get("[Content_Types].xml")
                workbook_rels = original_files.get("xl/_rels/workbook.xml.rels")
                if content_types and workbook_rels:
                    try:
                        # Patch [Content_Types].xml
                        ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
                        ct_root = etree.fromstring(content_types)
                        override_exists = any(
                            el.get("PartName") == "/xl/sharedStrings.xml"
                            for el in ct_root.findall(f".//{{{ct_ns}}}Override")
                        )
                        if not override_exists:
                            override = etree.SubElement(ct_root, f"{{{ct_ns}}}Override")
                            override.set("PartName", "/xl/sharedStrings.xml")
                            override.set(
                                "ContentType",
                                "application/vnd.openxmlformats-officedocument.spreadsheetml.sharedStrings+xml",
                            )
                        replacements["[Content_Types].xml"] = etree.tostring(
                            ct_root, xml_declaration=True, encoding="UTF-8", standalone=True
                        )

                        # Patch xl/_rels/workbook.xml.rels
                        rels_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
                        rels_root = etree.fromstring(workbook_rels)
                        rel_exists = any(
                            el.get("Target") == "sharedStrings.xml"
                            for el in rels_root.findall(f".//{{{rels_ns}}}Relationship")
                        )
                        if not rel_exists:
                            existing_ids = [
                                el.get("Id")
                                for el in rels_root.findall(f".//{{{rels_ns}}}Relationship")
                            ]
                            max_id = 0
                            for rid in existing_ids:
                                if rid and rid.startswith("rId"):
                                    try:
                                        max_id = max(max_id, int(rid[3:]))
                                    except ValueError:
                                        continue
                            rel = etree.SubElement(rels_root, f"{{{rels_ns}}}Relationship")
                            rel.set("Id", f"rId{max_id + 1}")
                            rel.set(
                                "Type",
                                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/sharedStrings",
                            )
                            rel.set("Target", "sharedStrings.xml")
                        replacements["xl/_rels/workbook.xml.rels"] = etree.tostring(
                            rels_root, xml_declaration=True, encoding="UTF-8", standalone=True
                        )
                    except Exception as e:
                        print(f"    ⚠️  Warning: could not register sharedStrings.xml: {e}")

            # Prefer system `zip` to preserve ZIP flags required by PowerPoint.
            self._update_xlsx_members_preserving_powerpoint_compat(
                xlsx_path=excel_path, replacements=replacements
            )

            return True

        except Exception as e:
            print(f"    ⚠️  Error updating Excel XML directly: {e}")
            import traceback

            traceback.print_exc()
            return False

    def _create_shared_strings_xml(
        self, path: str, strings: List[str], ns: str
    ) -> None:
        """Create or update sharedStrings.xml with the given strings."""
        # Create root element
        sst = etree.Element(f"{{{ns}}}sst", nsmap={None: ns})
        sst.set("count", str(len(strings)))
        sst.set("uniqueCount", str(len(strings)))

        for text in strings:
            si = etree.SubElement(sst, f"{{{ns}}}si")
            t = etree.SubElement(si, f"{{{ns}}}t")
            t.text = str(text) if text else ""

        # Write the file with correct double-quote declaration for Office
        tree = etree.ElementTree(sst)
        self._write_xml_with_office_declaration(tree, path)

    def _add_shared_strings_references(self, xlsx_temp_dir: str) -> None:
        """
        Add sharedStrings.xml references to [Content_Types].xml and workbook.xml.rels
        when creating a new sharedStrings.xml file.
        """
        content_types_ns = (
            "http://schemas.openxmlformats.org/package/2006/content-types"
        )
        rels_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
        shared_strings_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sharedStrings+xml"
        shared_strings_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/sharedStrings"

        try:
            # Update [Content_Types].xml
            content_types_path = os.path.join(xlsx_temp_dir, "[Content_Types].xml")
            if os.path.exists(content_types_path):
                tree = etree.parse(content_types_path)
                root = tree.getroot()

                # Check if sharedStrings override already exists
                override_exists = any(
                    elem.get("PartName") == "/xl/sharedStrings.xml"
                    for elem in root.findall(f".//{{{content_types_ns}}}Override")
                )

                if not override_exists:
                    override = etree.SubElement(root, f"{{{content_types_ns}}}Override")
                    override.set("PartName", "/xl/sharedStrings.xml")
                    override.set("ContentType", shared_strings_type)
                    self._write_xml_with_office_declaration(tree, content_types_path)

            # Update xl/_rels/workbook.xml.rels
            workbook_rels_path = os.path.join(
                xlsx_temp_dir, "xl", "_rels", "workbook.xml.rels"
            )
            if os.path.exists(workbook_rels_path):
                tree = etree.parse(workbook_rels_path)
                root = tree.getroot()

                # Check if sharedStrings relationship already exists
                rel_exists = any(
                    elem.get("Target") == "sharedStrings.xml"
                    for elem in root.findall(f".//{{{rels_ns}}}Relationship")
                )

                if not rel_exists:
                    # Find the next available rId
                    existing_ids = [
                        elem.get("Id")
                        for elem in root.findall(f".//{{{rels_ns}}}Relationship")
                    ]
                    max_id = 0
                    for rid in existing_ids:
                        if rid and rid.startswith("rId"):
                            try:
                                num = int(rid[3:])
                                max_id = max(max_id, num)
                            except ValueError:
                                pass
                    new_id = f"rId{max_id + 1}"

                    rel = etree.SubElement(root, f"{{{rels_ns}}}Relationship")
                    rel.set("Id", new_id)
                    rel.set("Type", shared_strings_rel_type)
                    rel.set("Target", "sharedStrings.xml")
                    self._write_xml_with_office_declaration(tree, workbook_rels_path)

        except Exception as e:
            print(f"    ⚠️  Warning: Could not add sharedStrings references: {e}")

    def _update_table_reference(
        self,
        table_path: str,
        categories: List[str],
        series_data: List[tuple],
        header_row: int,
        category_col: int,
        series_start_col: int,
        ns: str,
    ) -> None:
        """Update the table reference range in table1.xml to match new data size."""
        try:
            tree = etree.parse(table_path)
            root = tree.getroot()

            # Calculate new reference range
            num_rows = len(categories) + 1  # +1 for header
            num_cols = 1 + len(series_data)  # category col + series cols

            def col_to_letter(col_num: int) -> str:
                result = ""
                while col_num > 0:
                    col_num, remainder = divmod(col_num - 1, 26)
                    result = chr(65 + remainder) + result
                return result

            start_col_letter = col_to_letter(category_col)
            end_col_letter = col_to_letter(category_col + num_cols - 1)
            end_row = header_row + num_rows - 1

            new_ref = f"{start_col_letter}{header_row}:{end_col_letter}{end_row}"
            root.set("ref", new_ref)

            # Update tableColumns
            table_columns = root.find(f".//{{{ns}}}tableColumns")
            if table_columns is not None:
                # Clear existing columns
                for col in list(table_columns):
                    table_columns.remove(col)

                # Set new count
                table_columns.set("count", str(num_cols))

                # Add category column (with space as name)
                cat_col = etree.SubElement(table_columns, f"{{{ns}}}tableColumn")
                cat_col.set("id", "1")
                cat_col.set("name", " ")

                # Add series columns
                for i, (series_name, _) in enumerate(series_data):
                    series_col = etree.SubElement(table_columns, f"{{{ns}}}tableColumn")
                    series_col.set("id", str(i + 2))
                    series_col.set("name", series_name)

            self._write_xml_with_office_declaration(tree, table_path)

        except Exception as e:
            print(f"    ⚠️  Warning: Could not update table1.xml: {e}")

    def _repackage_xlsx(self, source_dir: str, output_path: str) -> None:
        """Repackage a directory structure into an XLSX file."""
        # Remove the old file
        if os.path.exists(output_path):
            os.remove(output_path)

        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as xlsx_zip:
            for root_dir, dirs, files in os.walk(source_dir):
                for file in files:
                    file_path = os.path.join(root_dir, file)
                    arcname = os.path.relpath(file_path, source_dir)
                    xlsx_zip.write(file_path, arcname)

    def _update_formula_references(
        self,
        root,
        ns: dict,
        num_categories: int,
        num_series: int,
        sheet_name: Optional[str] = None,
    ) -> bool:
        """
        Update the formula references (c:f elements) in chart XML to match
        the new data range.

        PowerPoint charts reference data in the embedded Excel using formulas
        like "Sheet1!$A$2:$A$10". When the data size changes, these references
        must be updated to reflect the new range.

        Args:
            root: XML root element of the chart
            ns: XML namespaces dict
            num_categories: Number of category values (rows of data)
            num_series: Number of data series (columns B, C, D, etc.)
            sheet_name: Name of the Excel sheet (if None, uses config default)

        Returns:
            True if successful, False otherwise
        """
        try:
            # Get configuration for Excel data layout
            from hello.utils.ppt_helpers_utils.services.template_config import (
                get_element_dimensions,
            )

            element_dims = get_element_dimensions()

            default_sheet_name = sheet_name or element_dims.excel_default_sheet_name
            header_row = element_dims.excel_header_row
            data_start_row = element_dims.excel_data_start_row
            category_col_letter = chr(
                ord("A") + element_dims.excel_category_column - 1
            )  # Convert 1 -> 'A'
            series_start_col_num = element_dims.excel_series_start_column

            # Calculate the last row number
            last_row = data_start_row + num_categories - 1

            # Find all series in the chart
            series_list = root.findall(".//c:ser", namespaces=ns)

            for series_idx, series in enumerate(series_list):
                # Calculate column letter for this series
                series_col = chr(ord("A") + series_start_col_num - 1 + series_idx)

                # Update category reference (c:cat/c:strRef/c:f or c:cat/c:numRef/c:f)
                cat = series.find(".//c:cat", namespaces=ns)
                if cat is not None:
                    # Try strRef first (text categories)
                    str_ref = cat.find(".//c:strRef", namespaces=ns)
                    if str_ref is not None:
                        f_elem = str_ref.find(".//c:f", namespaces=ns)
                        if f_elem is not None:
                            # Parse existing formula to preserve sheet name
                            existing_sheet = self._extract_sheet_name(
                                f_elem.text, default_sheet_name
                            )
                            f_elem.text = f"{existing_sheet}!${category_col_letter}${data_start_row}:${category_col_letter}${last_row}"

                    # Try numRef (numeric categories like years)
                    num_ref = cat.find(".//c:numRef", namespaces=ns)
                    if num_ref is not None:
                        f_elem = num_ref.find(".//c:f", namespaces=ns)
                        if f_elem is not None:
                            existing_sheet = self._extract_sheet_name(
                                f_elem.text, default_sheet_name
                            )
                            f_elem.text = f"{existing_sheet}!${category_col_letter}${data_start_row}:${category_col_letter}${last_row}"

                # Update value reference (c:val/c:numRef/c:f)
                val = series.find(".//c:val", namespaces=ns)
                if val is not None:
                    num_ref = val.find(".//c:numRef", namespaces=ns)
                    if num_ref is not None:
                        f_elem = num_ref.find(".//c:f", namespaces=ns)
                        if f_elem is not None:
                            existing_sheet = self._extract_sheet_name(
                                f_elem.text, default_sheet_name
                            )
                            f_elem.text = f"{existing_sheet}!${series_col}${data_start_row}:${series_col}${last_row}"

                # Update series name reference (c:tx/c:strRef/c:f) - points to header row
                tx = series.find(".//c:tx", namespaces=ns)
                if tx is not None:
                    str_ref = tx.find(".//c:strRef", namespaces=ns)
                    if str_ref is not None:
                        f_elem = str_ref.find(".//c:f", namespaces=ns)
                        if f_elem is not None:
                            existing_sheet = self._extract_sheet_name(
                                f_elem.text, default_sheet_name
                            )
                            f_elem.text = f"{existing_sheet}!${series_col}${header_row}"

            print(
                f"    ✓ Updated formula references: {category_col_letter}{data_start_row}:{category_col_letter}{last_row} (categories), {num_series} series columns"
            )
            return True

        except Exception as e:
            print(f"    ⚠️  Error updating formula references: {e}")
            import traceback

            traceback.print_exc()
            return False

    def _extract_sheet_name(self, formula: str, default: Optional[str] = None) -> str:
        """
        Extract the sheet name from an Excel formula reference.

        Examples:
            "Sheet1!$A$2:$A$10" -> "Sheet1"
            "'My Sheet'!$B$1" -> "'My Sheet'"

        Args:
            formula: Excel formula string
            default: Default sheet name if extraction fails (if None, uses config default)

        Returns:
            Extracted sheet name or default
        """
        # Get default from config if not provided
        if default is None:
            from hello.utils.ppt_helpers_utils.services.template_config import (
                get_element_dimensions,
            )

            element_dims = get_element_dimensions()
            default = element_dims.excel_default_sheet_name

        if not formula or "!" not in formula:
            return default

        try:
            # Split at '!' and take the part before it
            sheet_part = formula.split("!")[0]
            return sheet_part if sheet_part else default
        except Exception:
            return default

    def _create_pptx(self, extract_dir: str, output_path: str):
        """
        Repackage directory as PPTX.

        Prefer system `zip` for Office compatibility; fall back to Python zipfile if unavailable.
        """
        # Kept as a class method for backward-compat; shared implementation is module-level.
        create_pptx_from_dir(extract_dir, output_path)

    def _normalize_pptx_embedded_workbooks(self, pptx_path: str) -> None:
        """
        Normalize a PPTX ZIP package for Office compatibility:
        1. Embedded workbooks (ppt/embeddings/*.xlsx) are stored (not deflated)
        2. Deflated members inside xlsx have flag_bits=6
        3. All XML files have double-quote XML declarations

        python-pptx writes PPTX files using Python zipfile which typically:
        - Deflates all entries including embedded XLSX packages
        - Sets flag_bits=0 for deflated members inside xlsx files
        - Uses single-quote XML declarations

        PowerPoint expects:
        - Embedded xlsx to be STORED (not deflated) in the outer PPTX
        - Deflated members inside the xlsx to have flag_bits=6
        - Double-quote XML declarations: <?xml version="1.0" ...?>

        We fix all of these by:
        1. Extracting and repackaging via the system `zip` CLI with `-n .xlsx`
        2. Patching flag_bits inside each embedded xlsx
        3. Converting single-quote XML declarations to double-quote
        """
        zip_cli = shutil.which("zip")
        if not zip_cli:
            return

        try:
            with tempfile.TemporaryDirectory() as td:
                extract_dir = os.path.join(td, "pptx_extract")
                os.makedirs(extract_dir, exist_ok=True)

                with zipfile.ZipFile(pptx_path, "r") as z:
                    z.extractall(extract_dir)

                # Fix XML declarations in all XML files (single quotes -> double quotes)
                self._fix_xml_declarations_in_directory(extract_dir)

                # Fix flag_bits inside each embedded xlsx before repackaging
                embed_dir = os.path.join(extract_dir, "ppt", "embeddings")
                if os.path.isdir(embed_dir):
                    for fname in os.listdir(embed_dir):
                        if fname.endswith(".xlsx"):
                            xlsx_path = os.path.join(embed_dir, fname)
                            self._fix_xlsx_internal_flag_bits(xlsx_path)

                tmp_out = os.path.join(td, "normalized.pptx")
                self._create_pptx(extract_dir, tmp_out)

                # Patch flag_bits=6 for all deflated members in the PPTX itself.
                # PowerPoint/Office expect deflated XML files to have flag_bits=6.
                # The system `zip` command sets flag_bits=0, which can cause
                # "The linked file isn't available" errors.
                self._fix_pptx_deflated_flag_bits(tmp_out)

                shutil.copy2(tmp_out, pptx_path)
        except Exception as e:
            # Best-effort only; never fail chart generation because normalization failed.
            print(f"    ⚠️  PPTX normalization skipped: {e}")

    def _fix_pptx_deflated_flag_bits(self, pptx_path: str) -> None:
        """
        Patch flag_bits to 6 for all deflated members in a PPTX file.

        PowerPoint/Office expect deflated members (especially XML files) to have
        flag_bits=6 (general-purpose bit flag indicating UTF-8 encoding).
        Python's zipfile and the system `zip` command set flag_bits=0, which
        can cause "The linked file isn't available" errors when clicking
        "Edit Data in Excel" on charts.

        This patches the ZIP metadata in-place without re-compressing.
        """
        try:
            # Find all deflated members that need patching
            deflated_members = []
            with zipfile.ZipFile(pptx_path, "r") as z:
                for info in z.infolist():
                    if info.compress_type == zipfile.ZIP_DEFLATED and info.flag_bits != 6:
                        deflated_members.append(info.filename)

            if deflated_members:
                self._patch_zip_flag_bits(pptx_path, deflated_members, flag_bits=6)
        except Exception:
            pass  # Best-effort only

    def _fix_xml_declarations_in_directory(self, directory: str) -> None:
        """
        Fix XML declarations in all XML files in a directory tree.

        Changes single-quote declarations to double-quote:
        <?xml version='1.0' encoding='UTF-8' standalone='yes'?>
        becomes:
        <?xml version="1.0" encoding="UTF-8" standalone="yes"?>

        Office is picky about this format.
        """
        import re

        # Pattern to match single-quote XML declaration
        single_quote_pattern = re.compile(
            rb"<\?xml\s+version=['\"]([^'\"]+)['\"]\s+encoding=['\"]([^'\"]+)['\"]\s*"
            rb"(?:standalone=['\"]([^'\"]+)['\"])?\s*\?>"
        )

        for root_dir, _, files in os.walk(directory):
            for filename in files:
                if filename.endswith((".xml", ".rels")):
                    file_path = os.path.join(root_dir, filename)
                    try:
                        with open(file_path, "rb") as f:
                            content = f.read()

                        # Check if there's a single-quote declaration to fix
                        match = single_quote_pattern.match(content[:200])
                        if match:
                            version = match.group(1).decode("utf-8", errors="replace")
                            encoding = match.group(2).decode("utf-8", errors="replace")
                            standalone = match.group(3)

                            # Build new declaration with double quotes
                            if standalone:
                                standalone_str = standalone.decode("utf-8", errors="replace")
                                new_decl = f'<?xml version="{version}" encoding="{encoding}" standalone="{standalone_str}"?>'.encode(
                                    "utf-8"
                                )
                            else:
                                new_decl = f'<?xml version="{version}" encoding="{encoding}"?>'.encode(
                                    "utf-8"
                                )

                            # Replace declaration
                            new_content = single_quote_pattern.sub(new_decl, content, count=1)

                            with open(file_path, "wb") as f:
                                f.write(new_content)
                    except Exception:
                        pass  # Best-effort only

    def _fix_xlsx_internal_flag_bits(self, xlsx_path: str) -> None:
        """
        Patch flag_bits to 6 for all deflated members inside an xlsx file.

        PowerPoint/Excel expect deflated members in embedded xlsx files to have
        flag_bits=6 (compression options). Python's zipfile sets flag_bits=0.
        """
        try:
            # Find all deflated members
            deflated_members = []
            with zipfile.ZipFile(xlsx_path, "r") as z:
                for info in z.infolist():
                    if info.compress_type == zipfile.ZIP_DEFLATED:
                        deflated_members.append(info.filename)

            if deflated_members:
                self._patch_zip_flag_bits(xlsx_path, deflated_members, flag_bits=6)
        except Exception:
            pass  # Best-effort only


class TableDataPopulator:
    """
    Populates table data by modifying table cells
    """

    def __init__(self):
        self.namespaces = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        }

    def _get_fixed_row_height(self) -> float:
        """
        Get fixed row height for all table rows.
        Uses base row height with safety margin to allow for text wrapping without expanding beyond bounds.

        Returns:
            Row height in inches
        """
        from hello.utils.ppt_helpers_utils.services.template_config import (
            get_element_dimensions,
        )

        element_dims = get_element_dimensions()
        min_row_height = element_dims.table_min_row_height
        height_safety_margin = element_dims.table_height_safety_margin

        # Use fixed height with safety margin
        # This allows rows to handle wrapped text without exceeding constraints
        fixed_row_height = min_row_height * height_safety_margin
        return fixed_row_height

    # NOTE: Row height calculation is centralized in TableBlock.calculate_content_based_height()
    # This ensures a single source of truth. Cached heights are passed to populate_table_data().

    def populate_table_data(
        self,
        pptx_path: str,
        slide_index: int,
        table_index: int,
        data: List[Dict[str, Any]],
        output_path: Optional[str] = None,
        table_source: Optional[str] = None,
        cell_frame_info: Optional[dict] = None,
        cached_row_heights: Optional[List[float]] = None,
    ) -> str:
        """
        Populate table with new data

        Args:
            pptx_path: Path to PowerPoint file
            slide_index: Slide number (0-based)
            table_index: Table index on slide (0-based)
            data: List of data dictionaries
            table_source: Optional source label text to add after population
            cell_frame_info: Optional dict with cell frame boundaries (top, height, bottom, left, width)
            output_path: Optional output path
            cached_row_heights: Optional pre-calculated row heights from TableBlock (single source of truth)

        Returns:
            Path to updated PowerPoint file
        """
        if output_path is None:
            output_path = pptx_path

        try:
            # Use python-pptx for table manipulation (easier than raw XML)
            from pptx import Presentation
            from copy import deepcopy

            prs = Presentation(pptx_path)
            slide = prs.slides[slide_index]

            # Find table shapes
            tables = [shape for shape in slide.shapes if shape.has_table]

            if table_index >= len(tables):
                print(
                    f"Warning: Table {table_index} not found on slide {slide_index + 1}"
                )
                return pptx_path

            table_shape = tables[table_index]
            table = table_shape.table  # type: ignore[attr-defined]

            # Check if this is a continuation table (no header)
            skip_header = (
                cell_frame_info.get("skip_header", False) if cell_frame_info else False
            )

            # Keep header row and first data row (for styling reference)
            original_row_count = len(table.rows)

            # Store reference to a data row (row 1) before clearing, for styling
            reference_data_row = None
            if original_row_count > 1:
                reference_data_row = deepcopy(table.rows[1]._tr)

            # Clear existing data rows (keep header for now - we'll handle skip_header later)
            while len(table.rows) > 1:
                # Remove from the end
                table._tbl.remove(table._tbl.tr_lst[-1])

            # Add new rows and populate
            if data and len(data) > 0:
                columns = list(data[0].keys())
                num_data_columns = len(columns)
                num_template_columns = len(table.columns)

                print(f"    📊 Table has {num_template_columns} template columns")
                print(f"    📊 Data has {num_data_columns} columns: {columns}")

                # REMOVE extra columns that don't have data - like deleting columns in PowerPoint
                if num_template_columns > num_data_columns:
                    columns_to_remove = num_template_columns - num_data_columns
                    print(
                        f"    📊 Deleting {columns_to_remove} empty columns (keeping {num_data_columns} data columns)..."
                    )

                    # Step 1: Remove column definitions from table grid
                    # Remove from the END (rightmost columns) - MUST use parent.remove() to actually delete XML
                    tblGrid = table._tbl.tblGrid
                    for _ in range(columns_to_remove):
                        # Get fresh reference to gridCol_lst each time
                        gridCols = tblGrid.gridCol_lst
                        if len(gridCols) > 0:
                            last_col = gridCols[-1]
                            tblGrid.remove(last_col)  # Actually remove from XML tree

                    # Check final count
                    gridCols = tblGrid.gridCol_lst
                    print(f"    ✓ Table grid now has {len(gridCols)} columns")

                    # Step 2: Remove cells from header row
                    if len(table.rows) > 0:
                        header_row = table.rows[0]
                        header_tr = header_row._tr
                        initial_header_cells = len(header_tr.tc_lst)
                        for _ in range(columns_to_remove):
                            # Get fresh reference each time
                            cells = header_tr.tc_lst
                            if len(cells) > 0:
                                last_cell = cells[-1]
                                header_tr.remove(
                                    last_cell
                                )  # Actually remove from XML tree
                        print(
                            f"    ✓ Header row: {initial_header_cells} → {len(header_tr.tc_lst)} cells"
                        )

                    # Step 3: Remove cells from reference data row (for styling template)
                    if reference_data_row is not None:
                        initial_ref_cells = len(reference_data_row.tc_lst)
                        for _ in range(columns_to_remove):
                            # Get fresh reference each time
                            cells = reference_data_row.tc_lst
                            if len(cells) > 0:
                                last_cell = cells[-1]
                                reference_data_row.remove(
                                    last_cell
                                )  # Actually remove from XML tree
                        print(
                            f"    ✓ Reference row: {initial_ref_cells} → {len(reference_data_row.tc_lst)} cells"
                        )

                # Determine row heights to use
                # Prefer cached content-based row heights (accounts for text wrapping)
                # Fall back to fixed row height if cache not available
                fixed_row_height_inches = self._get_fixed_row_height()
                fixed_row_height_emus = int(fixed_row_height_inches * 914400)
                
                # Prepare row heights array: header + data rows + optional source
                use_cached_heights = False
                row_heights_to_use = []
                if cached_row_heights and len(cached_row_heights) > 0:
                    # Convert cached heights (inches) to EMUs
                    row_heights_to_use = [int(h * 914400) for h in cached_row_heights]
                    use_cached_heights = True
                    print(f"    📐 Using content-based row heights: {len(row_heights_to_use)} rows from cache")

                # Handle header row based on skip_header flag
                if skip_header and len(table.rows) > 0:
                    # CONTINUATION TABLE: Convert header row to first data row
                    # Remove the header row completely and we'll add all data as new rows
                    print(
                        "    📋 Continuation table - removing header row, all rows are data"
                    )
                    table._tbl.remove(table._tbl.tr_lst[0])

                    # IMPORTANT: Disable "firstRow" table style to prevent bold formatting
                    # PowerPoint's table styles apply special formatting to the first row
                    # For continuation tables, we need to turn this off
                    tblPr = table._tbl.tblPr
                    if tblPr is not None:
                        tblPr.set("firstRow", "0")  # Disable first row header styling
                        print("    📋 Disabled firstRow table style (no header band)")
                else:
                    # Update headers while preserving formatting (work directly with XML)
                    if len(table.rows) > 0:
                        header_row = table.rows[0]
                        # Access cells directly from XML to avoid caching issues
                        header_cells = header_row._tr.tc_lst
                        print(
                            f"    📋 Updating {min(len(header_cells), num_data_columns)} header cells with data column names..."
                        )

                        from pptx.oxml.xmlchemy import OxmlElement

                        for col_idx in range(min(len(header_cells), num_data_columns)):
                            # Column exists in data - update header text while preserving formatting
                            col_name = columns[col_idx]
                            tc = header_cells[col_idx]  # Get cell directly from XML

                            # Set vertical anchor to CENTER for balanced text alignment
                            # Template has anchor="b" (bottom) which causes top padding
                            a_ns = (
                                "http://schemas.openxmlformats.org/drawingml/2006/main"
                            )
                            tcPr = tc.find(f".//{{{a_ns}}}tcPr")
                            if tcPr is None:
                                # Create tcPr if it doesn't exist (insert at beginning of tc)
                                tcPr = OxmlElement("a:tcPr")
                                tc.insert(0, tcPr)
                            tcPr.set("anchor", "ctr")  # "ctr" = center/middle alignment

                            # Preserve formatting by updating text at XML level
                            txBody = tc.find(
                                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}txBody"
                            )
                            if txBody is not None:
                                t_elem = txBody.find(
                                    ".//{http://schemas.openxmlformats.org/drawingml/2006/main}t"
                                )
                                if t_elem is not None:
                                    t_elem.text = str(col_name)
                                    # t_elem.text = format_label(str(col_name))

                    # Set header row height and LOCK it
                    if len(table.rows) > 0:
                        header_row = table.rows[0]
                        # Use cached header height if available, otherwise fixed
                        if use_cached_heights and len(row_heights_to_use) > 0 and cached_row_heights is not None:
                            header_height_emus = row_heights_to_use[0]
                            header_height_inches = cached_row_heights[0]
                        else:
                            header_height_emus = fixed_row_height_emus
                            header_height_inches = fixed_row_height_inches
                        header_row._tr.set("h", str(header_height_emus))
                        # LOCK header row height to prevent auto-expansion
                        header_row._tr.set("minH", str(header_height_emus))
                        print(
                            f'    📏 Header row height: {header_height_inches:.3f}" ({header_height_emus} EMUs)'
                        )

                # Add data rows
                from pptx.oxml.xmlchemy import OxmlElement

                # Get table grid for truncation calculations
                tbl_grid = table._tbl.tblGrid
                grid_cols = (
                    tbl_grid.gridCol_lst if hasattr(tbl_grid, "gridCol_lst") else []
                )

                print(
                    f"    📝 Adding {len(data)} data rows with {num_data_columns} columns each..."
                )
                # Data rows start at index 1 in cached heights (index 0 is header)
                # For skip_header (continuation) tables, data starts at index 0 (no header)
                data_row_height_offset = 0 if skip_header else 1
                
                for row_idx, row_data in enumerate(data):
                    # Detect TOTAL rows based on the first column only. Must be exact match
                    # (case-insensitive + markdown-wrapped label allowed per config).
                    first_col_name = columns[0] if columns else None
                    first_cell_raw = (
                        row_data.get(first_col_name, "") if first_col_name else ""
                    )
                    is_total_row = is_total_label(first_cell_raw)

                    # Create new row element
                    tr = OxmlElement("a:tr")

                    # Determine row height - use cached content-based height if available
                    row_height_idx = row_idx + data_row_height_offset
                    if use_cached_heights and row_height_idx < len(row_heights_to_use):
                        row_height_emus = row_heights_to_use[row_height_idx]
                    else:
                        row_height_emus = fixed_row_height_emus
                    
                    # Set height as string (OxmlElement.set expects string)
                    tr.set("h", str(row_height_emus))

                    # LOCK the row height to prevent PowerPoint from auto-expanding
                    # Setting minH prevents desktop PowerPoint from overriding our height constraints
                    tr.set("minH", str(row_height_emus))

                    # Add cells to the row - ONLY for data columns
                    cells_added = 0
                    for col_idx in range(num_data_columns):
                        # Clone from DATA row (not header) to preserve data cell styling
                        if reference_data_row is not None and col_idx < len(
                            reference_data_row.tc_lst
                        ):
                            reference_cell = reference_data_row.tc_lst[col_idx]

                            # Clone the cell structure to preserve data cell styling
                            tc = deepcopy(reference_cell)

                            # Get the value for this column
                            col_name = columns[col_idx]
                            raw_value = row_data.get(col_name, "")

                            # If this is a TOTAL row and first cell, normalize display while preserving
                            # input casing (config-driven). This also strips markdown wrappers like **...**.
                            if is_total_row and col_idx == 0:
                                cell_value = total_display_text(first_cell_raw)
                            else:
                                # Data is already transformed by frontend_json_processor
                                # (format_table_cell_value applied during _transform_table_data)
                                cell_value = str(raw_value) if raw_value is not None else ""

                            # NOTE: We no longer truncate cell content here.
                            # Instead, we drop entire rows from the end when table doesn't fit (see _render_table in orchestrator_renderer.py)
                            # This provides better UX - complete rows are shown rather than truncated cell content

                            # DEBUG: Show what we're populating (first row only)
                            if row_idx == 0:
                                print(
                                    f"        Column {col_idx} ({col_name}): '{cell_value[:50]}{'...' if len(cell_value) > 50 else ''}'"
                                )

                            # Update the text content while keeping all formatting
                            txBody = tc.find(
                                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}txBody"
                            )
                            if txBody is not None:
                                # Find the text element
                                t_elem = txBody.find(
                                    ".//{http://schemas.openxmlformats.org/drawingml/2006/main}t"
                                )
                                if t_elem is not None:
                                    t_elem.text = cell_value
                                else:
                                    # If no text element, add one
                                    p = txBody.find(
                                        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}p"
                                    )
                                    if p is not None:
                                        r = p.find(
                                            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}r"
                                        )
                                        if r is None:
                                            r = OxmlElement("a:r")
                                            p.append(r)
                                        t = OxmlElement("a:t")
                                        t.text = cell_value
                                        r.append(t)
                        else:
                            # Fallback: create basic cell if no reference
                            tc = OxmlElement("a:tc")
                            txBody = OxmlElement("a:txBody")
                            bodyPr = OxmlElement("a:bodyPr")
                            lstStyle = OxmlElement("a:lstStyle")
                            p = OxmlElement("a:p")
                            r = OxmlElement("a:r")
                            t = OxmlElement("a:t")
                            col_name = columns[col_idx]
                            raw_value = row_data.get(col_name, "")
                            if is_total_row and col_idx == 0:
                                t.text = total_display_text(first_cell_raw)
                            else:
                                # Data is already transformed by frontend_json_processor
                                t.text = str(raw_value) if raw_value is not None else ""
                            r.append(t)
                            p.append(r)
                            txBody.append(bodyPr)
                            txBody.append(lstStyle)
                            txBody.append(p)
                            tc.append(txBody)
                            tcPr = OxmlElement("a:tcPr")
                            tc.append(tcPr)

                        tr.append(tc)
                        cells_added += 1

                    # If this is a TOTAL row, force the entire row to bold.
                    if is_total_row:
                        try:
                            a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
                            ns = {"a": a_ns}
                            for tc in tr.findall(".//a:tc", namespaces=ns):
                                # Set bold on all runs, and also end-paragraph properties.
                                for r in tc.findall(".//a:r", namespaces=ns):
                                    rPr = r.find("a:rPr", namespaces=ns)
                                    if rPr is None:
                                        rPr = OxmlElement("a:rPr")
                                        r.insert(0, rPr)
                                    rPr.set("b", "1")
                                for p in tc.findall(".//a:p", namespaces=ns):
                                    endPr = p.find("a:endParaRPr", namespaces=ns)
                                    if endPr is None:
                                        endPr = OxmlElement("a:endParaRPr")
                                        p.append(endPr)
                                    endPr.set("b", "1")
                        except Exception:
                            # Non-fatal: table will still render; bolding is best-effort if XML differs.
                            pass

                    # DEBUG: Show cell count for first row
                    if row_idx == 0:
                        print(f"        ✓ Added {cells_added} cells to row")

                    # Add row to table
                    table._tbl.append(tr)

                # Add source as a special row at the bottom of the table
                # This ensures source always stays with the table content
                if table_source:
                    # Get element dimensions from config for consistent source row height
                    from hello.utils.ppt_helpers_utils.services.template_config import (
                        get_element_dimensions,
                    )
                    source_element_dims = get_element_dimensions()
                    # Use config value for source row height (converted to EMUs) and LOCK it.
                    # PowerPoint may otherwise auto-adjust the source row to the table's minimum row height.
                    source_row_height_emus = int(
                        source_element_dims.table_source_row_height * 914400
                    )
                    source_tr = OxmlElement("a:tr")
                    source_tr.set("h", str(source_row_height_emus))
                    source_tr.set("minH", str(source_row_height_emus))

                    # Create first cell that spans all columns (merged cell)
                    source_tc = OxmlElement("a:tc")
                    source_tc.set(
                        "gridSpan", str(num_data_columns)
                    )  # Merge across all columns

                    # Create text body with source text
                    txBody = OxmlElement("a:txBody")
                    bodyPr = OxmlElement("a:bodyPr")
                    bodyPr.set("wrap", "square")
                    lstStyle = OxmlElement("a:lstStyle")
                    p = OxmlElement("a:p")
                    r = OxmlElement("a:r")

                    # Set font properties - 7pt, gray color
                    rPr = OxmlElement("a:rPr")
                    rPr.set("sz", "700")  # 7pt = 700 hundredths of a point
                    rPr.set("lang", "en-US")
                    solidFill = OxmlElement("a:solidFill")
                    srgbClr = OxmlElement("a:srgbClr")
                    srgbClr.set("val", "666666")  # Gray color
                    solidFill.append(srgbClr)
                    rPr.append(solidFill)
                    r.append(rPr)

                    t = OxmlElement("a:t")
                    t.text = table_source
                    r.append(t)
                    p.append(r)
                    txBody.append(bodyPr)
                    txBody.append(lstStyle)
                    txBody.append(p)
                    source_tc.append(txBody)

                    # Cell properties - NO borders, minimal margins
                    # Reduced top margin to minimize visual gap between tables
                    tcPr = OxmlElement("a:tcPr")
                    # Vertical alignment from config: "t" (top), "ctr" (center), "b" (bottom)
                    tcPr.set("anchor", source_element_dims.source_text_vertical_alignment)
                    tcPr.set("marL", "0")  # No left margin
                    tcPr.set("marR", "0")  # No right margin
                    tcPr.set("marT", "22860")  # Reduced top padding (~0.025" instead of 0.05")
                    tcPr.set("marB", "0")  # No bottom margin

                    # Remove all borders by setting them to no line
                    for border_name in ["lnL", "lnR", "lnT", "lnB"]:
                        ln = OxmlElement(f"a:{border_name}")
                        noFill = OxmlElement("a:noFill")
                        ln.append(noFill)
                        tcPr.append(ln)

                    # No fill (transparent background)
                    noFillBg = OxmlElement("a:noFill")
                    tcPr.append(noFillBg)

                    source_tc.append(tcPr)
                    source_tr.append(source_tc)

                    # Add placeholder cells for the merged columns (required by OOXML spec)
                    for _ in range(num_data_columns - 1):
                        placeholder_tc = OxmlElement("a:tc")
                        placeholder_tc.set(
                            "hMerge", "1"
                        )  # This cell is horizontally merged
                        # Minimal content for placeholder
                        ph_txBody = OxmlElement("a:txBody")
                        ph_bodyPr = OxmlElement("a:bodyPr")
                        ph_lstStyle = OxmlElement("a:lstStyle")
                        ph_p = OxmlElement("a:p")
                        ph_txBody.append(ph_bodyPr)
                        ph_txBody.append(ph_lstStyle)
                        ph_txBody.append(ph_p)
                        placeholder_tc.append(ph_txBody)
                        ph_tcPr = OxmlElement("a:tcPr")
                        placeholder_tc.append(ph_tcPr)
                        source_tr.append(placeholder_tc)

                    table._tbl.append(source_tr)
                    print(
                        f"    ✓ Added source as merged table row: '{table_source[:30]}...'"
                    )

            # Adjust table column widths to match the graphic frame width
            # Uses variable column widths if provided (first column expanded to fit content),
            # otherwise falls back to even distribution
            try:
                # Get the table's graphic frame width in EMUs (914400 EMUs = 1 inch)
                frame_width_emu = int(table_shape.width)
                EMU_PER_INCH = 914400

                # Get the table grid and columns
                tbl_grid = table._tbl.tblGrid
                grid_cols = tbl_grid.gridCol_lst

                if grid_cols and len(grid_cols) > 0:
                    # Count valid columns (columns that have width attribute)
                    valid_cols = [gc for gc in grid_cols if hasattr(gc, "w") and gc.w]
                    num_columns = len(valid_cols)

                    if num_columns > 0:
                        # Check if variable column widths are provided
                        column_widths_inches = (
                            cell_frame_info.get("column_widths")
                            if cell_frame_info
                            else None
                        )

                        if column_widths_inches and len(column_widths_inches) == num_columns:
                            # Use pre-calculated variable column widths
                            # First column is expanded to fit content, others share remaining space
                            print(
                                "    [DEBUG] Adjusting table columns (variable width - first col expanded)"
                            )
                            print(
                                f'    [DEBUG] Frame width: {frame_width_emu} EMUs ({table_shape.width / EMU_PER_INCH:.2f}")'
                            )
                            
                            # Convert inches to EMUs and apply
                            total_applied_emu = 0
                            for idx, grid_col in enumerate(valid_cols):
                                col_width_inches = column_widths_inches[idx]
                                col_width_emu = int(col_width_inches * EMU_PER_INCH)
                                
                                # For last column, absorb any rounding difference
                                if idx == len(valid_cols) - 1:
                                    col_width_emu = frame_width_emu - total_applied_emu
                                
                                grid_col.w = col_width_emu
                                total_applied_emu += col_width_emu
                                
                                print(
                                    f'    [DEBUG] Column {idx}: {col_width_inches:.2f}" = {col_width_emu} EMUs'
                                )
                            
                            print(
                                f'    ✓ Applied variable column widths (first col: {column_widths_inches[0]:.2f}")'
                            )
                        else:
                            # Fall back to equal width distribution
                            equal_width_per_col = frame_width_emu // num_columns
                            remainder = frame_width_emu % num_columns

                            print(
                                "    [DEBUG] Adjusting table columns (even distribution fallback)"
                            )
                            print(
                                f'    [DEBUG] Frame width: {frame_width_emu} EMUs ({table_shape.width / EMU_PER_INCH:.2f}")'
                            )
                            print(
                                f"    [DEBUG] Number of columns: {num_columns}, equal width per column: {equal_width_per_col} EMUs"
                            )

                            # Set all columns to equal width
                            for idx, grid_col in enumerate(valid_cols):
                                # Add remainder to the last column to handle rounding
                                if idx == len(valid_cols) - 1:
                                    grid_col.w = equal_width_per_col + remainder
                                else:
                                    grid_col.w = equal_width_per_col

                            print(
                                f"    ✓ Evenly distributed {num_columns} table columns to match graphic frame width"
                            )

                        # Verify final total
                        final_total = sum(
                            int(gc.w) for gc in grid_cols if hasattr(gc, "w") and gc.w
                        )
                        print(
                            f"    [DEBUG] Final column width total: {final_total} EMUs, target: {frame_width_emu} EMUs"
                        )
                    else:
                        print(
                            "    ⚠ Warning: No valid columns found with width attribute"
                        )
                else:
                    print("    ⚠ Warning: No gridCol elements found in table")
            except Exception as e:
                print(f"    ⚠ Warning: Could not adjust table column widths: {e}")
                import traceback

                traceback.print_exc()

            # CRITICAL: Enforce table shape height to prevent overflow and overlaps
            # After data population, PowerPoint may auto-expand rows, causing tables to overlap with elements below
            # We must strictly enforce the table height constraint at the shape level
            try:
                # Get element dimensions for font size enforcement
                from hello.utils.ppt_helpers_utils.services.template_config import (
                    get_element_dimensions,
                )

                element_dims_local = get_element_dimensions()

                # Use the table_shape that was already found at the beginning of this function
                # (no need to search again)
                if (
                    table_shape
                    and cell_frame_info
                    and "table_height" in cell_frame_info
                ):
                    # CRITICAL: Use the EXACT height that was pre-calculated and reserved
                    # This is the height we allocated when rendering the table
                    # DO NOT let PowerPoint change this, regardless of content
                    reserved_height_inches = cell_frame_info["table_height"]
                    enforced_height_emus = int(reserved_height_inches * 914400)

                    # Force the table to the EXACT reserved height
                    from pptx.util import Emu

                    table_shape.height = Emu(enforced_height_emus)

                    total_rows = len(table.rows)

                    print(
                        f'    🔒 Enforced table to reserved height: {reserved_height_inches:.3f}" ({total_rows} rows)'
                    )

                    # CRITICAL: LOCK table and row heights after data population
                    # PowerPoint will try to expand rows to fit content - we must prevent this

                    # Step 1: Set table shape height to EXACT reserved value (again, after population)
                    table_shape.height = Emu(enforced_height_emus)

                    # Step 2: Use CACHED row heights from TableBlock (single source of truth)
                    # Height calculation happens ONLY in TableBlock.calculate_content_based_height()
                    total_rows = len(table.rows)
                    if total_rows > 0:
                        from pptx.util import Pt
                        from pptx.enum.text import MSO_AUTO_SIZE

                        min_row_height = element_dims_local.table_min_row_height

                        # Use cached row heights from TableBlock if available
                        # IMPORTANT: cached_row_heights structure depends on the table type:
                        # - For normal tables:      [header, data1, data2, ..., dataN, source]
                        # - For continuation tables: [data1, data2, ..., dataN, source] (NO header)
                        # The source row (if present) is ALWAYS the LAST element
                        has_source_row = table_source is not None and len(str(table_source).strip()) > 0
                        
                        if cached_row_heights and len(cached_row_heights) > 0:
                            # Determine table structure
                            # total_rows = header (0 or 1) + data_rows + source (0 or 1)
                            num_header_rows = 0 if skip_header else 1
                            num_source_rows = 1 if has_source_row else 0
                            num_data_rows = total_rows - num_header_rows - num_source_rows
                            
                            # Build row_heights_inches with correct heights for each row type
                            row_heights_inches = []
                            
                            # For continuation tables (skip_header=True), cached_row_heights starts with data rows
                            # For normal tables, cached_row_heights starts with header
                            cache_has_header = not skip_header  # Cache structure matches table structure
                            
                            # Header row height (if present) - from start of cache
                            if num_header_rows > 0 and cache_has_header and len(cached_row_heights) > 0:
                                row_heights_inches.append(cached_row_heights[0])
                            
                            # Data row heights
                            # For normal tables: data starts at index 1 (after header)
                            # For continuation tables: data starts at index 0 (no header in cache)
                            data_start_idx = 1 if cache_has_header else 0
                            
                            # Calculate how many data rows are available in cache (excluding source)
                            cache_data_count = len(cached_row_heights) - (1 if cache_has_header else 0) - (1 if has_source_row else 0)
                            
                            # Take the first num_data_rows heights (matching the table's data rows)
                            for i in range(min(num_data_rows, cache_data_count)):
                                cache_idx = data_start_idx + i
                                if cache_idx < len(cached_row_heights):
                                    row_heights_inches.append(cached_row_heights[cache_idx])
                            
                            # Source row height (if present) - from END of cache (always last element)
                            if has_source_row and len(cached_row_heights) > 0:
                                # Source row height is always the last element in cached_row_heights
                                source_height = cached_row_heights[-1]
                                row_heights_inches.append(source_height)
                            
                            # Fill any missing heights with minimum
                            while len(row_heights_inches) < total_rows:
                                row_heights_inches.append(min_row_height)
                            
                            print(
                                f"    📐 Using cached row heights: {len(row_heights_inches)} rows "
                                f"(header={num_header_rows}, data={num_data_rows}, source={num_source_rows}, cache_has_header={cache_has_header})"
                            )
                        else:
                            # Fallback: use minimum row height for all rows
                            row_heights_inches = [min_row_height] * total_rows
                            print(
                                f'    📐 No cached heights - using minimum row height: {min_row_height}" × {total_rows} rows'
                            )

                        # Calculate total content-based height INCLUDING overhead (config-driven).
                        # For current templates, border_overhead and row_gap_padding are calibrated to 0.0,
                        # but we keep the formula for compatibility with other table styles/templates.
                        row_heights_sum = sum(row_heights_inches)
                        border_overhead = element_dims_local.table_border_overhead
                        row_gap_padding = element_dims_local.table_row_gap_padding * max(
                            0, total_rows - 1
                        )
                        total_content_height = row_heights_sum + border_overhead + row_gap_padding

                        # If content fits within reserved space, use content-based heights
                        # If content exceeds reserved space, scale proportionally
                        if total_content_height <= reserved_height_inches:
                            # Content fits - use content-based heights (rows only use what they need)
                            scale_factor = 1.0
                            print(
                                f'    📐 Content-based row heights: rows={row_heights_sum:.3f}" + overhead={border_overhead + row_gap_padding:.3f}" = {total_content_height:.3f}" fits in {reserved_height_inches:.3f}"'
                            )
                        else:
                            # Content exceeds space - scale down proportionally
                            # Only scale the row heights, not the overhead
                            available_for_rows = reserved_height_inches - border_overhead - row_gap_padding
                            scale_factor = available_for_rows / row_heights_sum if row_heights_sum > 0 else 1.0
                            print(
                                f'    📐 Scaling rows: {row_heights_sum:.3f}" → {available_for_rows:.3f}" (scale={scale_factor:.2f}, overhead={border_overhead + row_gap_padding:.3f}")'
                            )

                        # Apply calculated heights to each row
                        for row_idx, row in enumerate(table.rows):
                            row_height_inches = (
                                row_heights_inches[row_idx] * scale_factor
                            )
                            # Ensure minimum height
                            row_height_inches = max(
                                row_height_inches, min_row_height * scale_factor
                            )
                            row_height_emus = int(row_height_inches * 914400)

                            # Set row height
                            row._tr.set("h", str(row_height_emus))
                            # Set minimum height to lock it (prevent PowerPoint expansion)
                            row._tr.set("minH", str(row_height_emus))
                            row.height = row_height_emus

                            # Configure cell text frames using config values
                            for cell_idx, cell in enumerate(row.cells):
                                try:
                                    text_frame = cell.text_frame
                                    if text_frame:
                                        # Apply word wrap setting from config
                                        text_frame.word_wrap = (
                                            element_dims_local.table_word_wrap
                                        )

                                        # Prevent auto-resize to maintain row height
                                        if element_dims_local.table_auto_size_disabled:
                                            text_frame.auto_size = MSO_AUTO_SIZE.NONE

                                        # Set margins from config
                                        text_frame.margin_top = Pt(
                                            element_dims_local.table_cell_margin_top
                                        )
                                        text_frame.margin_bottom = Pt(
                                            element_dims_local.table_cell_margin_bottom
                                        )
                                        text_frame.margin_left = Pt(
                                            element_dims_local.table_cell_margin_left
                                        )
                                        text_frame.margin_right = Pt(
                                            element_dims_local.table_cell_margin_right
                                        )

                                        # Ensure font size is appropriate
                                        for paragraph in text_frame.paragraphs:
                                            for run in paragraph.runs:
                                                if (
                                                    run.font.size is None
                                                    or run.font.size
                                                    > Pt(
                                                        element_dims_local.table_font_size
                                                    )
                                                ):
                                                    run.font.size = Pt(
                                                        element_dims_local.table_font_size
                                                    )
                                                if not run.font.name:
                                                    run.font.name = element_dims_local.table_font_name
                                except Exception:
                                    pass

                        # Keep table shape at allocated height - don't shrink to content
                        # PowerPoint will expand rows visually, so source must be below allocated space
                        actual_content_height = sum(row_heights_inches) * scale_factor
                        print(
                            f"    ✓ Enforced {len(table.rows)} rows with content-based heights"
                        )
                        print(
                            f'       Content height: {actual_content_height:.3f}", Shape height: {reserved_height_inches:.3f}"'
                        )

            except Exception as e:
                print(f"    ⚠ Warning: Could not enforce table height constraint: {e}")
                import traceback

                traceback.print_exc()

            print(
                f"✓ Updated table data: {len(data)} rows × {num_data_columns} columns"
            )

            # Save the presentation (source is already added as a table row)
            prs.save(output_path)
            print(f"✓ Saved table with source row: {output_path}")

        except Exception as e:
            print(f"Error populating table: {e}")
            import traceback

            traceback.print_exc()

        return output_path


class HeroStatsPopulator:
    """
    Populates hero stats on the first slide with trend arrows and values
    """

    def __init__(self):
        self.namespaces = {
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }

        # Mapping between hero_fields keys and template identifiers
        self.hero_field_mapping = {
            "vacancy_rate": {
                "arrow": "slide_1_arrow_1_picture",
                "ticker": "slide_1_ticker_1_text",
                "title": "slide_1_ticker_1_title_text",
            },
            "sf_net_absorption": {
                "arrow": "slide_1_arrow_2_picture",
                "ticker": "slide_1_ticker_2_text",
                "title": "slide_1_ticker_2_title_text",
            },
            "sf_construction_delivered": {
                "arrow": "slide_1_arrow_3_picture",
                "ticker": "slide_1_ticker_3_text",
                "title": "slide_1_ticker_3_title_text",
            },
            "sf_under_construction": {
                "arrow": "slide_1_arrow_4_picture",
                "ticker": "slide_1_ticker_4_text",
                "title": "slide_1_ticker_4_title_text",
            },
            "lease_rate": {
                "arrow": "slide_1_arrow_5_picture",
                "ticker": "slide_1_ticker_5_text",
                "title": "slide_1_ticker_5_title_text",
            },
        }

        # Mapping between trend values and arrow image names
        self.trend_to_arrow = {
            "up": "up_arrow_picture",
            "down": "down_arrow_picture",
            "neutral": "neutral_arrow_picture",
        }

    def populate_hero_stats(
        self,
        pptx_path: str,
        hero_fields: Dict[str, Any],
        output_path: Optional[str] = None,
    ) -> str:
        """
        Populate hero stats on the first slide with trend arrows and values

        Args:
            pptx_path: Path to PowerPoint file
            hero_fields: Dictionary containing stats data with trend and value
            output_path: Optional output path

        Returns:
            Path to updated PowerPoint file
        """
        if output_path is None:
            output_path = pptx_path

        if not hero_fields or "stats" not in hero_fields:
            print("  No hero_fields stats found, skipping population")
            return pptx_path

        stats = hero_fields.get("stats", {})
        if not stats:
            print("  No stats data found in hero_fields, skipping population")
            return pptx_path

        print(f"  📊 Populating hero stats for {len(stats)} metrics...")

        # Create temp directory for extraction
        with tempfile.TemporaryDirectory() as temp_dir:
            extract_dir = os.path.join(temp_dir, "pptx_extract")

            # Extract PPTX
            with zipfile.ZipFile(pptx_path, "r") as zip_ref:
                zip_ref.extractall(extract_dir)

            # Process slide1.xml
            slide_path = os.path.join(extract_dir, "ppt", "slides", "slide1.xml")
            rels_path = os.path.join(
                extract_dir, "ppt", "slides", "_rels", "slide1.xml.rels"
            )

            if not os.path.exists(slide_path):
                print("  Warning: slide1.xml not found")
                return pptx_path

            # Update slide content
            success = self._update_slide_content(slide_path, rels_path, stats)

            if success:
                # Repackage PPTX
                create_pptx_from_dir(extract_dir, output_path)
                print("  ✓ Updated hero stats successfully")
            else:
                print("  Warning: Could not update hero stats")

        return output_path

    def _update_slide_content(
        self, slide_path: str, rels_path: str, stats: Dict[str, Any]
    ) -> bool:
        """Update slide content with hero stats data"""
        try:
            # Parse slide XML
            tree = etree.parse(slide_path)
            root = tree.getroot()

            # Load relationships
            rels_tree = etree.parse(rels_path)
            rels_root = rels_tree.getroot()

            # Process each stat
            for stat_key, stat_data in stats.items():
                if stat_key not in self.hero_field_mapping:
                    print(f"    Warning: Unknown stat key '{stat_key}', skipping")
                    continue

                mapping = self.hero_field_mapping[stat_key]
                print(
                    f"    Processing {stat_key}: {stat_data.get('value', 'N/A')} ({stat_data.get('trend', 'N/A')})"
                )

                # Update ticker title text if label provided
                label = stat_data.get("label")
                title_name = mapping.get("title")
                if title_name and label is not None:
                    self._update_ticker_text(root, title_name, label)

                # Update ticker text
                self._update_ticker_text(
                    root, mapping["ticker"], stat_data.get("value", "")
                )

                # Update arrow image
                trend = stat_data.get("trend", "neutral").lower()
                if trend not in self.trend_to_arrow:
                    print(
                        f"    Warning: Unknown trend '{trend}' for {stat_key}, using neutral"
                    )
                    trend = "neutral"

                arrow_image_name = self.trend_to_arrow[trend]
                self._update_arrow_image(
                    root, rels_root, mapping["arrow"], arrow_image_name
                )

            # Save modified XML with correct double-quote declaration for Office.
            write_xml_with_office_declaration(tree, slide_path)
            write_xml_with_office_declaration(rels_tree, rels_path)
            return True

        except Exception as e:
            print(f"    Error updating slide content: {e}")
            import traceback

            traceback.print_exc()
            return False

    def _update_ticker_text(self, root, ticker_name: str, value: str) -> bool:
        """Update ticker text value"""
        try:
            # Find the ticker text shape by name (requires full XPath support)
            ticker_shapes = root.xpath(
                f'.//p:sp[p:nvSpPr/p:cNvPr[@name="{ticker_name}"]]',
                namespaces=self.namespaces,
            )

            if not ticker_shapes:
                print(f"      Warning: Ticker shape '{ticker_name}' not found")
                return False

            ticker_shape = ticker_shapes[0]

            # Find the text element
            text_elem = ticker_shape.find(".//a:t", namespaces=self.namespaces)
            if text_elem is not None:
                text_elem.text = str(value)
                print(f"      ✓ Updated ticker text: {value}")
                return True
            else:
                print(
                    f"      Warning: Text element not found in ticker '{ticker_name}'"
                )
                return False

        except Exception as e:
            print(f"      Error updating ticker text: {e}")
            return False

    def _update_arrow_image(
        self, root, rels_root, arrow_name: str, image_name: str
    ) -> bool:
        """Update arrow image reference"""
        try:
            # Find the arrow picture shape by name (requires full XPath support)
            arrow_shapes = root.xpath(
                f'.//p:pic[p:nvPicPr/p:cNvPr[@name="{arrow_name}"]]',
                namespaces=self.namespaces,
            )

            if not arrow_shapes:
                print(f"      Warning: Arrow shape '{arrow_name}' not found")
                return False

            arrow_shape = arrow_shapes[0]

            # Find the reference picture (hidden template that has desired image)
            reference_shapes = root.xpath(
                f'.//p:pic[p:nvPicPr/p:cNvPr[@name="{image_name}"]]',
                namespaces=self.namespaces,
            )

            if not reference_shapes:
                print(f"      Warning: Reference arrow image '{image_name}' not found")
                return False

            reference_shape = reference_shapes[0]

            # Locate the blip elements that hold the r:embed relationship id
            ref_blip = reference_shape.find(".//a:blip", namespaces=self.namespaces)
            arrow_blip = arrow_shape.find(".//a:blip", namespaces=self.namespaces)

            if ref_blip is None or arrow_blip is None:
                print(
                    f"      Warning: Could not find blip elements for arrow update '{arrow_name}'"
                )
                return False

            embed_key = f"{{{self.namespaces['r']}}}embed"
            new_rid = ref_blip.get(embed_key)

            if not new_rid:
                print(
                    f"      Warning: Reference image '{image_name}' missing relationship id"
                )
                return False

            arrow_blip.set(embed_key, new_rid)
            print(f"      ✓ Arrow image updated using '{image_name}' (rId={new_rid})")
            return True

        except Exception as e:
            print(f"      Error updating arrow image: {e}")
            return False

# Convenience functions
def populate_chart(
    pptx_path: str,
    slide_index: int,
    chart_index: int,
    data: List[Dict],
    output_path: Optional[str] = None,
    primary_y_axis_title: Optional[str] = None,
    secondary_y_axis_title: Optional[str] = None,
    x_axis_title: Optional[str] = None,
    y_axis_keys: Optional[List[str]] = None,
    is_multi_axis: bool = True,
    primary_y_axis_format_code: Optional[str] = None,
    secondary_y_axis_format_code: Optional[str] = None,
) -> str:
    """Populate chart data and optionally update axis titles.

    Args:
        is_multi_axis: When False, removes secondary axis and plots all series on primary
    """
    populator = ChartDataPopulator()
    return populator.populate_chart_data(
        pptx_path,
        slide_index,
        chart_index,
        data,
        output_path,
        primary_y_axis_title=primary_y_axis_title,
        secondary_y_axis_title=secondary_y_axis_title,
        x_axis_title=x_axis_title,
        y_axis_keys=y_axis_keys,
        is_multi_axis=is_multi_axis,
        primary_y_axis_format_code=primary_y_axis_format_code,
        secondary_y_axis_format_code=secondary_y_axis_format_code,
    )


def populate_table(
    pptx_path: str,
    slide_index: int,
    table_index: int,
    data: List[Dict],
    output_path: Optional[str] = None,
) -> str:
    """Populate table data"""
    populator = TableDataPopulator()
    return populator.populate_table_data(
        pptx_path, slide_index, table_index, data, output_path
    )


def populate_hero_stats(
    pptx_path: str, hero_fields: Dict[str, Any], output_path: Optional[str] = None
) -> str:
    """Populate hero stats data"""
    populator = HeroStatsPopulator()
    return populator.populate_hero_stats(pptx_path, hero_fields, output_path)

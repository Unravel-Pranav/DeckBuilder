import os
import sys
import random
from pathlib import Path
from datetime import datetime

# Add the backend directory to the search path for imports
# This allows 'import app' to work
current_dir = Path(__file__).parent
# backend_dir points to this repository's backend/ directory
backend_node_dir = current_dir.parent.parent
sys.path.insert(0, str(backend_node_dir))

from app.ppt_engine.ppt_helpers_utils.services.frontend_json_processor import FrontendJSONProcessor
from app.ppt_engine.ppt_helpers_utils.services.presentation_generator import PresentationGenerator

def create_sample_json():
    """Create sample JSON: title-only first slide plus one section per chart/table type with distinct data."""
    return {
        "report": {
            "id": 1,
            "name": "Comprehensive Market Analysis 2025",
            "template_name": "first_slide_base",  # Randomized in verify_generation
            "property_type": "Industrial",
            "property_sub_type": "figures",
            "defined_markets": ["Chicago", "Dallas", "Empire State"],
            "quarter": "2025 Q1",
            # Reserve slide 1 for cover/title placeholders only (no charts); see slide_number_assigner
            "title_only_first_slide": True,
        },
        "sections": [
            {
                "id": 1,
                "key": "ExecSummary",
                "name": "Executive Summary",
                "display_order": 0,
                "selected": True,
                "layout_preference": "Full Width",
                "elements": [
                    {
                        "id": 101,
                        "element_type": "commentary",
                        "label": "Market Outlook",
                        "selected": True,
                        "display_order": 0,
                        "config": {
                            "commentary_text": "The industrial sector continues to exhibit resilience. Vacancy rates in core logistics hubs remain below historical averages, although new deliveries are beginning to impact supply-demand dynamics. Rent growth is moderating but remains positive across most primary markets."
                        }
                    },
                    {
                        "id": 102,
                        "element_type": "chart",
                        "label": "Absorption Trends",
                        "selected": True,
                        "display_order": 1,
                        "config": {
                            "chart_type": "Line - Single axis",
                            "chart_name": "Quarterly Net Absorption (MSF)",
                            "chart_data": [
                                {"quarter": "2023 Q3", "absorption": 3.1},
                                {"quarter": "2023 Q4", "absorption": 3.6},
                                {"quarter": "2024 Q1", "absorption": 4.2},
                                {"quarter": "2024 Q2", "absorption": 3.8},
                                {"quarter": "2024 Q3", "absorption": 4.5},
                                {"quarter": "2024 Q4", "absorption": 5.1},
                                {"quarter": "2025 Q1", "absorption": 4.7},
                            ],
                        }
                    }
                ]
            },
            {
                "id": 2,
                "key": "MarketOverviews",
                "name": "Market Overviews",
                "display_order": 1,
                "selected": True,
                "layout_preference": "Content (2x2 Grid)",
                "elements": [
                    {
                        "id": 201,
                        "element_type": "chart",
                        "label": "Market Composition",
                        "selected": True,
                        "display_order": 0,
                        "config": {
                            "chart_type": "Bar Chart",
                            "chart_name": "SF Mix by Use",
                            "chart_data": [
                                {"Category": "Bulk Warehouse", "Value": 38},
                                {"Category": "Last-mile / Urban", "Value": 22},
                                {"Category": "Distribution", "Value": 18},
                                {"Category": "Flex/R&D", "Value": 12},
                                {"Category": "Cold Storage", "Value": 10},
                            ],
                        }
                    },
                    {
                        "id": 202,
                        "element_type": "chart",
                        "label": "Asset Type Allocation",
                        "selected": True,
                        "display_order": 1,
                        "config": {
                            "chart_type": "Pie Chart",
                            "chart_name": "Inventory by Class",
                            "chart_data": [
                                {"Type": "Class A", "Value": 42},
                                {"Type": "Class B", "Value": 33},
                                {"Type": "Class C", "Value": 18},
                                {"Type": "Unclassified", "Value": 7},
                            ],
                        }
                    },
                    {
                        "id": 203,
                        "element_type": "chart",
                        "label": "Inventory Distribution",
                        "selected": True,
                        "display_order": 2,
                        "config": {
                            "chart_type": "Donut Chart",
                            "chart_name": "Leasing Volume by Segment",
                            "chart_data": [
                                {"Segment": "Logistics / 3PL", "Value": 48},
                                {"Segment": "Light Manufacturing", "Value": 27},
                                {"Segment": "R&D / Lab", "Value": 14},
                                {"Segment": "Other Industrial", "Value": 11},
                            ],
                        }
                    },
                    {
                        "id": 204,
                        "element_type": "chart",
                        "label": "Sublease Availability",
                        "selected": True,
                        "display_order": 3,
                        "config": {
                            "chart_type": "Stacked bar",
                            "chart_name": "Availability: Direct vs Sublease (MSF)",
                            "chart_data": [
                                {"Market": "North", "Direct": 820, "Sublease": 140},
                                {"Market": "South", "Direct": 610, "Sublease": 210},
                                {"Market": "East", "Direct": 705, "Sublease": 95},
                                {"Market": "West", "Direct": 890, "Sublease": 260},
                                {"Market": "Central", "Direct": 540, "Sublease": 120},
                            ],
                        }
                    }
                ]
            },
            {
                "id": 3,
                "key": "Trends",
                "name": "Market Trends",
                "display_order": 2,
                "selected": True,
                "layout_preference": "Full Width",
                "elements": [
                    {
                        "id": 301,
                        "element_type": "chart",
                        "label": "Regional Comparison",
                        "selected": True,
                        "display_order": 0,
                        "config": {
                            "chart_type": "Horizontal Bar",
                            "chart_name": "YoY Rent Growth by Region (%)",
                            "chart_data": [
                                {"Region": "Northeast", "Growth": 5.2},
                                {"Region": "Mid-Atlantic", "Growth": 4.4},
                                {"Region": "Midwest", "Growth": 3.8},
                                {"Region": "South", "Growth": 6.1},
                                {"Region": "Mountain", "Growth": 5.0},
                                {"Region": "West", "Growth": 4.9},
                            ],
                        }
                    },
                    {
                        "id": 302,
                        "element_type": "chart",
                        "label": "Dual Axis Metrics",
                        "selected": True,
                        "display_order": 1,
                        "config": {
                            "chart_type": "Line - Multi axis",
                            "chart_name": "Asking Rate vs Leasing Volume",
                            "chart_data": [
                                {"quarter": "2024 Q1", "Rate": 12.1, "Volume": 430},
                                {"quarter": "2024 Q2", "Rate": 12.5, "Volume": 450},
                                {"quarter": "2024 Q3", "Rate": 12.8, "Volume": 420},
                                {"quarter": "2024 Q4", "Rate": 13.2, "Volume": 480},
                                {"quarter": "2025 Q1", "Rate": 13.1, "Volume": 510},
                            ],
                        }
                    }
                ]
            },
            {
                "id": 4,
                "key": "Financials",
                "name": "Market Statistics",
                "display_order": 3,
                "selected": True,
                "layout_preference": "Full Width",
                "elements": [
                    {
                        "id": 401,
                        "element_type": "table",
                        "label": "Submarket Statistical Overview",
                        "selected": True,
                        "display_order": 0,
                        "config": {
                            "table_type": "market_stats_table",
                            "table_data": [
                                {"Submarket": "O'Hare", "Inventory": "120,450,000", "Vacancy": "4.2%", "Absorption": "245,000", "Under_Const": "1,200,000"},
                                {"Submarket": "I-88 Corridor", "Inventory": "85,200,000", "Vacancy": "5.1%", "Absorption": "(12,000)", "Under_Const": "450,000"},
                                {"Submarket": "South Suburbs", "Inventory": "92,100,000", "Vacancy": "6.8%", "Absorption": "115,000", "Under_Const": "0"},
                                {"Submarket": "Central City", "Inventory": "45,000,000", "Vacancy": "3.5%", "Absorption": "45,000", "Under_Const": "85,000"}
                            ]
                        }
                    },
                    {
                        "id": 402,
                        "element_type": "table",
                        "label": "Generic Pricing Trends",
                        "selected": True,
                        "display_order": 1,
                        "config": {
                            "table_data": [
                                {"Metric": "Asking Rent ($/SF)", "High": "15.50", "Low": "8.25", "Avg": "11.20"},
                                {"Metric": "Sale Price ($/SF)", "High": "210", "Low": "145", "Avg": "178"},
                                {"Metric": "Cap Rate (%)", "High": "6.5", "Low": "4.8", "Avg": "5.4"}
                            ]
                        }
                    }
                ]
            },
            {
                "id": 5,
                "key": "ComboShowcase",
                "name": "Advanced Visualizations",
                "display_order": 4,
                "selected": True,
                "layout_preference": "Content (2x2 Grid)",
                "elements": [
                    {
                        "id": 501,
                        "element_type": "chart",
                        "label": "Single Bar + Line",
                        "selected": True,
                        "display_order": 0,
                        "config": {
                            "chart_type": "Combo - Single Bar + Line",
                            "chart_name": "Leasing Volume vs Cap Rate",
                            "chart_data": [
                                {"category": "Chicago", "bar": 118, "line": 5.4},
                                {"category": "Dallas", "bar": 96, "line": 5.8},
                                {"category": "Atlanta", "bar": 104, "line": 5.6},
                                {"category": "Phoenix", "bar": 88, "line": 6.1},
                            ],
                        }
                    },
                    {
                        "id": 502,
                        "element_type": "chart",
                        "label": "Double Bar + Line",
                        "selected": True,
                        "display_order": 1,
                        "config": {
                            "chart_type": "Combo - Double Bar + Line",
                            "chart_name": "New Supply vs Net Absorption + Vacancy",
                            "chart_data": [
                                {"category": "2023", "bar1": 42, "bar2": 38, "line": 5.1},
                                {"category": "2024", "bar1": 55, "bar2": 44, "line": 5.4},
                                {"category": "2025 E", "bar1": 48, "bar2": 41, "line": 5.7},
                            ],
                        }
                    },
                    {
                        "id": 503,
                        "element_type": "chart",
                        "label": "Stacked Bar + Line",
                        "selected": True,
                        "display_order": 2,
                        "config": {
                            "chart_type": "Combo - Stacked Bar + Line",
                            "chart_name": "Direct vs Sublease Stack + Asking Rent Index",
                            "chart_data": [
                                {"category": "2024 Q1", "s1": 320, "s2": 90, "line": 100},
                                {"category": "2024 Q2", "s1": 340, "s2": 110, "line": 102},
                                {"category": "2024 Q3", "s1": 360, "s2": 95, "line": 105},
                                {"category": "2024 Q4", "s1": 375, "s2": 125, "line": 108},
                            ],
                        }
                    },
                    {
                        "id": 504,
                        "element_type": "chart",
                        "label": "Area + Bar Combo",
                        "selected": True,
                        "display_order": 3,
                        "config": {
                            "chart_type": "Combo - Area + Bar",
                            "chart_name": "Trailing Construction Pipeline vs Deliveries",
                            "chart_data": [
                                {"category": "Jan", "area": 72, "bar": 58},
                                {"category": "Feb", "area": 88, "bar": 64},
                                {"category": "Mar", "area": 91, "bar": 70},
                                {"category": "Apr", "area": 85, "bar": 90},
                            ],
                        }
                    }
                ]
            },
            {
                "id": 6,
                "key": "Proprietary",
                "name": "Custom Insights",
                "display_order": 5,
                "selected": True,
                "layout_preference": "Full Width",
                "elements": [
                    {
                        "id": 601,
                        "element_type": "chart",
                        "label": "Segmented Performance",
                        "selected": True,
                        "display_order": 0,
                        "config": {
                            "chart_type": "Single Column Stacked Chart",
                            "chart_name": "Allocation by Risk Profile",
                            "chart_data": [
                                {"Component": "Core / Core+", "Value": 48},
                                {"Component": "Value-Add", "Value": 28},
                                {"Component": "Opportunistic", "Value": 14},
                                {"Component": "Development", "Value": 10},
                            ],
                        }
                    },
                    {
                        "id": 602,
                        "element_type": "commentary",
                        "label": "Investment Theses",
                        "selected": True,
                        "display_order": 1,
                        "config": {
                            "commentary_text": "We recommend a defensive posture for the upcoming quarter. Focus on core assets with long-term credit tenants. Secondary markets may offer spread opportunities but carry higher liquidity risk in the current interest rate environment."
                        }
                    }
                ]
            },
            {
                "id": 7,
                "key": "MarketDetails",
                "name": "Detailed Market Stats",
                "display_order": 6,
                "selected": True,
                "layout_preference": "Full Width",
                "elements": [
                    {
                        "id": 701,
                        "element_type": "table",
                        "label": "Specific Market Breakdown",
                        "selected": True,
                        "display_order": 0,
                        "config": {
                            "table_type": "market_stats_sub_table",
                            "table_data": [
                                {"Item": "Leasing Activity", "Chicago": "1.2M", "Dallas": "0.9M", "NY": "0.5M"},
                                {"Item": "Net Absorption", "Chicago": "450K", "Dallas": "120K", "NY": "(50K)"},
                                {"Item": "Deliveries", "Chicago": "800K", "Dallas": "1.5M", "NY": "200K"}
                            ]
                        }
                    }
                ]
            },
            {
                "id": 8,
                "key": "IndustrialSpecific",
                "name": "Industrial Figures",
                "display_order": 7,
                "selected": True,
                "layout_preference": "Full Width",
                "elements": [
                    {
                        "id": 801,
                        "element_type": "table",
                        "label": "Industrial Figures Summary",
                        "selected": True,
                        "display_order": 0,
                        "config": {
                            "table_type": "industrial_figures_template",
                            "table_data": [
                                {"Market": "Hub A", "Total_SF": "10M", "Available_SF": "500K", "Pct": "5.0%"},
                                {"Market": "Hub B", "Total_SF": "15M", "Available_SF": "1.2M", "Pct": "8.0%"}
                            ]
                        }
                    }
                ]
            }
        ]
    }

def verify_generation():
    """Run verification: title-only cover slide, then content (random first-slide template)."""
    print("🚀 Starting PPT verification (title cover + full chart/table fixture)...")
    
    # 1. Setup paths
    templates_dir = current_dir / "ppt_helpers_utils" / "individual_templates"
    output_dir = backend_node_dir / "data" / "output_ppt"
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # 2. Randomly select a front-page template
    front_page_templates = [
        "first_slide_base",
        "snapshot_first_slide_base",
        "submarket_first_slide_base"
    ]
    selected_template = random.choice(front_page_templates)
    print(f"🎲 Randomly selected front-page template: {selected_template}")
    
    # 3. Initialize Processor and Generator
    processor = FrontendJSONProcessor(templates_dir=str(templates_dir))
    generator = PresentationGenerator(output_dir=str(output_dir))
    
    # 4. Create sample data and inject selected template
    json_data = create_sample_json()
    json_data["report"]["template_name"] = selected_template

    metadata = processor.extract_presentation_metadata(json_data)
    
    # 5. Parse JSON into Sections
    print("📄 Parsing complex JSON data...")
    sections = processor.parse_frontend_json(json_data)
    print(f"✅ Created {len(sections)} sections")
    
    # 6. Generate PPT
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_filename = f"FINAL_VERIFICATION_{timestamp}.pptx"
    output_path = output_dir / output_filename
    
    print(f"🛠 Generating presentation: {output_filename}...")
    try:
        final_path = generator.generate_presentation(
            sections=sections,
            title=metadata["title"],
            author=metadata.get("author", "CBRE Research"),
            output_filename=output_filename,
            metadata=metadata,
        )
        print(f"✅ PRESENTATION COMPLETE: {final_path}")
        
        # Verify slide count using python-pptx
        from pptx import Presentation
        prs = Presentation(final_path)
        slide_count = len(prs.slides)
        file_size = os.path.getsize(final_path)
        
        print(f"📊 ACTUAL SLIDE COUNT: {slide_count}")
        print(f"📊 FILE SIZE: {file_size / 1024:.2f} KB")

        # Title-only slide 1 + eight content sections (some sections share a slide; others split).
        expected_min_slides = 9
        if slide_count >= expected_min_slides:
            print(
                f"✨ SUCCESS! At least {expected_min_slides} slides (title cover + full fixture)."
            )
        else:
            print(
                f"❌ ERROR: {slide_count} slides; expected at least {expected_min_slides} "
                "(1 title cover + content)."
            )
            
    except Exception as e:
        print(f"❌ ERROR during generation: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    verify_generation()

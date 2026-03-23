from __future__ import annotations

"""
Utilities for updating PowerPoint templates in-place, adapted from the
at-report-generator project. These helpers let us update text, charts and
tables on existing slides by referencing shapes by name, while preserving
formatting and layout.

Only PPT-generation/manipulation logic is included here. No Snowflake or
external API dependencies are carried over.
"""

import math
from copy import deepcopy
from io import BytesIO
from typing import Any, Optional, List

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData
from app.utils.formatting import format_label
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape


def _round_to_nice_whole_number(value: float, direction: str = "round") -> int:
    """
    Round a value to a "nice" whole number for axis scaling.

    Args:
        value: The value to round
        direction: "floor" for min values, "ceil" for max values, "round" for nearest

    Returns:
        A nicely rounded whole number (integer)
    """
    if value == 0:
        return 0

    abs_value = abs(value)
    step: int

    if abs_value >= 10000:
        step = 1000 if abs_value >= 50000 else 500
    elif abs_value >= 1000:
        step = 100 if abs_value >= 5000 else 50
    elif abs_value >= 100:
        step = 10 if abs_value >= 500 else 5
    elif abs_value >= 10:
        step = 5 if abs_value >= 50 else 2
    elif abs_value >= 1:
        step = 1
    else:
        step = 1

    if direction == "floor":
        result = math.floor(value / step) * step
    elif direction == "ceil":
        result = math.ceil(value / step) * step
    else:
        result = round(value / step) * step

    return int(result)


# --- Generic helpers (local copies) -------------------------------------------------

def move_slide_to_index(prs: Presentation, old_index: int, new_index: int) -> None:
    """Move a slide from one position to another within a presentation."""
    slides = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slide = slides[old_index]
    slides.remove(slide)
    slides.insert(new_index, slide)


def identify_groups(dataframe: pd.DataFrame) -> List[List[int]]:
    """Group dataframe rows based on indentation in the first column.

    Rows with a leading space/nbsp in the first column are treated as
    belonging to the current group; a non-indented row starts a new group.
    Returns lists of row indices per group, preserving row order.
    """
    groups: List[List[int]] = []
    current: List[int] = []
    for idx, row in dataframe.iterrows():
        first_col_value = str(row.iloc[0])
        indented = first_col_value.startswith(" ") or first_col_value.startswith("\u00A0")
        if not indented:
            if current:
                groups.append(current)
            current = [idx]
        else:
            current.append(idx)
    if current:
        groups.append(current)
    return groups


# --- Text --------------------------------------------------------------------------

def update_slide_text(shape: BaseShape, new_text: str) -> None:
    """Replace text content while preserving run-level formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    words = str(new_text).split()
    para = tf.paragraphs[0]
    runs = para.runs
    if not runs:
        para.text = new_text
        return
    for r in runs:
        r.text = ""
    run_idx = 0
    for i, word in enumerate(words):
        space = " " if i > 0 else ""
        idx = min(run_idx, len(runs) - 1)
        runs[idx].text += space + word
        run_idx += 1

def update_shape(prs: Presentation, shape_name: str, update_data: Any, dynamic_y_axis: bool = False, category_cols: Optional[list[str]] = None, series_col: Optional[str] = None) -> None:
    """Updates a specific shape in a PowerPoint presentation based on its name.

    Searches through all slides in the presentation for a shape with the specified name
    and updates its content based on the shape type (text, chart, or table). If the shape 
    is not found in regular slides and its name starts with 'slide_master_', it searches 
    in the slide master's group shapes.

    Args:
        prs (Presentation): A python-pptx Presentation object containing the slides to update.
        shape_name (str): The name of the shape to update.
        update_data (Any): The data to update the shape with. Can be text, chart data, 
            or table data depending on the shape type.
        dynamic_y_axis (bool, optional): Whether to dynamically adjust the Y-axis scale 
            for charts based on data. Defaults to False. Only applies to chart shapes.
        category_cols (Optional[list[str]], optional): List of column names to use as 
            category levels for charts. If None (default), the first column will be used as the 
            single category level. For multi-level categories, provide column names in
            order from outermost to innermost category level.

    Returns:
        None

    Notes:
        For text frames, the update_data is converted to string.
        For charts and tables, the update_data format should match what's expected by
        update_chart() and update_table() respectively.
        For slide master shapes, only text frames are supported.
        When dynamic_y_axis is True, chart Y-axis min/max will be automatically 
        adjusted based on data.
        For multi-level chart categories, specify the column names in category_cols.

    Examples:
        >>> update_shape(presentation, "title_text", "New Title")
        >>> update_shape(presentation, "sales_chart", sales_data_df, dynamic_y_axis=True)
        >>> update_shape(presentation, "employee_table", employee_data_df)
        >>> update_shape(presentation, "hierarchical_chart", multi_level_df, category_cols=["year", "period"])
    """    
    # Find shape in slides    
    for slide in prs.slides:
        shape = next((s for s in slide.shapes if s.name == shape_name), None)
        if shape:
            if shape.has_text_frame:
                update_slide_text(shape, str(update_data))
            elif shape.has_chart:
                update_chart(shape.chart, update_data, dynamic_y_axis, category_cols, series_col)
            elif shape.has_table:
                update_table(prs, shape_name, update_data)
            return

    # Check slide master and all layouts if not found in slides
    if shape_name.startswith('slide_master_'):
        def check_master_shapes(shapes):
            for shape in shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # Check shapes within group
                    subshape = next((s for s in shape.shapes if s.name == shape_name), None)
                    if subshape and subshape.has_text_frame:
                        update_slide_text(subshape, str(update_data))
                        return True
                elif shape.name == shape_name and shape.has_text_frame:
                    update_slide_text(shape, str(update_data))
                    return True
            return False

        # Check main slide master
        if check_master_shapes(prs.slide_master.shapes):
            return

        # Check all slide layouts
        for layout in prs.slide_master.slide_layouts:
            if check_master_shapes(layout.shapes):
                return



# --- Charts ------------------------------------------------------------------------

def _analyze_dataframe_types(df: pd.DataFrame) -> tuple[list[str], list[str]]:
    """
    Analyze DataFrame to determine which columns are categorical vs numerical.
    
    Args:
        df: pandas DataFrame
        
    Returns:
        Tuple of (categorical_columns, numerical_columns)
    """
    categorical_columns = []
    numerical_columns = []
    
    for column in df.columns:
        # Check if column contains numeric values
        # First try to convert to numeric, see how many succeed
        numeric_converted = pd.to_numeric(df[column], errors='coerce')
        non_null_original = df[column].dropna()
        non_null_converted = numeric_converted.dropna()
        
        # If at least 70% of non-null values can be converted to numeric, treat as numerical
        if len(non_null_original) > 0:
            conversion_rate = len(non_null_converted) / len(non_null_original)
            if conversion_rate >= 0.7:
                numerical_columns.append(column)
            else:
                categorical_columns.append(column)
        else:
            # Empty column, treat as categorical
            categorical_columns.append(column)
    
    return categorical_columns, numerical_columns


def _select_dataframe_columns(df: pd.DataFrame, chart_type: str = 'auto') -> tuple[list[str], list[str]]:
    """
    Select appropriate columns for chart based on data types and chart type.
    
    Args:
        df: pandas DataFrame
        chart_type: Type of chart or 'auto' for intelligent detection
        
    Supported chart types:
    - 'line_single_axis': Single line - 1 categorical + 1 numerical
    - 'line_multi_axis': Multi-line - 1 categorical + multiple numerical  
    - 'combo_bar_line': Combo - 1 categorical + 2 numerical
    - 'combo_double_bar_line': Double combo - 1 categorical + 3 numerical
    - 'combo_stacked_bar_line': Stacked bar + line combo - 1 categorical + 2+ numerical
    - 'combo_area_bar': Area+bar combo - 1 categorical + 2 numerical
    - 'bar_chart': Bar chart - 1 categorical + numerical(s)
    - 'horizontal_bar': Horizontal bar - axes swapped
    - 'auto': Intelligent detection based on data
        
    Returns:
        Tuple of (category_columns, series_columns)
    """
    categorical_cols, numerical_cols = _analyze_dataframe_types(df)
    
    print(f"DataFrame analysis for {chart_type}:")
    print(f"  Categorical columns: {categorical_cols}")
    print(f"  Numerical columns: {numerical_cols}")
    
    # Select category columns (typically X-axis)
    if categorical_cols:
        category_columns = [categorical_cols[0]]  # Use first categorical
    else:
        category_columns = [df.columns[0]]  # Fallback to first column
        print(f"  No categorical columns found, using '{category_columns[0]}' as category")
    
    # Select series columns based on chart type
    if chart_type == 'line_single_axis':
        series_columns = [numerical_cols[0]] if numerical_cols else [df.columns[1]]
    elif chart_type == 'line_multi_axis':
        series_columns = numerical_cols if numerical_cols else df.columns[1:].tolist()
    elif chart_type in ['combo_bar_line', 'combo_area_bar']:
        series_columns = numerical_cols[:2] if len(numerical_cols) >= 2 else numerical_cols
        if len(series_columns) < 2:
            print(f"  Warning: {chart_type} needs 2 series, only found {len(series_columns)}")
    elif chart_type == 'combo_double_bar_line':
        series_columns = numerical_cols[:3] if len(numerical_cols) >= 3 else numerical_cols
        if len(series_columns) < 3:
            print(f"  Warning: {chart_type} needs 3 series, only found {len(series_columns)}")
    elif chart_type == 'combo_stacked_bar_line':
        # Stacked bar + line: use all available numerical columns
        series_columns = numerical_cols if numerical_cols else [col for col in df.columns if col not in category_columns]
        if len(series_columns) < 2:
            print(f"  Warning: {chart_type} needs at least 2 series, only found {len(series_columns)}")
    elif chart_type in ['bar_chart', 'horizontal_bar']:
        # Bar charts can be single or multi-series
        series_columns = numerical_cols if numerical_cols else [col for col in df.columns if col not in category_columns]
    elif chart_type in ['pie', 'donut', 'pie_chart', 'donut_chart']:
        # Pie/donut charts: 1 categorical (labels) + 1 numerical (values)
        series_columns = [numerical_cols[0]] if numerical_cols else [df.columns[1]]
        print(f"  {chart_type} chart: using single value column for slice sizes")
    else:  # 'auto' or unknown
        # Intelligent selection: use all numerical columns
        series_columns = numerical_cols if numerical_cols else [col for col in df.columns if col not in category_columns]
    
    print(f"  Selected category: {category_columns}")
    print(f"  Selected series: {series_columns}")
    
    return category_columns, series_columns


def update_chart(
    chart: Chart,
    update_data: pd.DataFrame,
    dynamic_y_axis: bool,
    category_cols: Optional[list[str]] = None,
    series_col: Optional[str] = None,
    chart_type: str = 'auto'
) -> None:
    """Replace chart data from a DataFrame and keep formatting.

    Five modes supported:
    1) series_col specified: distinct values in that column become series; other columns are categories
    2) category_cols specified: those columns are the categories; all other columns become series
    3) chart_type specified: intelligent selection based on specific chart type requirements
    4) intelligent auto-detection: analyze data types to determine categorical (x-axis) vs numerical (y-axis) columns
    5) fallback default: first column is categories; remaining columns are series
    """
    if not isinstance(update_data, pd.DataFrame):
        return

    df = update_data.copy()
    chart_data = CategoryChartData()

    # Mode 1: explicit series column
    if series_col is not None and series_col in df.columns:
        categories = [c for c in df.columns if c != series_col]
        chart_data.categories = categories
        for series_name in df[series_col].unique().tolist():
            row = df[df[series_col] == series_name]
            if row.empty:
                continue
            values: list[Optional[float]] = []
            for cat in categories:
                v = row[cat].iloc[0] if cat in row.columns else None
                if pd.isna(v) or (isinstance(v, float) and np.isinf(v)):
                    values.append(None)
                else:
                    try:
                        values.append(float(v))
                    except (ValueError, TypeError):
                        values.append(None)
            clean_name = str(series_name)
            # format_label(str(series_name))
            chart_data.add_series(clean_name, values)

    # Mode 2: explicit category columns
    elif category_cols:
        for col in category_cols:
            if col not in df.columns:
                return
        series_cols = [c for c in df.columns if c not in category_cols]
        if len(category_cols) == 1:
            cat = category_cols[0]
            df[cat] = df[cat].fillna("").astype(str)
            chart_data.categories = df[cat].tolist()
            for col in series_cols:
                cleaned: list[Any] = []
                for v in df[col].tolist():
                    try:
                        if pd.isna(v):
                            cleaned.append(None)
                        else:
                            vf = float(v)
                            cleaned.append(None if np.isinf(vf) else vf)
                    except (ValueError, TypeError):
                        cleaned.append(v)
                chart_data.add_series(col, cleaned)
                    # format_label(col), cleaned)
        else:
            # Multi-level categories: convert all to string and build tuples
            for c in category_cols:
                df[c] = df[c].fillna("").astype(str)
            categories = [tuple(r) for r in df[category_cols].itertuples(index=False, name=None)]
            chart_data.categories = categories
            for col in series_cols:
                cleaned = []
                for v in df[col].tolist():
                    try:
                        if pd.isna(v):
                            cleaned.append(None)
                        else:
                            vf = float(v)
                            cleaned.append(None if np.isinf(vf) else vf)
                    except (ValueError, TypeError):
                        cleaned.append(v)
                chart_data.add_series(col, cleaned)
                    # format_label(col), cleaned)

    # Mode 3: chart_type specified - intelligent selection based on chart type
    elif chart_type != 'auto':
        category_columns, series_columns = _select_dataframe_columns(df, chart_type)
        
        # Set categories
        if len(category_columns) == 1:
            chart_data.categories = df[category_columns[0]].tolist()
        else:
            # Multi-level categories
            chart_data.categories = [tuple(r) for r in df[category_columns].itertuples(index=False, name=None)]
        
        # Add series data
        for col in series_columns:
            cleaned: list[Any] = []
            for v in df[col].tolist():
                try:
                    if pd.isna(v):
                        cleaned.append(None)
                        continue
                    vf = float(v)
                    cleaned.append(None if np.isinf(vf) else vf)
                except (ValueError, TypeError):
                    cleaned.append(v)
            chart_data.add_series(col, cleaned)
                # format_label(col), cleaned)

    # Mode 4: intelligent auto-detection based on data types
    elif chart_type == 'auto':
        # Analyze data types
        categorical_cols, numerical_cols = _analyze_dataframe_types(df)
        
        print(f"Chart data analysis:")
        print(f"  Categorical columns: {categorical_cols}")
        print(f"  Numerical columns: {numerical_cols}")
        
        # Select category column (x-axis) - prefer first categorical column
        if categorical_cols:
            category_col = categorical_cols[0]
            chart_data.categories = df[category_col].tolist()
            
            # Use numerical columns as series (y-axis)
            series_columns = numerical_cols if numerical_cols else [col for col in df.columns if col != category_col]
        else:
            # Fallback to Mode 5: traditional default (first column = categories)
            print("  No categorical columns found, using traditional default mode")
            category_col = df.columns[0]
            chart_data.categories = df.iloc[:, 0].tolist()
            series_columns = df.columns[1:].tolist()
        
        print(f"  Selected category (X-axis): {category_col}")
        print(f"  Selected series (Y-axis): {series_columns}")
        
        # Add series data
        for col in series_columns:
            cleaned: list[Any] = []
            for v in df[col].tolist():
                try:
                    if pd.isna(v):
                        cleaned.append(None)
                        continue
                    vf = float(v)
                    cleaned.append(None if np.isinf(vf) else vf)
                except (ValueError, TypeError):
                    cleaned.append(v)
            chart_data.add_series(col, cleaned)
                # format_label(col), cleaned)

    chart.replace_data(chart_data)

    # Optional dynamic Y-axis scaling
    if dynamic_y_axis and hasattr(chart, "value_axis") and chart.value_axis is not None:
        numeric_values: list[float] = []
        for c in update_data.columns[1:]:
            for v in update_data[c]:
                try:
                    if pd.notnull(v) and not np.isinf(float(v)):
                        numeric_values.append(float(v))
                except (ValueError, TypeError):
                    pass
        if numeric_values:
            min_val, max_val = min(numeric_values), max(numeric_values)

            # Calculate padding (10% of the range)
            data_range = max_val - min_val
            padding = data_range * 0.1 if data_range > 0 else (abs(min_val) * 0.1 if min_val != 0 else 1.0)

            # Dynamic scaling - apply padding and round to nice whole number
            scaled_min = _round_to_nice_whole_number(min_val - padding, direction="floor")
            if scaled_min > min_val:
                scaled_min = _round_to_nice_whole_number(min_val, direction="floor")

            # If all data is non-negative, don't let the axis go negative
            if min_val >= 0 and scaled_min < 0:
                scaled_min = 0

            # Calculate max with padding, round to nice whole number, add 1 unit top padding
            scaled_max = _round_to_nice_whole_number(max_val + padding, direction="ceil")
            if scaled_max < max_val:
                scaled_max = _round_to_nice_whole_number(max_val, direction="ceil")
            scaled_max = scaled_max + 1  # Add 1 unit padding on top

            chart.value_axis.minimum_scale = scaled_min
            chart.value_axis.maximum_scale = scaled_max
            try:
                chart.value_axis.tick_labels.number_format = "#,##0"
                chart.value_axis.tick_labels.number_format_is_linked = False
                for plot in chart.plots:
                    if plot.has_data_labels:
                        dl = plot.data_labels
                        dl.number_format = "#,##0"
                        dl.number_format_is_linked = False
            except Exception:
                pass

    # Bar charts: reverse order and set label formatting
    if chart.chart_type in (
        XL_CHART_TYPE.BAR_CLUSTERED,
        XL_CHART_TYPE.BAR_STACKED,
        XL_CHART_TYPE.BAR_STACKED_100,
    ):
        try:
            chart.category_axis.reverse_order = True
            for plot in chart.plots:
                if plot.has_data_labels:
                    dl = plot.data_labels
                    dl.number_format = "#,##0"
                    dl.number_format_is_linked = False
                    dl.position = XL_LABEL_POSITION.INSIDE_END
        except Exception:
            pass


# --- Tables ------------------------------------------------------------------------

def update_table(prs: Presentation, shape_name: str, dataframe: pd.DataFrame) -> None:
    """Populate a named table shape with DataFrame content.

    - Preserves table formatting
    - Splits across cloned slides if data exceeds table capacity
    - Maintains indentation groupings across slides
    """
    if not isinstance(dataframe, pd.DataFrame):
        raise TypeError("dataframe must be a pandas DataFrame")

    df = dataframe.copy()
    df["_original_index"] = range(len(df))

    # Locate the table shape
    target_slide = None
    table_shape = None
    slide_idx = 0
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if getattr(shape, "name", None) == shape_name and getattr(shape, "has_table", False):
                target_slide = slide
                table_shape = shape
                slide_idx = idx
                break
        if target_slide:
            break
    if not table_shape:
        raise ValueError(f"Table shape {shape_name} not found in any slide")

    # Original dims and table
    original_left = table_shape.left
    original_top = table_shape.top
    original_width = table_shape.width
    original_height = table_shape.height
    original_table = table_shape.table

    rows_needed = len(df) + 1  # header row
    rows_per_table = len(original_table.rows)
    num_tables = (rows_needed + rows_per_table - 1) // rows_per_table
    tables = [original_table]

    # Create additional slides/tables if required
    if num_tables > 1:
        for i in range(num_tables - 1):
            new_slide = prs.slides.add_slide(target_slide.slide_layout)
            # Shallow copy of shapes on the original slide
            for sh in target_slide.shapes:
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_bytes = sh.image.blob
                    pic = new_slide.shapes.add_picture(BytesIO(img_bytes), sh.left, sh.top, sh.width, sh.height)
                    pic.name = sh.name
                else:
                    el = sh.element
                    new_el = deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

            # Drop common content placeholders that some templates carry
            to_remove = []
            for sh in new_slide.shapes:
                if getattr(sh, "name", None) in {"Content Placeholder 1", "Content Placeholder 2", "Text Placeholder 1"}:
                    to_remove.append(sh)
            for sh in to_remove:
                sp = sh.element
                sp.getparent().remove(sp)

            # Move immediately after the original slide
            move_slide_to_index(prs, len(prs.slides) - 1, slide_idx + i + 1)

            # Restore table dims on the new slide and collect table handle
            for sh in new_slide.shapes:
                if getattr(sh, "name", None) == shape_name and getattr(sh, "has_table", False):
                    sh.left, sh.top, sh.width, sh.height = original_left, original_top, original_width, original_height
                    tables.append(sh.table)
                    break

    # Build indentation-aware row groups and distribute across tables
    indent_groups = [sorted(g) for g in identify_groups(df)]
    table_groups: list[list[int]] = []
    current: list[int] = []
    used = 0
    capacity = rows_per_table - 1  # minus header
    for grp in indent_groups:
        if used + len(grp) > capacity:
            if current:
                table_groups.append(current)
            current = grp
            used = len(grp)
        else:
            current.extend(grp)
            used += len(grp)
    if current:
        table_groups.append(current)

    # Fill tables
    current_group_idx = 0
    for table in tables:
        visible_cols = [c for c in df.columns if c != "_original_index"]
        for col_idx, col_name in enumerate(visible_cols):
            cell = table.cell(0, col_idx)
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = str(col_name)
            else:
                para.text = str(col_name)

        current_row = 1
        while current_group_idx < len(table_groups):
            grp = table_groups[current_group_idx]
            if current_row + len(grp) <= len(table.rows):
                for row_idx in grp:
                    for col_idx, col in enumerate(visible_cols):
                        cell = table.cell(current_row, col_idx)
                        para = cell.text_frame.paragraphs[0]
                        val = df.iloc[row_idx][col]
                        val_str = "" if pd.isna(val) else str(val)
                        if not para.runs:
                            # copy formatting from row 1 if available
                            first_para = table.cell(1, col_idx).text_frame.paragraphs[0]
                            run = para.add_run()
                            if first_para.runs:
                                run.font.name = first_para.runs[0].font.name
                                run.font.size = first_para.runs[0].font.size
                        para.runs[0].text = val_str
                    current_row += 1
                current_group_idx += 1
            else:
                break

        # Trim any unused rows
        while current_row < len(table.rows):
            tr = table._tbl.tr_lst[current_row]
            table._tbl.remove(tr)

        # Restore original geometry (protects from placeholder copying quirks)
        table_shape = table._graphic_frame  # type: ignore[attr-defined]
        table_shape.left = original_left
        table_shape.top = original_top
        table_shape.width = original_width
        table_shape.height = original_height


# --- Simple directional indicator ---------------------------------------------------

def update_ticker_arrows(prs: Presentation, shape_name: str, previous_quarter: float, current_quarter: float) -> None:
    """Replace a placeholder arrow by another picture based on the sign/change."""
    if previous_quarter < 0 and current_quarter >= 0:
        new_shape_name = "up_arrow_picture"
    elif previous_quarter > 0 and current_quarter <= 0:
        new_shape_name = "down_arrow_picture"
    elif current_quarter > previous_quarter:
        new_shape_name = "up_arrow_picture"
    elif current_quarter < previous_quarter:
        new_shape_name = "down_arrow_picture"
    else:
        new_shape_name = "neutral_arrow_picture"

    for slide in prs.slides:
        old_shape = next((s for s in slide.shapes if s.name == shape_name), None)
        new_shape = next((s for s in slide.shapes if s.name == new_shape_name), None)
        if old_shape and new_shape:
            left, top, width, height = old_shape.left, old_shape.top, old_shape.width, old_shape.height
            image_bytes = new_shape.image.blob
            slide.shapes.add_picture(BytesIO(image_bytes), left, top, width, height)
            sp = old_shape._element
            sp.getparent().remove(sp)
            return


#!/usr/bin/env python3
"""
Presentation Generator
Uses SlideOrchestrator and OrchestratorRenderer to create PowerPoint presentations
"""

import os
from typing import List, Optional
from datetime import datetime

from app.ppt_engine.ppt_helpers_utils.ppt_helpers import (
    SlideOrchestrator,
    Section,
    SlideConstraints,
    export_layouts_to_json,
    OrchestratorRenderer,
)
from app.ppt_engine.ppt_helpers_utils.ppt_helpers.data_populator import populate_hero_stats
from app.utils.formatting import is_total_label, total_display_text
from app.ppt_engine.ppt_helpers_utils.services.template_config import (
    get_ppt_template_config,
    get_slide_constraint_profile,
    get_header_format_config,
)


class PresentationGenerator:
    """
    High-level presentation generator that orchestrates the entire workflow
    """

    def __init__(
        self, output_dir: str = "output", constraints: Optional[SlideConstraints] = None
    ):
        """
        Initialize presentation generator

        Args:
            output_dir: Directory for output files
            constraints: Optional custom slide constraints
        """
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

        # Initialize orchestrator with template-driven constraints
        default_constraints = constraints or get_slide_constraint_profile(None)

        self.orchestrator = SlideOrchestrator(constraints=default_constraints)
        self.renderer = OrchestratorRenderer()
        self.constraints = default_constraints

    @staticmethod
    def _write_xml_with_office_declaration(tree_or_root, file_path: str) -> None:
        """
        Write XML with double-quote XML declaration that Office expects.

        lxml's tree.write(xml_declaration=True) produces single quotes:
            <?xml version='1.0' encoding='UTF-8' standalone='yes'?>

        But Office expects double quotes:
            <?xml version="1.0" encoding="UTF-8" standalone="yes"?>

        This helper serializes without declaration and prepends the correct one.
        """
        from lxml import etree

        # Get the root element (handle both ElementTree and Element)
        if hasattr(tree_or_root, "getroot"):
            root = tree_or_root.getroot()
        else:
            root = tree_or_root

        # Serialize without declaration
        xml_bytes = etree.tostring(root, encoding="UTF-8")

        # Prepend correct double-quote declaration
        xml_decl = b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'

        with open(file_path, "wb") as f:
            f.write(xml_decl + xml_bytes)

    def generate_presentation(
        self,
        sections: List[Section],
        title: str = "Market Report",
        author: str = "CBRE",
        output_filename: Optional[str] = None,
        metadata: dict = None,
    ) -> str:
        """
        Generate complete PowerPoint presentation

        Args:
            sections: List of Section objects with content blocks
            title: Presentation title
            author: Presentation author
            output_filename: Optional custom output filename
            metadata: Optional metadata dict for placeholder replacement

        Returns:
            Path to generated PowerPoint file
        """
        print("\n" + "=" * 60)
        print(f"🚀 GENERATING PRESENTATION: {title}")
        print("=" * 60)

        # Generate output filename
        if output_filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = f"presentation_{timestamp}.pptx"

        output_path = os.path.join(self.output_dir, output_filename)

        # Handle empty sections - generate empty PPT with first slide template only
        if not sections or len(sections) == 0:
            print("\n⚠️  No processable sections found - generating empty presentation")
            return self._generate_empty_presentation(
                output_path, title, author, metadata
            )

        # CRITICAL: Reset renderer state for clean generation
        # When the same PresentationGenerator instance is reused (e.g., for multi-market reports),
        # the renderer's state must be reset to prevent issues:
        # - _first_block_rendered_per_section: Tracks which sections have rendered their first block
        #   (for section title placement). Without reset, section titles disappear in subsequent PPTs.
        # - pending_data_updates: Queue of chart/table data to populate. Must be cleared between runs.
        # This fix addresses bugs where:
        # 1. Section titles disappeared for 2nd+ market in multi-market Figures reports
        # 2. Sections could be skipped due to stale tracking state
        self.renderer._first_block_rendered_per_section = {}
        self.renderer.pending_data_updates = []

        # Configure orchestrator policy for this run
        property_sub_type = metadata.get("property_sub_type") if metadata else None
        self.orchestrator.set_property_sub_type(property_sub_type)

        # STEP 3: Group blocks by pre-assigned slide numbers
        print("\n📐 Step 3: Organizing blocks by slide numbers...")
        grouped_sections = self._group_blocks_by_slide_number(sections)

        # Slide rendering plan
        print(f"\n{'=' * 60}")
        print("🎨 SLIDE RENDERING PLAN")
        print(f"{'=' * 60}")

        for slide_idx, slide_section in enumerate(grouped_sections):
            slide_num = slide_section.get("slide_number", slide_idx + 1)
            block_count = len(slide_section["blocks"])
            section_names = slide_section.get("section_names", [])

            print(f"\nSlide {slide_num}:")
            print(f"  Blocks: {block_count}")
            print(f"  Sections: {', '.join(section_names)}")
            print(f"  Layout: {slide_section.get('layout_preference', 'auto')}")

            for block in slide_section["blocks"]:
                block_type = block.type.value if hasattr(block, "type") else "unknown"
                print(f"    - {block_type}")

        print(f"{'=' * 60}\n")

        # Process each slide (now represented as a section)
        all_layouts = []
        total_slides = len(
            grouped_sections
        )  # Calculate total slides from grouped sections

        for slide_idx, slide_section in enumerate(grouped_sections):
            slide_num = slide_section.get("slide_number", slide_idx + 1)
            print(f"\n📊 Processing Slide {slide_num}")
            print(f"   Blocks: {len(slide_section['blocks'])}")
            print(f"   Original sections: {slide_section.get('section_names', [])}")

            # VERIFY DATA IS POPULATED BEFORE ORCHESTRATION
            print(
                "\n   🔍 Verifying data availability in blocks (BEFORE orchestration):"
            )
            for block in slide_section["blocks"]:
                if hasattr(block, "type"):
                    if block.type.value == "table":
                        data_count = (
                            len(block.data)
                            if hasattr(block, "data") and block.data
                            else 0
                        )
                        if data_count > 0:
                            print(
                                f"      ✅ Block {block.id} (table): {data_count} data rows"
                            )
                        else:
                            print(
                                f"      ⚠️  Block {block.id} (table): NO DATA - size estimation will use fallback"
                            )
                    elif block.type.value == "chart":
                        data_count = (
                            len(block.data)
                            if hasattr(block, "data") and block.data
                            else 0
                        )
                        if data_count > 0:
                            print(
                                f"      ✅ Block {block.id} (chart): {data_count} data rows"
                            )
                        else:
                            print(
                                f"      ⚠️  Block {block.id} (chart): NO DATA - size estimation will use fallback"
                            )

            # Create a Section object for this slide
            # Use the section_style from the original section (preserved from frontend_json_processor)
            section_obj = Section(
                id=f"slide_{slide_num}",
                title=f"Slide {slide_num}",
                blocks=slide_section["blocks"],
                layout_preference=slide_section.get("layout_preference"),
                style=slide_section.get(
                    "section_style", {}
                ),  # Use original section style
            )

            # Orchestrate layout for this slide (pass total_slides and slide_num for first/last slide detection)
            # NOTE: Blocks should have data populated at this point for accurate size estimation
            layouts = self.orchestrator.orchestrate_section(
                section_obj, total_slides=total_slides, start_slide_number=slide_num
            )

            print(f"   Generated: {len(layouts)} layout(s)")
            for layout in layouts:
                print(f"   - Layout type: {layout.layout_type.value}")
                layout.slide_number = slide_num  # Use pre-assigned number

            all_layouts.extend(layouts)

            # Export layout JSON for debugging
            layout_json_path = os.path.join(
                self.output_dir, f"layout_slide_{slide_num}.json"
            )
            export_layouts_to_json(layouts, layout_json_path)
            print(f"   Layout spec: {layout_json_path}")

        print(f"\n📄 Total slides: {len(all_layouts)}")

        # Render all slides into single presentation
        print("\n🎨 Rendering to PowerPoint...")
        self._render_all_sections(
            grouped_sections, all_layouts, output_path, title, author, metadata
        )

        # Final normalization pass: various steps inside rendering use python-pptx saves,
        # which can deflate ppt/embeddings/*.xlsx again. Normalize once at the end so
        # embedded workbooks remain editable via PowerPoint's "Edit Data in Excel".
        try:
            from app.ppt_engine.ppt_helpers_utils.ppt_helpers.data_populator import (
                ChartDataPopulator,
            )

            ChartDataPopulator()._normalize_pptx_embedded_workbooks(output_path)
        except Exception:
            # Best-effort only; never fail generation due to normalization.
            pass

        print(f"\n✅ PRESENTATION COMPLETE: {output_path}")
        print("=" * 60 + "\n")

        return output_path

    def _generate_empty_presentation(
        self, output_path: str, title: str, author: str, metadata: dict = None
    ) -> str:
        """
        Generate an empty presentation when there are no processable sections.

        Creates a PPT with just the first slide template (if available) or a blank presentation.
        This ensures the workflow doesn't fail when all sections have empty data.

        Args:
            output_path: Path to save the presentation
            title: Presentation title
            author: Presentation author
            metadata: Optional metadata dict for placeholder replacement

        Returns:
            Path to generated PowerPoint file
        """
        from pptx import Presentation

        property_sub_type = metadata.get("property_sub_type") if metadata else None
        template_config = get_ppt_template_config(property_sub_type)

        # Template paths
        template_dir = os.path.join(
            os.path.dirname(__file__), "..", "individual_templates"
        )
        template_dir = os.path.abspath(template_dir)

        first_slide_template = (
            os.path.join(template_dir, template_config.first_slide)
            if template_config.first_slide
            else None
        )
        base_template_path = os.path.join(template_dir, template_config.base_slide)

        first_slide_available = first_slide_template and os.path.exists(
            first_slide_template
        )

        if first_slide_available:
            print(
                f"   Using first slide template for empty presentation: {first_slide_template}"
            )
            prs = Presentation(first_slide_template)
        elif os.path.exists(base_template_path):
            print(
                f"   Using base template for empty presentation: {base_template_path}"
            )
            prs = Presentation(base_template_path)
            # Remove existing slides from template to create empty PPT
            while len(prs.slides) > 0:
                xml_slides = prs.slides._sldIdLst
                xml_slides.remove(xml_slides[0])
        else:
            print("   No templates found, creating blank empty presentation")
            prs = Presentation()

        # Populate template placeholders if metadata is provided and we have a first slide
        if first_slide_available and metadata and len(prs.slides) > 0:
            self._populate_template_placeholders(prs, title, author, metadata)

        prs.core_properties.title = title
        prs.core_properties.author = author

        # Save the empty presentation
        prs.save(output_path)

        # Populate hero stats if available (must be done after saving the PPTX)
        should_populate_hero = (
            first_slide_available and metadata and metadata.get("hero_fields")
        )
        if should_populate_hero:
            print("   🎯 Populating hero stats...")
            self._populate_hero_stats(output_path, metadata["hero_fields"])

        print(f"\n✅ EMPTY PRESENTATION GENERATED: {output_path}")
        print("   Note: No sections had processable data")
        print("=" * 60 + "\n")

        # Best-effort normalization for consistency with full pipeline.
        try:
            from app.ppt_engine.ppt_helpers_utils.ppt_helpers.data_populator import (
                ChartDataPopulator,
            )

            ChartDataPopulator()._normalize_pptx_embedded_workbooks(output_path)
        except Exception:
            pass

        return output_path

    def _group_blocks_by_slide_number(self, sections: List[Section]) -> List[dict]:
        """
        Group content blocks by their pre-assigned slide numbers.

        Args:
            sections: List of Section objects with blocks that have slide numbers

        Returns:
            List of dicts, one per slide, containing grouped blocks
        """
        from collections import defaultdict

        slides_dict = defaultdict(
            lambda: {
                "blocks": [],
                "section_names": [],
                "layout_preference": None,
                "slide_number": None,
                "section_style": None,  # Track original section style
            }
        )

        # Debug logging before grouping
        print(f"\n{'=' * 60}")
        print("📑 GROUPING BLOCKS BY SLIDE NUMBER")
        print(f"{'=' * 60}")

        section_block_counts = {}
        for section in sections:
            section_name = section.title
            block_count = len(section.blocks)
            section_block_counts[section_name] = block_count
            print(f"\nSection: '{section_name}' has {block_count} blocks")

            for block_idx, block in enumerate(section.blocks):
                slide_num = None
                if hasattr(block, "metadata") and isinstance(block.metadata, dict):
                    slide_num = block.metadata.get("slide_number")

                block_type = block.type.value if hasattr(block, "type") else "unknown"
                block_id = block.id if hasattr(block, "id") else "?"

                if slide_num is None:
                    print(
                        f"   ⚠️  Block {block_idx} ({block_type}, id={block_id}): NO SLIDE NUMBER - will assign to slide 1"
                    )
                else:
                    print(
                        f"   ✓  Block {block_idx} ({block_type}, id={block_id}): Slide {slide_num}"
                    )

        print(f"{'=' * 60}\n")

        # Collect all blocks and group by slide number
        for section in sections:
            section_name = section.title
            section_id = section.id  # Get section ID for tracking

            for block in section.blocks:
                # Get slide number from block metadata
                slide_num = None
                if hasattr(block, "metadata") and isinstance(block.metadata, dict):
                    slide_num = block.metadata.get("slide_number")

                if slide_num is None:
                    # Fallback: assign to slide 1 if no number
                    slide_num = 1
                    print(
                        f"   ⚠️ Block {block.id} has no slide number, assigning to slide 1"
                    )

                # CRITICAL: Set section_id, section_name, and section_style on block for section boundary constraint
                # This allows the renderer to detect section changes and add section titles appropriately
                block.section_id = section_id
                block.section_name = section_name
                block.section_style = section.style if hasattr(section, "style") else {}
                print(
                    f"   ✓ Block {block.id}: section_id='{section_id}', section_name='{section_name}', slide={slide_num}"
                )

                slides_dict[slide_num]["blocks"].append(block)
                slides_dict[slide_num]["slide_number"] = slide_num

                # Track which sections contributed to this slide
                if section_name not in slides_dict[slide_num]["section_names"]:
                    slides_dict[slide_num]["section_names"].append(section_name)

                # Set layout preference from section if not set
                if slides_dict[slide_num]["layout_preference"] is None:
                    slides_dict[slide_num]["layout_preference"] = (
                        section.layout_preference
                    )

                # Set section style from first section that contributes to this slide
                if slides_dict[slide_num]["section_style"] is None:
                    slides_dict[slide_num]["section_style"] = (
                        section.style if hasattr(section, "style") else {}
                    )

        # Convert to sorted list by slide number
        sorted_slides = [
            slides_dict[slide_num] for slide_num in sorted(slides_dict.keys())
        ]

        print(f"   Grouped into {len(sorted_slides)} slides")
        for slide_data in sorted_slides:
            print(
                f"     Slide {slide_data['slide_number']}: {len(slide_data['blocks'])} blocks from {slide_data['section_names']}"
            )

        return sorted_slides

    def _select_content_layout(
        self, prs, is_last_slide: bool = False, property_sub_type: str = None
    ):
        """
        Select appropriate layout based on slide position and property_sub_type.

        For the last slide with "submarket" property_sub_type, searches for a "back"
        layout first (e.g., "F&LR - Back"). For all other cases (including Figures
        and Snapshot), uses "blank" or "body" layouts.

        Args:
            prs: Presentation object
            is_last_slide: Whether this is the last content slide in the presentation
            property_sub_type: The property sub type (e.g., "Figures", "Snapshot", "Submarket")

        Returns:
            Selected layout object
        """
        selected_layout = None

        # Only use "back" layout for submarket property_sub_type on the last slide
        use_back_layout = (
            is_last_slide and (property_sub_type or "").strip().lower() == "submarket"
        )

        # For the last slide with submarket, try to find a "back" layout first
        if use_back_layout:
            for layout in prs.slide_layouts:
                layout_name_lower = layout.name.lower()
                if "back" in layout_name_lower:
                    selected_layout = layout
                    print(
                        f"     Selected layout: {layout.name} (contains 'back' - last slide, submarket)"
                    )
                    break

        # If no "back" layout found (or not submarket/last slide), look for "blank" or "body"
        if selected_layout is None:
            for layout in prs.slide_layouts:
                layout_name_lower = layout.name.lower()
                if "blank" in layout_name_lower or "body" in layout_name_lower:
                    selected_layout = layout
                    print(
                        f"     Selected layout: {layout.name} (contains 'blank' or 'body')"
                    )
                    break

        # Fallback to last layout if criteria not met
        if selected_layout is None:
            selected_layout = prs.slide_layouts[-1]
            print(
                f"     Using fallback layout: {selected_layout.name} (no layout found with 'blank' or 'body')"
            )

        return selected_layout

    def _process_table_continuations(
        self,
        section,  # Section object
        layout,  # SlideLayout object
        current_pptx: str,
        slide_constraints,  # SlideConstraints object
        content_layout_selector,
        metadata: dict = None,
    ) -> str:
        """
        Process table continuation data and create continuation slides.

        When a table is too large to fit on a single slide in full_width layout,
        the renderer stores the remaining rows in _continuation_data. This method:
        1. Collects blocks with continuation data
        2. Creates new slides for the continuation rows
        3. Recursively handles multi-page continuations

        Args:
            section: The Section containing blocks
            layout: The SlideLayout used for the original slide
            current_pptx: Path to the current presentation file
            slide_constraints: Slide constraint profile
            content_layout_selector: Function to select content layout
            metadata: Optional metadata dict

        Returns:
            Path to the updated presentation file
        """
        from pptx import Presentation
        from app.ppt_engine.ppt_helpers_utils.ppt_helpers.slide_orchestrator import (
            TableBlock,
            SlideLayout,
            LayoutType,
            CellFrame,
        )

        # Collect blocks with continuation data
        continuation_blocks = []
        for block in section.blocks:
            if hasattr(block, "_continuation_data") and block._continuation_data:
                continuation_blocks.append(block)

        if not continuation_blocks:
            return current_pptx

        print(f"\n   📄 Processing {len(continuation_blocks)} table continuation(s)...")

        for block in continuation_blocks:
            remaining_rows = block._continuation_data  # type: ignore
            continuation_id = getattr(block, "_continuation_id", f"table_{id(block)}")
            continuation_label = getattr(block, "_continuation_label", "Table")
            continuation_section_name = getattr(
                block, "_continuation_section", "Continued"
            )

            # Get the original block's template path and source for rendering
            original_template_path = getattr(block, "template_path", None)
            original_table_source = getattr(block, "table_source", None)

            # Track continuation number for labeling
            continuation_num = 1

            while remaining_rows:
                print(
                    f"      Creating continuation slide {continuation_num} for '{continuation_label}' ({len(remaining_rows)} rows remaining)"
                )

                # Load presentation
                prs = Presentation(current_pptx)

                # Add new slide
                content_layout = content_layout_selector(prs)
                slide = prs.slides.add_slide(content_layout)

                # Remove placeholders
                self._clean_slide_placeholders(slide)

                # Calculate how many rows fit on this continuation slide
                # Use regular slide content height (not first slide)
                from app.ppt_engine.ppt_helpers_utils.services.template_config import (
                    get_slide_layout_config,
                )

                property_sub_type = (
                    metadata.get("property_sub_type") if metadata else None
                )
                layout_config = get_slide_layout_config(property_sub_type)

                # Calculate available height for continuation slide
                regular_constraints = layout_config.get_constraints(
                    is_first_slide=False
                )
                available_height = regular_constraints.content_height

                # Create a continuation TableBlock with the same template as original
                # Continuation tables SHOULD repeat header row on every slide for readability.
                continuation_block = TableBlock(
                    id=f"{continuation_id}_cont_{continuation_num}",
                    label=f"{continuation_label} (continued)",
                    data=remaining_rows,
                    template_path=original_template_path,  # Use original block's template
                    has_header=True,  # Continuation - repeat header row on each continuation slide
                    table_source=original_table_source,  # Source label (will appear on last slide)
                )
                # Mark as continuation slide content (used for label/heading suppression, not header behavior)
                continuation_block.is_continuation = True  # type: ignore[attr-defined]
                continuation_block.section_name = (
                    f"{continuation_section_name} (continued)"
                )

                # Create a CellFrame for the continuation table
                continuation_cell = CellFrame(
                    left=regular_constraints.margin_left,
                    top=regular_constraints.margin_top,
                    width=regular_constraints.content_width,
                    height=available_height,
                )

                # Create SlideLayout with required parameters
                continuation_layout = SlideLayout(
                    slide_number=len(prs.slides),
                    section_id=f"continuation_{continuation_num}",
                    layout_type=LayoutType.FULL_WIDTH,
                    cell_frames=[continuation_cell],
                    assigned_blocks=[(0, continuation_block.id)],
                    scale_factors={continuation_block.id: 1.0},
                )

                # Create a Section for rendering
                continuation_section_obj = Section(
                    id=f"continuation_{continuation_num}",
                    title=f"{continuation_section_name} (continued)",
                    blocks=[continuation_block],
                    layout_preference="full_width",
                    style={
                        "show_title": False
                    },  # Don't show section title on continuation
                    constraints=slide_constraints,
                )

                # Save temp file
                temp_path = f"{current_pptx}.cont_{continuation_num}"
                prs.save(temp_path)

                # Render the continuation table
                render_slide_index = len(prs.slides) - 1
                next_path = self.renderer._render_slide_content(
                    section=continuation_section_obj,
                    layout=continuation_layout,
                    current_file=temp_path,
                    slide_index=render_slide_index,
                    constraints=slide_constraints,
                )

                # Update current file
                current_pptx = next_path if next_path != temp_path else temp_path

                # Check if this continuation block also has remaining rows
                if (
                    hasattr(continuation_block, "_continuation_data")
                    and continuation_block._continuation_data
                ):
                    remaining_rows = continuation_block._continuation_data  # type: ignore
                    continuation_num += 1
                else:
                    remaining_rows = None

            # Clear the original block's continuation data
            block._continuation_data = None  # type: ignore
            print(
                f"      ✓ Created {continuation_num} continuation slide(s) for '{continuation_label}'"
            )

        return current_pptx

    def _render_all_sections(
        self,
        sections: List,  # Now list of dicts with grouped blocks
        all_layouts: List,
        output_path: str,
        title: str,
        author: str,
        metadata: dict = None,
    ):
        """
        Render all sections into a single presentation
        """
        from pptx import Presentation

        property_sub_type = metadata.get("property_sub_type") if metadata else None
        template_config = get_ppt_template_config(property_sub_type)
        slide_constraints = get_slide_constraint_profile(property_sub_type)

        # Template paths - using individual_templates for chart-type-based templates
        template_dir = os.path.join(
            os.path.dirname(__file__), "..", "individual_templates"
        )
        template_dir = os.path.abspath(template_dir)

        base_template_path = os.path.join(template_dir, template_config.base_slide)
        first_slide_template = (
            os.path.join(template_dir, template_config.first_slide)
            if template_config.first_slide
            else None
        )

        # Check if we have any BASE_SLIDE layouts
        has_base_slide = any(
            layout.layout_type.value == "base_slide" for layout in all_layouts
        )
        first_slide_available = first_slide_template and os.path.exists(
            first_slide_template
        )

        if first_slide_available:
            print(f"   Using first slide template: {first_slide_template}")
            prs = Presentation(first_slide_template)
            if has_base_slide:
                print("   Preserving template first slide for BASE_SLIDE layout")
            else:
                print("   Keeping template first slide as intro cover")
        elif os.path.exists(base_template_path):
            print(f"   Using base template: {base_template_path}")
            prs = Presentation(base_template_path)
            # Remove existing slides from template
            while len(prs.slides) > 0:
                xml_slides = prs.slides._sldIdLst
                xml_slides.remove(xml_slides[0])
        else:
            print("   No templates found, using blank presentation")
            prs = Presentation()

        # Populate template placeholders if metadata is provided and we have a first slide
        if first_slide_available and metadata and len(prs.slides) > 0:
            self._populate_template_placeholders(prs, title, author, metadata)

        should_populate_hero = (
            first_slide_available and metadata and metadata.get("hero_fields")
        )

        prs.core_properties.title = title
        prs.core_properties.author = author

        # Track which blocks and layout each slide index corresponds to
        # IMPORTANT: all_layouts may have more items than sections when a single grouped section
        # generates multiple layouts (e.g., when blocks don't all fit and are split)
        section_map = {}

        # Build a lookup of all blocks by their ID for quick access
        all_blocks_by_id = {}
        for section in sections:
            for block in section["blocks"]:
                all_blocks_by_id[block.id] = block

        # Create section map from all layouts
        for layout_idx, layout in enumerate(all_layouts):
            # Get the blocks assigned to this specific layout using assigned_blocks
            layout_blocks = []
            for cell_idx, block_id in layout.assigned_blocks:
                if block_id in all_blocks_by_id:
                    layout_blocks.append(all_blocks_by_id[block_id])

            # Determine section name and style from first block (blocks carry their original section metadata)
            # If multiple sections on same slide, first block's section name and style are used
            section_title = f"Slide {layout_idx + 1}"  # Default fallback
            section_style = {}  # Default fallback
            if layout_blocks:
                first_block_section_name = getattr(
                    layout_blocks[0], "section_name", None
                )
                if first_block_section_name:
                    section_title = first_block_section_name
                # Get section style from first block (preserves show_title config)
                first_block_section_style = getattr(
                    layout_blocks[0], "section_style", None
                )
                if first_block_section_style:
                    section_style = first_block_section_style

            # Create a Section object with blocks from this specific layout
            section_obj = Section(
                id=f"slide_{layout_idx + 1}",
                title=section_title,  # Use section name from first block
                blocks=layout_blocks,  # Use only blocks assigned to this layout
                layout_preference=None,  # Layout is already determined
                style=section_style,  # Use section style from first block (preserves show_title config)
                constraints=slide_constraints,
            )
            section_map[layout_idx] = (section_obj, layout)

        # Create slides and render content
        temp_files = []
        current_pptx = output_path

        # Save initial presentation
        prs.save(current_pptx)

        # Populate hero stats once the PPTX exists on disk, so updates persist
        if should_populate_hero:
            self._populate_hero_stats(current_pptx, metadata["hero_fields"])

        # Determine the last slide index for back layout selection
        sorted_slide_indices = sorted(section_map.keys())
        last_slide_idx = sorted_slide_indices[-1] if sorted_slide_indices else -1

        for slide_idx in sorted_slide_indices:
            section, layout = section_map[slide_idx]

            # Detect if this is the last slide in the presentation
            is_last_slide = slide_idx == last_slide_idx

            print(
                f"   Rendering slide {slide_idx + 1}/{len(section_map)}: {layout.layout_type.value}"
            )

            # Load current presentation
            prs = Presentation(current_pptx)

            # Select appropriate layout for this slide
            # For submarket: last slide uses "back" layout if available (e.g., "F&LR - Back")
            # For figures/snapshot: always use "blank" or "body" layouts (no back slide)
            content_layout = self._select_content_layout(
                prs, is_last_slide=is_last_slide, property_sub_type=property_sub_type
            )

            # Handle different slide types
            # Use existing first slide ONLY if:
            # 1. This is the first slide to render (slide_idx == 0)
            # 2. First slide template is available
            # 3. The layout's slide_number is 1 (content intended for first slide, not overflow)
            # CRITICAL: If slide_number > 1, content was moved to a continuation slide
            # and should NOT use the first slide template
            using_existing_first_slide = False
            is_first_slide_content = (
                layout.slide_number == 1
            )  # Check actual slide number, not render order

            if slide_idx == 0 and first_slide_available and is_first_slide_content:
                if len(prs.slides) > 0:
                    slide = prs.slides[0]
                    using_existing_first_slide = True
                    print(
                        f"     Using existing first slide from template (layout: {layout.layout_type.value})"
                    )
                else:
                    slide = prs.slides.add_slide(content_layout)
                    print(
                        "     Added new slide for first slide (template had no slides)"
                    )
            else:
                slide = prs.slides.add_slide(content_layout)
                if is_first_slide_content:
                    print(f"     Added new slide for {layout.layout_type.value}")
                else:
                    print(
                        f"     Added new slide for {layout.layout_type.value} (content on slide {layout.slide_number})"
                    )

            # Remove template placeholders for new slides or when using existing first slide with non-base_slide layouts
            # Preserve placeholders only for base_slide layouts on first slide
            # IMPORTANT: Never clean first slide if hero stats were populated (they're already on that slide)
            should_clean_placeholders = not (
                using_existing_first_slide
                and (layout.layout_type.value == "base_slide" or should_populate_hero)
            )
            if should_clean_placeholders:
                self._clean_slide_placeholders(slide)

            # Save current state
            temp_path = f"{current_pptx}.temp{slide_idx}"
            prs.save(temp_path)
            temp_files.append(temp_path)

            # Determine correct slide index for rendering
            # Use 0 when using existing first slide, otherwise use the last slide index
            render_slide_index = (
                0 if using_existing_first_slide else len(prs.slides) - 1
            )

            # Render content blocks for this slide
            next_path = self.renderer._render_slide_content(
                section=section,
                layout=layout,
                current_file=temp_path,
                slide_index=render_slide_index,
                constraints=slide_constraints,
            )

            # Update current file
            if next_path != temp_path:
                current_pptx = next_path
            else:
                current_pptx = temp_path

            # Check for table continuation data and create continuation slides
            # Note: Continuation slides never use back layout (is_last_slide=False)
            current_pptx = self._process_table_continuations(
                section=section,
                layout=layout,
                current_pptx=current_pptx,
                slide_constraints=slide_constraints,
                content_layout_selector=lambda prs,
                pst=property_sub_type: self._select_content_layout(
                    prs, property_sub_type=pst
                ),
                metadata=metadata,
            )

        # Copy final result to output path
        if current_pptx != output_path:
            import shutil

            shutil.copy2(current_pptx, output_path)

        # Populate chart and table data
        if self.renderer.pending_data_updates:
            print(f"\n{'=' * 60}")
            print("POPULATING DATA")
            print(f"{'=' * 60}\n")

            for update_info in self.renderer.pending_data_updates:
                try:
                    slide_idx, element_type, element_idx = (
                        update_info[0],
                        update_info[1],
                        update_info[2],
                    )
                    data = update_info[3]

                    if element_type == "chart":
                        print(
                            f"Populating chart on slide {slide_idx + 1}, chart index {element_idx}..."
                        )
                        # Extract axis titles/keys and formatting from update_info (indices 4..10)
                        primary_y_axis_title = (
                            update_info[4] if len(update_info) > 4 else None
                        )
                        secondary_y_axis_title = (
                            update_info[5] if len(update_info) > 5 else None
                        )
                        x_axis_title = update_info[6] if len(update_info) > 6 else None
                        y_axis_keys = update_info[7] if len(update_info) > 7 else []
                        is_multi_axis = update_info[8] if len(update_info) > 8 else True
                        primary_y_axis_format_code = (
                            update_info[9] if len(update_info) > 9 else None
                        )
                        secondary_y_axis_format_code = (
                            update_info[10] if len(update_info) > 10 else None
                        )
                        output_path = self.renderer.chart_populator.populate_chart_data(
                            output_path,
                            slide_idx,
                            element_idx,
                            data,
                            output_path,
                            primary_y_axis_title=primary_y_axis_title,
                            secondary_y_axis_title=secondary_y_axis_title,
                            x_axis_title=x_axis_title,
                            y_axis_keys=y_axis_keys,
                            is_multi_axis=is_multi_axis,
                            primary_y_axis_format_code=primary_y_axis_format_code,
                            secondary_y_axis_format_code=secondary_y_axis_format_code,
                        )
                    elif element_type == "table":
                        # Tables may have optional table_source (5th), cell_frame_info (6th), and cached_row_heights (7th) parameters
                        table_source = update_info[4] if len(update_info) > 4 else None
                        cell_frame_info = (
                            update_info[5] if len(update_info) > 5 else None
                        )
                        cached_row_heights = (
                            update_info[6] if len(update_info) > 6 else None
                        )
                        print(
                            f"Populating table on slide {slide_idx + 1}, table index {element_idx}..."
                        )
                        output_path = self.renderer.table_populator.populate_table_data(
                            output_path,
                            slide_idx,
                            element_idx,
                            data,
                            output_path,
                            table_source=table_source,
                            cell_frame_info=cell_frame_info,
                            cached_row_heights=cached_row_heights,
                        )
                except Exception as e:
                    print(f"Warning: Error populating {element_type}: {e}")

            print("\nData population completed")

            # Clear pending updates
            self.renderer.pending_data_updates = []

        # Populate headers on all slides if metadata provided
        if metadata:
            self._populate_first_slide_headers(output_path, metadata)
            self._populate_subsequent_slide_headers(output_path, metadata)

        # Add custom template slides at the end (market stats, sub table, last slide)
        template_dir = os.path.join(
            os.path.dirname(__file__), "..", "individual_templates"
        )
        template_dir = os.path.abspath(template_dir)

        # Conditionally build template files list based on property_type
        custom_template_files = []

        # Get property_type from metadata
        property_type = (
            metadata.get("property_type", metadata.get("sector_type", ""))
            if metadata
            else ""
        )

        # If Industrial, include market stats tables
        # if property_type.lower() == "industrial":
        #     custom_template_files.extend([
        #         os.path.join(template_dir, "market_stats_table.pptx"),
        #         os.path.join(template_dir, "market_stats_sub_table.pptx")
        #     ])
        #     print(f"   Including market stats tables for Industrial property type")

        # Add configured last slide template if present
        will_add_last_slide = False
        if template_config.last_slide:
            custom_template_files.append(
                os.path.join(template_dir, template_config.last_slide)
            )
            will_add_last_slide = True

        if custom_template_files:
            self._add_custom_template_slides(output_path, custom_template_files)

            # If we added a last slide template, change its layout to "back" if available
            # if will_add_last_slide:
            #     self._apply_back_layout_to_last_slide(output_path)

        # Populate industrial custom tables (market stats) if industrial property type
        # try:
        #     if metadata and property_type.lower() == "industrial":
        #         self._populate_industrial_custom_tables(output_path, metadata)
        # except Exception as e:
        #     print(f"   ⚠️ Error populating industrial custom tables: {e}")

        # Cleanup temp files
        for temp_file in temp_files:
            if os.path.exists(temp_file) and temp_file != output_path:
                try:
                    os.remove(temp_file)
                except:
                    pass

    def preview_layouts(self, sections: List[Section]) -> List:
        """Preview layouts without generating PowerPoint"""
        all_layouts = []
        for section in sections:
            layouts = self.orchestrator.orchestrate_section(section)
            all_layouts.extend(layouts)
        return all_layouts

    def export_layout_json(self, sections: List[Section], output_path: str):
        """Export layout specifications to JSON"""
        all_layouts = self.preview_layouts(sections)
        export_layouts_to_json(all_layouts, output_path)
        print(f"✅ Layout exported to {output_path}")

    def _clean_slide_placeholders(self, slide):
        """Remove template placeholders that might conflict with our content"""
        try:
            shapes_to_remove = []

            for shape in slide.shapes:
                # Never remove hero field shapes (they contain important data)
                if hasattr(shape, "name") and shape.name:
                    shape_name = shape.name.lower()
                    if any(pattern in shape_name for pattern in ["ticker_", "arrow_"]):
                        continue  # Skip hero field shapes

                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    if hasattr(shape, "placeholder_format"):
                        placeholder_type = shape.placeholder_format.type
                        # Keep header/footer/slide number placeholders
                        if placeholder_type not in [15, 16, 17]:
                            shapes_to_remove.append(shape)

                elif hasattr(shape, "text_frame") and shape.text_frame:
                    text_content = shape.text_frame.text.lower()
                    if any(
                        phrase in text_content
                        for phrase in ["click to add", "stat", "figure"]
                    ):
                        shapes_to_remove.append(shape)

            # Remove identified shapes
            for shape in shapes_to_remove:
                try:
                    slide.shapes._spTree.remove(shape._element)
                    print(
                        f"     Removed template placeholder: {getattr(shape, 'name', 'unnamed')}"
                    )
                except:
                    pass

        except Exception as e:
            print(f"     Warning: Could not clean placeholders: {e}")
            pass

    def _populate_template_placeholders(
        self, prs, title: str, author: str, metadata: dict = None
    ) -> None:
        """
        Populate placeholder text in PowerPoint template with data from JSON

        Args:
            prs: PowerPoint presentation object
            title: Report title
            author: Author name
            metadata: Metadata from JSON containing header_prefix, market_name, etc.
        """
        try:
            if len(prs.slides) == 0:
                return

            first_slide = prs.slides[0]

            # Extract values from metadata - use separate components
            header_prefix = (
                metadata.get("header_prefix", "Figures") if metadata else "Figures"
            )
            market_name = (
                metadata.get("market_name", title.split()[0] if title else "Market")
                if metadata
                else "Market"
            )
            sector_type = (
                metadata.get(
                    "sector_type",
                    metadata.get("property_type", "Industrial")
                    if metadata
                    else "Industrial",
                )
                if metadata
                else "Industrial"
            )
            quarter = (
                " ".join(metadata.get("quarter", "2024 Q3").split()[::-1])
                if metadata
                else "Q3 2024"
            )

            print("   🔄 Populating template placeholders...")
            print(f"      Header prefix: '{header_prefix}'")
            print(f"      Market name: '{market_name}'")
            print(f"      Sector type: '{sector_type}'")
            print(f"      Quarter: '{quarter}'")

            # Dynamic replacement using selection pane approach
            # Build replacement mappings from metadata dynamically
            property_sub_type = (
                metadata.get("property_sub_type", "") if metadata else ""
            )
            replacement_values = {
                "title": title,
                "header_prefix": header_prefix,
                "market_name": market_name,
                "sector_type": sector_type,
                "quarter": quarter,
                "author": author,
                "property_sub_type": property_sub_type,
                "replacements": metadata.get("replacements", {}) if metadata else {},
            }

            print(
                f"      Available replacement values: {list(replacement_values.keys())}"
            )

            replacements_made = 0

            # Process all shapes on the first slide
            for shape_idx, shape in enumerate(first_slide.shapes):
                # Skip hero field shapes to avoid interference with populate_hero_stats
                # Hero field shapes follow patterns like: slide_1_ticker_1_text, slide_1_arrow_1_picture, etc.
                if hasattr(shape, "name") and shape.name:
                    shape_name = shape.name.lower()
                    if any(pattern in shape_name for pattern in ["ticker_", "arrow_"]):
                        continue  # Skip hero field shapes

                if hasattr(shape, "text_frame") and shape.text_frame:
                    # Get the complete text from all paragraphs and runs
                    complete_text = ""
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            complete_text += run.text

                    if complete_text.strip():
                        print(f"      Processing Shape {shape_idx}: '{complete_text}'")

                        # Apply dynamic replacements based on template content
                        updated_text = self._apply_dynamic_replacements(
                            complete_text, replacement_values
                        )

                        if updated_text != complete_text:
                            print(f"         Text changed from: '{complete_text}'")
                            print(f"         Text changed to: '{updated_text}'")

                        # If text was changed, update the entire text frame
                        if updated_text != complete_text:
                            # Store original formatting before clearing
                            original_font_size = None
                            original_font_name = None
                            original_font_bold = None

                            # Get formatting from the first run if it exists
                            if (
                                shape.text_frame.paragraphs
                                and shape.text_frame.paragraphs[0].runs
                            ):
                                first_run = shape.text_frame.paragraphs[0].runs[0]
                                if (
                                    hasattr(first_run.font, "size")
                                    and first_run.font.size
                                ):
                                    original_font_size = first_run.font.size
                                if hasattr(first_run.font, "name"):
                                    original_font_name = first_run.font.name
                                if hasattr(first_run.font, "bold"):
                                    original_font_bold = first_run.font.bold

                            # Clear all existing content and set fresh text
                            shape.text_frame.clear()

                            # Remove text frame margins to start from very top
                            shape.text_frame.margin_top = 0
                            shape.text_frame.margin_bottom = 0
                            shape.text_frame.margin_left = 0
                            shape.text_frame.margin_right = 0

                            # Set vertical anchor to top
                            from pptx.enum.text import MSO_ANCHOR

                            shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP

                            # Ensure we have a paragraph to work with
                            if len(shape.text_frame.paragraphs) == 0:
                                p = shape.text_frame.paragraphs.add()
                            else:
                                p = shape.text_frame.paragraphs[0]

                            # Clear any existing runs and add fresh text
                            p.clear()
                            run = p.add_run()
                            run.text = updated_text

                            # Restore original formatting to preserve template design
                            if original_font_size:
                                run.font.size = original_font_size
                                print(
                                    f"         Preserved font size: {original_font_size}"
                                )
                            if original_font_name:
                                run.font.name = original_font_name
                                print(f"         Preserved font: {original_font_name}")
                            if original_font_bold is not None:
                                run.font.bold = original_font_bold

                            # Remove any extra spacing/margins that might push text down
                            p.space_before = 0
                            p.space_after = 0

                            replacements_made += 1
                            print(
                                f"         Updated shape text: '{updated_text}' (margins set to 0)"
                            )

            print(f"   ✅ Made {replacements_made} placeholder replacements")

        except Exception as e:
            print(f"   ⚠️ Error populating placeholders: {e}")

    def _apply_dynamic_replacements(self, text: str, replacement_values: dict) -> str:
        """
        Apply truly dynamic replacements based on what's actually in the template
        No hardcoded patterns - let the JSON define what should be replaced
        """
        updated_text = text

        # Get replacement mappings from JSON metadata
        # The JSON should specify what text to find and what to replace it with
        replacement_mappings = replacement_values.get("replacements", {})

        if replacement_mappings:
            # Use mappings defined in JSON
            for find_text, replace_with in replacement_mappings.items():
                if find_text in updated_text:
                    old_text = updated_text
                    updated_text = updated_text.replace(find_text, replace_with)
                    if old_text != updated_text:
                        print(
                            f"         ✓ JSON Replacement: '{find_text}' → '{replace_with}'"
                        )
        else:
            # Fallback: Use auto-detection for common placeholder patterns
            # This is temporary until all templates use JSON-defined replacements
            import re

            # Auto-detect bracket patterns [anything]
            bracket_patterns = re.findall(r"\[([^\]]+)\]", text)
            for pattern_content in bracket_patterns:
                full_pattern = f"[{pattern_content}]"

                # Try to map common bracket patterns to our data
                replacement = None
                pattern_lower = pattern_content.lower()

                if "market" in pattern_lower and "submarket" not in pattern_lower:
                    replacement = replacement_values.get("market_name", "")
                elif "sector" in pattern_lower:
                    replacement = replacement_values.get("sector_type", "")
                elif "period" in pattern_lower or "quarter" in pattern_lower:
                    replacement = replacement_values.get("quarter", "")
                elif "prefix" in pattern_lower or "header" in pattern_lower:
                    replacement = replacement_values.get("header_prefix", "")

                if replacement and full_pattern in updated_text:
                    updated_text = updated_text.replace(full_pattern, replacement)
                    print(
                        f"         ✓ Auto-detected: '{full_pattern}' → '{replacement}'"
                    )

            # Handle common title placeholders
            title_patterns = [
                "Document title goes here",
                "lorem ipsum dolor sit amet",
                "Lorem ipsum dolor sit amet",
                "LOREM IPSUM DOLOR SIT AMET",
            ]

            for pattern in title_patterns:
                if pattern in updated_text:
                    replacement = (
                        replacement_values.get("title", "")
                        if "document title" in pattern.lower()
                        else ""
                    )
                    updated_text = updated_text.replace(pattern, replacement)
                    print(
                        f"         ✓ Title replacement: '{pattern}' → '{replacement}'"
                    )

            # Handle submarket-specific template placeholder
            from app.ppt_engine.ppt_helpers_utils.services.template_config import (
                get_title_strategy,
            )

            property_sub_type = replacement_values.get("property_sub_type", "")
            if property_sub_type:
                strategy = get_title_strategy(property_sub_type)
                template_placeholder = strategy.get("template_placeholder")
                if template_placeholder and template_placeholder in updated_text:
                    title_value = replacement_values.get("title", "")
                    updated_text = updated_text.replace(
                        template_placeholder, title_value
                    )
                    print(
                        f"         ✓ Title replacement: '{template_placeholder}' → '{title_value}'"
                    )

        return updated_text

    def _populate_first_slide_headers(self, pptx_path: str, metadata: dict) -> None:
        """
        Populate headers on all slides with market, sector, and quarter information

        Args:
            pptx_path: Path to the PowerPoint file
            metadata: Metadata dict with market_name, sector_type, quarter, etc.
        """
        from pptx import Presentation

        try:
            print(f"\n{'=' * 60}")
            print("🔄 POPULATING HEADERS ON ALL SLIDES")
            print(f"{'=' * 60}\n")

            # Load presentation
            prs = Presentation(pptx_path)

            # Extract values from metadata
            market_name = metadata.get("market_name", "")
            sector_type = metadata.get(
                "sector_type", metadata.get("property_type", "Industrial")
            )
            quarter = " ".join(metadata.get("quarter", "").split()[::-1])
            header_prefix = metadata.get("header_prefix", "")

            print(f"   Market: '{market_name}'")
            print(f"   Sector: '{sector_type}'")
            print(f"   Quarter: '{quarter}'")
            print(f"   Header Prefix: '{header_prefix}'")
            print()

            # Build replacement values
            replacement_values = {
                "market_name": market_name,
                "sector_type": sector_type,
                "quarter": quarter,
                "header_prefix": header_prefix,
                "replacements": metadata.get("replacements", {}),
            }

            total_replacements = 0

            # Iterate through all slides (skip first slide as it's already handled)
            for slide_idx, slide in enumerate(prs.slides):
                if slide_idx == 0:
                    continue  # Skip first slide - already populated in _populate_template_placeholders

                slide_replacements = 0

                # Process all shapes on the slide looking for header areas
                for shape in slide.shapes:
                    if not hasattr(shape, "text_frame") or not shape.text_frame:
                        continue

                    # Check if this shape is likely a header (small height, top position)
                    # Headers are typically in the top 1 inch of the slide
                    from pptx.util import Inches

                    is_likely_header = (
                        shape.top < Inches(1.0)  # Top 1 inch
                        and shape.height < Inches(0.8)  # Small height
                    )

                    # Also check for placeholder types that indicate headers
                    is_header_placeholder = False
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                        if hasattr(shape, "placeholder_format"):
                            placeholder_type = shape.placeholder_format.type
                            # Type 15 = Header, Type 16 = Footer
                            if placeholder_type == 15:
                                is_header_placeholder = True

                    # Process if it's a header
                    if is_likely_header or is_header_placeholder:
                        # Get current text
                        complete_text = ""
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                complete_text += run.text

                        if not complete_text.strip():
                            continue

                        # Try to apply replacements
                        updated_text = self._apply_dynamic_replacements(
                            complete_text, replacement_values
                        )

                        # If text changed, update the shape
                        if updated_text != complete_text:
                            # Store original formatting
                            original_font_size = None
                            original_font_name = None
                            original_font_bold = None
                            original_font_color = None

                            if (
                                shape.text_frame.paragraphs
                                and shape.text_frame.paragraphs[0].runs
                            ):
                                first_run = shape.text_frame.paragraphs[0].runs[0]
                                if (
                                    hasattr(first_run.font, "size")
                                    and first_run.font.size
                                ):
                                    original_font_size = first_run.font.size
                                if hasattr(first_run.font, "name"):
                                    original_font_name = first_run.font.name
                                if hasattr(first_run.font, "bold"):
                                    original_font_bold = first_run.font.bold
                                if hasattr(first_run.font, "color") and hasattr(
                                    first_run.font.color, "rgb"
                                ):
                                    original_font_color = first_run.font.color.rgb

                            # Clear and update text
                            shape.text_frame.clear()

                            if len(shape.text_frame.paragraphs) == 0:
                                p = shape.text_frame.paragraphs.add()
                            else:
                                p = shape.text_frame.paragraphs[0]

                            p.clear()
                            run = p.add_run()
                            run.text = updated_text

                            # Restore formatting
                            if original_font_size:
                                run.font.size = original_font_size
                            if original_font_name:
                                run.font.name = original_font_name
                            if original_font_bold is not None:
                                run.font.bold = original_font_bold
                            if original_font_color:
                                run.font.color.rgb = original_font_color

                            slide_replacements += 1
                            print(
                                f"   Slide {slide_idx + 1}: Updated header '{complete_text}' → '{updated_text}'"
                            )

                total_replacements += slide_replacements

            # Save the updated presentation
            prs.save(pptx_path)

            print(
                f"\n   ✅ Updated headers on {total_replacements} shapes across all slides"
            )
            print(f"{'=' * 60}\n")

        except Exception as e:
            print(f"   ⚠️ Error populating slide headers: {e}")
            import traceback

            traceback.print_exc()

    def _populate_subsequent_slide_headers(
        self, pptx_path: str, metadata: dict
    ) -> None:
        """
        Populate headers on all slides (including first slide) using update_shape.

        This method updates ALL instances of the slide_master_market_sector_period_header_text shape
        found on any slide (including the first slide) and in master/layouts with the formatted header text.

        Args:
            pptx_path: Path to the PowerPoint file
            metadata: Metadata dict with market_name, sector_type, quarter, header_prefix, property_sub_type
        """
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from app.ppt_engine.pptx_update import update_slide_text

        try:
            prs = Presentation(pptx_path)

            # Get property_sub_type from metadata
            property_sub_type = metadata.get("property_sub_type") if metadata else None

            # Get header format template from config
            header_format_template = get_header_format_config(property_sub_type)

            # Extract values from metadata
            market_name = metadata.get("market_name", "")
            quarter = " ".join(metadata.get("quarter", "").split()[::-1])
            header_prefix = ""
            # metadata.get("header_prefix", "")
            sector_type = metadata.get("sector_type", metadata.get("property_type", ""))
            title = metadata.get("title", "")

            # Format header text based on template
            # For figures: uses market_name and header_prefix (both uppercased)
            # For snapshot: uses market_name and sector_type (both uppercased)
            # For submarket: uses dynamic title (geography-based, uppercased)
            pst_lower = (property_sub_type or "").strip().lower()
            if pst_lower == "submarket":
                header_text = header_format_template.format(
                    title=title.upper(), quarter=quarter
                )
            elif pst_lower == "snapshot":
                header_text = header_format_template.format(
                    market_name=market_name.upper(),
                    sector_type=sector_type.upper(),
                    quarter=quarter,
                )
            else:
                # Default to figures format
                header_text = header_format_template.format(
                    market_name=market_name.upper(),
                    header_prefix=header_prefix.upper(),
                    quarter=quarter,
                )

            shape_name = "slide_master_market_sector_period_header_text"
            updated_count = 0

            # First, update all slide-specific instances (including first slide)
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.name == shape_name and shape.has_text_frame:
                        update_slide_text(shape, header_text)
                        updated_count += 1
                        print(f"   Updated header on slide {slide_idx + 1}")

            # Then, update master shape instances
            def check_master_shapes(shapes):
                found = False
                for shape in shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        # Check shapes within group
                        subshape = next(
                            (s for s in shape.shapes if s.name == shape_name), None
                        )
                        if subshape and subshape.has_text_frame:
                            update_slide_text(subshape, header_text)
                            found = True
                    elif shape.name == shape_name and shape.has_text_frame:
                        update_slide_text(shape, header_text)
                        found = True
                return found

            # Check main slide master
            if check_master_shapes(prs.slide_master.shapes):
                updated_count += 1
                print("   Updated master header shape")

            # Check all slide layouts
            for layout in prs.slide_master.slide_layouts:
                if check_master_shapes(layout.shapes):
                    updated_count += 1
                    print(f"   Updated layout header shape: {layout.name}")

            # Save the updated presentation
            prs.save(pptx_path)

            print(
                f"Subsequent slide headers updated successfully ({updated_count} instance(s) updated)"
            )

        except Exception as e:
            print(f"Error populating subsequent slide headers: {e}")
            import traceback

            traceback.print_exc()

    def _add_custom_template_slides(
        self, pptx_path: str, template_files: List[str]
    ) -> None:
        """
        Add slides from template PPTX files to the end of the presentation by directly
        importing the slide XML at the ZIP level. This preserves all formatting, backgrounds,
        images, and layouts perfectly.

        Args:
            pptx_path: Path to the main presentation file
            template_files: List of template PPTX file paths to import slides from
        """
        from pptx import Presentation

        try:
            print(f"\n{'=' * 60}")
            print("📎 ADDING CUSTOM TEMPLATE SLIDES")
            print(f"{'=' * 60}\n")

            slides_added = 0

            # Process each template file
            for template_file in template_files:
                if not os.path.exists(template_file):
                    print(f"   ⚠️ Template file not found: {template_file}")
                    continue

                template_name = os.path.basename(template_file)
                print(f"   Processing template: {template_name}")

                # Import slides from this template file
                pptx_path = self._import_template_slides_to_file(
                    pptx_path, template_file
                )

                # Count slides that were added
                from pptx import Presentation

                template_prs = Presentation(template_file)
                slides_added += len(template_prs.slides)

            print(
                f"\n   ✅ Imported {slides_added} custom slide(s) from {len(template_files)} template(s)"
            )
            print(f"{'=' * 60}\n")

        except Exception as e:
            print(f"   ⚠️ Error adding custom template slides: {e}")
            import traceback

            traceback.print_exc()

    def _import_template_slides_to_file(
        self, target_pptx_path: str, template_pptx_path: str
    ) -> str:
        """
        Import all slides from a template PPTX file to the target PPTX file using ZIP-level copy.
        This preserves ALL formatting, backgrounds, shapes, images, and relationships perfectly
        by copying the complete slide XML structure and all related parts at the PPTX ZIP level.

        Args:
            target_pptx_path: Path to the target presentation file
            template_pptx_path: Path to the template presentation file

        Returns:
            Path to the updated presentation file (same as input)
        """
        import zipfile
        import tempfile
        import shutil
        import re
        from lxml import etree
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as temp_dir:
            # Step 1: Determine how many slides are in target and template
            target_prs = Presentation(target_pptx_path)
            template_prs = Presentation(template_pptx_path)

            target_slide_count = len(target_prs.slides)
            template_slide_count = len(template_prs.slides)

            print(
                f"      Target has {target_slide_count} slides, template has {template_slide_count} slides"
            )

            # Step 2: Add blank slides to target for each template slide
            blank_layout = None
            for layout in target_prs.slide_layouts:
                if "blank" in layout.name.lower():
                    blank_layout = layout
                    break
            if blank_layout is None:
                blank_layout = target_prs.slide_layouts[-1]

            for i in range(template_slide_count):
                target_prs.slides.add_slide(blank_layout)

            # Save the target with new blank slides
            target_temp = os.path.join(temp_dir, "target_with_blanks.pptx")
            target_prs.save(target_temp)

            # Step 3: Extract both presentations as ZIP files
            source_extract = os.path.join(temp_dir, "source_extract")
            target_extract = os.path.join(temp_dir, "target_extract")

            with zipfile.ZipFile(template_pptx_path, "r") as source_zip:
                source_zip.extractall(source_extract)

            with zipfile.ZipFile(target_temp, "r") as target_zip:
                target_zip.extractall(target_extract)

            # Step 4: Copy each slide from template to target
            media_files_copied = set()

            for template_slide_idx in range(1, template_slide_count + 1):
                target_slide_idx = target_slide_count + template_slide_idx

                print(
                    f"      Importing slide {template_slide_idx} from template as slide {target_slide_idx} in target"
                )

                # Copy slide XML file
                source_slide_file = f"ppt/slides/slide{template_slide_idx}.xml"
                target_slide_file = f"ppt/slides/slide{target_slide_idx}.xml"

                source_slide_path = os.path.join(source_extract, source_slide_file)
                target_slide_path = os.path.join(target_extract, target_slide_file)

                if os.path.exists(source_slide_path):
                    shutil.copy2(source_slide_path, target_slide_path)
                    print("         ✓ Copied slide XML")
                else:
                    print(
                        f"         Warning: Source slide XML not found: {source_slide_file}"
                    )
                    continue

                # Copy slide relationships file
                source_rels_file = (
                    f"ppt/slides/_rels/slide{template_slide_idx}.xml.rels"
                )
                target_rels_file = f"ppt/slides/_rels/slide{target_slide_idx}.xml.rels"

                source_rels_path = os.path.join(source_extract, source_rels_file)
                target_rels_path = os.path.join(target_extract, target_rels_file)

                if os.path.exists(source_rels_path):
                    os.makedirs(os.path.dirname(target_rels_path), exist_ok=True)

                    # Clean notesSlide references from slide relationships before copying
                    # Multiple slides cannot share the same notesSlide - this causes corruption
                    self._clean_notes_slide_references(
                        source_rels_path, target_rels_path
                    )
                    print("         ✓ Copied and cleaned slide relationships")

                    # Copy all referenced media files and ensure chart dependencies
                    # (charts, chart rels, embedded workbooks, chart styles/colors) are present.
                    try:
                        tree = etree.parse(source_rels_path)
                        root = tree.getroot()
                        ns = {
                            "r": "http://schemas.openxmlformats.org/package/2006/relationships"
                        }

                        # Helpers to allocate unique part names in the target package.
                        def _next_chart_num(pptx_dir: str) -> int:
                            charts_dir = os.path.join(pptx_dir, "ppt", "charts")
                            max_num = 0
                            if os.path.exists(charts_dir):
                                for fn in os.listdir(charts_dir):
                                    m = re.match(r"chart(\d+)\.xml$", fn)
                                    if m:
                                        max_num = max(max_num, int(m.group(1)))
                            return max_num + 1

                        def _next_embed_num(pptx_dir: str) -> int:
                            emb_dir = os.path.join(pptx_dir, "ppt", "embeddings")
                            max_num = 0
                            if os.path.exists(emb_dir):
                                for fn in os.listdir(emb_dir):
                                    m = re.match(r"Microsoft_Excel_Worksheet(\d+)\.xlsx$", fn)
                                    if m:
                                        max_num = max(max_num, int(m.group(1)))
                            return max_num + 1

                        # Copy chart parts referenced by the slide rels, allocating unique chart
                        # and embedding names in the target to avoid collisions across templates.
                        chart_rel_updated = False
                        for rel in root.findall(".//r:Relationship", ns):
                            target_rel_path = rel.get("Target")
                            if target_rel_path and "../media/" in target_rel_path:
                                media_filename = os.path.basename(target_rel_path)

                                # Skip if already copied
                                if media_filename in media_files_copied:
                                    continue

                                source_media = os.path.join(
                                    source_extract, "ppt/media", media_filename
                                )
                                target_media_dir = os.path.join(
                                    target_extract, "ppt/media"
                                )
                                target_media = os.path.join(
                                    target_media_dir, media_filename
                                )

                                if os.path.exists(source_media):
                                    os.makedirs(target_media_dir, exist_ok=True)
                                    if not os.path.exists(target_media):
                                        shutil.copy2(source_media, target_media)
                                        media_files_copied.add(media_filename)
                                        print(
                                            f"         ✓ Copied media: {media_filename}"
                                        )

                            # Charts: slide rels point to ../charts/chartN.xml
                            if (
                                target_rel_path
                                and "../charts/" in target_rel_path
                                and os.path.basename(target_rel_path).startswith("chart")
                                and target_rel_path.endswith(".xml")
                            ):
                                src_chart_rel = target_rel_path.replace("../", "")
                                src_chart_path = os.path.join(source_extract, "ppt", src_chart_rel)

                                m = re.match(r"chart(\d+)\.xml$", os.path.basename(src_chart_rel))
                                if not m or not os.path.exists(src_chart_path):
                                    continue

                                new_chart_num = _next_chart_num(target_extract)
                                dst_chart_rel = f"charts/chart{new_chart_num}.xml"
                                dst_chart_path = os.path.join(
                                    target_extract, "ppt", dst_chart_rel
                                )
                                os.makedirs(os.path.dirname(dst_chart_path), exist_ok=True)
                                shutil.copy2(src_chart_path, dst_chart_path)

                                # Copy chart relationships (needed to resolve embedded workbook)
                                src_chart_rels_path = os.path.join(
                                    source_extract,
                                    "ppt",
                                    "charts",
                                    "_rels",
                                    os.path.basename(src_chart_rel) + ".rels",
                                )
                                if os.path.exists(src_chart_rels_path):
                                    os.makedirs(
                                        os.path.join(target_extract, "ppt", "charts", "_rels"),
                                        exist_ok=True,
                                    )
                                    rels_tree = etree.parse(src_chart_rels_path)
                                    rels_root = rels_tree.getroot()

                                    for chart_rel in rels_root.findall(".//r:Relationship", ns):
                                        chart_target = chart_rel.get("Target") or ""

                                        # Embedded workbook: rename to avoid collisions
                                        if chart_target.endswith(".xlsx") and "embeddings" in chart_target:
                                            old_embed = os.path.basename(chart_target)
                                            new_embed_num = _next_embed_num(target_extract)
                                            new_embed = f"Microsoft_Excel_Worksheet{new_embed_num}.xlsx"

                                            src_embed_path = os.path.join(
                                                source_extract, "ppt", "embeddings", old_embed
                                            )
                                            dst_embed_path = os.path.join(
                                                target_extract, "ppt", "embeddings", new_embed
                                            )
                                            os.makedirs(os.path.dirname(dst_embed_path), exist_ok=True)
                                            if os.path.exists(src_embed_path):
                                                shutil.copy2(src_embed_path, dst_embed_path)
                                                chart_rel.set(
                                                    "Target", f"../embeddings/{new_embed}"
                                                )

                                        # Chart styles/colors are in ppt/charts/ and can be shared.
                                        # Copy them if missing.
                                        if chart_target and not chart_target.startswith("..") and chart_target.endswith(".xml"):
                                            src_style = os.path.join(
                                                source_extract, "ppt", "charts", chart_target
                                            )
                                            dst_style = os.path.join(
                                                target_extract, "ppt", "charts", chart_target
                                            )
                                            if os.path.exists(src_style) and not os.path.exists(dst_style):
                                                os.makedirs(os.path.dirname(dst_style), exist_ok=True)
                                                shutil.copy2(src_style, dst_style)

                                    dst_chart_rels_path = os.path.join(
                                        target_extract,
                                        "ppt",
                                        "charts",
                                        "_rels",
                                        f"chart{new_chart_num}.xml.rels",
                                    )
                                    self._write_xml_with_office_declaration(
                                        rels_tree, dst_chart_rels_path
                                    )

                                # Update slide relationship target to the new chart part
                                rel.set("Target", f"../{dst_chart_rel}")
                                chart_rel_updated = True

                        # If we updated any slide relationship targets (e.g., charts),
                        # persist the modified relationships XML to the target slide rels file.
                        if chart_rel_updated:
                            self._write_xml_with_office_declaration(
                                tree, target_rels_path
                            )
                    except Exception as e:
                        print(f"         Warning: Error copying media files: {e}")

            # Step 4.5: Copy customXml files and docMetadata from template
            print("\n      Copying additional resources from template...")
            self._copy_template_resources(source_extract, target_extract)

            # Step 4.6: Merge content types from template to target
            print("      Merging content types...")
            self._merge_content_types(source_extract, target_extract)

            # Step 5: Repackage the presentation
            output_path = os.path.join(temp_dir, "merged.pptx")
            zip_cli = shutil.which("zip")
            if zip_cli:
                import subprocess

                if os.path.exists(output_path):
                    os.remove(output_path)

                # Prefer OS zip for better interoperability with Office.
                # -r: recurse, -X: strip extra file attrs, -q: quiet
                # -n .xlsx: store embedded workbooks without compression (they are already zip files)
                proc = subprocess.run(
                    [zip_cli, "-q", "-r", "-X", "-n", ".xlsx", output_path, "."],
                    cwd=target_extract,
                    capture_output=True,
                    text=True,
                )
                if proc.returncode != 0 or not os.path.exists(output_path):
                    stderr = (proc.stderr or "").strip()
                    print(
                        f"      Warning: zip CLI repack failed (rc={proc.returncode}); falling back to python zipfile. {stderr}"
                    )
                    zip_cli = None

            if not zip_cli:
                # Fallback: Python zipfile (best-effort)
                with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zipf:
                    for root_dir, _dirs, files in os.walk(target_extract):
                        for file in files:
                            file_path = os.path.join(root_dir, file)
                            arcname = os.path.relpath(file_path, target_extract)
                            zipf.write(file_path, arcname)

            print("      ✓ Repackaged presentation")

            # Step 6: Replace the original target file with the merged one
            shutil.copy2(output_path, target_pptx_path)
            print("      ✓ Updated target presentation")

        return target_pptx_path

    def _clean_notes_slide_references(
        self, source_rels_path: str, target_rels_path: str
    ) -> None:
        """
        Remove notesSlide references from slide relationships to prevent corruption.

        Multiple slides cannot share the same notesSlide file. When importing template slides,
        we need to remove notesSlide references since the template's notesSlide1.xml would be
        shared by multiple slides, causing PowerPoint to mark the file as corrupted.

        Args:
            source_rels_path: Path to source slide relationships file
            target_rels_path: Path to target slide relationships file (output)
        """
        from lxml import etree

        try:
            # Parse the source relationships
            tree = etree.parse(source_rels_path)
            root = tree.getroot()
            ns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}

            # Find and remove notesSlide relationships
            notes_rels = root.findall(
                './/r:Relationship[@Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"]',
                ns,
            )

            for notes_rel in notes_rels:
                root.remove(notes_rel)
                print(
                    f"         ✓ Removed notesSlide reference: {notes_rel.get('Target')}"
                )

            # Write cleaned relationships to target
            self._write_xml_with_office_declaration(tree, target_rels_path)

        except Exception as e:
            print(f"         Warning: Error cleaning notesSlide references: {e}")
            # Fallback: just copy the file as-is
            import shutil

            shutil.copy2(source_rels_path, target_rels_path)

    def _copy_template_resources(
        self, source_extract: str, target_extract: str
    ) -> None:
        """
        Copy customXml files, docMetadata, and other resources from template to target.
        These files are often referenced by slides with complex tables and need to be preserved.

        Args:
            source_extract: Path to extracted source template
            target_extract: Path to extracted target presentation
        """
        import shutil

        try:
            # Copy customXml directory (contains metadata referenced by tables/slides)
            source_customxml = os.path.join(source_extract, "customXml")
            target_customxml = os.path.join(target_extract, "customXml")

            if os.path.exists(source_customxml):
                # If target already has customXml, we need to merge intelligently
                if os.path.exists(target_customxml):
                    # Copy files that don't exist in target
                    for filename in os.listdir(source_customxml):
                        source_file = os.path.join(source_customxml, filename)
                        target_file = os.path.join(target_customxml, filename)

                        if os.path.isfile(source_file) and not os.path.exists(
                            target_file
                        ):
                            shutil.copy2(source_file, target_file)
                            print(f"         ✓ Copied customXml/{filename}")

                    # Also copy _rels subdirectory
                    source_rels = os.path.join(source_customxml, "_rels")
                    target_rels = os.path.join(target_customxml, "_rels")

                    if os.path.exists(source_rels):
                        os.makedirs(target_rels, exist_ok=True)
                        for filename in os.listdir(source_rels):
                            source_file = os.path.join(source_rels, filename)
                            target_file = os.path.join(target_rels, filename)

                            if os.path.isfile(source_file) and not os.path.exists(
                                target_file
                            ):
                                shutil.copy2(source_file, target_file)
                                print(f"         ✓ Copied customXml/_rels/{filename}")
                else:
                    # No existing customXml, copy entire directory
                    shutil.copytree(source_customxml, target_customxml)
                    print("         ✓ Copied entire customXml directory")

            # Copy docMetadata directory (contains classification labels)
            source_docmetadata = os.path.join(source_extract, "docMetadata")
            target_docmetadata = os.path.join(target_extract, "docMetadata")

            if os.path.exists(source_docmetadata):
                if not os.path.exists(target_docmetadata):
                    shutil.copytree(source_docmetadata, target_docmetadata)
                    print("         ✓ Copied docMetadata directory")
                else:
                    # Merge docMetadata files
                    for filename in os.listdir(source_docmetadata):
                        source_file = os.path.join(source_docmetadata, filename)
                        target_file = os.path.join(target_docmetadata, filename)

                        if os.path.isfile(source_file) and not os.path.exists(
                            target_file
                        ):
                            shutil.copy2(source_file, target_file)
                            print(f"         ✓ Copied docMetadata/{filename}")

        except Exception as e:
            print(f"         Warning: Error copying template resources: {e}")
            import traceback

            traceback.print_exc()

    def _merge_content_types(self, source_extract: str, target_extract: str) -> None:
        """
        Merge [Content_Types].xml from template into target to register all content types.
        This is critical - PowerPoint will mark the file as corrupted if content types are missing.

        Args:
            source_extract: Path to extracted source template
            target_extract: Path to extracted target presentation
        """
        from lxml import etree

        try:
            target_content_types_path = os.path.join(
                target_extract, "[Content_Types].xml"
            )
            source_content_types_path = os.path.join(
                source_extract, "[Content_Types].xml"
            )

            if not os.path.exists(target_content_types_path):
                print("         Warning: [Content_Types].xml not found in target")
                return

            if not os.path.exists(source_content_types_path):
                print("         Warning: [Content_Types].xml not found in source")
                return

            # Parse both files
            target_tree = etree.parse(target_content_types_path)
            target_root = target_tree.getroot()

            source_tree = etree.parse(source_content_types_path)
            source_root = source_tree.getroot()

            ns = {"ct": "http://schemas.openxmlformats.org/package/2006/content-types"}

            # Get existing content types from target
            existing_extensions = set()
            existing_part_names = set()

            for default in target_root.findall(".//ct:Default", ns):
                ext = default.get("Extension")
                if ext:
                    existing_extensions.add(ext)

            for override in target_root.findall(".//ct:Override", ns):
                part_name = override.get("PartName")
                if part_name:
                    existing_part_names.add(part_name)

            added_count = 0

            # Add missing Default elements (file extensions) from source
            for default in source_root.findall(".//ct:Default", ns):
                ext = default.get("Extension")
                content_type = default.get("ContentType")
                if ext and ext not in existing_extensions:
                    new_default = etree.SubElement(
                        target_root,
                        "{http://schemas.openxmlformats.org/package/2006/content-types}Default",
                    )
                    new_default.set("Extension", ext)
                    new_default.set("ContentType", content_type)
                    existing_extensions.add(ext)
                    added_count += 1
                    print(f"         Added content type for .{ext}")

            # Add missing Override elements (specific part names) from source
            for override in source_root.findall(".//ct:Override", ns):
                part_name = override.get("PartName")
                content_type = override.get("ContentType")
                if part_name and part_name not in existing_part_names:
                    # Don't add overrides for slide-specific paths (they're renumbered)
                    if (
                        "/slides/slide" in part_name
                        or "/notesSlides/notesSlide" in part_name
                    ):
                        continue

                    new_override = etree.SubElement(
                        target_root,
                        "{http://schemas.openxmlformats.org/package/2006/content-types}Override",
                    )
                    new_override.set("PartName", part_name)
                    new_override.set("ContentType", content_type)
                    existing_part_names.add(part_name)
                    added_count += 1
                    print(f"         Added content type: {part_name}")

            # Write updated content types back
            if added_count > 0:
                self._write_xml_with_office_declaration(
                    target_tree, target_content_types_path
                )
                print(f"         ✓ Merged {added_count} content type entries")
            else:
                print("         ✓ Content types already up to date")

        except Exception as e:
            print(f"         Warning: Error merging content types: {e}")
            import traceback

            traceback.print_exc()

    def populate_custom_slides(self, pptx_path: str, slide_data: dict) -> None:
        """
        Populate data in the custom template slide that was added at the end.

        Args:
            pptx_path: Path to the presentation file
            slide_data: Dictionary with slide-specific data
                Format: {
                    'last_slide': {
                        'table_data': [[row1_data], [row2_data], ...],
                        'text_replacements': {'placeholder_text': 'replacement_value'}
                    }
                }
        """
        from pptx import Presentation

        try:
            print(f"\n{'=' * 60}")
            print("📝 POPULATING CUSTOM SLIDES DATA")
            print(f"{'=' * 60}\n")

            # Load presentation
            prs = Presentation(pptx_path)
            total_slides = len(prs.slides)

            # The custom slide is the last slide
            if total_slides < 1:
                print(
                    f"   ⚠️ Not enough slides to populate custom data (found {total_slides})"
                )
                return

            # Map slide name to index (last slide)
            slide_mapping = {"last_slide": total_slides - 1}

            # Process each custom slide
            for slide_name, slide_idx in slide_mapping.items():
                if slide_name not in slide_data:
                    print(f"   ⚠️ No data provided for {slide_name}, skipping...")
                    continue

                print(f"   Processing slide: {slide_name} (index {slide_idx})")
                slide = prs.slides[slide_idx]
                data = slide_data[slide_name]

                # Handle table data
                if "table_data" in data and data["table_data"]:
                    self._populate_slide_table(slide, data["table_data"])

                # Handle text replacements
                if "text_replacements" in data and data["text_replacements"]:
                    self._populate_slide_text(slide, data["text_replacements"])

                print("      ✓ Slide populated successfully")

            # Save updated presentation
            prs.save(pptx_path)

            print("\n   ✅ Custom slide populated successfully")
            print(f"{'=' * 60}\n")

        except Exception as e:
            print(f"   ⚠️ Error populating custom slide: {e}")
            import traceback

            traceback.print_exc()

    def _populate_slide_table(self, slide, table_data):
        """
        Populate table data in a slide.

        Args:
            slide: Slide object
            table_data: List of lists containing table data [[row1], [row2], ...]
        """

        # Find tables in the slide
        tables = [shape for shape in slide.shapes if shape.has_table]

        if not tables:
            print("         Warning: No tables found in slide")
            return

        # Use the first table
        table_shape = tables[0]
        table = table_shape.table

        print(f"         Populating table with {len(table_data)} rows")

        # Clear existing data rows (keep header if exists)
        while len(table.rows) > 1:
            table._tbl.remove(table._tbl.tr_lst[-1])

        # Add new rows with data
        for row_data in table_data:
            row = table.rows.add()
            row_idx = len(table.rows) - 1
            
            # Detect if this is a TOTAL row (first cell matches "total" case-insensitively,
            # with markdown variations like **TOTAL**, __total__, etc.)
            first_cell_value = row_data[0] if row_data else ""
            is_total_row = is_total_label(first_cell_value)
            
            for col_idx, cell_value in enumerate(row_data):
                if col_idx < len(table.columns):
                    cell = table.cell(row_idx, col_idx)
                    
                    # For TOTAL row first cell, use normalized display text
                    if is_total_row and col_idx == 0:
                        cell.text = total_display_text(cell_value)
                    else:
                        # Data is already transformed by frontend_json_processor
                        cell.text = str(cell_value) if cell_value is not None else ""
                    
                    # Apply bold formatting to all cells in TOTAL rows
                    if is_total_row:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.bold = True

    def _populate_slide_text(self, slide, text_replacements):
        """
        Replace placeholder text in slide shapes.

        Args:
            slide: Slide object
            text_replacements: Dictionary of {find_text: replace_text}
        """
        print(f"         Applying {len(text_replacements)} text replacements")

        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue

            # Get current text
            complete_text = ""
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    complete_text += run.text

            # Check if any replacement applies
            new_text = complete_text
            for find_text, replace_text in text_replacements.items():
                if find_text in new_text:
                    new_text = new_text.replace(find_text, replace_text)

            # Update text if changed
            if new_text != complete_text:
                # Store formatting
                original_font_size = None
                original_font_name = None
                original_font_bold = None

                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    first_run = shape.text_frame.paragraphs[0].runs[0]
                    if hasattr(first_run.font, "size") and first_run.font.size:
                        original_font_size = first_run.font.size
                    if hasattr(first_run.font, "name"):
                        original_font_name = first_run.font.name
                    if hasattr(first_run.font, "bold"):
                        original_font_bold = first_run.font.bold

                # Update text
                shape.text_frame.clear()
                if len(shape.text_frame.paragraphs) == 0:
                    p = shape.text_frame.paragraphs.add()
                else:
                    p = shape.text_frame.paragraphs[0]

                p.clear()
                run = p.add_run()
                run.text = new_text

                # Restore formatting
                if original_font_size:
                    run.font.size = original_font_size
                if original_font_name:
                    run.font.name = original_font_name
                if original_font_bold is not None:
                    run.font.bold = original_font_bold

                print(f"            ✓ Replaced '{find_text}' with '{replace_text}'")

    def _populate_hero_stats(self, pptx_path: str, hero_fields: dict) -> None:
        """
        Populate hero stats on the first slide with trend arrows and values

        Args:
            pptx_path: Path to PowerPoint file
            hero_fields: Dictionary containing stats data with trend and value
        """
        try:
            print(f"\n{'=' * 60}")
            print("🔄 POPULATING HERO STATS")
            print(f"{'=' * 60}\n")

            # Use the hero stats populator
            populate_hero_stats(pptx_path, hero_fields, pptx_path)

            print("\n   ✅ Hero stats population completed")
            print(f"{'=' * 60}\n")

        except Exception as e:
            print(f"   ⚠️ Error populating hero stats: {e}")
            import traceback

            traceback.print_exc()

    # # -------------------------------------------------------------
    # # Industrial Custom Template Table Population
    # # -------------------------------------------------------------
    # def _populate_industrial_custom_tables(self, pptx_path: str, metadata: dict) -> None:
    #     """Populate data tables on imported industrial custom template slides.

    #     Mapping of shape names (selection pane names) to data functions:
    #       slide_5_table_1 -> fetch_industrial_size_bucket_data
    #       slide_5_table_2 -> fetch_industrial_product_type_data
    #       slide_5_table_3 -> fetch_industrial_class_data
    #       slide_6_table_1 -> fetch_industrial_submarket_data

    #     IMPORTANT: Numeric fragments in shape names do not correspond to real slide indices
    #     post-import; we therefore scan all slides' shapes and match by shape.name.
    #     """
    #     from pptx import Presentation
    #     import pandas as pd
    #     print(f"\n{'='*60}")
    #     print("📊 POPULATING INDUSTRIAL TEMPLATE TABLES")
    #     print(f"{'='*60}\n")

    #     prs = Presentation(pptx_path)
    #     # Inline config creation using metadata (functions access attributes via dot notation)
    #     from types import SimpleNamespace
    #     md = dict(metadata)  # shallow copy
    #     # Ensure defined_markets list
    #     if not md.get('defined_markets'):
    #         md['defined_markets'] = [md.get('market_name') or md.get('defined_market_name') or 'Unknown']
    #     # Singular convenience
    #     md.setdefault('defined_market_name', md['defined_markets'][0])
    #     # Align quarter/current_quarter
    #     if 'quarter' not in md and 'current_quarter' in md:
    #         md['quarter'] = md['current_quarter']
    #     if 'current_quarter' not in md and 'quarter' in md:
    #         md['current_quarter'] = md['quarter']
    #     # Monthly/yearly select fallback from asking_rate_frequency
    #     md.setdefault('monthly_yearly_select', md.get('asking_rate_frequency', 'yearly'))
    #     md.setdefault('asking_rate_field', 'DIRECT_NET_ASKING_RATE_NUMERATOR / DIRECT_NET_ASKING_RATE_DENOMINATOR')
    #     md.setdefault('publishing_group', 'US')
    #     md.setdefault('asking_rate_type', md.get('asking_rate_type', 'average'))
    #     md.setdefault('absorption_calculation', 'standard')
    #     md.setdefault('total_vs_direct_absorption', 'direct')
    #     cfg = SimpleNamespace(**md)

    #     # Import data functions
    #     from app.ppt_engine.ppt_helpers_utils.static_data.industrial_figures import (
    #         fetch_industrial_size_bucket_data,
    #         fetch_industrial_product_type_data,
    #         fetch_industrial_class_data,
    #         fetch_industrial_submarket_data,
    #     )

    #     mapping: Dict[str, Any] = {
    #         "slide_5_table_1": fetch_industrial_size_bucket_data,
    #         "slide_5_table_2": fetch_industrial_product_type_data,
    #         "slide_5_table_3": fetch_industrial_class_data,
    #         "slide_6_table_1": fetch_industrial_submarket_data,
    #     }

    #     dataset_cache: Dict[str, pd.DataFrame] = {}
    #     populated = 0

    #     # Use the shared table_populator's populate_table_data for robust table updates
    #     table_populator = getattr(self.renderer, 'table_populator', None)
    #     if table_populator is None:
    #         from app.ppt_engine.ppt_helpers_utils.ppt_helpers.data_populator import TableDataPopulator
    #         table_populator = TableDataPopulator()

    #     for slide_idx, slide in enumerate(prs.slides):
    #         for shape_idx, shape in enumerate(slide.shapes):
    #             if not getattr(shape, 'has_table', False):
    #                 continue
    #             shape_name = getattr(shape, 'name', '')
    #             if shape_name in mapping:
    #                 func = mapping[shape_name]
    #                 print(f"   • Found target table shape '{shape_name}' on slide {slide_idx + 1}")
    #                 if shape_name not in dataset_cache:
    #                     try:
    #                         df = func(cfg)
    #                         if not isinstance(df, pd.DataFrame):
    #                             raise ValueError("Returned object is not a DataFrame")
    #                         dataset_cache[shape_name] = df
    #                         print(f"     ✓ Retrieved dataset rows={len(df)} cols={len(df.columns)}")
    #                     except Exception as e:
    #                         print(f"     ⚠️ Data fetch failed for {shape_name}: {e}")
    #                         continue
    #                 df = dataset_cache[shape_name]
    #                 try:
    #                     # Convert DataFrame to list of dicts for populate_table_data
    #                     data_dicts = df.to_dict(orient='records')
    #                     # Find the table index on this slide
    #                     table_shapes = [s for s in slide.shapes if getattr(s, 'has_table', False)]
    #                     table_index = table_shapes.index(shape)
    #                     # Use the robust populator
    #                     table_populator.populate_table_data(
    #                         pptx_path,
    #                         slide_idx,
    #                         table_index,
    #                         data_dicts,
    #                         pptx_path
    #                     )
    #                     populated += 1
    #                     print(f"     ✅ Populated table '{shape_name}' using shared populator")
    #                 except Exception as e:
    #                     print(f"     ⚠️ Population error for '{shape_name}': {e}")

    #     print(f"\n   ✅ Industrial table population complete. Tables updated: {populated}")
    #     print(f"{'='*60}\n")

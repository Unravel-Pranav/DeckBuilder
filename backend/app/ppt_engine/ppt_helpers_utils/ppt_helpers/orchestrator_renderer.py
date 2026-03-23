"""
Orchestrator Renderer - Bridge between layout specs and PowerPoint generation

This module takes the JSON layout specifications from the orchestrator
and renders them into actual PowerPoint presentations using test_open_xml.py
"""

from app.ppt_engine.ppt_helpers_utils.ppt_helpers.slide_orchestrator import (
    SlideOrchestrator,
    Section,
    TextBlock,
    ChartBlock,
    TableBlock,
    SlideLayout,
    ContentType,
)
from app.ppt_engine.ppt_helpers_utils.ppt_helpers.test_open_xml import (
    clone_chart,
    clone_table,
)
from app.ppt_engine.ppt_helpers_utils.ppt_helpers.data_populator import (
    ChartDataPopulator,
    TableDataPopulator,
)
from pptx import Presentation
from pptx.util import Pt
from typing import List, Optional
import os

from app.ppt_engine.ppt_helpers_utils.services.template_config import (
    SlideConstraints,
    get_ppt_template_config,
    get_slide_constraint_profile,
)


class OrchestratorRenderer:
    """
    Renders layout specifications into PowerPoint presentations.

    Takes the pure layout output from SlideOrchestrator and creates
    actual PPTX files using existing chart/table cloning tools.
    """

    def __init__(self, temp_dir=None):
        """Initialize renderer."""
        self.temp_counter = 0
        # Use provided temp_dir or default to output_ppt
        if temp_dir is None:
            temp_dir = os.path.join(os.path.dirname(__file__), "..", "output_ppt")
            temp_dir = os.path.abspath(temp_dir)
            os.makedirs(temp_dir, exist_ok=True)
        self.temp_dir = temp_dir

        # Set templates directory to individual_templates
        self.templates_dir = os.path.join(
            os.path.dirname(__file__), "..", "individual_templates"
        )
        self.templates_dir = os.path.abspath(self.templates_dir)

        # Initialize data populators
        self.chart_populator = ChartDataPopulator()
        self.table_populator = TableDataPopulator()

        # Track chart/table indices for data population
        self.pending_data_updates = []  # List of (slide_idx, element_type, element_idx, data)

        # Track if first block has been rendered (for section title placement)
        self._first_block_rendered_per_section = {}  # Track by section_id which section has had its first block rendered

        # Track which sections have rendered their first block (for section title rendering)
        # This prevents duplicate section titles within the same section

    def render_to_pptx(
        self,
        section: Section,
        layouts: List[SlideLayout],
        output_path: str,
        template_pptx: str = None,
        property_sub_type: Optional[str] = None,
    ) -> str:
        """
        Render layout specifications to PowerPoint file.

        Note: This method expects that blocks already have data populated.
        Layout orchestration should happen AFTER data is available in blocks
        to ensure accurate size estimation.

        Args:
            section: The section with content blocks (blocks should have data populated)
            layouts: Layout specifications from orchestrator
            output_path: Path for output PPTX file
            template_pptx: Optional base template to start from
            property_sub_type: Determines which template set should be used

        Returns:
            Path to generated PPTX file
        """
        print(f"\n{'=' * 60}")
        print(f"RENDERING: {section.title}")
        print(f"{'=' * 60}\n")

        # Reset trackers for clean state (in case renderer instance is reused)
        self._first_block_rendered_per_section = {}

        # Verify data availability in blocks (for logging/debugging)
        for block in section.blocks:
            if block.type == ContentType.TABLE:
                data_count = (
                    len(block.data) if hasattr(block, "data") and block.data else 0
                )
                if data_count == 0:
                    print(
                        f"  ⚠ Warning: TableBlock {block.id} has no data - using estimated dimensions"
                    )
            elif block.type == ContentType.CHART:
                data_count = (
                    len(block.data) if hasattr(block, "data") and block.data else 0
                )
                if data_count == 0:
                    print(
                        f"  ⚠ Warning: ChartBlock {block.id} has no data - using estimated dimensions"
                    )

        # Create base presentation using appropriate template
        # Get template directory - using individual_templates for chart-type-based templates
        template_dir = os.path.join(
            os.path.dirname(__file__), "..", "individual_templates"
        )
        template_dir = os.path.abspath(template_dir)
        template_config = get_ppt_template_config(property_sub_type)
        slide_constraints = get_slide_constraint_profile(property_sub_type)
        base_template_path = os.path.join(template_dir, template_config.base_slide)
        first_slide_template = (
            os.path.join(template_dir, template_config.first_slide)
            if template_config.first_slide
            else None
        )

        # Check if we need the first slide base template
        use_first_slide_base = any(
            layout.layout_type.value == "base_slide" for layout in layouts
        )

        if (
            use_first_slide_base
            and first_slide_template
            and os.path.exists(first_slide_template)
        ):
            print(f"Using first slide base template: {first_slide_template}")
            prs = Presentation(first_slide_template)
        elif template_pptx and os.path.exists(template_pptx):
            prs = Presentation(template_pptx)
        elif os.path.exists(base_template_path):
            print(f"Using base template: {base_template_path}")
            prs = Presentation(base_template_path)
        else:
            prs = Presentation()

        # Find the best layout - prefer blank layout with headers/footers
        content_layout = None
        for i, layout in enumerate(prs.slide_layouts):
            print(f"Available Layout {i}: {layout.name}")
            # Look for blank or minimal content layouts (avoid layouts with placeholders)
            if "blank" in layout.name.lower():
                content_layout = layout
                print(f"Using blank layout: {layout.name}")
                break

        # If no blank found, use the last layout (often blank) or fallback
        if content_layout is None:
            if len(prs.slide_layouts) > 6:
                content_layout = prs.slide_layouts[6]  # Standard blank layout
                print("Using standard blank layout")
            else:
                content_layout = prs.slide_layouts[-1]  # Last layout
                print(f"Using last layout: {content_layout.name}")

        # Track current file as we build it
        current_file = output_path

        # Render each slide
        for i, layout in enumerate(layouts):
            print(
                f"Rendering Slide {layout.slide_number} ({layout.layout_type.value})..."
            )

            # Handle BASE_SLIDE differently - use existing slide from template
            if (
                layout.layout_type.value == "base_slide"
                and i == 0
                and use_first_slide_base
            ):
                # Use the existing first slide from the template (already has KPIs and title)
                if len(prs.slides) > 0:
                    slide = prs.slides[0]  # Use existing first slide
                    print("  Using existing first slide from base template")
                else:
                    slide = prs.slides.add_slide(content_layout)
                    print("  Added new slide (base template had no slides)")
            else:
                # Add new slide using base template layout
                slide = prs.slides.add_slide(content_layout)

            # Save current state
            temp_file = os.path.join(self.temp_dir, f"temp_render_{i}.pptx")
            prs.save(temp_file)

            # Add content to this slide
            next_file = self._render_slide_content(
                section=section,
                layout=layout,
                current_file=temp_file,
                slide_index=len(prs.slides) - 1,
                constraints=slide_constraints,
            )

            # Load the updated presentation
            if next_file != temp_file:
                prs = Presentation(next_file)
                current_file = next_file
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            else:
                current_file = temp_file

        # Save final presentation
        if current_file != output_path:
            if os.path.exists(output_path):
                os.remove(output_path)
            os.rename(current_file, output_path)

        # Populate chart and table data into PPTX file
        # NOTE: This step writes the data (that was already used for layout calculation) into the actual PowerPoint file.
        # The data was already available in blocks during orchestration, so layouts were calculated using actual dimensions.
        # This step is still necessary to populate the rendered PPTX with the actual data values.
        if self.pending_data_updates:
            print(f"\n{'=' * 60}")
            print("POPULATING DATA INTO PPTX")
            print(f"{'=' * 60}\n")

            for update_info in self.pending_data_updates:
                try:
                    slide_idx, element_type, element_idx = (
                        update_info[0],
                        update_info[1],
                        update_info[2],
                    )
                    data = update_info[3]

                    if element_type == "chart":
                        print(
                            f"Populating chart on slide {slide_idx + 1}, chart index {element_idx}..."
                        )
                        # Extract axis titles/keys and formatting from update_info (indices 4..10)
                        primary_y_axis_title = (
                            update_info[4] if len(update_info) > 4 else None
                        )
                        secondary_y_axis_title = (
                            update_info[5] if len(update_info) > 5 else None
                        )
                        x_axis_title = update_info[6] if len(update_info) > 6 else None
                        y_axis_keys = update_info[7] if len(update_info) > 7 else []
                        is_multi_axis = update_info[8] if len(update_info) > 8 else True
                        primary_y_axis_format_code = (
                            update_info[9] if len(update_info) > 9 else None
                        )
                        secondary_y_axis_format_code = (
                            update_info[10] if len(update_info) > 10 else None
                        )
                        output_path = self.chart_populator.populate_chart_data(
                            output_path,
                            slide_idx,
                            element_idx,
                            data,
                            output_path,
                            primary_y_axis_title=primary_y_axis_title,
                            secondary_y_axis_title=secondary_y_axis_title,
                            x_axis_title=x_axis_title,
                            y_axis_keys=y_axis_keys,
                            is_multi_axis=is_multi_axis,
                            primary_y_axis_format_code=primary_y_axis_format_code,
                            secondary_y_axis_format_code=secondary_y_axis_format_code,
                        )
                    elif element_type == "table":
                        # Tables may have optional table_source (5th), cell_frame_info (6th), and cached_row_heights (7th) parameters
                        table_source = update_info[4] if len(update_info) > 4 else None
                        cell_frame_info = (
                            update_info[5] if len(update_info) > 5 else None
                        )
                        cached_row_heights = (
                            update_info[6] if len(update_info) > 6 else None
                        )
                        print(
                            f"Populating table on slide {slide_idx + 1}, table index {element_idx}..."
                        )
                        output_path = self.table_populator.populate_table_data(
                            output_path,
                            slide_idx,
                            element_idx,
                            data,
                            output_path,
                            table_source=table_source,
                            cell_frame_info=cell_frame_info,
                            cached_row_heights=cached_row_heights,
                        )
                except Exception as e:
                    print(f"⚠ Error populating {element_type}: {e}")

            print("\n✓ Data population completed")

            # Clear pending updates for next section
            self.pending_data_updates = []

        # Clean up any temp files
        self._cleanup_temp_files()

        print(f"\n✓ Presentation saved: {output_path}")
        print(f"  Total slides: {len(layouts)}")
        print(f"{'=' * 60}\n")

        return output_path

    def _render_slide_content(
        self,
        section: Section,
        layout: SlideLayout,
        current_file: str,
        slide_index: int,
        constraints: Optional[SlideConstraints] = None,
    ) -> str:
        """
        Render content for a single slide.

        Args:
            section: Section containing blocks
            layout: Layout specification
            current_file: Current PPTX file path
            slide_index: Which slide to add content to (0-based)

        Returns:
            Path to updated PPTX file
        """
        output_file = current_file

        # NOTE: Do NOT reset _first_block_rendered_per_section here!
        # Section titles should only appear on the FIRST element of each section
        # across the entire presentation, not on every slide.

        # Store section title for ALL text blocks (don't clear it)
        self._section_title_for_all = section.title if section.title else None
        resolved_constraints = constraints or getattr(section, "constraints", None)

        # Get blocks from section
        block_dict = {block.id: block for block in section.blocks}

        # Group blocks by their section_id to determine first block per logical section
        blocks_by_section = {}
        for block in section.blocks:
            # Get section_id from block (set by _group_blocks_by_slide_number)
            block_section_id = getattr(block, "section_id", None) or section.id
            if block_section_id not in blocks_by_section:
                blocks_by_section[block_section_id] = []
            blocks_by_section[block_section_id].append(block)

        # Determine the first block ID for each logical section using display_order
        first_block_per_section = {}
        for section_id, blocks in blocks_by_section.items():
            # Try to find block with smallest display_order
            blocks_with_order = [
                b
                for b in blocks
                if hasattr(b, "display_order") and b.display_order is not None
            ]
            if blocks_with_order:
                first_block = min(blocks_with_order, key=lambda b: b.display_order)
                first_block_per_section[section_id] = first_block.id
            else:
                # Fallback to first block in list
                first_block_per_section[section_id] = blocks[0].id

        # Render each assigned block
        for cell_idx, block_id in layout.assigned_blocks:
            block = block_dict.get(block_id)
            if not block:
                print(f"  ⚠ Block {block_id} not found in section")
                continue

            cell = layout.cell_frames[cell_idx]

            print(f"  Adding {block.type.value}: {block_id}")
            print(f'    Position: ({cell.left:.2f}", {cell.top:.2f}")')
            print(f'    Size: {cell.width:.2f}" × {cell.height:.2f}"')

            # Determine if we should pass section title to this block
            # Get the logical section this block belongs to
            block_section_id = getattr(block, "section_id", None) or section.id

            # Determine section title for this block (use section_name from block, fallback to section.title)
            # When blocks are from different logical sections on same slide, they each get their own section title
            block_section_title = getattr(block, "section_name", None) or section.title

            # Pass section title to first block of each logical section (regardless of type)
            section_title_for_block = None
            is_first_block_of_section = block_id == first_block_per_section.get(
                block_section_id
            )
            already_rendered = self._first_block_rendered_per_section.get(
                block_section_id, False
            )

            # Check if section titles should be shown (from section style config)
            show_section_title = (
                section.style.get("show_title", True)
                if hasattr(section, "style") and section.style
                else True
            )

            # Only pass section title to the FIRST block of each section, if configured to show titles
            if (
                is_first_block_of_section
                and not already_rendered
                and block_section_title
                and show_section_title
            ):
                section_title_for_block = block_section_title
                self._first_block_rendered_per_section[block_section_id] = True
                print(
                    f"    [SECTION TITLE] Adding to first element: '{block_section_title}' (type={block.type.value})"
                )

            # Render based on type
            if block.type == ContentType.CHART:
                output_file = self._render_chart(
                    block=block,
                    cell=cell,
                    current_file=output_file,
                    slide_index=slide_index,
                    constraints=resolved_constraints,
                    layout_type=layout.layout_type,
                    section_title=section_title_for_block,
                )
            elif block.type == ContentType.TABLE:
                output_file = self._render_table(
                    block=block,
                    cell=cell,
                    current_file=output_file,
                    slide_index=slide_index,
                    constraints=resolved_constraints,
                    layout_type=layout.layout_type,
                    section_title=section_title_for_block,
                )
            elif block.type == ContentType.TEXT:
                output_file = self._render_text(
                    block=block,
                    cell=cell,
                    current_file=output_file,
                    slide_index=slide_index,
                    section_title=section_title_for_block,
                )

        return output_file

    def _render_chart(
        self,
        block: ChartBlock,
        cell,
        current_file: str,
        slide_index: int,
        constraints=None,
        layout_type=None,
        section_title: Optional[str] = None,
    ) -> str:
        """Render a chart block."""
        if not os.path.exists(block.template_path):
            print(f"    ⚠ Template not found: {block.template_path}")
            return current_file

        # Derive slide geometry from template constraints
        constraint_profile = constraints or SlideConstraints()
        slide_width = constraint_profile.slide_width
        margin_right = constraint_profile.margin_right
        right_boundary = slide_width - margin_right

        # Calculate if chart would exceed the REAL boundaries
        chart_right_edge = cell.left + cell.width
        if chart_right_edge > right_boundary:
            available_width = max(right_boundary - cell.left, 0.0)

            # CRITICAL: Prevent negative or zero widths (would corrupt file!)
            if available_width <= 0.5:
                print(
                    f'    ❌ ERROR: Chart position ({cell.left:.2f}") too far right! Cannot fit on slide.'
                )
                print("       Skipping this chart to prevent file corruption.")
                return current_file  # Skip this chart

            print(
                f'    ⚠ Chart too wide ({cell.width:.2f}"), constraining to {available_width:.2f}"'
            )
            cell_width = available_width
            # Maintain aspect ratio
            scale_factor = available_width / cell.width
            cell_height = cell.height * scale_factor
        else:
            # Use the full cell size - no artificial constraints
            cell_width = cell.width
            cell_height = cell.height

        # Get label dimensions from config
        from app.ppt_engine.ppt_helpers_utils.services.template_config import (
            get_element_dimensions,
        )

        element_dims = get_element_dimensions()
        FIGURE_LABEL_HEIGHT = element_dims.figure_label_height
        FIGURE_GAP = element_dims.figure_gap
        SOURCE_LABEL_HEIGHT = element_dims.source_label_height
        SOURCE_GAP = element_dims.source_gap

        # Always reserve space for labels within cell (labels shown/hidden based on config)
        # Cell height already includes this space from estimate_intrinsic_size()
        figure_label_space = FIGURE_LABEL_HEIGHT + FIGURE_GAP

        # Add section title space if provided
        if section_title:
            SECTION_TITLE_HEIGHT = FIGURE_LABEL_HEIGHT  # Same height as figure label
            SECTION_TITLE_GAP = FIGURE_GAP  # Same gap as figure label
            figure_label_space += SECTION_TITLE_HEIGHT + SECTION_TITLE_GAP

        source_label_space = SOURCE_GAP + SOURCE_LABEL_HEIGHT
        total_label_space = figure_label_space + source_label_space

        # Chart should start below FIGURE label space (always reserved)
        chart_top = cell.top + figure_label_space

        # Chart height fits in the space between figure label and source label
        # The cell.height already includes label space (from estimate_intrinsic_size)
        available_chart_height = cell.height - total_label_space
        chart_height = min(cell_height, available_chart_height)

        # Ensure chart_height is positive
        if chart_height <= 0:
            chart_height = max(
                0.1, cell.height - total_label_space
            )  # Minimum 0.1" if labels take all space

        # Centering logic for single column stacked charts
        # These charts should be narrower and centered horizontally
        chart_left = cell.left
        chart_type = getattr(block, "chart_type", "").lower()
        is_single_column_stacked = "single_column_stacked" in chart_type or "single column stacked" in chart_type
        
        if is_single_column_stacked:
            # Single column stacked charts should be centered
            # Use 50% of the available cell width for the chart
            centered_width_ratio = 0.50
            centered_chart_width = cell_width * centered_width_ratio
            # Ensure minimum width of 2.5 inches
            centered_chart_width = max(centered_chart_width, 2.5)
            # Don't exceed original width
            centered_chart_width = min(centered_chart_width, cell_width)
            
            # Calculate left position to center the chart within the cell
            horizontal_padding = (cell_width - centered_chart_width) / 2
            chart_left = cell.left + horizontal_padding
            cell_width = centered_chart_width
            
            print(f"    📊 Single column stacked chart: centering with width {cell_width:.2f}\" at left {chart_left:.2f}\"")

        temp_output = os.path.join(
            self.temp_dir, f"temp_with_chart_{self.temp_counter}.pptx"
        )
        self.temp_counter += 1

        try:
            success = clone_chart(
                template_path=block.template_path,
                target_path=current_file,
                output_path=temp_output,
                template_slide_idx=0,
                target_slide_idx=slide_index,
                left_inches=chart_left,  # Use centered position for single column stacked
                top_inches=chart_top,  # Start below FIGURE label
                width_inches=cell_width,  # Use constrained/centered width
                height_inches=chart_height,  # Reduced height to fit below FIGURE label
            )

            if success and os.path.exists(temp_output):
                # Add figure name and source if available
                if (
                    hasattr(block, "figure_name")
                    or hasattr(block, "figure_source")
                    or section_title
                ):
                    temp_output = self._add_chart_labels(
                        temp_output,
                        slide_index,
                        cell,
                        getattr(block, "figure_name", ""),
                        getattr(block, "figure_source", ""),
                        section_title=section_title,
                    )

                # Track data for population
                if hasattr(block, "data") and block.data:
                    # Count existing charts on this slide to determine chart index
                    chart_index = len(
                        [
                            u
                            for u in self.pending_data_updates
                            if u[0] == slide_index and u[1] == "chart"
                        ]
                    )
                    # Include axis titles and column keys from block for chart XML update
                    primary_y_axis_title = getattr(block, "primary_y_axis_title", None)
                    secondary_y_axis_title = getattr(
                        block, "secondary_y_axis_title", None
                    )
                    x_axis_title = getattr(block, "x_axis_title", None)
                    y_axis_keys = getattr(block, "y_axis_keys", [])
                    is_multi_axis = getattr(block, "is_multi_axis", True)
                    primary_y_axis_format_code = getattr(
                        block, "primary_y_axis_format_code", None
                    )
                    secondary_y_axis_format_code = getattr(
                        block, "secondary_y_axis_format_code", None
                    )
                    self.pending_data_updates.append(
                        (
                            slide_index,
                            "chart",
                            chart_index,
                            block.data,
                            primary_y_axis_title,
                            secondary_y_axis_title,
                            x_axis_title,
                            y_axis_keys,
                            is_multi_axis,
                            primary_y_axis_format_code,
                            secondary_y_axis_format_code,
                        )
                    )
                    print(
                        f"    [DATA] Tracked chart data for population: {len(block.data)} rows, is_multi_axis={is_multi_axis}"
                    )
                    if primary_y_axis_title or secondary_y_axis_title or x_axis_title:
                        print(
                            f"    [DATA] Axis titles - Primary Y: '{primary_y_axis_title}', Secondary Y: '{secondary_y_axis_title}', X: '{x_axis_title}'"
                        )
                    if y_axis_keys:
                        print(f"    [DATA] Y-axis column keys: {y_axis_keys}")
                else:
                    print(
                        f"    [DATA] Chart has no data (hasattr={hasattr(block, 'data')}, data={getattr(block, 'data', [])})"
                    )

                # Clean up old file
                if current_file.startswith("temp_"):
                    os.remove(current_file)
                return temp_output
            else:
                print("    ✗ Failed to clone chart")
                return current_file

        except Exception as e:
            print(f"    ✗ Error cloning chart: {e}")
            return current_file

    def _calculate_actual_table_height(
        self, block: TableBlock, table_width: float
    ) -> float:
        """
        Calculate the actual table height - DELEGATES to TableBlock.calculate_content_based_height.

        This is a thin wrapper to maintain API compatibility.
        The SINGLE SOURCE OF TRUTH is TableBlock.calculate_content_based_height().

        Args:
            block: TableBlock with data
            table_width: Available width for the table in inches

        Returns:
            Table height in inches (excluding labels)
        """
        # Delegate to the single source of truth
        total_height, _ = block.calculate_content_based_height(table_width)
        return total_height

    def _render_table(
        self,
        block: TableBlock,
        cell,
        current_file: str,
        slide_index: int,
        constraints: Optional[SlideConstraints] = None,
        layout_type=None,
        section_title: Optional[str] = None,
    ) -> str:
        """Render a table block."""
        if not block.template_path or not os.path.exists(block.template_path):
            print(f"    ⚠ Template not found: {block.template_path}")
            return current_file

        # Get label dimensions from config
        from app.ppt_engine.ppt_helpers_utils.services.template_config import (
            get_element_dimensions,
        )

        element_dims = get_element_dimensions()
        # Use figure_label_height for both table headings and figure labels (they're the same)
        LABEL_HEIGHT = element_dims.figure_label_height
        # Use table-specific gap for tables (larger than chart gap to prevent overlapping)
        LABEL_GAP = element_dims.table_figure_gap
        SOURCE_LABEL_HEIGHT = element_dims.source_label_height
        SOURCE_GAP = element_dims.source_gap

        # Check if table has a heading label (figure_name) or table source
        has_heading = (hasattr(block, "figure_name") and block.figure_name) or (
            hasattr(block, "label") and block.label
        )
        has_table_source = hasattr(block, "table_source") and block.table_source

        # Track whether rows are trimmed (for height allocation decision)
        rows_trimmed = False

        # Calculate space needed at top (table heading/label)
        top_space = 0.0

        # Add section title space if provided
        if section_title:
            SECTION_TITLE_HEIGHT = LABEL_HEIGHT  # Same height as figure label
            SECTION_TITLE_GAP = LABEL_GAP  # Same gap as figure label for tables
            top_space += SECTION_TITLE_HEIGHT + SECTION_TITLE_GAP

        if has_heading:
            top_space += LABEL_HEIGHT + LABEL_GAP

        # Determine if this is a fixed-height layout (grid cells) vs expandable (full-width)
        # Grid layouts (GRID_2x2, SPLIT, etc.) have fixed dimensions and should trim rows to fit
        # Full-width layouts can expand vertically and should reserve extra space for overflow

        is_fixed_layout = False
        if layout_type:
            # Check if it's a grid or split layout (fixed dimensions)
            layout_name = (
                layout_type.value if hasattr(layout_type, "value") else str(layout_type)
            )
            is_fixed_layout = (
                "grid" in layout_name.lower() or "split" in layout_name.lower()
            )

        print(f"    📐 Layout: {layout_type}, Fixed: {is_fixed_layout}")

        # Calculate space needed at bottom (table source label)
        bottom_space = 0.0
        rows_trimmed = False

        # Calculate table height and allocated height (needed for all paths)
        predicted_height = self._calculate_actual_table_height(block, cell.width)

        # Source is now a row INSIDE the table, so we don't subtract it from available space
        # The source row height is already included in cell.height via estimate_intrinsic_size,
        # and it will be rendered as part of the table content, not as external space
        # bottom_space should be 0 for tables since source is internal
        source_space = 0.0  # Source is inside table, no external space needed
        allocated_height = cell.height - top_space - source_space

        if has_table_source:
            # Reserve space based on predicted height vs available height
            # If table is predicted to overflow, reserve extra space
            if predicted_height > allocated_height:
                overflow = predicted_height - allocated_height

                # SMART TRIMMING LOGIC:
                # - For FULL-WIDTH layouts: Only trim if table exceeds slide bounds
                # - For FIXED layouts (grid/split): Trim if table exceeds allocated cell

                if not is_fixed_layout:
                    # FULL-WIDTH LAYOUT: Check if table fits within the CELL's allocated space
                    # The cell was already positioned by slide orchestrator accounting for KPIs, margins, etc.
                    # Source is INSIDE the table (as a row), so no external source space needed

                    # Available for table content (including source row) = cell - top_space
                    cell_available_for_table = cell.height - top_space

                    # Use small tolerance for floating point comparison (0.01" = ~0.25mm)
                    TOLERANCE = 0.01
                    if predicted_height <= cell_available_for_table + TOLERANCE:
                        # Table fits in cell, no trimming needed
                        # Source is inside table, no external bottom space
                        bottom_space = 0.0
                        print(
                            f'    ✓ FULL-WIDTH: Table fits in cell ({predicted_height:.2f}" ≤ {cell_available_for_table:.2f}"), no trimming'
                        )
                        rows_trimmed = False  # Explicitly track no trimming
                    else:
                        # Table exceeds cell height - check if continuation (split) is enabled
                        # Available for table = cell height - top labels - source
                        available_for_table = cell_available_for_table

                        # Check if continuation (split across slides) is enabled
                        from app.ppt_engine.ppt_helpers_utils.services.template_config import (
                            get_element_dimensions,
                        )

                        element_dims = get_element_dimensions()
                        allow_continuation = element_dims.allow_full_width_overflow

                        # Use TableBlock's centralized method for accurate row fitting calculation
                        # This is the SINGLE SOURCE OF TRUTH for height calculations
                        if (
                            hasattr(block, "data")
                            and block.data
                            and hasattr(block, "get_max_rows_for_available_height")
                        ):
                            total_data_rows = len(block.data)

                            # Get accurate count using content-based row heights (cached)
                            max_rows = block.get_max_rows_for_available_height(
                                available_for_table, cell.width
                            )

                            if max_rows >= total_data_rows:
                                # All rows fit - don't split
                                max_rows = total_data_rows
                                cached = block.get_cached_height()
                                if cached:
                                    total_height, _ = cached
                                    print(
                                        f'    ✓ All {total_data_rows} rows fit (content-based: {total_height:.2f}" ≤ {available_for_table:.2f}")'
                                    )
                            else:
                                print(
                                    f'    📐 Content-based: {max_rows} data rows fit in {available_for_table:.2f}"'
                                )
                        else:
                            # Fallback if method not available
                            max_rows = self._calculate_max_rows_for_height(
                                block, cell.width, available_for_table
                            )

                        # Check minimum reasonable row count
                        if max_rows < 3:
                            max_rows = max(
                                2, max_rows
                            )  # At minimum, header + 2 data rows
                            print(
                                f'    ⚠️  WARNING: Table severely compressed - only {max_rows} rows fit in {available_for_table:.2f}"!'
                            )
                            print(
                                "       Consider reassigning this table to next slide in layout planning phase"
                            )

                        # Trim rows and handle continuation
                        if (
                            hasattr(block, "data")
                            and block.data
                            and len(block.data) > max_rows
                        ):
                            original_count = len(block.data)

                            if allow_continuation:
                                # Store remaining rows for continuation slide
                                remaining_rows = block.data[max_rows:]
                                block.data = block.data[:max_rows]

                                # Store continuation data on the block for later processing
                                if not hasattr(block, "_continuation_data"):
                                    block._continuation_data = []
                                block._continuation_data = remaining_rows
                                block._continuation_id = getattr(
                                    block, "id", f"table_{id(block)}"
                                )
                                block._continuation_label = getattr(
                                    block, "label", "Table"
                                )
                                block._continuation_section = getattr(
                                    block, "section_name", "Continued"
                                )

                                # IMPORTANT: Don't show source on this slide - it continues to next slide
                                # Source label should only appear on the LAST slide with table data
                                block._skip_source = True

                                print(
                                    f"    📄 FULL-WIDTH CONTINUATION: {original_count} rows split → {max_rows} on this slide, {len(remaining_rows)} continue to next slide"
                                )
                            else:
                                # No continuation - just trim
                                block.data = block.data[:max_rows]
                                print(
                                    f"    ✂️  FULL-WIDTH: Table exceeds slide, trimming {original_count} → {max_rows} rows"
                                )

                            rows_trimmed = True  # Track that trimming occurred

                            # Recalculate predicted height after trimming
                            predicted_height = self._calculate_actual_table_height(
                                block=block, table_width=cell.width
                            )

                        # Source is inside table, no external bottom space needed
                        bottom_space = 0.0
                else:
                    # FIXED LAYOUT (grid/split): Must fit in allocated cell
                    # Source is inside table, so no external source space needed
                    available_for_table = allocated_height

                    # Use centralized height calculation from TableBlock
                    if hasattr(block, "get_max_rows_for_available_height"):
                        max_rows = block.get_max_rows_for_available_height(
                            available_for_table, cell.width
                        )
                    else:
                        # Fallback for blocks without the method
                        max_rows = self._calculate_max_rows_for_height(
                            block, cell.width, available_for_table
                        )

                    # Check minimum reasonable row count
                    if max_rows < 3:
                        max_rows = max(2, max_rows)
                        print(
                            f'    ⚠️  WARNING: Table severely compressed - only {max_rows} rows fit in {available_for_table:.2f}"!'
                        )
                        print(
                            "       Consider reassigning this table to next slide in layout planning phase"
                        )

                    # Trim rows
                    if (
                        hasattr(block, "data")
                        and block.data
                        and len(block.data) > max_rows
                    ):
                        original_count = len(block.data)
                        block.data = block.data[:max_rows]
                        print(
                            f'    ✂️  FIXED LAYOUT: Trimming {original_count} → {max_rows} rows to fit {allocated_height:.2f}" cell'
                        )
                        rows_trimmed = True  # Track that trimming occurred

                        # Recalculate predicted height after trimming
                        predicted_height = self._calculate_actual_table_height(
                            block=block, table_width=cell.width
                        )

                    # Source is inside table, no external bottom space needed
                    bottom_space = 0.0
            else:
                # Table fits, source is inside table
                bottom_space = 0.0
        else:
            # Table fits within allocated height, source is inside table
            bottom_space = 0.0
            rows_trimmed = False  # No trimming needed

        print(
            f'    📏 Height: predicted={predicted_height:.3f}", allocated={allocated_height:.3f}", bottom_space={bottom_space:.3f}" (fixed={is_fixed_layout})'
        )

        # Adjust table position to account for top space
        table_top = cell.top + top_space

        # Use slide constraints to determine available width
        constraint_profile = constraints or SlideConstraints()
        slide_width = constraint_profile.slide_width
        margin_right = constraint_profile.margin_right
        right_boundary = slide_width - margin_right

        # Calculate if table would exceed the REAL boundaries
        table_right_edge = cell.left + cell.width
        if table_right_edge > right_boundary:
            # Table would exceed slide, scale it down
            available_width = max(right_boundary - cell.left, 0.0)
            scale_factor = available_width / cell.width

            print(
                f'    ⚠ Table too wide ({cell.width:.2f}"), scaling to {available_width:.2f}" (scale: {scale_factor:.1%})'
            )

            # Adjust width to fit
            cell_width = available_width
        else:
            cell_width = cell.width

        # Calculate available height for the table (excluding labels)
        # Note: Table borders add visual height beyond programmatic height,
        # which is accounted for in source label positioning
        available_table_height = cell.height - top_space - bottom_space

        # Calculate actual table height based on rows, row height, and padding
        # This gives us the exact height the table will be when rendered
        # Use cell width for accurate wrapping calculation
        actual_table_height = self._calculate_actual_table_height(
            block=block, table_width=cell.width
        )

        from app.ppt_engine.ppt_helpers_utils.services.template_config import (
            get_element_dimensions,
        )

        element_dims = get_element_dimensions()

        # CRITICAL FIX: For full-width tables where rows were NOT trimmed,
        # use the actual calculated height instead of constraining to allocated height.
        # This prevents the source label from appearing inside the table.

        # For full-width layouts where no trimming occurred, use actual height
        if (
            not is_fixed_layout
            and not rows_trimmed
            and actual_table_height > available_table_height
        ):
            # Table needs more space than allocated, but we decided not to trim
            # So give it the full height it needs
            cell_height = actual_table_height
            print(
                f'    ✓ FULL-WIDTH (no trim): Using actual height {actual_table_height:.3f}" (exceeds allocated {available_table_height:.3f}")'
            )
        else:
            # Fixed layout OR trimming occurred - use available height
            cell_height = available_table_height

            # Ensure minimum height
            cell_height = max(0.1, cell_height)

            # Log table sizing
            if actual_table_height > cell_height:
                print(
                    f'    ⚠ Table constrained: calculated={actual_table_height:.3f}", allocated={available_table_height:.3f}"'
                )
            else:
                print(
                    f'    ✓ Table fits: calculated={actual_table_height:.3f}", allocated={available_table_height:.3f}"'
                )

        temp_output = os.path.join(
            self.temp_dir, f"temp_with_table_{self.temp_counter}.pptx"
        )
        self.temp_counter += 1

        try:
            success = clone_table(
                template_path=block.template_path,
                target_path=current_file,
                output_path=temp_output,
                template_slide_idx=0,
                target_slide_idx=slide_index,
                left_inches=cell.left,
                top_inches=table_top,  # Use adjusted position
                width_inches=cell_width,  # Use constrained width
                height_inches=cell_height,
            )

            # CRITICAL: After cloning, ensure the table shape height matches our constraint
            # The template may have its own row heights, so we need to enforce ours
            if success:
                from pptx import Presentation
                from pptx.enum.text import MSO_AUTO_SIZE

                prs = Presentation(temp_output)
                if slide_index < len(prs.slides):
                    slide = prs.slides[slide_index]
                    # Find the table shape (should be the last added shape)
                    for shape in reversed(slide.shapes):
                        if shape.has_table:
                            # Set table shape height to our calculated cell_height
                            # Convert from inches to EMUs
                            shape.height = int(cell_height * 914400)
                            print(
                                f'    ✓ Enforced table shape height: {cell_height:.3f}" ({shape.height} EMUs)'
                            )

                            # Configure table cells to prevent auto-expansion
                            # This is a precautionary measure before data population
                            try:
                                # Get table configuration from config
                                from app.ppt_engine.ppt_helpers_utils.services.template_config import (
                                    get_element_dimensions,
                                )

                                elem_dims = get_element_dimensions()

                                table = shape.table
                                for row in table.rows:
                                    for table_cell in row.cells:
                                        text_frame = table_cell.text_frame
                                        # Apply word wrap setting from config
                                        text_frame.word_wrap = elem_dims.table_word_wrap
                                        # Apply auto-size setting from config
                                        if elem_dims.table_auto_size_disabled:
                                            text_frame.auto_size = MSO_AUTO_SIZE.NONE
                                        # Set cell margins from config
                                        text_frame.margin_top = Pt(
                                            elem_dims.table_cell_margin_top
                                        )
                                        text_frame.margin_bottom = Pt(
                                            elem_dims.table_cell_margin_bottom
                                        )
                                        text_frame.margin_left = Pt(
                                            elem_dims.table_cell_margin_left
                                        )
                                        text_frame.margin_right = Pt(
                                            elem_dims.table_cell_margin_right
                                        )
                                print(
                                    "    ✓ Configured table cells for fixed dimensions (from config)"
                                )
                            except Exception as cell_ex:
                                print(
                                    f"    ⚠ Could not pre-configure table cells: {cell_ex}"
                                )

                            prs.save(temp_output)
                            break

            if success and os.path.exists(temp_output):
                # Check if this is a continuation table.
                # IMPORTANT: Continuation slides SHOULD repeat the table header row.
                # Continuation status is about slide semantics (hide headings/labels), not header presence.
                is_continuation = bool(getattr(block, "is_continuation", False)) or (
                    hasattr(block, "has_header") and not block.has_header
                )

                # Add table heading if figure_name exists (preferred) or label exists
                # SKIP heading for continuation tables - they don't get labels
                heading_text = None
                if not is_continuation:
                    if hasattr(block, "figure_name") and block.figure_name:
                        heading_text = block.figure_name
                    elif hasattr(block, "label") and block.label:
                        heading_text = block.label

                if heading_text or (section_title and not is_continuation):
                    temp_output = self._add_table_heading(
                        pptx_file=temp_output,
                        slide_index=slide_index,
                        cell_frame=cell,
                        heading_text=heading_text,
                        section_title=section_title if not is_continuation else None,
                    )
                elif is_continuation:
                    print("    📋 Continuation table - skipping heading/labels")

                # NOTE: Source label will be added AFTER data population in populate_table_data()
                # to ensure proper Z-order (source on top of table) and accurate positioning
                # based on the actual enforced table height

                # Track data for population
                if hasattr(block, "data") and block.data:
                    # Count existing tables on this slide to determine table index
                    table_index = len(
                        [
                            u
                            for u in self.pending_data_updates
                            if u[0] == slide_index and u[1] == "table"
                        ]
                    )
                    # Include source label info and cell frame boundaries for post-population addition
                    # Don't add source if table continues to next slide (_skip_source flag)
                    skip_source = getattr(block, "_skip_source", False)
                    table_source = (
                        None if skip_source else getattr(block, "table_source", None)
                    )
                    if skip_source:
                        print(
                            "    📋 Source label skipped - table continues on next slide"
                        )
                    # Position source BELOW the table content (not relative to cell boundary)
                    # This ensures source follows the actual table, even if table renders taller
                    # table_top + cell_height = table bottom
                    # + SOURCE_GAP = gap between table and source
                    reserved_source_top = table_top + cell_height + SOURCE_GAP
                    # Header rows must be preserved on continuation slides.
                    # `skip_header` is only for cases where we explicitly want to remove the header row.
                    # Continuation is controlled via `is_continuation`, not by removing headers.
                    skip_header = bool(getattr(block, "skip_header", False))
                    if (
                        not skip_header
                        and hasattr(block, "has_header")
                        and block.has_header is False
                    ):
                        skip_header = True

                    # Get cached column widths from TableBlock (for variable column width rendering)
                    cached_column_widths = block.get_cached_column_widths()
                    cell_frame_info = {
                        "top": cell.top,
                        "height": cell.height,
                        "bottom": cell.top + cell.height,
                        "left": cell.left,
                        "width": cell.width,
                        "reserved_source_top": reserved_source_top,  # Pre-calculated source position
                        "table_top": table_top,
                        "table_height": cell_height,  # The constrained table height
                        "skip_header": skip_header,  # True only when header should be removed
                        "column_widths": cached_column_widths,  # Variable column widths (first col expanded)
                    }
                    # Get cached row heights from TableBlock (single source of truth for height calculation)
                    cached_heights = block.get_cached_height()
                    cached_row_heights = cached_heights[1] if cached_heights else None
                    self.pending_data_updates.append(
                        (
                            slide_index,
                            "table",
                            table_index,
                            block.data,
                            table_source,
                            cell_frame_info,
                            cached_row_heights,
                        )
                    )
                    if skip_header:
                        print("    📋 Table configured to skip header row")
                    print(
                        f"    [DATA] Tracked table data for population: {len(block.data)} rows"
                    )
                else:
                    print(
                        f"    [DATA] Table has no data (hasattr={hasattr(block, 'data')}, data={getattr(block, 'data', [])})"
                    )

                if current_file.startswith("temp_"):
                    os.remove(current_file)
                return temp_output
            else:
                print("    ✗ Failed to clone table")
                return current_file

        except Exception as e:
            print(f"    ✗ Error cloning table: {e}")
            return current_file

    def _render_text(
        self,
        block: TextBlock,
        cell,
        current_file: str,
        slide_index: int,
        section_title: Optional[str] = None,
    ) -> str:
        """Render a text block with heading on top if present.

        Args:
            block: TextBlock to render
            cell: Cell frame for positioning
            current_file: Current PPTX file path
            slide_index: Index of slide being rendered
            section_title: Optional section title to use as heading (overrides block.text)
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.dml.color import RGBColor

            prs = Presentation(current_file)
            slide = prs.slides[slide_index]

            # Determine heading text
            # 1. If section_title is provided (first element), use it
            # 2. If section_title is None (non-first element), check if block.text matches section name
            #    - If it matches, skip heading (section title already shown on first element)
            #    - If it's different, render block.text as heading
            block_section_id = getattr(block, "section_id", None)
            block_section_name = getattr(block, "section_name", None)

            heading_text = None
            if section_title:
                # First element: use section title
                heading_text = section_title
            elif block.text and block.text.strip():
                # Non-first element: check if block.text is same as section name
                if block_section_name and block.text.strip() == block_section_name:
                    # Skip heading - section title already rendered on first element
                    # (whether that was a chart or commentary)
                    heading_text = None
                    print(
                        f"    📝 Skipping duplicate heading (matches section name): '{block.text.strip()}'"
                    )
                else:
                    # Different text, render as heading
                    heading_text = block.text.strip()

            has_heading = heading_text is not None

            print(
                f"    📝 Commentary: has_heading={has_heading}, section_title={section_title}"
            )
            if has_heading:
                print(f"       Heading text: '{heading_text}'")
            print(f"       Description paragraphs: {len(block.bullet_points)}")

            # Get heading dimensions from config (use section title gap for commentary headings)
            from app.ppt_engine.ppt_helpers_utils.services.template_config import (
                get_element_dimensions,
            )

            element_dims = get_element_dimensions()
            HEADING_HEIGHT = element_dims.figure_label_height
            HEADING_GAP = (
                element_dims.section_title_gap
            )  # Use larger gap for section titles

            # Use cell's left and width directly - no changes!
            padding = 0.05
            left = Inches(cell.left + padding)
            width = Inches(cell.width - 2 * padding)

            # Calculate positions first
            if has_heading:
                # Heading at top (no padding at top), description below
                heading_top = cell.top  # No padding at top to stick to top
                desc_top = cell.top + HEADING_HEIGHT + HEADING_GAP
                desc_height = (
                    cell.height - HEADING_HEIGHT - HEADING_GAP - padding
                )  # Only bottom padding
                print(
                    f'    📐 Positions: heading_top={heading_top:.3f}", desc_top={desc_top:.3f}", desc_height={desc_height:.3f}"'
                )
            else:
                # No heading, use full cell for description
                desc_top = cell.top + padding
                desc_height = cell.height - 2 * padding
                print(
                    f'    📐 No heading: desc_top={desc_top:.3f}", desc_height={desc_height:.3f}"'
                )

            # Create description text box FIRST
            top = Inches(desc_top)
            height = Inches(desc_height)

            print(
                f'    📄 Creating description box: top={desc_top:.3f}" ({top}), height={desc_height:.3f}" ({height})'
            )

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame

            # Configure text frame for professional appearance
            tf.word_wrap = True
            tf.auto_size = None
            tf.margin_left = Inches(0.00)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(
                0
            )  # No top margin - description sticks to top of box
            tf.margin_bottom = Inches(0.02)

            # Calculate appropriate font size based on cell size for better readability
            # Use more conservative scaling to maintain readability
            scale_factor = min(
                cell.width / 3.5, cell.height / 2.5
            )  # More conservative scaling
            adjusted_font_size = max(
                11, min(16, block.font_size * scale_factor)
            )  # Better font size range

            # Clear default paragraph
            tf.clear()

            # Add description text (bullet_points) - heading is already added separately at top
            # No need for section title inside description box anymore
            for i, text_line in enumerate(block.bullet_points):
                p = (
                    tf.add_paragraph()
                    if i > 0
                    else (
                        tf.paragraphs[0]
                        if len(tf.paragraphs) > 0
                        else tf.add_paragraph()
                    )
                )
                clean_text = text_line.strip()

                p.text = clean_text
                p.font.size = Pt(10.5)  # User specified: 10.5pt
                p.font.name = "Calibre (Body)"  # Calibri (Body)
                p.font.color.rgb = RGBColor(
                    69, 79, 81
                )  # User specified: RGB(69, 79, 81)
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(0)  # No space before
                p.space_after = Pt(0)  # No space after - tight spacing
                p.line_spacing = 1.0  # Single line spacing

            # Set vertical alignment to TOP - text starts from top
            tf.vertical_anchor = MSO_ANCHOR.TOP

            # Now add heading box if it exists (AFTER description, so it appears on top)
            if has_heading:
                heading_top_inches = Inches(heading_top)
                heading_height_inches = Inches(HEADING_HEIGHT)

                print(
                    f'    ✏️  Creating heading: top={heading_top:.3f}" ({heading_top_inches}), height={HEADING_HEIGHT:.3f}" ({heading_height_inches})'
                )
                print(f"    ✏️  Heading text: '{heading_text}'")

                # Create heading text box
                heading_box = slide.shapes.add_textbox(
                    left, heading_top_inches, width, heading_height_inches
                )
                htf = heading_box.text_frame
                htf.word_wrap = True
                htf.auto_size = None
                htf.margin_left = Inches(0)
                htf.margin_right = Inches(0)
                htf.margin_top = Inches(0)
                htf.margin_bottom = Inches(0)

                # Add heading text
                hp = htf.paragraphs[0]
                hp.text = heading_text
                hp.font.size = Pt(12) if slide_index == 0 else Pt(16)
                hp.font.bold = False
                # Use Calibri for first page, Financier Display for other pages
                hp.font.name = (
                    "Calibre" if slide_index == 0 else "Financier Display (Headings)"
                )
                hp.font.color.rgb = RGBColor(70, 82, 84)
                hp.alignment = PP_ALIGN.LEFT

                print(
                    f'    ✓ Heading box created at position {heading_top:.3f}" (should be at TOP)'
                )

            # Save with proper temp file naming - add unique timestamp to avoid conflicts
            import time

            timestamp = int(time.time() * 1000)  # milliseconds for uniqueness
            temp_output = os.path.join(
                self.temp_dir, f"temp_with_text_{self.temp_counter}_{timestamp}.pptx"
            )
            self.temp_counter += 1

            # Try to save, with retry logic for permission issues
            max_retries = 3
            for attempt in range(max_retries):
                try:
                    prs.save(temp_output)
                    break
                except PermissionError as e:
                    if attempt < max_retries - 1:
                        import time

                        time.sleep(0.1)  # Wait 100ms before retry
                        temp_output = os.path.join(
                            self.temp_dir,
                            f"temp_with_text_{self.temp_counter}_{timestamp}_{attempt}.pptx",
                        )
                    else:
                        print(
                            f"    ⚠ Permission error saving text, using current file: {e}"
                        )
                        temp_output = current_file  # Fallback to current file
                        break

            # Cleanup previous temp file
            if current_file.startswith("temp_"):
                try:
                    os.remove(current_file)
                except:
                    pass

            heading_info = (
                f" with heading '{block.text.strip()}'" if has_heading else ""
            )
            print(
                f"    ✓ Commentary added: {len(block.bullet_points)} paragraphs{heading_info}"
            )
            return temp_output

        except Exception as e:
            print(f"    ✗ Error adding text: {e}")
            import traceback

            traceback.print_exc()
            return current_file

    def _cleanup_temp_files(self):
        """Remove temporary files."""
        for i in range(self.temp_counter + 10):  # Extra margin
            for prefix in [
                "temp_render_",
                "temp_with_chart_",
                "temp_with_table_",
                "temp_with_text_",
            ]:
                temp_file = os.path.join(self.temp_dir, f"{prefix}{i}.pptx")
                if os.path.exists(temp_file):
                    try:
                        os.remove(temp_file)
                    except:
                        pass

    def _add_chart_labels(
        self,
        pptx_file: str,
        slide_index: int,
        cell,
        figure_name: str,
        figure_source: str,
        section_title: Optional[str] = None,
    ) -> str:
        """Add figure name and source labels to chart"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            # Get label dimensions from config
            from app.ppt_engine.ppt_helpers_utils.services.template_config import (
                get_element_dimensions,
            )

            element_dims = get_element_dimensions()

            prs = Presentation(pptx_file)
            slide = prs.slides[slide_index]

            FIGURE_HEIGHT = element_dims.figure_label_height
            FIGURE_GAP = element_dims.figure_gap
            SECTION_TITLE_GAP = element_dims.section_title_gap

            # Track current vertical position for labels
            current_top = cell.top

            # Add section title above everything (if provided)
            if section_title:
                section_box = slide.shapes.add_textbox(
                    Inches(cell.left),
                    Inches(current_top),
                    Inches(cell.width),
                    Inches(FIGURE_HEIGHT),
                )
                section_frame = section_box.text_frame
                section_frame.text = section_title
                section_frame.margin_left = Inches(0)
                section_frame.margin_right = Inches(0)
                section_frame.margin_top = Inches(0)
                section_frame.margin_bottom = Inches(0)
                p = section_frame.paragraphs[0]
                # Match styling from _render_text method (lines 866-872)
                p.font.size = Pt(12) if slide_index == 0 else Pt(16)
                p.font.bold = False
                p.font.name = (
                    "Calibre" if slide_index == 0 else "Financier Display (Headings)"
                )
                p.font.color.rgb = RGBColor(70, 82, 84)
                p.alignment = PP_ALIGN.LEFT

                # Move down for next label - use section_title_gap for space after section title
                current_top += FIGURE_HEIGHT + SECTION_TITLE_GAP

            # Add figure name above chart (if provided)
            if figure_name:
                name_box = slide.shapes.add_textbox(
                    Inches(cell.left),
                    Inches(current_top),
                    Inches(cell.width),
                    Inches(FIGURE_HEIGHT),
                )
                name_frame = name_box.text_frame
                name_frame.text = figure_name
                name_frame.margin_left = Inches(0)
                name_frame.margin_right = Inches(0)
                name_frame.margin_top = Inches(0)
                name_frame.margin_bottom = Inches(0)
                p = name_frame.paragraphs[0]
                p.font.size = Pt(element_dims.figure_label_font_size)
                p.font.bold = False
                p.font.name = element_dims.figure_label_font_name
                p.font.color.rgb = RGBColor(101, 112, 113)
                p.alignment = PP_ALIGN.LEFT

            # Add source at bottom of cell (if provided) - positioned within cell bounds
            if figure_source:
                SOURCE_WIDTH = 5.969  # 15.16 cm in inches
                SOURCE_HEIGHT = element_dims.source_label_height
                SOURCE_GAP = element_dims.source_gap

                # Position source at bottom of cell (within cell bounds)
                source_top = cell.top + cell.height - SOURCE_HEIGHT

                source_box = slide.shapes.add_textbox(
                    Inches(cell.left),
                    Inches(source_top),  # Within cell bounds at bottom
                    Inches(SOURCE_WIDTH),
                    Inches(SOURCE_HEIGHT),
                )
                source_frame = source_box.text_frame
                source_frame.text = figure_source
                source_frame.margin_left = Inches(0)
                source_frame.margin_right = Inches(0)
                source_frame.margin_top = Inches(0)
                source_frame.margin_bottom = Inches(0)
                p = source_frame.paragraphs[0]
                p.font.size = Pt(element_dims.source_label_font_size)
                p.font.name = element_dims.source_label_font_name
                p.font.color.rgb = RGBColor(88, 89, 91)  # CBRE Gray
                p.alignment = PP_ALIGN.LEFT

            # Save updated presentation
            temp_labeled = os.path.join(
                self.temp_dir, f"temp_labeled_{self.temp_counter}.pptx"
            )
            self.temp_counter += 1
            prs.save(temp_labeled)

            # Clean up
            if pptx_file.startswith("temp_"):
                try:
                    os.remove(pptx_file)
                except:
                    pass

            label_info = []
            if section_title:
                label_info.append(f"section: {section_title}")
            if figure_name:
                label_info.append(f"figure: {figure_name[:30]}")
            if figure_source:
                label_info.append("source")
            print(f"    ✓ Added chart labels: {', '.join(label_info)}")
            return temp_labeled

        except Exception as e:
            print(f"    ⚠ Could not add chart labels: {e}")
            return pptx_file

    def _add_table_heading(
        self,
        pptx_file: str,
        slide_index: int,
        cell_frame,
        heading_text: str = None,
        section_title: Optional[str] = None,
    ) -> str:
        """Add heading label above table

        Args:
            pptx_file: Path to PPT file
            slide_index: Slide index
            cell_frame: CellFrame object with layout position (left, top, width, height)
            heading_text: Text for the heading
            section_title: Optional section title to add above heading
        """
        if not heading_text and not section_title:
            return pptx_file

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            prs = Presentation(pptx_file)
            slide = prs.slides[slide_index]

            # Get table heading dimensions from config
            from app.ppt_engine.ppt_helpers_utils.services.template_config import (
                get_element_dimensions,
            )

            element_dims = get_element_dimensions()
            # Use figure_label_height for table headings (same as figure labels)
            HEADING_HEIGHT = element_dims.figure_label_height
            # Use table-specific gap for tables (larger than chart gap to prevent overlapping)
            HEADING_GAP = element_dims.table_figure_gap
            SECTION_TITLE_GAP = element_dims.section_title_gap

            # Track current vertical position for labels
            current_top = cell_frame.top

            # Add section title above everything (if provided)
            if section_title:
                section_box = slide.shapes.add_textbox(
                    Inches(cell_frame.left),
                    Inches(current_top),
                    Inches(cell_frame.width),
                    Inches(HEADING_HEIGHT),
                )
                section_frame = section_box.text_frame
                section_frame.text = section_title
                section_frame.margin_left = Inches(0)
                section_frame.margin_right = Inches(0)
                section_frame.margin_top = Inches(0)
                section_frame.margin_bottom = Inches(0)
                section_frame.word_wrap = False
                section_frame.vertical_anchor = 1  # Top alignment

                p = section_frame.paragraphs[0]
                # Match styling from _render_text method (lines 866-872)
                p.font.size = Pt(12) if slide_index == 0 else Pt(16)
                p.font.bold = False
                p.font.name = (
                    "Calibre" if slide_index == 0 else "Financier Display (Headings)"
                )
                p.font.color.rgb = RGBColor(70, 82, 84)
                p.alignment = PP_ALIGN.LEFT

                # Move down for next label - use section_title_gap for space after section title
                current_top += HEADING_HEIGHT + SECTION_TITLE_GAP

            # Add heading above table (if provided)
            if heading_text:
                heading_box = slide.shapes.add_textbox(
                    Inches(cell_frame.left),
                    Inches(current_top),
                    Inches(cell_frame.width),
                    Inches(HEADING_HEIGHT),
                )
                heading_frame = heading_box.text_frame
                heading_frame.text = heading_text
                heading_frame.margin_left = Inches(0)
                heading_frame.margin_right = Inches(0)
                heading_frame.margin_top = Inches(0)
                heading_frame.margin_bottom = Inches(0)
                heading_frame.word_wrap = False
                heading_frame.vertical_anchor = 1  # Top alignment

                p = heading_frame.paragraphs[0]
                # Match chart label style (use config for font settings)
                p.font.size = Pt(element_dims.table_heading_font_size)
                p.font.bold = False
                p.font.name = element_dims.table_heading_font_name
                p.font.color.rgb = RGBColor(101, 112, 113)
                p.alignment = PP_ALIGN.LEFT

            # Save updated presentation
            temp_labeled = os.path.join(
                self.temp_dir, f"temp_table_labeled_{self.temp_counter}.pptx"
            )
            self.temp_counter += 1
            prs.save(temp_labeled)

            # Clean up
            if pptx_file.startswith("temp_"):
                try:
                    os.remove(pptx_file)
                except:
                    pass

            label_info = []
            if section_title:
                label_info.append(f"section: {section_title}")
            if heading_text:
                label_info.append(f"heading: {heading_text}")
            print(f"    ✓ Added table labels: {', '.join(label_info)}")
            return temp_labeled

        except Exception as e:
            print(f"    ⚠ Could not add table heading/section title: {e}")
            return pptx_file

    def _add_table_labels(
        self,
        pptx_file: str,
        slide_index: int,
        cell_frame,
        figure_name: str,
        table_source: str,
        has_heading: bool = False,
        table_top: float = None,
        table_height: float = None,
    ) -> str:
        """Add figure name and source labels to table (for submarket/snapshot tables)

        Args:
            pptx_file: Path to PPT file
            slide_index: Slide index
            cell_frame: CellFrame object with layout position (left, top, width, height)
            figure_name: Figure name text
            table_source: Source label text
            has_heading: Whether table already has a heading
            table_top: Calculated table top position
            table_height: Calculated table height
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            # Get label dimensions from config
            from app.ppt_engine.ppt_helpers_utils.services.template_config import (
                get_element_dimensions,
            )

            element_dims = get_element_dimensions()
            FIGURE_HEIGHT = element_dims.figure_label_height
            # Use table-specific gap for tables (larger than chart gap to prevent overlapping)
            FIGURE_GAP = element_dims.table_figure_gap
            # Use figure_label_height for table headings (same as figure labels)
            HEADING_HEIGHT = element_dims.figure_label_height
            HEADING_GAP = element_dims.table_figure_gap

            prs = Presentation(pptx_file)
            slide = prs.slides[slide_index]

            # Add figure name above table (if provided)
            # If table has heading, place figure name below heading; otherwise at cell_frame.top
            if figure_name:
                if has_heading:
                    # Place figure name below heading
                    figure_top = cell_frame.top + HEADING_HEIGHT + HEADING_GAP
                else:
                    # Place figure name at cell_frame.top
                    figure_top = cell_frame.top

                name_box = slide.shapes.add_textbox(
                    Inches(cell_frame.left),
                    Inches(figure_top),
                    Inches(cell_frame.width),
                    Inches(FIGURE_HEIGHT),
                )
                name_frame = name_box.text_frame
                name_frame.text = figure_name
                name_frame.margin_left = Inches(0)
                name_frame.margin_right = Inches(0)
                name_frame.margin_top = Inches(0)
                name_frame.margin_bottom = Inches(0)
                p = name_frame.paragraphs[0]
                p.font.size = Pt(element_dims.figure_label_font_size)
                p.font.bold = False
                p.font.name = element_dims.figure_label_font_name
                p.font.color.rgb = RGBColor(101, 112, 113)
                p.alignment = PP_ALIGN.LEFT

            # Add source at bottom of cell (if provided) - positioned below table with proper spacing
            if table_source:
                SOURCE_WIDTH = 5.969  # 15.16 cm in inches
                SOURCE_HEIGHT = element_dims.source_label_height
                SOURCE_GAP = element_dims.source_gap

                # ALWAYS use actual table shape position after rendering
                # This ensures source is positioned correctly based on the ACTUAL constrained table height,
                # not the predicted height (which may be larger before trimming/constraining)
                table_shape = None
                for shape in slide.shapes:
                    if shape.has_table:
                        table_shape = shape
                        break

                if table_shape:
                    # Use actual table position (the constrained/rendered height)
                    table_bottom = table_shape.top + table_shape.height
                    source_top = table_bottom + SOURCE_GAP

                    # Ensure source doesn't exceed cell_frame bounds
                    max_source_top = cell_frame.top + cell_frame.height - SOURCE_HEIGHT
                    if source_top + SOURCE_HEIGHT > cell_frame.top + cell_frame.height:
                        source_top = max_source_top
                        print(
                            "    ⚠ Table extends beyond allocated space, positioning source at cell bottom"
                        )

                    print("    ✓ Source positioned from ACTUAL table bottom")
                    print(
                        f'       Table bottom: {table_bottom:.3f}", Source at: {source_top:.3f}", Gap: {SOURCE_GAP:.3f}"'
                    )
                else:
                    # Fallback: position at bottom of cell_frame if no table found
                    source_top = cell_frame.top + cell_frame.height - SOURCE_HEIGHT
                    print("    ⚠ No table found, positioning source at cell bottom")

                source_box = slide.shapes.add_textbox(
                    Inches(cell_frame.left),
                    Inches(source_top),  # Within cell bounds at bottom
                    Inches(SOURCE_WIDTH),
                    Inches(SOURCE_HEIGHT),
                )
                source_frame = source_box.text_frame
                source_frame.text = table_source
                source_frame.margin_left = Inches(0)
                source_frame.margin_right = Inches(0)
                source_frame.margin_top = Inches(0)
                source_frame.margin_bottom = Inches(0)
                p = source_frame.paragraphs[0]
                p.font.size = Pt(element_dims.source_label_font_size)
                p.font.name = element_dims.source_label_font_name
                p.font.color.rgb = RGBColor(88, 89, 91)  # CBRE Gray
                p.alignment = PP_ALIGN.LEFT

            # Save updated presentation
            temp_labeled = os.path.join(
                self.temp_dir, f"temp_table_labeled_{self.temp_counter}.pptx"
            )
            self.temp_counter += 1
            prs.save(temp_labeled)

            # Clean up
            if pptx_file.startswith("temp_"):
                try:
                    os.remove(pptx_file)
                except:
                    pass

            print(
                f"    ✓ Added table labels: {figure_name[:30] if figure_name else 'source only'}..."
            )
            return temp_labeled

        except Exception as e:
            print(f"    ⚠ Could not add table labels: {e}")
            return pptx_file

    def _calculate_max_rows_for_height(
        self, block: TableBlock, table_width: float, max_height: float
    ) -> int:
        """
        Calculate maximum number of rows that can fit in a given height.

        Uses the shared content-based height calculation to ensure consistency
        with size estimation. Progressively adds rows until max_height is exceeded.

        Args:
            block: TableBlock with data
            table_width: Available width for the table in inches
            max_height: Maximum height available in inches

        Returns:
            Maximum number of data rows that fit (excluding header)
        """
        from app.ppt_engine.ppt_helpers_utils.services.template_config import (
            get_element_dimensions,
        )

        element_dims = get_element_dimensions()

        # Get table data
        data = block.data if hasattr(block, "data") and block.data else []
        if not data:
            return 0

        # Get row heights using the shared calculation method
        _, row_heights = block.calculate_content_based_height(table_width)

        if not row_heights:
            return 0

        # Account for row gap padding and border overhead - consistent with content_height_calculator
        row_gap_padding = element_dims.table_row_gap_padding
        border_overhead = element_dims.table_border_overhead

        # Calculate available height for rows (excluding borders)
        # Formula from content_height_calculator: total = sum(row_heights) + border_overhead + (num_rows - 1) * row_gap
        available_height = max_height - border_overhead

        # Header row (if present)
        header_idx = 0
        if block.has_header and len(row_heights) > 0:
            header_height = row_heights[header_idx]
            available_height -= header_height
            available_height -= row_gap_padding  # Gap after header
            data_row_start = 1
        else:
            data_row_start = 0

        if available_height <= 0:
            return 0

        # Check if source row is present (always the last row if present)
        # Source row height is 0.21" (0.18" row + 0.025" margin)
        has_source = hasattr(block, "table_source") and block.table_source
        source_row_idx = len(row_heights) - 1 if has_source and len(row_heights) > 0 else None
        
        # Subtract source row height from available height if present
        if source_row_idx is not None and source_row_idx >= data_row_start:
            source_row_height = row_heights[source_row_idx]
            available_height -= source_row_height
            # Exclude source row from iteration
            max_idx = source_row_idx
        else:
            max_idx = len(row_heights)
        
        # Add data rows one by one until we exceed available height
        # (excluding source row which is always last)
        rows_that_fit = 0
        for idx in range(data_row_start, max_idx):
            row_height = row_heights[idx]
            # Include row gap padding after this row (except for last row)
            is_last_possible_row = (idx == max_idx - 1)
            needed_height = row_height + (0 if is_last_possible_row else row_gap_padding)

            if available_height >= needed_height:
                rows_that_fit += 1
                available_height -= needed_height
            else:
                break

        return max(1, rows_that_fit)  # Always keep at least 1 row


# ============================================================================
# CONVENIENCE FUNCTION
# ============================================================================


def orchestrate_and_render(
    section: Section,
    output_path: str,
    constraints=None,
    property_sub_type: Optional[str] = None,
) -> str:
    """
    One-shot function: orchestrate layout AND render to PowerPoint.

    Args:
        section: Section with content blocks
        output_path: Path for output PPTX file
        constraints: Optional custom SlideConstraints
        property_sub_type: Template selection hint

    Returns:
        Path to generated PPTX file
    """
    # Orchestrate (pure function, generates layout spec)
    orchestrator = SlideOrchestrator(constraints=constraints)
    layouts = orchestrator.orchestrate_section(section)

    # Render (side effect, creates PPTX file)
    renderer = OrchestratorRenderer()
    return renderer.render_to_pptx(
        section, layouts, output_path, property_sub_type=property_sub_type
    )


# ============================================================================
# COMPLETE EXAMPLES
# ============================================================================


def example_complete_workflow():
    """Complete example: orchestrate + render to actual PPTX."""
    print("\n" + "=" * 70)
    print(" " * 15 + "COMPLETE WORKFLOW EXAMPLE")
    print("=" * 70)

    # Create content blocks with REAL templates
    blocks = [
        ChartBlock(
            id="chart_asking_rents",
            template_path="templates/Asking Rents.pptx",
            intrinsic_aspect_ratio=16 / 9,
            priority=1,
        ),
        ChartBlock(
            id="chart_leasing",
            template_path="templates/Leasing Activity Trend.pptx",
            intrinsic_aspect_ratio=16 / 9,
            priority=1,
        ),
        TextBlock(
            id="summary_text",
            text="Q3 2024 Market Summary",
            bullet_points=[
                "Strong performance across all metrics",
                "Leasing activity exceeded expectations",
                "Asking rents continue upward trend",
            ],
            font_size=14.0,
        ),
        TableBlock(
            id="market_stats",
            template_path="templates/table_Market_stats.pptx",
            rows=8,
            columns=6,
            can_split_rows=True,
            priority=2,
        ),
    ]

    # Create section
    section = Section(id="q3_analysis", title="Q3 2024 Analysis", blocks=blocks)

    # Orchestrate and render in one call
    output_file = orchestrate_and_render(
        section=section, output_path="orchestrated_presentation.pptx"
    )

    print(f"\n{'=' * 70}")
    print(f"✓ COMPLETE! Open: {output_file}")
    print(f"{'=' * 70}\n")

    return output_file


def example_two_sections():
    """Example with multiple sections."""
    print("\n" + "=" * 70)
    print(" " * 15 + "MULTI-SECTION EXAMPLE")
    print("=" * 70)

    # Section 1: Market Overview
    section1_blocks = [
        ChartBlock(
            id="chart1", template_path="templates/Asking Rents.pptx", priority=2
        ),
        ChartBlock(
            id="chart2", template_path="templates/Availability Rates.pptx", priority=2
        ),
    ]

    section1 = Section(id="overview", title="Market Overview", blocks=section1_blocks)

    # Section 2: Detailed Analysis
    section2_blocks = [
        TextBlock(
            id="intro",
            text="Detailed Market Analysis",
            bullet_points=["Q3 performance metrics", "Year-over-year comparison"],
            font_size=14.0,
        ),
        ChartBlock(
            id="chart3",
            template_path="templates/Leasing Activity Trend.pptx",
            priority=3,
        ),
        TableBlock(
            id="table1",
            template_path="templates/table_Market_stats.pptx",
            rows=6,
            columns=5,
            priority=1,
        ),
    ]

    section2 = Section(id="analysis", title="Detailed Analysis", blocks=section2_blocks)

    # Render section 1
    output1 = orchestrate_and_render(section=section1, output_path="temp_section1.pptx")

    # Render section 2 (starts fresh)
    output2 = orchestrate_and_render(section=section2, output_path="temp_section2.pptx")

    # Merge sections (simple version - just use section2 as final)
    # In production, you'd merge the two PPTX files
    final_output = "orchestrated_multi_section.pptx"
    os.rename(output2, final_output)

    if os.path.exists(output1):
        os.remove(output1)

    print(f"\n{'=' * 70}")
    print(f"✓ COMPLETE! Open: {final_output}")
    print(f"{'=' * 70}\n")

    return final_output


if __name__ == "__main__":
    print("\n" + "=" * 70)
    print(" " * 10 + "ORCHESTRATOR RENDERER - COMPLETE EXAMPLES")
    print("=" * 70)

    # Run complete workflow
    example_complete_workflow()

    # Run multi-section example
    # example_two_sections()

    print("\n" + "=" * 70)
    print("✓ All examples complete!")
    print("✓ Check 'orchestrated_*.pptx' files")
    print("=" * 70 + "\n")
